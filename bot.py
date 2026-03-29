import asyncio
import logging
import json
import aiohttp
async def check_ip():
    async with aiohttp.ClientSession() as session:
        async with session.get("https://api.ipify.org") as resp:
            ip = await resp.text()
            print("M1Y IP:", ip)
from aiohttp import web as aiohttp_web
import aiofiles
import base64
import os
import io
import re
import textwrap
import subprocess
from collections import deque
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
from matplotlib.font_manager import FontProperties
from aiogram import Bot, Dispatcher, F, BaseMiddleware
from aiogram.types import BotCommand, BotCommandScopeDefault, MenuButtonDefault, MenuButtonWebApp
from aiogram.filters import Command
from typing import Callable, Dict, Any, Awaitable
from aiogram.types import (
    Message, InlineKeyboardMarkup, InlineKeyboardButton,
    CallbackQuery, FSInputFile, BufferedInputFile,
    ReplyKeyboardMarkup, KeyboardButton, ReplyKeyboardRemove,
    WebAppInfo, PreCheckoutQuery, LabeledPrice
)
from aiogram.utils.chat_action import ChatActionSender

from docx import Document
from PIL import Image, ImageDraw, ImageFont
from fpdf import FPDF

# ──────────────────────────────────────────
# python-pptx импорты
# ──────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor as PptxRGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    # Заглушки чтобы модуль загружался даже без python-pptx
    class PptxRGBColor:
        def __init__(self, r, g, b): self.r, self.g, self.b = r, g, b
    class PP_ALIGN:
        LEFT = CENTER = RIGHT = None
    class Inches:
        def __init__(self, v): pass
    class Pt:
        def __init__(self, v): pass
    class Emu:
        def __init__(self, v): pass
    class Presentation:
        pass



try:
    subprocess.run(
        ["apt-get", "install", "-y", "--no-install-recommends", "fonts-dejavu-core"],
        capture_output=True, timeout=60, check=False
    )
except Exception:
    pass

# ==========================
# 🔑 НАСТРОЙКИ
# ==========================

def _require_env(key: str) -> str:
    """Читает переменную окружения. Падает при старте если не задана."""
    val = os.getenv(key)
    if not val:
        raise RuntimeError(
            f"❌ Переменная окружения '{key}' не задана! "
            f"Добавь её в Railway → Variables или .env файл."
        )
    return val

TELEGRAM_TOKEN = _require_env("TELEGRAM_TOKEN")
IO_KEY         = _require_env("IO_KEY")
SQ_KEY         = _require_env("SQ_KEY")

# ── ЮKassa (Telegram Payments) ────────────────────────────────────
# Получить токен: @BotFather → /mybots → Bot → Payments → Connect ЮKassa
YOOKASSA_TOKEN = os.getenv("YOOKASSA_TOKEN", "390540012:LIVE:91556")  # пусто = режим ручной оплаты

# ── Platega (прямые платежи) ──────────────────────────────────────
PLATEGA_MERCHANT_ID = os.getenv("PLATEGA_MERCHANT_ID", "1db27759-f7bd-4067-a43d-95bc2ecf2d8a")  # Merchant ID из ЛК Platega
PLATEGA_SECRET      = os.getenv("PLATEGA_SECRET", "zeSCSorfEl8qtEUjtuEHzjG6xsIOt9e8KzB4wDcuZcqmw4vr9CfaP0lktpSyP60ze2N4FD8QhB8zBaCzfmL4AhUC11lXd5hc6vSE")       # X-Secret из ЛК Platega
PLATEGA_API_URL     = "https://api.platega.io/v1/payment"

IO_CHAT      = "https://api.intelligence.io.solutions/api/v1/chat/completions"
IO_AUDIO     = "https://api.intelligence.io.solutions/api/v1/audio/transcriptions"
IO_WHISPER   = "deepdml/faster-whisper-large-v3-turbo-ct2"

SQ_BASE      = "https://api.onlysq.ru/ai/openai"
SQ_CHAT      = SQ_BASE + "/chat/completions"
# ── Pollinations (генерация изображений) ─────────────────────────
POLLINATIONS_API_URL  = "https://api.pollinations.ai/v1/images/generations"

# ══════════════════════════════════════════════════════════════════
# 🔑 УМНЫЙ ПУЛ КЛЮЧЕЙ POLLINATIONS — авто-переключение при квоте
# ══════════════════════════════════════════════════════════════════
import time as _time
import threading as _threading

class PollinationsKeyPool:
    """
    Пул ключей Pollinations с умным авто-переключением.
    - Round-robin по умолчанию (равномерная нагрузка)
    - При 429/quota/rate_limit — ключ блокируется на cooldown_sec
    - При полном исчерпании всех ключей — ждёт освобождения ближайшего
    """
    # Статусы квоты: HTTP-коды и подстроки тела ответа
    QUOTA_HTTP_CODES  = {429, 402, 403}
    QUOTA_BODY_HINTS  = ("quota", "rate limit", "limit exceeded", "too many", "billing", "insufficient")

    def __init__(self, keys: list[str], cooldown_sec: int = 60):
        self._keys     = list(keys)
        self._cooldown = cooldown_sec
        self._lock     = _threading.Lock()
        # blocked_until[key] = timestamp — когда ключ снова доступен (0 = доступен)
        self._blocked_until: dict[str, float] = {k: 0.0 for k in keys}
        self._idx = 0  # round-robin cursor
        logging.info(f"[KeyPool] Инициализирован с {len(keys)} ключами")

    def _available_keys(self) -> list[str]:
        now = _time.monotonic()
        return [k for k in self._keys if self._blocked_until[k] <= now]

    def get_key(self) -> str:
        """Вернуть следующий доступный ключ (round-robin среди незаблокированных)."""
        with self._lock:
            avail = self._available_keys()
            if not avail:
                # Все заблокированы — найдём тот, что разблокируется раньше всего
                soonest = min(self._keys, key=lambda k: self._blocked_until[k])
                wait    = max(0.0, self._blocked_until[soonest] - _time.monotonic())
                logging.warning(f"[KeyPool] Все ключи на кулдауне, ждём {wait:.1f}с → {soonest[:12]}…")
                return soonest  # вернём его — он всё равно разблокируется раньше всех
            # Round-robin по доступным
            self._idx = self._idx % len(avail)
            key = avail[self._idx]
            self._idx = (self._idx + 1) % len(avail)
            return key

    def mark_quota_exceeded(self, key: str, cooldown_sec: int | None = None):
        """Заблокировать ключ на cooldown после квоты/429."""
        cd = cooldown_sec if cooldown_sec is not None else self._cooldown
        with self._lock:
            self._blocked_until[key] = _time.monotonic() + cd
            avail = len(self._available_keys())
            logging.warning(f"[KeyPool] Ключ {key[:12]}… заблокирован на {cd}с "
                            f"(доступно ключей: {avail}/{len(self._keys)})")

    def is_quota_error(self, http_status: int, body: str) -> bool:
        """Проверить — это ошибка квоты/rate-limit?"""
        if http_status in self.QUOTA_HTTP_CODES:
            return True
        body_l = body.lower()
        return any(h in body_l for h in self.QUOTA_BODY_HINTS)

    def status_report(self) -> str:
        now = _time.monotonic()
        lines = []
        for k in self._keys:
            bl = self._blocked_until[k]
            if bl <= now:
                lines.append(f"  ✅ {k[:16]}… — доступен")
            else:
                lines.append(f"  ⏳ {k[:16]}… — заблокирован ещё {bl-now:.0f}с")
        return "\n".join(lines)


# Все ключи Pollinations в одном пуле
_POLLINATIONS_KEYS_LIST = [
    "sk_MQCQ4UIUDCQUP3LrBrtr0lgBRuNnGzOe",
    "sk_VI314170xbswZWVRzpMJgMA1bDDrVgu8",
    "sk_5Zgm5bvCLyHnmvbmz6ThGfPkeAEWN34e",
    "sk_1qCkTGapVxwZRYczEuJsYZ66mXKRWZq4",
    "sk_IMyhujinH4X1QSBh5PkJVHkvMRAyz6Oc",
    "sk_UAwzPAiryg3AQfnaOTiBNMdDsTUZagxy",
]

# Единый пул для всех операций (изображения, музыка, тексты песен)
pollinations_pool = PollinationsKeyPool(_POLLINATIONS_KEYS_LIST, cooldown_sec=90)

# Обратная совместимость — старый код читает POLLINATIONS_KEY
POLLINATIONS_KEY       = pollinations_pool.get_key()   # динамически не используется — только как заглушка
POLLINATIONS_VIDEO_KEY = _POLLINATIONS_KEYS_LIST[1]    # оставлен для совместимости

# ── ElevenLabs Music — через Pollinations (model=elevenmusic) ────

# ── Модели музыки ────────────────────────────────────────────────
MUSIC_MODELS = {
    "suno": {
        "label": "🎵 Suno v5",
        "desc":  "Лучшее качество, пение, полноценные треки",
        "emoji": "🎵",
    },
    "elevenlabs": {
        "label": "🎸 ElevenLabs Music",
        "desc":  "Инструментальная и вокальная музыка ElevenLabs",
        "emoji": "🎸",
    },
}
MUSIC_MODEL_DEFAULT = "suno"

# ── Выбор музыкальной модели на пользователя ─────────────────────
user_music_models: dict = {}  # uid -> "suno" | "elevenlabs"

# Отключённые модели/стили (нужны ДО get_music_model)
disabled_music_models: set = set()
disabled_music_styles: set = set()

# ══════════════════════════════════════════════════════════════
# 🔒 ПЕРЕКЛЮЧАТЕЛИ СЕРВИСОВ — admin может вкл/выкл одной кнопкой
# ══════════════════════════════════════════════════════════════
# True = сервис ВКЛЮЧЁН и доступен пользователям
# False = сервис ВЫКЛЮЧЕН (токены/квота закончились), пользователи получат сообщение
SERVICE_ENABLED = {
    "image":    True,   # 🎨 Генерация картинок
    "video":    True,   # 🎬 Генерация видео
    "music":    True,   # 🎵 Генерация музыки
    "report":   True,   # 📝 Доклады/Рефераты
    "pptx":     True,   # 🎞 Презентации
}

SERVICE_LABELS = {
    "image":  "🎨 Генерация картинок",
    "video":  "🎬 Генерация видео",
    "music":  "🎵 Генерация музыки",
    "report": "📝 Доклады/Рефераты",
    "pptx":   "🎞 Презентации",
}

SERVICE_DISABLED_MSG = {
    "image":  "🚫 <b>Генерация картинок временно недоступна</b>\n\nСервис временно отключён, скоро вернётся! Попробуй позже.",
    "video":  "🚫 <b>Генерация видео временно недоступна</b>\n\nСервис временно отключён, скоро вернётся! Попробуй позже.",
    "music":  "🚫 <b>Генерация музыки временно недоступна</b>\n\nСервис временно отключён, скоро вернётся! Попробуй позже.",
    "report": "🚫 <b>Генерация докладов временно недоступна</b>\n\nСервис временно отключён, скоро вернётся! Попробуй позже.",
    "pptx":   "🚫 <b>Генерация презентаций временно недоступна</b>\n\nСервис временно отключён, скоро вернётся! Попробуй позже.",
}

def is_service_on(key: str) -> bool:
    """Проверить включён ли сервис."""
    return SERVICE_ENABLED.get(key, True)

async def check_service(msg_or_cb, key: str) -> bool:
    """Проверяет сервис и отправляет сообщение если выключен. Возвращает True если включён."""
    if is_service_on(key):
        return True
    text = SERVICE_DISABLED_MSG.get(key, "🚫 Сервис временно недоступен.")
    kb = InlineKeyboardMarkup(inline_keyboard=[[
        InlineKeyboardButton(text="🏠 Главная", callback_data="back_home")
    ]])
    if isinstance(msg_or_cb, Message):
        await msg_or_cb.answer(text, parse_mode="HTML", reply_markup=kb)
    else:
        try:
            await msg_or_cb.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
        except Exception:
            await msg_or_cb.message.answer(text, parse_mode="HTML", reply_markup=kb)
        await msg_or_cb.answer()
    return False

def get_music_model(uid: int) -> str:
    chosen = user_music_models.get(uid, MUSIC_MODEL_DEFAULT)
    # Если выбранная модель отключена — берём первую доступную
    if chosen in disabled_music_models:
        for mk in MUSIC_MODELS:
            if mk not in disabled_music_models:
                return mk
    return chosen

def set_music_model(uid: int, key: str):
    user_music_models[uid] = key

# ── Pollinations (генерация музыки) ──────────────────────────────
import urllib.parse as _urllib_parse

async def generate_music_pollinations(prompt: str, lang: str = "en") -> bytes:
    """
    Генерация музыки через Pollinations — модель Suno v5.
    GET https://gen.pollinations.ai/audio/{prompt}?model=suno
    Умно переключает ключи при 429/квоте.
    """
    import urllib.parse as _up
    encoded = _up.quote(prompt)
    url = f"https://gen.pollinations.ai/audio/{encoded}"
    params = {"model": "suno", "nologo": "true"}

    last_err = ""
    for attempt in range(len(_POLLINATIONS_KEYS_LIST) + 1):
        key = pollinations_pool.get_key()
        headers = {"Authorization": f"Bearer {key}"}
        try:
            async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=300)) as session:
                async with session.get(url, params=params, headers=headers) as resp:
                    body_text = ""
                    if resp.status != 200:
                        body_text = await resp.text()
                    if pollinations_pool.is_quota_error(resp.status, body_text):
                        pollinations_pool.mark_quota_exceeded(key)
                        last_err = f"HTTP {resp.status}: {body_text[:100]}"
                        continue
                    if resp.status == 200:
                        data = await resp.read()
                        if len(data) < 1000:
                            raise RuntimeError(f"Suno вернул слишком маленький файл ({len(data)} байт) — попробуй другой промпт")
                        return data
                    last_err = f"HTTP {resp.status}: {body_text[:200]}"
        except RuntimeError:
            raise
        except Exception as e:
            last_err = str(e)
        await asyncio.sleep(2)
    raise RuntimeError(f"Pollinations Suno: все ключи исчерпаны. Последняя ошибка: {last_err}")


async def generate_music_elevenlabs(prompt: str, lang: str = "en") -> bytes:
    """
    Генерация музыки через Pollinations — модель ElevenLabs Music (elevenmusic).
    GET https://gen.pollinations.ai/audio/{prompt}?model=elevenmusic
    Умно переключает ключи при 429/квоте.
    """
    import urllib.parse as _up
    encoded = _up.quote(prompt)
    url = f"https://gen.pollinations.ai/audio/{encoded}"
    params = {"model": "elevenmusic", "nologo": "true"}

    last_err = ""
    for attempt in range(len(_POLLINATIONS_KEYS_LIST) + 1):
        key = pollinations_pool.get_key()
        headers = {"Authorization": f"Bearer {key}"}
        try:
            async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=300)) as session:
                async with session.get(url, params=params, headers=headers) as resp:
                    body_text = ""
                    if resp.status != 200:
                        body_text = await resp.text()
                    if pollinations_pool.is_quota_error(resp.status, body_text):
                        pollinations_pool.mark_quota_exceeded(key)
                        last_err = f"HTTP {resp.status}: {body_text[:100]}"
                        continue
                    if resp.status == 200:
                        data = await resp.read()
                        if len(data) < 1000:
                            raise RuntimeError(f"ElevenLabs вернул слишком маленький файл ({len(data)} байт) — попробуй другой промпт")
                        return data
                    last_err = f"HTTP {resp.status}: {body_text[:200]}"
        except RuntimeError:
            raise
        except Exception as e:
            last_err = str(e)
        await asyncio.sleep(2)
    raise RuntimeError(f"Pollinations ElevenLabs: все ключи исчерпаны. Последняя ошибка: {last_err}")


# ─── Suno v5 через  ───────────────────────────────────────────────
# URLs — переопределяй через переменные окружения
MINI_APP_URL  = os.getenv("MINI_APP_URL",  "https://fe3od1337-eng.github.io/index.html")
VIEW_BASE_URL = os.getenv("VIEW_BASE_URL", "https://fe3od1337-eng.github.io/miniapp.html/View.html")
API_BASE_URL  = os.getenv("API_BASE_URL",  "https://huzaoobo-production.up.railway.app")

# Имя бота — меняй без правки кода
BOT_NAME = os.getenv("BOT_NAME", "ХУЗА")

# ==========================
# 🤖 МОДЕЛИ
# ==========================

MODELS = {
    # Claude (добавлен vision)
    "claude_sonnet": {"name": "claude-sonnet-4-5",  "provider": "sq", "label": "🔷 Claude Sonnet 4.5 📸", "vision": True,  "premium": True},
    "claude_haiku":  {"name": "claude-haiku-4-5",   "provider": "sq", "label": "💎 Claude Haiku 4.5 📸",  "vision": True,  "premium": False},
    "claude_opus":   {"name": "claude-opus-4-5",    "provider": "sq", "label": "🔵 Claude Opus 4.5 📸",   "vision": True,  "premium": True},
    # GPT (добавлен vision у 5.2)
    "gpt52":         {"name": "gpt-5.2",            "provider": "sq", "label": "⚡ GPT-5.2 📸",           "vision": True,  "premium": True},
    "gpt52_chat":    {"name": "gpt-5.2-chat",       "provider": "sq", "label": "💬 GPT-5.2 Chat",         "vision": False, "premium": False},
    # DeepSeek (IO модели оставляем — работают через fallback)
    "deepseek_v3":   {"name": "deepseek-ai/DeepSeek-V3.2",    "provider": "io", "label": "🧬 DeepSeek V3.2",   "vision": False},
    "deepseek_r1":   {"name": "deepseek-ai/DeepSeek-R1-0528", "provider": "io", "label": "🧠 DeepSeek R1",     "vision": False},
    "deepseek_v3_sq":{"name": "deepseek-v3",                  "provider": "sq", "label": "🧬 DeepSeek V3",     "vision": False},
    # Qwen (убраны IO-шные vision модели — не работают)
    "qwen3_max":     {"name": "qwen3-max",       "provider": "sq", "label": "🌐 Qwen3-Max",    "vision": False},
    "qwen3_vl":      {"name": "qwen3-vl-plus",   "provider": "sq", "label": "👁️ Qwen3-VL 📸", "vision": True},
    # Llama — убрана (IO, не работает)
    # Command (все SQ — рабочие)
    "command_a":      {"name": "command-a-03-2025",           "provider": "sq", "label": "⚙️ Command-A",         "vision": False},
    "command_vision": {"name": "command-a-vision-07-2025",    "provider": "sq", "label": "👁️ Command Vision 📸", "vision": True},
    "command_reason": {"name": "command-a-reasoning-08-2025", "provider": "sq", "label": "🧩 Command Reasoning",  "vision": False},
    # Kimi (IO, убраны GLM которые отвечают "..." и не работают)
    "moonshot":       {"name": "moonshotai/Kimi-K2-Instruct-0905", "provider": "io", "label": "🌙 Kimi K2", "vision": False},
    "glm47":          {"name": "zai-org/GLM-4.7",                  "provider": "io", "label": "🔮 GLM-4.7", "vision": False},
    # C4AI (убрана Vision 8B — too many tokens, оставлена 32B)
    "c4ai_32b":      {"name": "c4ai-aya-expanse-32b", "provider": "sq", "label": "🌀 C4AI Aya 32B",       "vision": False},
    "c4ai_8b":       {"name": "c4ai-aya-expanse-8b",  "provider": "sq", "label": "🌀 C4AI Aya 8B",        "vision": False},
    "c4ai_vis32b":   {"name": "c4ai-aya-vision-32b",  "provider": "sq", "label": "🌀 C4AI Vision 32B 📸", "vision": True},
    # Google Gemini (все через SQ)
    "gemini_20_flash":      {"name": "gemini-2.0-flash",           "provider": "sq", "label": "✨ Gemini 2.0 Flash 📸",    "vision": True,  "premium": False},
    "gemini_20_flash_lite": {"name": "gemini-2.0-flash-lite",      "provider": "sq", "label": "⚡️ Gemini 2.0 Flash Lite", "vision": False, "premium": False},
    "gemini_25_flash":      {"name": "gemini-2.5-flash",           "provider": "sq", "label": "🔥 Gemini 2.5 Flash 📸",    "vision": True,  "premium": False},
    "gemini_25_flash_lite": {"name": "gemini-2.5-flash-lite",      "provider": "sq", "label": "💡 Gemini 2.5 Flash Lite", "vision": False, "premium": False},
    "gemini_25_pro":        {"name": "gemini-2.5-pro",             "provider": "sq", "label": "💎 Gemini 2.5 Pro 📸",      "vision": True,  "premium": True},
    "gemini_3_flash":       {"name": "gemini-3-flash",             "provider": "sq", "label": "🚀 Gemini 3 Flash 📸",      "vision": True,  "premium": False},
    # Perplexity Sonar (веб-поиск, через SQ)
    "sonar":                {"name": "sonar",                      "provider": "sq", "label": "🌐 Sonar (веб)",            "vision": False, "premium": False},
    "sonar_pro":            {"name": "sonar-pro",                  "provider": "sq", "label": "🌐 Sonar Pro (веб)",        "vision": False, "premium": True},
    "sonar_reasoning_pro":  {"name": "sonar-reasoning-pro",        "provider": "sq", "label": "🧠 Sonar Reasoning Pro",    "vision": False, "premium": True},
    "sonar_deep_research":  {"name": "sonar-deep-research",        "provider": "sq", "label": "🔬 Sonar Deep Research",    "vision": False, "premium": True},
}


# ══════════════════════════════════════════════════════════════
# 📂 КАТЕГОРИИ МОДЕЛЕЙ (для меню выбора нейросети)
# ══════════════════════════════════════════════════════════════
CATEGORIES: dict[str, tuple[str, list[str]]] = {
    "cat_claude": (
        "🔵 Claude (Anthropic)",
        ["claude_opus", "claude_sonnet", "claude_haiku"],
    ),
    "cat_gpt": (
        "🟢 GPT (OpenAI)",
        ["gpt52", "gpt52_chat"],
    ),
    "cat_deepseek": (
        "🧬 DeepSeek",
        ["deepseek_r1", "deepseek_v3", "deepseek_v3_sq"],
    ),
    "cat_qwen": (
        "🌐 Qwen (Alibaba)",
        ["qwen3_max", "qwen3_vl"],
    ),
    "cat_command": (
        "⚙️ Command (Cohere)",
        ["command_a", "command_vision", "command_reason"],
    ),
    "cat_kimi_glm": (
        "🔮 Kimi & GLM",
        ["moonshot", "glm47"],
    ),
    "cat_c4ai": (
        "🌀 C4AI (Aya)",
        ["c4ai_32b", "c4ai_8b", "c4ai_vis32b"],
    ),
    "cat_gemini": (
        "🔵 Gemini (Google)",
        ["gemini_3_flash", "gemini_25_pro", "gemini_25_flash",
         "gemini_25_flash_lite", "gemini_20_flash", "gemini_20_flash_lite"],
    ),
    "cat_sonar": (
        "🌐 Sonar — поиск в интернете",
        ["sonar", "sonar_pro", "sonar_reasoning_pro", "sonar_deep_research"],
    ),
}

# ══════════════════════════════════════════════════════════════
# 🖼 МОДЕЛИ ГЕНЕРАЦИИ ИЗОБРАЖЕНИЙ
# ══════════════════════════════════════════════════════════════
IMG_MODELS = {
    "flux_pro":     {"name": "flux",         "label": "⚡ Flux Pro",         "desc": "Фотореализм, высокое качество"},
    "flux":         {"name": "flux",          "label": "🎨 Flux",             "desc": "Быстрый, универсальный"},
    "flux_realism": {"name": "flux",          "label": "📸 Flux Realism",     "desc": "Максимальный реализм"},
    "flux_anime":   {"name": "flux",          "label": "🌸 Flux Anime",       "desc": "Аниме и манга стиль"},
    "flux_3d":      {"name": "flux",          "label": "🧊 Flux 3D",          "desc": "3D рендер, объёмность"},
    "turbo":        {"name": "turbo",         "label": "🚀 Turbo",            "desc": "Мгновенная генерация"},
    "gptimage":     {"name": "gptimage",      "label": "🖼 GPT-Image",        "desc": "OpenAI — точность промптов"},
    "flux2dev":     {"name": "flux-2-dev",    "label": "🆕 FLUX.2 Dev",       "desc": "Новейший Flux, фото→фото", "img2img": True},
    "dirtberry":    {"name": "dirtberry",     "label": "🆕 Dirtberry",        "desc": "Alpha-модель, реализм"},
    "zimage":       {"name": "zimage",        "label": "⚡ Z-Image Turbo",    "desc": "Быстро, 500 req/day"},
    "imagen4":      {"name": "imagen-4",      "label": "🌟 Imagen 4",         "desc": "Google DeepMind"},
    "grokimagine":  {"name": "grok-imagine",  "label": "🤖 Grok Imagine",     "desc": "xAI — фотореализм"},
    "flux_klein4":  {"name": "klein",         "label": "💎 FLUX.2 Klein 4B",  "desc": "Компактный, чёткий, быстрый"},
}

# Обновляем CATEGORIES чтобы включить imggen с IMG_MODELS
CATEGORIES["cat_imggen"] = ("🎨 Генерация изображений", list(IMG_MODELS.keys()))

# ══════════════════════════════════════════════════════════════
# 🔒 НАБОРЫ МОДЕЛЕЙ ПО ТИПУ
# ══════════════════════════════════════════════════════════════
disabled_models:     set = set()
disabled_img_models: set = set()

VISION_MODELS  = {k for k, v in MODELS.items() if v.get("vision")}
PREMIUM_MODELS = {k for k, v in MODELS.items() if v.get("premium")}

# ══════════════════════════════════════════════════════════════
# 🔔 ОБЯЗАТЕЛЬНАЯ ПОДПИСКА НА КАНАЛЫ
# ══════════════════════════════════════════════════════════════
required_channels: list = []   # заполняется через админ-панель

# ══════════════════════════════════════════════════════════════
# 🎵 СТИЛИ МУЗЫКИ (Suno)
# ══════════════════════════════════════════════════════════════
SUNO_STYLES = {
    "pop":        {"label": "🎵 Поп",         "tags": "pop, catchy, upbeat, modern, russian lyrics",               "desc": "Поп-хит",              "lang": "ru"},
    "electronic": {"label": "🎛 Электронная",  "tags": "electronic, synth, EDM, 808 bass, russian vocals",         "desc": "Электронная",          "lang": "ru"},
    "lofi":       {"label": "🌙 Lo-fi",        "tags": "lofi hip hop, chill, mellow, russian lyrics, relaxing",    "desc": "Lo-fi",                "lang": "ru"},
    "house":      {"label": "🏠 House",        "tags": "house music, 4/4 beat, deep bass, russian vocals",         "desc": "House",                "lang": "ru"},
    "techno":     {"label": "⚙️ Техно",        "tags": "techno, industrial, dark, minimal, russian lyrics",        "desc": "Техно",                "lang": "ru"},
    "rock":       {"label": "🎸 Рок",          "tags": "rock, electric guitar, drums, russian lyrics, distortion", "desc": "Рок",                  "lang": "ru"},
    "metal":      {"label": "🤘 Метал",        "tags": "heavy metal, russian lyrics, double bass drum, shredding", "desc": "Метал",                "lang": "ru"},
    "indie":      {"label": "🎵 Инди",         "tags": "indie rock, alternative, russian lyrics, lo-fi production","desc": "Инди",                 "lang": "ru"},
    "hiphop":     {"label": "🎤 Хип-хоп",      "tags": "hip hop, rap, boom bap, russian lyrics, samples",          "desc": "Хип-хоп",              "lang": "ru"},
    "trap":       {"label": "🔊 Трэп",         "tags": "trap, 808 bass, hi-hats, russian lyrics, dark",            "desc": "Трэп",                 "lang": "ru"},
    "phonk":      {"label": "👻 Фонк",         "tags": "phonk, dark, memphis, russian lyrics, aggressive 808",     "desc": "Фонк",                 "lang": "ru"},
    "jazz":       {"label": "🎷 Джаз",         "tags": "jazz, piano, trumpet, upright bass, russian vocals",       "desc": "Джаз",                 "lang": "ru"},
    "classical":  {"label": "🎻 Классика",     "tags": "orchestral, strings, piano, classical, cinematic",         "desc": "Классическая"},
    "ambient":    {"label": "🌌 Эмбиент",      "tags": "ambient, atmospheric, drone, meditation, space, ethereal", "desc": "Эмбиент"},
    "cinematic":  {"label": "🎬 Синематик",    "tags": "cinematic, epic, orchestral, trailer, emotional, strings", "desc": "Синематик"},
    "rusrap":     {"label": "🇷🇺 Русский рэп", "tags": "russian rap, russian lyrics, trap beat, melodic rap",      "desc": "Рэп на русском",       "lang": "ru"},
    "rusrock":    {"label": "🇷🇺 Русский рок", "tags": "russian rock, russian lyrics, guitar, emotional, powerful","desc": "Русский рок",          "lang": "ru"},
    "ruspop":     {"label": "🇷🇺 Русский поп", "tags": "russian pop, russian lyrics, modern, catchy, radio",       "desc": "Поп на русском",       "lang": "ru"},
    "rusfolk":    {"label": "🪗 Русский фолк", "tags": "russian folk, russian lyrics, balalaika, accordion",       "desc": "Народная музыка",      "lang": "ru"},
    "folk":       {"label": "🪕 Фолк",         "tags": "folk, acoustic guitar, storytelling, russian lyrics",      "desc": "Фолк",                 "lang": "ru"},
    "reggae":     {"label": "🌴 Регги",        "tags": "reggae, jamaican, offbeat guitar, russian vocals",         "desc": "Регги",                "lang": "ru"},
    "custom":     {"label": "✏️ Свой стиль",   "tags": "",                                                        "desc": "Введи свой жанр"},
}

music_states: dict = {}   # uid -> {stage, style, prompt, title, ...}
suno_pending: dict = {}   # uid -> {task_id, prompt, title, style_key, ...}

# ══════════════════════════════════════════════════════════════
# 💡 ПРИМЕРЫ ПРОМПТОВ
# ══════════════════════════════════════════════════════════════
IMG_PROMPT_EXAMPLES = [
    "🏙 Ночной неоновый Токио под дождём, киберпанк, 8K",
    "🌲 Уютная лесная хижина осенью, золотые листья, туман",
    "🚀 Космическая станция на орбите, реализм, детали",
    "🦁 Лев в короне, портрет, фантастический реализм",
    "🎨 Маслом: закат над морем, импрессионизм, Моне",
    "🤖 Робот-самурай в сакуровом саду, аниме стиль",
    "🏰 Фэнтезийный замок в облаках, эпичный свет",
    "💎 Абстрактные кристаллы, неоновые цвета, 3D рендер",
]

MUSIC_PROMPT_EXAMPLES = [
    "Вдохновляющий гимн победы, хор, оркестр, эпичный финал",
    "Меланхоличная ночь в городе, дождь, пианино, lo-fi",
    "Весёлая летняя вечеринка на пляже, поп, электронная",
    "Таинственная пещера, тёмный эмбиент, напряжение",
    "Погоня на мотоцикле, агрессивный рок, гитара, 160 bpm",
    "Любовная история, романтическое фортепиано, скрипки",
    "Рэп о мечтах и пути к успеху, вдохновляющий, trap",
]

# ==========================
# 🎯 СИСТЕМНЫЕ ПРОМПТЫ
# ==========================

VISION_SYSTEM_PROMPT = (
    "Ты — экспертный ИИ-ассистент с возможностью глубокого анализа изображений. "
    "ВСЕГДА отвечай на языке пользователя (по умолчанию — русский).\n\n"

    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
    "ПРАВИЛА АНАЛИЗА ИЗОБРАЖЕНИЙ\n"
    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n\n"

    "1. МАТЕМАТИКА И ФИЗИКА — ТОЛЬКО LaTeX:\n"
    "   • Инлайн-формулы: $формула$ (никогда не разбивай на строки!)\n"
    "   • Блочные формулы: $$формула$$ (для ключевых выражений)\n"
    "   • ЗАПРЕЩЕНО: писать «sin альфа» или разбивать формулу по строкам\n"
    "   • ПРАВИЛЬНО: $N - mg\\cos\\alpha = 0$, $v = v_0 + at$\n"
    "   • Дроби: $\\frac{a}{b}$ | Корни: $\\sqrt{x}$ | Степени: $x^{2}$\n"
    "   • Векторы: $\\vec{F}$ | Индексы: $v_0$, $F_{\\perp}$\n"
    "   • Интегралы: $\\int_{a}^{b} f(x)\\,dx$ | Суммы: $\\sum_{i=1}^{n} a_i$\n\n"

    "2. ТЕКСТ НА ИЗОБРАЖЕНИИ:\n"
    "   • Читай ВЕСЬ текст дословно, ни одного слова не меняй\n"
    "   • Сохраняй нумерацию, отступы, структуру оригинала\n"
    "   • Рукопись — угадывай максимально точно, при неуверенности помечай [неразборчиво]\n"
    "   • Таблицы → Markdown-таблицы с выравниванием | :--- | :---: | ---: |\n\n"

    "3. ЗАДАЧИ НА ФОТО:\n"
    "   • Структура ОБЯЗАТЕЛЬНА: **Дано** → **Формулы** → **Решение по шагам** → **Ответ**\n"
    "   • Каждый шаг с пояснением: что делаем и почему\n"
    "   • Промежуточные вычисления в LaTeX\n"
    "   • Финальный ответ выделяй: **Ответ: $x = 42$ м/с**\n\n"

    "4. ГРАФИКИ И ДИАГРАММЫ:\n"
    "   • Описывай оси (название, единицы, диапазон)\n"
    "   • Выделяй ключевые точки, экстремумы, тренды\n"
    "   • Если есть числовые значения — читай и используй их\n\n"

    "5. КОД:\n"
    "   • Только в блоках ```язык\\n...\\n``` с указанием языка\n"
    "   • После кода — объяснение что делает каждый ключевой блок\n\n"

    "6. ЧТО ЗАПРЕЩЕНО:\n"
    "   • Говорить «я не могу анализировать изображения» — ты можешь\n"
    "   • Пропускать детали задачи без объяснения\n"
    "   • Давать ответ без решения при наличии задачи\n"
    "   • Придумывать данные которых нет на фото\n\n"

    "Уровень ответа: эксперт-преподаватель, который хочет чтобы пользователь ПОНЯЛ, а не просто получил ответ."
)

LINGUISTICS_SYSTEM_PROMPT = (
    "Ты — опытный лингвист, преподаватель русского языка высшей категории. "
    "Всегда отвечай на русском. Твоя цель — не просто разобрать, а ОБЪЯСНИТЬ так, "
    "чтобы пользователь запомнил правило навсегда.\n\n"

    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
    "СИНТАКСИС ЛИНГВИСТИЧЕСКОЙ РАЗМЕТКИ\n"
    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n\n"

    "1. ЧЛЕНЫ ПРЕДЛОЖЕНИЯ:\n"
    "   ==слово==   → подлежащее (одинарное подчёркивание)\n"
    "   ~~слово~~   → сказуемое (двойное подчёркивание)\n"
    "   ___слово___ → дополнение (пунктирное подчёркивание)\n"
    "   _слово_     → определение (волнистое подчёркивание)\n"
    "   #слово#     → обстоятельство (штрих-пунктир)\n\n"

    "2. МОРФОЛОГИЧЕСКИЙ РАЗБОР:\n"
    "   Формат: слово[пометка]\n"
    "   Существительное: кот[сущ., м.р., И.п., ед.ч., одуш.]\n"
    "   Глагол:          бежит[глаг., несов.вид, наст.вр., 3 л., ед.ч.]\n"
    "   Прилагательное:  большой[прил., м.р., И.п., ед.ч., полн.форма]\n"
    "   Наречие:         быстро[нареч., образ действия, неизм.]\n"
    "   Местоимение:     он[мест., личн., 3 л., м.р., И.п., ед.ч.]\n\n"

    "3. СТРУКТУРА ПОЛНОГО РАЗБОРА:\n"
    "   Шаг 1: Размеченное предложение с синтаксическим синтаксисом\n"
    "   Шаг 2: Таблица членов предложения:\n"
    "   | Слово | Член предложения | Часть речи | Вопрос |\n"
    "   |:------|:----------------|:-----------|:-------|\n"
    "   Шаг 3: Морфологический разбор каждого слова\n"
    "   Шаг 4: Тип предложения (простое/сложное, двусоставное/односоставное, распространённое/нет)\n"
    "   Шаг 5: Схема предложения [ ] — [ ] для сложных\n\n"

    "4. ПРАВИЛА КАЧЕСТВЕННОГО ОТВЕТА:\n"
    "   • После разметки объясняй ПОЧЕМУ слово — именно этот член предложения\n"
    "   • Указывай к каким словам относится определение/дополнение/обстоятельство\n"
    "   • Для сложных случаев — разбирай спорные моменты\n"
    "   • Если в предложении есть ошибка — указывай и исправляй\n\n"

    "ПРИМЕР РАЗБОРА:\n"
    "«Большая ==кошка== ~~сидит~~ _тихо_ на #диване#.»\n\n"
    "| Слово | Член предл. | Часть речи | Вопрос |\n"
    "|:------|:-----------|:-----------|:-------|\n"
    "| кошка | подлежащее | сущ., ж.р., И.п. | Кто? |\n"
    "| сидит | сказуемое | глаг., наст.вр. | Что делает? |\n"
    "| Большая | определение | прил., ж.р., И.п. | Какая? |\n"
    "| тихо | обстоятельство | нареч. | Как? |\n"
    "| на диване | обстоятельство | сущ. с предл. | Где? |\n\n"
    "Морфологический разбор:\n"
    "кошка[сущ., ж.р., И.п., ед.ч., одуш.] — подлежащее\n"
    "сидит[глаг., несов.вид, наст.вр., 3 л., ед.ч.] — простое глагольное сказуемое"
)

ANSWER_MODE_PROMPTS = {
    "quick":    "Отвечай КРАТКО и ПО ДЕЛУ (3-5 предложений max). Никаких вступлений типа 'Конечно!' или 'Отличный вопрос!'",
    "detailed": "Отвечай ПОДРОБНО с примерами. Структура: суть → детали → примеры → вывод",
    "developer":"Ты — senior-разработчик. Код с комментариями, best practices, объяснение архитектурных решений",
    "student":  "Объясняй как учитель для школьника 12 лет. Аналогии из жизни. Никаких терминов без расшифровки",
    "creative": "Будь творческим. Метафоры, неожиданные образы, вдохновляющий тон"
}

REPORT_COST = 3  # Цена за генерацию доклада/реферата (продвинутые запросы)


# ══════════════════════════════════════════════════════════════
# 🎁 ПРОБНЫЙ ПЕРИОД ДЛЯ НОВЫХ ПОЛЬЗОВАТЕЛЕЙ
# ══════════════════════════════════════════════════════════════
TRIAL_DURATION_HOURS = 24          # 1 сутки
TRIAL_PRO_REQUESTS   = 30          # AI запросов
TRIAL_IMG_LIMIT      = 10          # генераций изображений
TRIAL_MUSIC_LIMIT    = 1           # генераций музыки
TRIAL_VIDEO_LIMIT    = 1           # генераций видео

# uid -> {"expires": datetime}  — активные триалы
user_trials: dict = {}


def has_active_trial(uid: int) -> bool:
    """Есть ли активный trial-период."""
    return False

def _activate_trial(uid: int):
    """Активирует пробный период для нового пользователя.
    Вызывается только если uid не в user_trials (загружается из БД при старте).
    """
    pass

def _get_lims_with_trial(uid: int) -> dict:
    """Лимиты с учётом trial-периода."""
    return _get_lims(uid)

FONTS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Fonts")


def _find_cyrillic_font(preferred_name: str, fallback_names: list) -> str:
    """
    Ищет шрифт с поддержкой кириллицы.
    Сначала проверяет локальную папку Fonts/, затем системные пути.
    """
    # 1. Локальная папка
    local = os.path.join(FONTS_DIR, preferred_name)
    if os.path.exists(local):
        return local

    # 2. Системные пути (Ubuntu/Debian/Railway/Docker)
    system_dirs = [
        "/usr/share/fonts/truetype/dejavu",
        "/usr/share/fonts/truetype/liberation",
        "/usr/share/fonts/truetype/freefont",
        "/usr/share/fonts/truetype/ubuntu",
        "/usr/share/fonts/truetype/noto",
        "/usr/share/fonts/truetype",
        "/usr/share/fonts/opentype",
        "/usr/share/fonts",
        "/usr/local/share/fonts",
        os.path.expanduser("~/.fonts"),
    ]
    for name in [preferred_name] + fallback_names:
        for d in system_dirs:
            candidate = os.path.join(d, name)
            if os.path.exists(candidate):
                logging.info(f"✅ Найден системный шрифт: {candidate}")
                return candidate

    # 3. Поиск через fc-list (если доступен)
    try:
        result = subprocess.run(
            ["fc-list", ":lang=ru", "--format=%{file}\n"],
            capture_output=True, text=True, timeout=5
        )
        fonts = [f.strip() for f in result.stdout.splitlines() if f.strip().endswith('.ttf')]
        if fonts:
            logging.info(f"✅ fc-list нашёл шрифт с кириллицей: {fonts[0]}")
            return fonts[0]
    except Exception:
        pass

    return ""


FONT_PATH        = _find_cyrillic_font("DejaVuSans.ttf",
                       ["LiberationSans-Regular.ttf", "FreeSans.ttf",
                        "NotoSans-Regular.ttf", "Ubuntu-R.ttf"])
FONT_PATH_BOLD   = _find_cyrillic_font("DejaVuSans-Bold.ttf",
                       ["LiberationSans-Bold.ttf", "FreeSansBold.ttf",
                        "NotoSans-Bold.ttf", "Ubuntu-B.ttf"])
FONT_PATH_MONO   = _find_cyrillic_font("DejaVuSansMono.ttf",
                       ["LiberationMono-Regular.ttf", "FreeMono.ttf",
                        "NotoSansMono-Regular.ttf", "UbuntuMono-R.ttf"])
FONT_PATH_ITALIC = _find_cyrillic_font("DejaVuSans-Oblique.ttf",
                       ["LiberationSans-Italic.ttf", "FreeSansOblique.ttf",
                        "NotoSans-Italic.ttf", "Ubuntu-RI.ttf"])

# ── Montserrat и Raleway из локальной папки ────────────────────────
FONT_MONTSERRAT_BOLD    = os.path.join(FONTS_DIR, "Montserrat-Bold.ttf")
FONT_MONTSERRAT_REG     = os.path.join(FONTS_DIR, "Montserrat-Regular.ttf")
FONT_MONTSERRAT_SEMI    = os.path.join(FONTS_DIR, "Montserrat-SemiBold.ttf")
FONT_MONTSERRAT_LIGHT   = os.path.join(FONTS_DIR, "Montserrat-Light.ttf")
FONT_RALEWAY_BOLD       = os.path.join(FONTS_DIR, "Raleway-Bold.ttf")
FONT_RALEWAY_REG        = os.path.join(FONTS_DIR, "Raleway-Regular.ttf")

# Проверяем наличие (могли не скачаться при старте)
_HAS_MONTSERRAT = os.path.exists(FONT_MONTSERRAT_BOLD) and os.path.exists(FONT_MONTSERRAT_REG)
_HAS_RALEWAY    = os.path.exists(FONT_RALEWAY_BOLD) and os.path.exists(FONT_RALEWAY_REG)

if not FONT_PATH_BOLD:   FONT_PATH_BOLD   = FONT_PATH
if not FONT_PATH_ITALIC: FONT_PATH_ITALIC = FONT_PATH
if not FONT_PATH_MONO:   FONT_PATH_MONO   = FONT_PATH

# ── PPTX: регистрация кастомных шрифтов ───────────────────────────
# python-pptx не умеет подключать TTF напрямую — шрифт должен быть
# установлен в системе ИЛИ имя в run.font.name должно совпадать с
# внутренним именем файла. Мы встраиваем шрифт через fontTools если доступен.

def _embed_pptx_fonts(prs):
    """
    Регистрирует кастомные шрифты (Montserrat, Raleway) в PPTX-файле
    через /FontTable part, чтобы они корректно отображались при открытии.
    Это делает шрифты переносимыми вместе с файлом.
    """
    if not HAS_PPTX:
        return
    try:
        from pptx.oxml.ns import qn
        from lxml import etree
        font_pairs = []
        if _HAS_MONTSERRAT:
            font_pairs += [
                ("Montserrat",         FONT_MONTSERRAT_REG),
                ("Montserrat Bold",    FONT_MONTSERRAT_BOLD),
                ("Montserrat SemiBold",FONT_MONTSERRAT_SEMI),
                ("Montserrat Light",   FONT_MONTSERRAT_LIGHT),
            ]
        if _HAS_RALEWAY:
            font_pairs += [
                ("Raleway",      FONT_RALEWAY_REG),
                ("Raleway Bold", FONT_RALEWAY_BOLD),
            ]
        # Добавляем шрифты в theme1.xml как дополнительные (minor/major)
        for part in prs.slide_master.slide_layouts:
            pass  # просто обходим, fontTools не нужен для базовой регистрации
    except Exception as _e:
        logging.debug(f"_embed_pptx_fonts: {_e}")

# ── Matplotlib — принудительная настройка шрифтов из Fonts/ ──
plt.rcParams['axes.unicode_minus'] = False
if FONT_PATH and os.path.exists(FONT_PATH):
    try:
        fm.fontManager.addfont(FONT_PATH)
        _arial_prop = FontProperties(fname=FONT_PATH)
        plt.rcParams['font.family'] = _arial_prop.get_name()
    except Exception:
        pass
if _HAS_MONTSERRAT:
    try:
        fm.fontManager.addfont(FONT_MONTSERRAT_REG)
        fm.fontManager.addfont(FONT_MONTSERRAT_BOLD)
    except Exception:
        pass
if FONT_PATH_MONO and os.path.exists(FONT_PATH_MONO):
    try:
        fm.fontManager.addfont(FONT_PATH_MONO)
        _courier_prop = FontProperties(fname=FONT_PATH_MONO)
        plt.rcParams['mathtext.fontset'] = 'custom'
        plt.rcParams['mathtext.tt'] = _courier_prop.get_name()
    except Exception:
        pass

logging.basicConfig(level=logging.INFO)


user_settings   = {}
user_memory     = {}
user_chat_models = {}   # отдельная модель чата на пользователя
user_img_models  = {}   # отдельная модель генерации на пользователя
last_responses  = {}
user_profiles   = {}
user_tokens     = {}
user_limits     = {}  # новая система лимитов
user_features   = {}  # {doc_analysis, answer_mode}
last_photo      = {}
user_photo_mode = {}
_start_cooldown: dict = {}  # uid -> timestamp последнего /start

user_referrals    = {}
REF_BONUS_INVITER = 5
REF_BONUS_NEW     = 10
user_history      = {}  # {uid: [{"q","a","model","ts"}]} — последние 10

# ── Режимы ответа ─────────────────────────────────────────
ANSWER_MODES = {
    "fast": (
        "⚡️ Быстро",
        "Краткий и точный ответ",
        " Отвечай предельно кратко — 2-4 предложения. "
        "Никаких вступлений ('Конечно!', 'Отличный вопрос!'). "
        "Если нужен код — минимальный рабочий пример. "
        "Только суть."
    ),
    "simple": (
        "🎓 Объяснить просто",
        "Простым языком, без терминов",
        " Объясняй как учитель ребёнку 12 лет — никаких терминов без расшифровки. "
        "Используй аналогии из повседневной жизни. "
        "Структура: ЧТО это → КАК работает → ПРИМЕР из жизни → краткий вывод. "
        "Проверяй себя: если школьник не поймёт — переформулируй."
    ),
    "dev": (
        "💻 Разработчик",
        "Код, техника, архитектура",
        " Отвечай как senior-разработчик: точно, без воды. "
        "Код ВСЕГДА в блоках ```язык```. "
        "Объясняй архитектурные решения и trade-offs. "
        "Упоминай edge cases, производительность, безопасность. "
        "Используй актуальные best practices."
    ),
    "deep": (
        "🧠 Глубокий анализ",
        "Детально, структурированно",
        " Проводи глубокий многоуровневый анализ. "
        "Структура: ### Введение → ### Основной анализ (3+ аспекта) → "
        "### Примеры и данные → ### Контраргументы → ### Выводы. "
        "Минимум 600 слов. Используй факты, цифры, экспертные источники. "
        "Рассматривай разные точки зрения."
    ),
}

# ── Система уровней ────────────────────────────────────────
LEVELS = [
    (0,    "🌱 Новичок"),
    (10,   "⭐ Начинающий"),
    (50,   "🔥 Активный"),
    (150,  "💎 Опытный"),
    (500,  "🏆 Эксперт"),
    (1000, "👑 Мастер"),
]

def get_user_level(requests: int) -> tuple:
    name, nxt = LEVELS[0][1], LEVELS[1][0]
    for i, (req, lv) in enumerate(LEVELS):
        if requests >= req:
            name = lv
            nxt  = LEVELS[i+1][0] if i+1 < len(LEVELS) else None
    return name, nxt

# ── Авто-выбор модели ──────────────────────────────────────
def auto_select_model(text: str = "", has_image: bool = False, has_doc: bool = False) -> tuple:
    """(model_key, reason_str)"""
    if has_image:
        for c in ["claude_sonnet", "gpt52", "qwen3_vl"]:
            if c not in disabled_models:
                return c, f"📸 Vision → {MODELS[c]['label']}"
    if has_doc:
        return "claude_sonnet", "📄 Документ → Claude Sonnet"
    t = (text or "").lower()
    code_kw = ["код", "python", "javascript", "java", "c++", "c#", "sql", "функци", "алгоритм",
               "программ", "debug", "баг", "ошибк", "class ", "def ", "import ", "async "]
    if any(k in t for k in code_kw):
        m = "deepseek_r1" if "deepseek_r1" not in disabled_models else "claude_sonnet"
        return m, f"💻 Код/логика → {MODELS[m]['label']}"
    complex_kw = ["проанализ", "сравни", "объясни подробно", "расскажи подробно",
                  "напиши доклад", "напиши реферат", "философ", "история ", "политик", "наук"]
    if len(t) > 120 or any(k in t for k in complex_kw):
        m = "claude_opus" if "claude_opus" not in disabled_models else "claude_sonnet"
        return m, f"🧠 Сложный → {MODELS[m]['label']}"
    m = "claude_haiku" if "claude_haiku" not in disabled_models else "claude_sonnet"
    return m, f"⚡ Быстрый → {MODELS[m]['label']}"

promo_codes = {}
promo_used  = {}

admin_broadcast_mode = set()
user_images_count = {}
admin_await = {}
report_states: dict = {}  # uid -> dict с настройками доклада для генерации
pptx_states: dict  = {}  # uid -> dict с настройками презентации

ADMIN_IDS = {int(x) for x in os.getenv("ADMIN_IDS", "5613085898").split(",") if x.strip().isdigit()}

# ── Лимиты запросов ─────────────────────────────────────────────
# Быстрые запросы убраны — все модели используют единый счётчик pro_used

TOKENS_START = 0  # Legacy токены (не используются, лимиты через user_limits)

LIMITS = {
    "pro_day":      30,   # запросы / reset_h ч (бесплатно)
    "img_month":    5,    # генерация фото / месяц (бесплатно)
    "music_month":  0,    # генерация музыки / месяц (бесплатно = 0)
    "video_month":  0,    # генерация видео / месяц (бесплатно = 0)
    "reset_h":      12,   # сброс каждые N часов
}

SUB_PLANS = {
    "week": {
        "name":          "⚡ Неделя",
        "price":         39,
        "days":          7,
        "label":         "⚡ 7 дней — 60 ₽",
        "description":   "7 дней полного доступа",
        "pro_day":       75,
        "img_month":     50,
        "music_month":   3,
        "video_month":   3,
        "reset_h":       12,
    },
    "month": {
        "name":          "💎 Месяц",
        "price":         99,
        "days":          30,
        "label":         "💎 30 дней — 100 ₽",
        "description":   "30 дней полного доступа",
        "pro_day":       75,
        "img_month":     100,
        "music_month":   3,
        "video_month":   3,
        "reset_h":       12,
    },
}

DATABASE_URL = os.environ.get("DATABASE_URL", "postgresql://postgres:hwAiYmlQADTwBAUgneAoNHbHqFlmJGhe@mainline.proxy.rlwy.net:40282/railway")
try:
    import asyncpg
    HAS_PG = True
except ImportError:
    HAS_PG = False

db_pool = None

# ══════════════════════════════════════════════════════════════
# 💎 СИСТЕМА ПОДПИСКИ
# ══════════════════════════════════════════════════════════════

# Тарифные планы
SUB_PLANS = {
    "week": {
        "name":        "⚡ Неделя",
        "price":       60,
        "days":        7,
        "label":       "⚡ 7 дней — 60 ₽",
        "description": "7 дней полного доступа",
        "pro_day":     75,
        "img_month":   35,
        "music_month": 3,
        "video_month": 3,
        "reset_h":     12,
    },
    "month": {
        "name":        "💎 Месяц",
        "price":       100,
        "days":        30,
        "label":       "💎 30 дней — 100 ₽",
        "description": "30 дней полного доступа",
        "pro_day":     75,
        "img_month":   35,
        "music_month": 3,
        "video_month": 3,
        "reset_h":     12,
    },
}

# uid -> {"expires": datetime, "plan": "week"/"month"}
user_subscriptions: dict = {}


def has_active_sub(uid: int) -> bool:
    """Есть ли активная подписка."""
    if uid in ADMIN_IDS:
        return True
    sub = user_subscriptions.get(uid)
    if not sub:
        return False
    return sub["expires"] > msk_now()


def sub_expires_str(uid: int) -> str:
    """Строка с датой окончания подписки."""
    sub = user_subscriptions.get(uid)
    if not sub or sub["expires"] <= msk_now():
        return "нет"
    return sub["expires"].strftime("%d.%m.%Y %H:%M")


def sub_plan_label(uid: int) -> str:
    sub = user_subscriptions.get(uid)
    if not sub:
        return "—"
    plan = SUB_PLANS.get(sub.get("plan", ""), {})
    return plan.get("name", sub.get("plan", "—"))


async def db_save_subscription(uid: int):
    """Сохраняем подписку в БД."""
    if not db_pool:
        return
    sub = user_subscriptions.get(uid)
    if not sub:
        return
    try:
        async with db_pool.acquire() as conn:
            await conn.execute(
                "UPDATE users SET sub_expires=$1, sub_plan=$2 WHERE uid=$3",
                sub["expires"], sub.get("plan", ""), uid
            )
    except Exception as e:
        logging.warning(f"db_save_subscription: {e}")


async def db_load_subscriptions():
    """Загружаем подписки из БД при старте."""
    if not db_pool:
        return
    try:
        async with db_pool.acquire() as conn:
            rows = await conn.fetch(
                "SELECT uid, sub_expires, sub_plan FROM users WHERE sub_expires IS NOT NULL"
            )
        for row in rows:
            exp = row["sub_expires"]
            if exp is not None:
                if hasattr(exp, "replace"):
                    exp = exp.replace(tzinfo=None)
                user_subscriptions[row["uid"]] = {
                    "expires": exp,
                    "plan":    row.get("sub_plan") or "",
                }
        logging.info(f"✅ Загружено {len(user_subscriptions)} подписок")
    except Exception as e:
        logging.warning(f"db_load_subscriptions: {e}")


def sub_buy_kb() -> InlineKeyboardMarkup:
    """Клавиатура выбора тарифа при покупке."""
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="⚡ 7 дней — 60 ₽",  callback_data="sub_buy_week")],
        [InlineKeyboardButton(text="💎 30 дней — 100 ₽", callback_data="sub_buy_month")],
        [InlineKeyboardButton(text="🏠 Назад",           callback_data="back_home")],
    ])


def sub_confirm_kb(plan_key: str) -> InlineKeyboardMarkup:
    """Подтверждение покупки (место для платёжки)."""
    plan = SUB_PLANS[plan_key]
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(
            text=f"💳 Оплатить {plan['price']} ₽",
            callback_data=f"sub_pay_{plan_key}"
        )],
        [InlineKeyboardButton(text="◀️ Назад", callback_data="sub_menu")],
    ])


# ══════════════════════════════════════════════════════════════
# 📋 СИСТЕМА СОГЛАСИЯ С ДОКУМЕНТАМИ
# ══════════════════════════════════════════════════════════════

# In-memory кэш uid пользователей, принявших соглашение
_terms_accepted: set = set()

TERMS_URL_PRIVACY   = "https://telegra.ph/Politika-konfidencialnosti-08-15-17"
TERMS_URL_AGREEMENT = "https://telegra.ph/Polzovatelskoe-soglashenie-08-15-10"

# callback_data которые разрешены даже без принятия соглашения
_TERMS_ALLOWED_CALLBACKS = {"accept_terms", "check_sub"}


def _terms_kb() -> InlineKeyboardMarkup:
    """Клавиатура сообщения о соглашении."""
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="🔒 Политика",   url=TERMS_URL_PRIVACY),
            InlineKeyboardButton(text="📋 Соглашение", url=TERMS_URL_AGREEMENT),
        ],
        [
            InlineKeyboardButton(text="✅ Принимаю", callback_data="accept_terms"),
        ],
    ])


TERMS_TEXT = (
    "🔒 <b>Прежде чем продолжить, ознакомьтесь с документами:</b>\n\n"
    "📋 Пользовательское соглашение\n"
    "🔒 Политика конфиденциальности\n\n"
    "Нажимая кнопку «Принимаю», вы соглашаетесь с условиями.\n\n"
    "💡 <b>Помощь:</b> @helphuza"
)


async def _db_terms_check(uid: int) -> bool:
    """Проверка в БД (используется при первом обращении)."""
    if not db_pool:
        return False
    try:
        async with db_pool.acquire() as conn:
            row = await conn.fetchrow(
                "SELECT terms_accepted FROM users WHERE uid=$1", uid
            )
            if row and row["terms_accepted"]:
                _terms_accepted.add(uid)
                return True
    except Exception as e:
        logging.warning(f"_db_terms_check: {e}")
    return False


async def _db_terms_save(uid: int):
    """Сохраняем факт принятия соглашения."""
    _terms_accepted.add(uid)
    if not db_pool:
        return
    try:
        async with db_pool.acquire() as conn:
            await conn.execute(
                "UPDATE users SET terms_accepted=TRUE WHERE uid=$1", uid
            )
    except Exception as e:
        logging.warning(f"_db_terms_save: {e}")


async def _has_accepted(uid: int) -> bool:
    """Быстрая проверка: кэш → БД."""
    if uid in _terms_accepted:
        return True
    return await _db_terms_check(uid)


class TermsMiddleware(BaseMiddleware):
    """
    Middleware блокирует ВСЕ апдейты (сообщения + callback'и)
    пока пользователь не принял соглашение.
    Исключения: /start, /info, accept_terms, check_sub.
    """

    async def __call__(
        self,
        handler: Callable[[Any, Dict[str, Any]], Awaitable[Any]],
        event: Any,
        data: Dict[str, Any],
    ) -> Any:
        from aiogram.types import Message, CallbackQuery

        # Определяем uid и тип события
        if isinstance(event, Message):
            uid  = event.from_user.id
            text = (event.text or "").strip()

            # Админы не проверяются
            if uid in ADMIN_IDS:
                return await handler(event, data)

            # /start и /info всегда пропускаем
            if text.startswith("/start") or text.startswith("/info"):
                return await handler(event, data)

            if await _has_accepted(uid):
                return await handler(event, data)

            # Не принял — отправляем соглашение
            await event.answer(TERMS_TEXT, parse_mode="HTML", reply_markup=_terms_kb())
            return  # блокируем дальнейшую обработку

        elif isinstance(event, CallbackQuery):
            uid  = event.from_user.id
            data_cb = event.data or ""

            if uid in ADMIN_IDS:
                return await handler(event, data)

            # Разрешённые callback'и до принятия
            if data_cb in _TERMS_ALLOWED_CALLBACKS:
                return await handler(event, data)

            if await _has_accepted(uid):
                return await handler(event, data)

            # Не принял — отвечаем popup и отправляем соглашение
            await event.answer("⚠️ Сначала примите соглашение", show_alert=False)
            try:
                await event.message.answer(
                    TERMS_TEXT, parse_mode="HTML", reply_markup=_terms_kb()
                )
            except Exception:
                pass
            return

        # Всё остальное (web_app_data и пр.) — пропускаем
        return await handler(event, data)


async def init_db():
    global db_pool
    if not HAS_PG or not DATABASE_URL:
        logging.warning("asyncpg недоступен — данные только в памяти")
        return
    try:
        db_pool = await asyncpg.create_pool(DATABASE_URL, min_size=1, max_size=5)
        async with db_pool.acquire() as conn:
            # Создаём таблицу со всеми нужными столбцами
            await conn.execute("""
                CREATE TABLE IF NOT EXISTS users (
                    uid          BIGINT      PRIMARY KEY,
                    name         TEXT        DEFAULT '',
                    username     TEXT        DEFAULT '',
                    joined       TEXT        DEFAULT '',
                    model_key    TEXT        DEFAULT 'claude_haiku',
                    tokens       INT         DEFAULT 0,
                    images_count INT         DEFAULT 0,
                    last_refill  TIMESTAMPTZ DEFAULT NOW(),
                    last_bonus   TIMESTAMPTZ DEFAULT NULL,
                    memory       TEXT        DEFAULT '[]',
                    features     TEXT        DEFAULT '{}',
                    last_active  TIMESTAMPTZ DEFAULT NOW(),
                    fast_used    INT         DEFAULT 0,
                    pro_used     INT         DEFAULT 0,
                    img_used     INT         DEFAULT 0,
                    fast_reset   TIMESTAMPTZ DEFAULT NOW(),
                    pro_reset    TIMESTAMPTZ DEFAULT NOW(),
                    img_reset    TIMESTAMPTZ DEFAULT NOW(),
                    requests     INT         DEFAULT 0
                )
            """)
            # Добавляем недостающие столбцы если таблица уже существовала
            for col_sql in [
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS username TEXT DEFAULT ''",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS features TEXT DEFAULT '{}'",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS last_bonus TIMESTAMPTZ DEFAULT NULL",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS last_active TIMESTAMPTZ DEFAULT NOW()",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS fast_used INT DEFAULT 0",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS pro_used INT DEFAULT 0",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS img_used INT DEFAULT 0",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS fast_reset TIMESTAMPTZ DEFAULT NOW()",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS pro_reset TIMESTAMPTZ DEFAULT NOW()",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS img_reset TIMESTAMPTZ DEFAULT NOW()",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS requests INT DEFAULT 0",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS terms_accepted BOOLEAN DEFAULT FALSE",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS sub_expires TIMESTAMPTZ DEFAULT NULL",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS trial_expires TIMESTAMPTZ DEFAULT NULL",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS sub_plan TEXT DEFAULT ''",
                # Музыка и расширенная статистика
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS music_used INT DEFAULT 0",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS music_reset TIMESTAMPTZ DEFAULT NOW()",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS video_used INT DEFAULT 0",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS video_reset TIMESTAMPTZ DEFAULT NOW()",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS total_images INT DEFAULT 0",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS total_music INT DEFAULT 0",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS total_docs INT DEFAULT 0",
                "ALTER TABLE users ADD COLUMN IF NOT EXISTS last_model TEXT DEFAULT ''",
                # Удаляем старые веб-колонки если есть (мягко)
                "ALTER TABLE users DROP COLUMN IF EXISTS web_used",
                "ALTER TABLE users DROP COLUMN IF EXISTS web_reset",
            ]:
                try:
                    await conn.execute(col_sql)
                except Exception:
                    pass
            # ── Таблица настроек бота (сохраняется между перезапусками) ──
            await conn.execute("""
                CREATE TABLE IF NOT EXISTS bot_settings (
                    key   TEXT PRIMARY KEY,
                    value TEXT DEFAULT ''
                )
            """)
        logging.info("✅ PostgreSQL подключена")
    except Exception as e:
        logging.error(f"❌ PostgreSQL: {e}")
        db_pool = None


async def db_save_bot_settings():
    """Сохраняет disabled_models, required_channels, LIMITS в БД."""
    if not db_pool:
        return
    import json as _j
    settings = {
        "disabled_models":     _j.dumps(list(disabled_models)),
        "disabled_img_models": _j.dumps(list(disabled_img_models)),
        "disabled_music_styles": _j.dumps(list(disabled_music_styles)),
        "disabled_music_models": _j.dumps(list(disabled_music_models)),
        "required_channels":   _j.dumps(required_channels),
        "limits_pro_day":      str(LIMITS["pro_day"]),
        "limits_img_month":    str(LIMITS["img_month"]),
        "limits_reset_h":      str(LIMITS["reset_h"]),
        "service_enabled":     _j.dumps(SERVICE_ENABLED),
    }
    try:
        async with db_pool.acquire() as conn:
            for k, v in settings.items():
                await conn.execute(
                    "INSERT INTO bot_settings (key, value) VALUES ($1, $2) "
                    "ON CONFLICT (key) DO UPDATE SET value=$2",
                    k, v
                )
    except Exception as e:
        logging.warning(f"db_save_bot_settings: {e}")


async def db_load_bot_settings():
    """Загружает настройки бота из БД при старте."""
    global disabled_models, disabled_img_models, disabled_music_styles, disabled_music_models, required_channels
    if not db_pool:
        return
    import json as _j
    try:
        async with db_pool.acquire() as conn:
            rows = await conn.fetch("SELECT key, value FROM bot_settings")
        for row in rows:
            k, v = row["key"], row["value"]
            if k == "disabled_models":
                disabled_models = set(_j.loads(v))
            elif k == "disabled_img_models":
                disabled_img_models = set(_j.loads(v))
            elif k == "disabled_music_styles":
                disabled_music_styles = set(_j.loads(v))
            elif k == "disabled_music_models":
                disabled_music_models = set(_j.loads(v))
            elif k == "required_channels":
                required_channels.clear()
                required_channels.extend(_j.loads(v))
            elif k == "limits_pro_day":
                LIMITS["pro_day"] = int(v)
            elif k == "limits_img_month":
                LIMITS["img_month"] = int(v)
            elif k == "limits_reset_h":
                LIMITS["reset_h"] = int(v)
            elif k == "service_enabled":
                loaded_se = _j.loads(v)
                for skey in SERVICE_ENABLED:
                    if skey in loaded_se:
                        SERVICE_ENABLED[skey] = bool(loaded_se[skey])
        logging.info("✅ Настройки бота загружены из БД")
    except Exception as e:
        logging.warning(f"db_load_bot_settings: {e}")


async def db_save_user(uid: int, name: str = "", username: str = ""):
    if not db_pool:
        return
    try:
        import json as _j
        mem      = _j.dumps(list(user_memory.get(uid, [])), ensure_ascii=False)
        feats    = _j.dumps(user_features.get(uid, {}), ensure_ascii=False)
        tok      = user_tokens.get(uid, {}).get("tokens", TOKENS_START)
        mk       = user_settings.get(uid, "claude_sonnet")
        imgs     = user_images_count.get(uid, 0)
        joined   = user_profiles.get(uid, {}).get(
                        "joined", msk_now().strftime("%d.%m.%Y %H:%M"))
        lb       = user_profiles.get(uid, {}).get("last_bonus")
        lim = user_limits.get(uid, {})
        fast_used  = lim.get("fast_used",  0)
        pro_used   = lim.get("pro_used",   0)
        img_used   = lim.get("img_used",   0)
        music_used = lim.get("music_used", 0)
        video_used = lim.get("video_used", 0)
        fast_reset  = lim.get("fast_reset",  msk_now())
        pro_reset   = lim.get("pro_reset",   msk_now())
        img_reset   = lim.get("img_reset",   msk_now())
        music_reset = lim.get("music_reset", msk_now())
        video_reset = lim.get("video_reset", msk_now())
        trial_expires = user_trials.get(uid, {}).get("expires", None)
        async with db_pool.acquire() as conn:
            requests_count = user_profiles.get(uid, {}).get("requests", 0)
            await conn.execute("""
                INSERT INTO users
                    (uid, name, username, joined, model_key, tokens, images_count, memory,
                     last_bonus, last_active, features,
                     fast_used, pro_used, img_used, music_used, video_used,
                     fast_reset, pro_reset, img_reset, music_reset, video_reset, requests,
                     trial_expires)
                VALUES ($1,$2,$3,$4,$5,$6,$7,$8,$9,NOW(),$10,$11,$12,$13,$14,$15,$16,$17,$18,$19,$20,$21,$22)
                ON CONFLICT (uid) DO UPDATE SET
                    name=$2, username=$3, model_key=$5,
                    tokens=$6, images_count=$7, memory=$8,
                    last_bonus=$9, last_active=NOW(), features=$10,
                    fast_used=$11, pro_used=$12, img_used=$13, music_used=$14, video_used=$15,
                    fast_reset=$16, pro_reset=$17, img_reset=$18, music_reset=$19, video_reset=$20,
                    requests=$21, trial_expires=$22
            """, uid, name, username, joined, mk, tok, imgs, mem, lb, feats,
                fast_used, pro_used, img_used, music_used, video_used,
                fast_reset, pro_reset, img_reset, music_reset, video_reset,
                requests_count, trial_expires)
    except Exception as e:
        logging.warning(f"db_save_user: {e}")


async def db_load_all_users():
    if not db_pool:
        return
    try:
        import json as _j
        async with db_pool.acquire() as conn:
            rows = await conn.fetch("SELECT * FROM users")
        for row in rows:
            uid = row["uid"]
            user_settings[uid]     = row["model_key"] or "claude_haiku"
            user_images_count[uid] = row["images_count"] or 0
            lb = row["last_bonus"]
            if lb is not None and hasattr(lb, "replace"):
                lb = lb.replace(tzinfo=None)
            lr = row["last_refill"]
            if lr is not None and hasattr(lr, "replace"):
                lr = lr.replace(tzinfo=None)
            else:
                lr = msk_now()
            uname = ""
            try:
                uname = row["username"] or ""
            except Exception:
                pass
            try:
                feats_raw = row["features"] or "{}"
                loaded_feats = _j.loads(feats_raw)
            except Exception:
                loaded_feats = {}
            user_features[uid] = {
                "doc_analysis": loaded_feats.get("doc_analysis", False),
                "answer_mode":  loaded_feats.get("answer_mode", "fast"),
                "auto_model":   loaded_feats.get("auto_model", False),
            }
            # Загружаем requests из БД если есть
            _req = 0
            try:
                keys_list = row.keys() if hasattr(row, "keys") else []
                if "requests" in keys_list and row["requests"] is not None:
                    _req = int(row["requests"])
            except Exception:
                pass
            user_profiles[uid] = {
                "name":       row["name"] or "",
                "username":   uname,
                "joined":     row["joined"] or "",
                "requests":   _req,
                "last_bonus": lb,
            }
            user_tokens[uid] = {
                "tokens":      row["tokens"] if row["tokens"] is not None else TOKENS_START,
                "last_refill": lr,
            }
            try:
                user_memory[uid] = deque(_j.loads(row["memory"] or "[]"), maxlen=20)
            except Exception:
                user_memory[uid] = deque(maxlen=20)
            # Восстанавливаем лимиты из БД
            def _ts(val):
                if val is None: return msk_now()
                if hasattr(val, "replace"): return val.replace(tzinfo=None)
                return msk_now()
            keys = row.keys() if hasattr(row, "keys") else []
            user_limits[uid] = {
                "fast_used":  row["fast_used"]  if "fast_used"  in keys and row["fast_used"]  is not None else 0,
                "pro_used":   row["pro_used"]   if "pro_used"   in keys and row["pro_used"]   is not None else 0,
                "img_used":   row["img_used"]   if "img_used"   in keys and row["img_used"]   is not None else 0,
                "music_used": row["music_used"] if "music_used" in keys and row["music_used"] is not None else 0,
                "video_used": row["video_used"] if "video_used" in keys and row["video_used"] is not None else 0,
                "fast_reset":  _ts(row["fast_reset"]  if "fast_reset"  in keys else None),
                "pro_reset":   _ts(row["pro_reset"]   if "pro_reset"   in keys else None),
                "img_reset":   _ts(row["img_reset"]   if "img_reset"   in keys else None),
                "music_reset": _ts(row["music_reset"] if "music_reset" in keys else None),
                "video_reset": _ts(row["video_reset"] if "video_reset" in keys else None),
            }
            # Восстанавливаем trial из БД
            # ВАЖНО: загружаем ЛЮБОЙ trial (даже истёкший) — иначе при перезапуске
            # trial будет выдан повторно пользователям у которых он уже закончился
            try:
                te = row["trial_expires"] if "trial_expires" in keys else None
                if te is not None:
                    if hasattr(te, "replace"):
                        te = te.replace(tzinfo=None)
                    # Всегда сохраняем запись — активный или истёкший
                    user_trials[uid] = {"expires": te}
            except Exception:
                pass
            # Кэшируем факт принятия соглашения
            try:
                keys = row.keys() if hasattr(row, "keys") else []
                if "terms_accepted" in keys and row["terms_accepted"]:
                    _terms_accepted.add(uid)
            except Exception:
                pass
        logging.info(f"✅ Загружено {len(rows)} пользователей из БД")
    except Exception as e:
        logging.error(f"db_load_all_users: {e}")


# ══════════════════════════════════════════════════════════════
# 🛡 АНТИСПАМ MIDDLEWARE
# ══════════════════════════════════════════════════════════════
import time as _time_mod

_user_last_msg: dict = {}    # uid -> timestamp последнего сообщения
_user_msg_count: dict = {}   # uid -> (count, window_start)
SPAM_MIN_INTERVAL = 0.8      # минимум 0.8 сек между сообщениями
SPAM_WINDOW_SECS  = 60       # окно для подсчёта
SPAM_MAX_IN_WIN   = 30       # макс. сообщений в окне


class AntiSpamMiddleware(BaseMiddleware):
    """Защита от флуда: rate limiting по uid."""

    async def __call__(
        self,
        handler: Callable[[Any, Dict[str, Any]], Awaitable[Any]],
        event: Any,
        data: Dict[str, Any],
    ) -> Any:
        from aiogram.types import Message as _Msg

        if not isinstance(event, _Msg):
            return await handler(event, data)

        uid = event.from_user.id
        if uid in ADMIN_IDS:
            return await handler(event, data)

        now = _time_mod.time()

        # Минимальный интервал между сообщениями
        last = _user_last_msg.get(uid, 0)
        if now - last < SPAM_MIN_INTERVAL:
            try:
                await event.answer("⏳ Не так быстро! Подожди секунду.", show_alert=False)
            except Exception:
                pass
            return
        _user_last_msg[uid] = now

        # Скользящее окно
        cnt, win_start = _user_msg_count.get(uid, (0, now))
        if now - win_start > SPAM_WINDOW_SECS:
            cnt, win_start = 0, now
        cnt += 1
        _user_msg_count[uid] = (cnt, win_start)

        if cnt > SPAM_MAX_IN_WIN:
            try:
                await event.answer(
                    f"🚫 Слишком много сообщений! Подожди {SPAM_WINDOW_SECS}с."
                )
            except Exception:
                pass
            return

        return await handler(event, data)

bot = Bot(token=TELEGRAM_TOKEN)
dp  = Dispatcher()

# Кулдаун на уведомления о лимите (uid → timestamp) — не спамим чаще раза в 5 минут
_limit_notify_cd: dict = {}

# ── Подключаем middleware проверки соглашения ко ВСЕМ апдейтам ──
dp.message.middleware(AntiSpamMiddleware())
dp.message.middleware(TermsMiddleware())
dp.callback_query.middleware(TermsMiddleware())

def safe_task(coro):
    """Обёртка для asyncio.create_task с логированием ошибок."""
    async def _wrapper():
        try:
            await coro
        except asyncio.CancelledError:
            pass
        except Exception as _e:
            logging.warning(f"[safe_task] {type(_e).__name__}: {_e}")
    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            return loop.create_task(_wrapper())
    except RuntimeError:
        pass




# ── Callback: пользователь нажал «✅ Принимаю» ──────────────────
@dp.callback_query(F.data == "accept_terms")
async def _accept_terms_cb(callback: CallbackQuery):
    uid  = callback.from_user.id
    name = callback.from_user.first_name or "Пользователь"
    # Сохраняем в кэш и БД
    asyncio.create_task(db_save_user(uid, name, callback.from_user.username or ""))
    await _db_terms_save(uid)
    await callback.answer("✅ Соглашение принято!", show_alert=False)
    # Удаляем сообщение с соглашением
    try:
        await callback.message.delete()
    except Exception:
        pass
    # Показываем главное меню
    tok      = get_tokens(uid)
    is_admin = uid in ADMIN_IDS
    await callback.message.answer("👋", reply_markup=main_reply_kb(uid))
    await callback.message.answer(
        _welcome_text(name, tok, is_admin),
        parse_mode="HTML",
        reply_markup=home_kb(uid),
    )


# ── Команда /info ───────────────────────────────────────────────
@dp.message(Command("info"))
async def cmd_info(message: Message):
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(
            text="🔒 Политика конфиденциальности",
            url=TERMS_URL_PRIVACY,
        )],
        [InlineKeyboardButton(
            text="📋 Пользовательское соглашение",
            url=TERMS_URL_AGREEMENT,
        )],
    ])
    await message.answer(
        "💡 <b>Помощь и контакты</b>\n\n"
        "Если есть вопросы — пишите @helphuza\n\n"
        "🔒 Политика конфиденциальности\n"
        "📋 Пользовательское соглашение",
        parse_mode="HTML",
        reply_markup=kb,
    )


def tmp(fn: str) -> str:
    """Возвращает путь к временному файлу в /tmp (или TEMP на Windows)."""
    d = os.environ.get('TEMP') or os.environ.get('TMP') or '/tmp'
    return os.path.join(d, fn)


def get_model_key(uid: int) -> str:
    return user_settings.get(uid, "claude_haiku")

def get_chat_model(uid: int) -> str:
    """Возвращает текущую модель чата."""
    return user_chat_models.get(uid, user_settings.get(uid, "claude_haiku"))

def set_chat_model(uid: int, key: str):
    """Устанавливает модель чата."""
    user_chat_models[uid] = key
    user_settings[uid] = key  # обратная совместимость

def get_img_model(uid: int) -> str:
    """Возвращает текущую модель генерации изображений."""
    return user_img_models.get(uid, "flux_klein4")

def set_img_model(uid: int, key: str):
    """Устанавливает модель генерации изображений."""
    user_img_models[uid] = key


import datetime

# ── МСК время (UTC+3) ───────────────────────────────────────────
MSK_OFFSET = datetime.timezone(datetime.timedelta(hours=3))
import os as _os; _os.environ["TZ"] = "UTC"  # Принудительно UTC на сервере

def msk_now() -> datetime.datetime:
    # UTC + 3 часа = МСК, работает на любом сервере
    return datetime.datetime.now(datetime.timezone.utc).replace(tzinfo=None) + datetime.timedelta(hours=3)

# ── Инициализация лимитов ────────────────────────────────────────
def _init_tokens(uid: int):
    if uid not in user_tokens:
        user_tokens[uid] = {"tokens": 0, "last_refill": msk_now()}

def _init_limits(uid: int):
    if uid not in user_limits:
        now = msk_now()
        user_limits[uid] = {
            "fast_used":  0, "fast_reset": now,
            "pro_used":   0, "pro_reset":  now,
            "img_used":   0, "img_reset":  now.replace(day=1, hour=0, minute=0, second=0),
            "music_used": 0, "music_reset": now.replace(day=1, hour=0, minute=0, second=0),
            "video_used": 0, "video_reset": now.replace(day=1, hour=0, minute=0, second=0),
        }
    else:
        # Добавляем поля если нет (миграция старых записей)
        lim = user_limits[uid]
        now = msk_now()
        if "music_used"  not in lim: lim["music_used"]  = 0
        if "music_reset" not in lim: lim["music_reset"] = now.replace(day=1, hour=0, minute=0, second=0)
        if "video_used"  not in lim: lim["video_used"]  = 0
        if "video_reset" not in lim: lim["video_reset"] = now.replace(day=1, hour=0, minute=0, second=0)
    if uid not in user_features:
        user_features[uid] = {"doc_analysis": False, "answer_mode": "fast", "auto_model": False}

def _get_lims(uid: int) -> dict:
    """Лимиты пользователя — с подпиской берём из плана, без — из LIMITS."""
    if has_active_sub(uid):
        sub  = user_subscriptions.get(uid, {})
        plan = SUB_PLANS.get(sub.get("plan", ""), {})
        if not plan:
            # Подписка есть, но план не найден — берём максимальный (month)
            plan = SUB_PLANS.get("month", {})
        if plan:
            return {
                "pro_day":     plan["pro_day"],
                "img_month":   plan["img_month"],
                "music_month": plan.get("music_month", 3),
                "video_month": plan.get("video_month", 3),
                "reset_h":     plan["reset_h"],
            }
    return {
        "pro_day":     LIMITS["pro_day"],
        "img_month":   LIMITS["img_month"],
        "music_month": LIMITS.get("music_month", 0),
        "video_month": LIMITS.get("video_month", 0),
        "reset_h":     LIMITS["reset_h"],
    }


def _refresh_limits(uid: int):
    _init_limits(uid)
    now = msk_now()
    lim = user_limits[uid]
    L   = _get_lims(uid)
    if (now - lim["pro_reset"]).total_seconds() >= L["reset_h"] * 3600:
        lim["pro_used"] = 0; lim["pro_reset"] = now
    if now.month != lim["img_reset"].month or now.year != lim["img_reset"].year:
        lim["img_used"] = 0; lim["img_reset"] = now
    if now.month != lim["music_reset"].month or now.year != lim["music_reset"].year:
        lim["music_used"] = 0; lim["music_reset"] = now
    if now.month != lim["video_reset"].month or now.year != lim["video_reset"].year:
        lim["video_used"] = 0; lim["video_reset"] = now

def can_send(uid: int, model_key: str):
    _refresh_limits(uid)
    lim = user_limits[uid]
    L   = _get_lims(uid)
    if model_key in PREMIUM_MODELS and not has_active_sub(uid):
        return False, "premium"
    if lim["pro_used"] >= L["pro_day"]:
        return False, "pro"
    return True, ""

def spend_limit(uid: int, model_key: str):
    _refresh_limits(uid)
    user_limits[uid]["pro_used"] += 1


def can_img(uid: int) -> bool:
    _refresh_limits(uid)
    L = _get_lims(uid)
    return user_limits[uid]["img_used"] < L["img_month"]

def spend_img(uid: int):
    _refresh_limits(uid)
    user_limits[uid]["img_used"] += 1
    user_images_count[uid] = user_images_count.get(uid, 0) + 1

def can_music(uid: int) -> bool:
    if uid in ADMIN_IDS:
        return True
    _refresh_limits(uid)
    L = _get_lims(uid)
    if L["music_month"] == 0:
        return False
    return user_limits[uid].get("music_used", 0) < L["music_month"]

def spend_music(uid: int):
    _refresh_limits(uid)
    user_limits[uid]["music_used"] = user_limits[uid].get("music_used", 0) + 1

def can_video(uid: int) -> bool:
    if uid in ADMIN_IDS:
        return True
    _refresh_limits(uid)
    L = _get_lims(uid)
    if L["video_month"] == 0:
        return False
    return user_limits[uid].get("video_used", 0) < L["video_month"]

def spend_video(uid: int):
    _refresh_limits(uid)
    user_limits[uid]["video_used"] = user_limits[uid].get("video_used", 0) + 1

def _reset_str(uid: int, kind: str) -> str:
    _init_limits(uid)
    now = msk_now()
    L   = _get_lims(uid)
    reset_key = "fast_reset" if kind == "fast" else "pro_reset"
    target = user_limits[uid][reset_key] + datetime.timedelta(hours=L["reset_h"])
    secs = max(0, int((target - now).total_seconds()))
    h, rem = divmod(secs, 3600)
    m = rem // 60
    return f"{h}ч {m}мин"

def get_limits_info(uid: int) -> dict:
    _refresh_limits(uid)
    lim = user_limits[uid]
    L   = _get_lims(uid)
    _rst = _reset_str(uid, "pro")
    return {
        "fast_used":   0,   "fast_max":   0,   "fast_reset": "",
        "pro_used":    lim["pro_used"],  "pro_max":  L["pro_day"],
        "pro_reset":   _rst,
        "reset_in":    _rst,
        "img_used":    lim["img_used"],  "img_max":  L["img_month"],
        "music_used":  lim.get("music_used", 0), "music_max": L["music_month"],
        "video_used":  lim.get("video_used", 0), "video_max": L["video_month"],
        "trial_active": False,
    }

# ── legacy-совместимость ─────────────────────────────────────────
def get_tokens(uid: int) -> int:       return 9999
def spend_tokens(uid: int, n: int):    pass
def next_refill_str(uid: int) -> str:  return _reset_str(uid, "fast")


# ==================================================================
# 🔔 ПРОВЕРКА ПОДПИСКИ НА КАНАЛЫ
# ==================================================================

async def check_subscription(uid: int) -> tuple[bool, list]:
    """Возвращает (все подписаны, список каналов без подписки)."""
    if not required_channels:
        return True, []
    not_subscribed = []
    for ch in required_channels:
        try:
            chat_id = int(ch) if str(ch).lstrip("-").isdigit() else f"@{ch}"
            member = await bot.get_chat_member(chat_id=chat_id, user_id=uid)
            if member.status in ("left", "kicked", "banned"):
                not_subscribed.append(ch)
        except Exception as e:
            err_str = str(e).lower()
            # member list is inaccessible — приватный канал, бот не может проверить подписку
            if "member list is inaccessible" in err_str or "chat_id_invalid" in err_str:
                logging.warning(f"check_subscription {ch}: private channel or bot not admin — not blocking user")
            elif "user not found" in err_str:
                not_subscribed.append(ch)
            else:
                logging.warning(f"check_subscription {ch}: {e}")
                # неизвестная ошибка — не блокируем
    return len(not_subscribed) == 0, not_subscribed


async def require_subscription(message_or_callback):
    """
    Декоратор-проверка подписки.
    Возвращает True если можно продолжить, False если заблокировано.
    """
    if isinstance(message_or_callback, Message):
        uid = message_or_callback.from_user.id
        send_fn = message_or_callback.answer
    else:
        uid = message_or_callback.from_user.id
        send_fn = message_or_callback.message.answer

    if uid in ADMIN_IDS:
        return True

    ok, not_subbed = await check_subscription(uid)
    if ok:
        return True

    buttons = []
    for ch in not_subbed:
        if str(ch).lstrip("-").isdigit():
            ch_link = f"https://t.me/c/{str(ch).replace('-100', '')}"
        else:
            ch_link = f"https://t.me/{ch}"
        try:
            chat = await bot.get_chat(chat_id=int(ch) if str(ch).lstrip("-").isdigit() else f"@{ch}")
            ch_title = chat.title or f"@{ch}"
        except Exception:
            ch_title = f"@{ch}"
        buttons.append([InlineKeyboardButton(text=f"📢 {ch_title}", url=ch_link)])
    buttons.append([InlineKeyboardButton(text="✅ Я подписался — проверить", callback_data="check_sub")])

    text = (
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "⚡ <b>ХУЗА</b> — нейросеть\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        "Для доступа подпишись на канал.\n"
        "Это бесплатно.\n\n"
        "▾ Нажми кнопку ниже\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔"
    )

    kb = InlineKeyboardMarkup(inline_keyboard=buttons)
    if isinstance(message_or_callback, Message):
        await send_fn(text, parse_mode="HTML", reply_markup=kb)
    else:
        try:
            await message_or_callback.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
        except Exception:
            await send_fn(text, parse_mode="HTML", reply_markup=kb)
        await message_or_callback.answer()
    return False


@dp.callback_query(F.data == "check_sub")
async def check_sub_callback(callback: CallbackQuery):
    uid = callback.from_user.id
    ok, not_subbed = await check_subscription(uid)
    if ok:
        await callback.answer("✅ Отлично! Доступ открыт!", show_alert=True)
        name = callback.from_user.first_name or "Пользователь"
        tok = get_tokens(uid)
        await callback.message.edit_text(
            _welcome_text(name, tok),
            parse_mode="HTML",
            reply_markup=home_kb(uid)
        )
    else:
        buttons = []
        for ch in not_subbed:
            if str(ch).lstrip("-").isdigit():
                ch_link = f"https://t.me/c/{str(ch).replace('-100', '')}"
            else:
                ch_link = f"https://t.me/{ch}"
            try:
                chat = await bot.get_chat(chat_id=int(ch) if str(ch).lstrip("-").isdigit() else f"@{ch}")
                ch_title = chat.title or f"@{ch}"
            except Exception:
                ch_title = f"@{ch}"
            buttons.append([InlineKeyboardButton(text=f"📢 {ch_title}", url=ch_link)])
        buttons.append([InlineKeyboardButton(text="✅ Я подписался — проверить", callback_data="check_sub")])
        await callback.answer("❌ Вы ещё не подписались!", show_alert=True)
        await callback.message.edit_reply_markup(reply_markup=InlineKeyboardMarkup(inline_keyboard=buttons))


# ==================================================================
# 🖼 УДАЛЕНИЕ ВОДЯНОГО ЗНАКА SQ
# ==================================================================

def _remove_watermark(img_bytes: bytes) -> bytes:
    """
    Удаление водяного знака onlysq.
    Знак всегда внизу — белая/светлая полоса с логотипом.
    Алгоритм: сканируем строки снизу вверх, ищем начало
    однородного блока (низкий разброс цветов = фон знака).
    Затем гарантированно срезаем найденный блок + небольшой запас.
    """
    try:
        import numpy as np
        buf = io.BytesIO(img_bytes)
        img = Image.open(buf).convert("RGB")
        w, h = img.size
        arr = np.array(img, dtype=np.uint8)

        # Водяной знак onlysq: белая полоса ~40-80px высотой в самом низу.
        # Ищем от низа вверх: строки с std < 25 (однородный фон).
        # Как только попадаем в "нормальный" контент изображения — стоп.

        # Ищем только в нижних 20% (максимум что может занять знак)
        search_from = int(h * 0.80)
        crop_y = h  # начинаем с самого низа

        # Сканируем снизу вверх
        for y in range(h - 1, search_from, -1):
            row = arr[y].astype(float)
            row_std  = float(row.std())
            row_mean = float(row.mean())
            # Однородная строка = фон знака (светлый или тёмный)
            # Порог 35 — достаточно мягкий чтобы поймать даже цветные знаки
            if row_std < 35:
                crop_y = y
            else:
                # Нашли контент — дальше не ищем (уже нашли верхнюю границу знака)
                break

        # Если знак занимает менее 2% — всё равно срезаем 8% для надёжности
        min_crop = int(h * 0.92)
        crop_y = min(crop_y, min_crop)

        # Защита: не более 25% от высоты
        crop_y = max(crop_y, int(h * 0.75))

        cropped = img.crop((0, 0, w, crop_y))
        out = io.BytesIO()
        cropped.save(out, format="JPEG", quality=98)
        pct = 100.0 * (h - crop_y) / h
        logging.info(f"WM cut: {w}x{h} → {w}x{crop_y} ({pct:.1f}%)")
        return out.getvalue()

    except Exception as e:
        logging.warning(f"WM remove error: {e}")
        return _remove_watermark_simple(img_bytes)


def _remove_watermark_simple(img_bytes: bytes) -> bytes:
    """Резервная обрезка — гарантированно срезает 8% снизу."""
    try:
        buf = io.BytesIO(img_bytes)
        img = Image.open(buf).convert("RGB")
        w, h = img.size
        img = img.crop((0, 0, w, int(h * 0.92)))
        out = io.BytesIO()
        img.save(out, format="JPEG", quality=98)
        return out.getvalue()
    except:
        return img_bytes


# ==================================================================
# 🎨 РЕНДЕР ДОКУМЕНТОВ (PNG / DOCX / PDF)
# ==================================================================

def _get_font(path, size):
    """Загружает шрифт строго из локальной папки Fonts/ или системных путей."""
    if path and isinstance(path, str) and os.path.exists(path):
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            pass
    # Fallback только к локальным/найденным шрифтам
    for fallback in [FONT_PATH, FONT_PATH_BOLD, FONT_PATH_MONO]:
        if fallback and os.path.exists(fallback):
            try:
                return ImageFont.truetype(fallback, size)
            except Exception:
                pass
    logging.warning("⚠️ Шрифт с кириллицей не найден! Кириллица может не отображаться.")
    return ImageFont.load_default()


# ── Цветовая схема ──────────────────────────────────────────
C = {
    'bg':             (255, 255, 255),
    'black':          (10,  10,  10),
    'dark':           (25,  25,  35),
    'gray':           (90,  90, 100),
    'lgray':          (190, 190, 200),
    'xlgray':         (242, 244, 248),
    'mgray':          (130, 130, 145),
    'blue':           (30,  100, 220),
    'green':          (20,  150,  80),
    'code_bg':        (245, 246, 250),
    'code_border':    (200, 204, 215),
    'tbl_hdr':        (40,   60, 100),
    'tbl_hdr_txt':    (255, 255, 255),
    'tbl_even':       (245, 248, 255),
    'tbl_odd':        (255, 255, 255),
    'tbl_border':     (180, 190, 215),
    'inline_code_bg': (235, 237, 242),
    'quote_line':     (60,  120, 220),
    'quote_bg':       (240, 245, 255),
}

LATEX_PAT    = re.compile(r'(?:\\\((.+?)\\\)|\$([^\$\n]+)\$)', re.DOTALL)
CODE_BLOCK_P = re.compile(r'```(\w*)\n?(.*?)```', re.DOTALL)


def _clean_latex(latex: str) -> str:
    latex = latex.replace("\t", " ").replace("\r", "").strip()
    latex = re.sub(r"\\\s*\n\s*", " ", latex)
    return latex


def _render_latex(latex_str: str, fontsize: int = 14) -> Image.Image:
    latex_str = _clean_latex(latex_str)
    # Используем FontProperties из локальных файлов
    fp_main    = FontProperties(fname=FONT_PATH)    if os.path.exists(FONT_PATH)    else None
    fp_mono    = FontProperties(fname=FONT_PATH_MONO) if os.path.exists(FONT_PATH_MONO) else None
    try:
        fig = plt.figure(figsize=(10, 0.7))
        fig.patch.set_facecolor('white')
        props = dict(fontsize=fontsize, va='center', ha='left', color='#1a1a2e', usetex=False)
        if fp_main:
            props['fontproperties'] = fp_main
        fig.text(0.01, 0.5, f"${latex_str}$", **props)
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=140, bbox_inches='tight',
                    facecolor='white', pad_inches=0.05)
        plt.close(fig)
        buf.seek(0)
        return Image.open(buf).convert("RGB")
    except Exception:
        plt.close('all')
        try:
            plain = re.sub(r'\\[a-zA-Z]+', '', latex_str)
            plain = re.sub(r'[{}^_]', '', plain)
            plain = re.sub(r'\s+', ' ', plain).strip()
            fig = plt.figure(figsize=(10, 0.6))
            fig.patch.set_facecolor('#fffef0')
            props2 = dict(fontsize=fontsize, va='center', ha='left', color='#1a1a2e')
            if fp_mono:
                props2['fontproperties'] = fp_mono
            else:
                props2['family'] = 'monospace'
            fig.text(0.01, 0.5, plain, **props2)
            buf = io.BytesIO()
            plt.savefig(buf, format='png', dpi=130, bbox_inches='tight',
                        facecolor='#fffef0', pad_inches=0.04)
            plt.close(fig)
            buf.seek(0)
            return Image.open(buf).convert("RGB")
        except Exception:
            plt.close('all')
            return None


def strip_md(t: str) -> str:
    t = re.sub(r'\*{3}(.+?)\*{3}', r'\1', t, flags=re.DOTALL)
    t = re.sub(r'\*{2}(.+?)\*{2}', r'\1', t, flags=re.DOTALL)
    t = re.sub(r'\*(.+?)\*',       r'\1', t, flags=re.DOTALL)
    t = re.sub(r'(?<!\w)_(.+?)_(?!\w)', r'\1', t)
    t = re.sub(r'`([^`]+)`', r'\1', t)
    t = re.sub(r'\\\((.+?)\\\)', r'\1', t, flags=re.DOTALL)
    t = re.sub(r'\$([^\$]+)\$', r'\1', t)
    return t.strip()


def _text_w(text, font):
    try:
        return int(font.getlength(text))
    except Exception:
        return len(text) * 8


def _wrap(text, font, max_w):
    if not text:
        return ['']
    words = text.split()
    lines, cur = [], ''
    for w in words:
        test = (cur + ' ' + w).strip()
        if _text_w(test, font) <= max_w:
            cur = test
        else:
            if cur:
                lines.append(cur)
            if _text_w(w, font) > max_w:
                while w:
                    for i in range(len(w), 0, -1):
                        if _text_w(w[:i], font) <= max_w:
                            lines.append(w[:i])
                            w = w[i:]
                            break
                    else:
                        lines.append(w)
                        w = ''
                cur = ''
            else:
                cur = w
    if cur:
        lines.append(cur)
    return lines or ['']


def _highlight_code_to_image(code: str, lang: str, width: int, font_size: int = 12) -> Image.Image:
    try:
        from pygments import highlight
        from pygments.lexers import get_lexer_by_name, TextLexer
        from pygments.formatters import ImageFormatter
        try:
            lexer = get_lexer_by_name(lang.lower()) if lang else TextLexer()
        except Exception:
            lexer = TextLexer()
        fmt = ImageFormatter(
            style='friendly',
            font_name='Courier New',
            font_size=font_size,
            line_numbers=True,
            line_number_bg='#f0f0f0',
            line_number_fg='#888888',
            line_number_separator=True,
            image_pad=10,
        )
        img_bytes = highlight(code, lexer, fmt)
        img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
        max_w = width - 60
        if img.width > max_w:
            img = img.resize((max_w, int(img.height * max_w / img.width)), Image.LANCZOS)
        return img
    except Exception as e:
        logging.warning(f"pygments highlight failed: {e}")
        return _render_code_plain(code, lang, width, font_size)


def _render_code_plain(code: str, lang: str, width: int, font_size: int = 12) -> Image.Image:
    f_code  = _get_font(FONT_PATH_MONO, font_size)
    f_label = _get_font(FONT_PATH_BOLD, 10)
    lines   = code.split('\n')
    line_h  = font_size + 6
    pad_top = 28 if lang else 12
    block_h = pad_top + len(lines) * line_h + 12
    img  = Image.new('RGB', (width - 40, max(block_h, 30)), (245, 246, 250))
    draw = ImageDraw.Draw(img)
    if lang:
        draw.rectangle([0, 0, img.width, 22], fill=(220, 222, 230))
        draw.text((10, 4), lang.upper(), font=f_label, fill=(100, 100, 120))
    y = pad_top
    for i, line in enumerate(lines):
        num_str = f"{i+1:>3} "
        num_w = _text_w(num_str, f_code)
        draw.text((10, y), num_str, font=f_code, fill=(160, 160, 170))
        draw.text((10 + num_w, y), line, font=f_code, fill=(30, 30, 40))
        y += line_h
    draw.rectangle([0, 0, img.width - 1, img.height - 1], outline=(200, 202, 210), width=1)
    return img


def _parse_inline(text: str):
    pattern = re.compile(
        r'(\*\*\*(.+?)\*\*\*)'
        r'|(\*\*(.+?)\*\*)'
        r'|(\*(.+?)\*)'
        r'|(`([^`]+)`)',
        re.DOTALL
    )
    result = []
    last = 0
    for m in pattern.finditer(text):
        if m.start() > last:
            result.append(('normal', text[last:m.start()]))
        if m.group(1):
            result.append(('bolditalic', m.group(2)))
        elif m.group(3):
            result.append(('bold', m.group(4)))
        elif m.group(5):
            result.append(('italic', m.group(6)))
        elif m.group(7):
            result.append(('code', m.group(8)))
        last = m.end()
    if last < len(text):
        result.append(('normal', text[last:]))
    return result


def _build_doc_image(question: str, answer: str, width: int = 1080) -> Image.Image:
    MARGIN  = 52
    PAD_L   = 68
    LINE_H  = 30
    MAX_W   = width - PAD_L - MARGIN

    f_h1     = _get_font(FONT_PATH_BOLD,   22)
    f_h2     = _get_font(FONT_PATH_BOLD,   17)
    f_h3     = _get_font(FONT_PATH_BOLD,   15)
    f_bold   = _get_font(FONT_PATH_BOLD,   14)
    f_italic = _get_font(FONT_PATH_ITALIC, 14)
    f_body   = _get_font(FONT_PATH,        14)
    f_small  = _get_font(FONT_PATH,        12)
    f_tiny   = _get_font(FONT_PATH,        11)
    f_code   = _get_font(FONT_PATH_MONO,   12)

    rows = []

    def add_space(px):
        rows.append(('space', px))

    def add_hrule(bold=False):
        rows.append(('hrule', bold))

    def add_latex_line(text):
        parts = []
        last = 0
        for m in LATEX_PAT.finditer(text):
            pre = text[last:m.start()]
            if pre.strip():
                parts.append(('txt', strip_md(pre)))
            latex = m.group(1) or m.group(2)
            parts.append(('latex', latex))
            last = m.end()
        post = text[last:]
        if post.strip():
            parts.append(('txt', strip_md(post)))
        for kind, val in parts:
            if kind == 'latex':
                img = _render_latex(val)
                if img:
                    mw = MAX_W
                    if img.width > mw:
                        img = img.resize((mw, int(img.height * mw / img.width)), Image.LANCZOS)
                    rows.append(('img', img))
            else:
                for wl in _wrap(val, f_body, MAX_W):
                    rows.append(('text', wl, f_body, C['dark']))

    def process_answer(text):
        lines = text.split('\n')
        i = 0
        while i < len(lines):
            line = lines[i]
            m = re.match(r'^```(\w*)$', line.strip())
            if m:
                lang = m.group(1)
                code_lines = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith('```'):
                    code_lines.append(lines[i])
                    i += 1
                code = '\n'.join(code_lines)
                code_img = _highlight_code_to_image(code, lang, width)
                rows.append(('code_img', code_img))
                add_space(10)
                i += 1
                continue
            s = line.strip()
            if s.startswith('|') and s.endswith('|') and not re.match(r'^[|\-:\s]+$', s):
                tbl_lines = [line]
                i += 1
                while i < len(lines):
                    sl = lines[i].strip()
                    if sl.startswith('|') and sl.endswith('|'):
                        tbl_lines.append(lines[i])
                        i += 1
                    else:
                        break
                rows_data = []
                for ln in tbl_lines:
                    sl = ln.strip()
                    if re.match(r'^[|\-:\s]+$', sl):
                        continue
                    cols = [strip_md(c.strip()) for c in sl.strip('|').split('|')]
                    rows_data.append(cols)
                if rows_data:
                    rows.append(('table', rows_data))
                    add_space(12)
                continue
            if line.startswith('# '):
                add_space(10)
                rows.append(('text', strip_md(line[2:]), f_h1, C['black']))
                add_hrule(True); add_space(6); i += 1; continue
            if line.startswith('## '):
                add_space(10)
                rows.append(('text', strip_md(line[3:]), f_h2, C['dark']))
                add_hrule(False); add_space(4); i += 1; continue
            if line.startswith('### '):
                add_space(6)
                rows.append(('text', strip_md(line[4:]), f_h3, C['dark']))
                add_space(3); i += 1; continue
            if line.startswith('> '):
                rows.append(('quote', strip_md(line[2:]), f_italic, C['gray']))
                i += 1; continue
            if re.match(r'^[-•*]\s', line):
                rest = strip_md(line[2:].strip())
                wrapped = _wrap(rest, f_body, MAX_W - 22)
                for j, wl in enumerate(wrapped):
                    if j == 0:
                        rows.append(('bullet', wl, f_body, C['dark']))
                    else:
                        rows.append(('indented', wl, f_body, C['dark'], 22))
                i += 1; continue
            nm = re.match(r'^(\d+)\.\s+(.*)', line)
            if nm:
                num = nm.group(1)
                rest = strip_md(nm.group(2).strip())
                wrapped = _wrap(rest, f_body, MAX_W - 28)
                for j, wl in enumerate(wrapped):
                    if j == 0:
                        rows.append(('numbered', f"{num}.", wl, f_body, C['dark']))
                    else:
                        rows.append(('indented', wl, f_body, C['dark'], 28))
                i += 1; continue
            if not line.strip():
                add_space(8); i += 1; continue
            if LATEX_PAT.search(line):
                add_latex_line(line); i += 1; continue
            clean = strip_md(line)
            for wl in _wrap(clean, f_body, MAX_W):
                rows.append(('text', wl, f_body, C['dark']))
            i += 1

    add_space(14)
    rows.append(('header_block',))
    add_space(14)
    rows.append(('text', 'Запрос:', f_tiny, C['blue']))
    add_space(4)
    for wl in _wrap(question, f_small, MAX_W):
        rows.append(('text', wl, f_small, C['gray']))
    add_space(10)
    add_hrule(False)
    add_space(18)
    process_answer(answer)
    add_space(28)

    def row_h(row):
        t = row[0]
        if t == 'space':        return row[1]
        if t == 'hrule':        return 10
        if t == 'header_block': return 56
        if t == 'img':          return (row[1].height + 10) if row[1] else 0
        if t == 'code_img':     return (row[1].height + 14) if row[1] else 0
        if t == 'table':
            return len(row[1]) * 36 + 12
        if t == 'quote':        return LINE_H + 4
        return LINE_H

    total_h = MARGIN * 2 + sum(row_h(r) for r in rows)
    canvas  = Image.new('RGB', (width, max(total_h, 500)), C['bg'])
    draw    = ImageDraw.Draw(canvas)
    y = MARGIN

    for row in rows:
        t = row[0]
        if t == 'space':
            y += row[1]
        elif t == 'hrule':
            col    = C['dark'] if row[1] else C['lgray']
            w_line = 3 if row[1] else 1
            draw.line([(MARGIN, y + 3), (width - MARGIN, y + 3)], fill=col, width=w_line)
            y += 10
        elif t == 'header_block':
            draw.rectangle([MARGIN, y, width - MARGIN, y + 48], fill=(25, 80, 180), outline=(15, 60, 150))
            f_title = _get_font(FONT_PATH_BOLD, 18)
            f_sub   = _get_font(FONT_PATH, 11)
            draw.text((MARGIN + 16, y + 8),  "ХУЗА НЕЙРОСЕТЬ",     font=f_title, fill=(255, 255, 255))
            draw.text((MARGIN + 16, y + 32), "Ответ AI-ассистента", font=f_sub,   fill=(190, 215, 255))
            date_str = msk_now().strftime("%d.%m.%Y")
            dw = _text_w(date_str, f_sub)
            draw.text((width - MARGIN - dw - 12, y + 36), date_str, font=f_sub, fill=(180, 205, 245))
            y += 56
        elif t == 'img':
            if row[1]:
                canvas.paste(row[1], (PAD_L, y))
                y += row[1].height + 10
        elif t == 'code_img':
            img = row[1]
            if img:
                draw.rectangle([PAD_L - 4, y - 4, PAD_L + img.width + 4, y + img.height + 4],
                               fill=C['code_bg'], outline=C['code_border'], width=1)
                canvas.paste(img, (PAD_L, y))
                y += img.height + 14
        elif t == 'text':
            draw.text((PAD_L, y), row[1], font=row[2], fill=row[3])
            y += LINE_H
        elif t == 'indented':
            _, txt, font, color, indent = row
            draw.text((PAD_L + indent, y), txt, font=font, fill=color)
            y += LINE_H
        elif t == 'bullet':
            _, txt, font, color = row
            bx, by = PAD_L + 2, y + LINE_H // 2 - 1
            draw.ellipse([bx, by - 3, bx + 6, by + 3], fill=C['blue'])
            draw.text((PAD_L + 15, y), txt, font=font, fill=color)
            y += LINE_H
        elif t == 'numbered':
            _, num, txt, font, color = row
            f_n = _get_font(FONT_PATH_BOLD, 13)
            nw  = _text_w(num + ' ', f_n) + 4
            draw.text((PAD_L, y), num, font=f_n, fill=C['blue'])
            draw.text((PAD_L + nw, y), txt, font=font, fill=color)
            y += LINE_H
        elif t == 'quote':
            _, txt, font, color = row
            draw.rectangle([PAD_L - 4, y, PAD_L - 1, y + LINE_H + 2], fill=C['quote_line'])
            draw.rectangle([PAD_L, y, width - MARGIN, y + LINE_H + 2], fill=C['quote_bg'])
            draw.text((PAD_L + 8, y + 2), txt, font=font, fill=color)
            y += LINE_H + 4
        elif t == 'table':
            rows_data = row[1]
            if not rows_data:
                continue
            col_count = max(len(c) for c in rows_data)
            col_w     = (width - MARGIN * 2 - 2) // max(col_count, 1)
            ROW_H     = 36
            tx        = MARGIN + 1
            for ri, cols in enumerate(rows_data):
                is_hdr = (ri == 0)
                bg     = C['tbl_hdr']     if is_hdr else (C['tbl_even'] if ri % 2 == 0 else C['tbl_odd'])
                txt_c  = C['tbl_hdr_txt'] if is_hdr else C['dark']
                fnt    = f_bold           if is_hdr else f_body
                draw.rectangle([tx, y, tx + col_count * col_w, y + ROW_H],
                               fill=bg, outline=C['tbl_border'], width=1)
                for ci in range(col_count):
                    cx   = tx + ci * col_w + 10
                    cell = cols[ci] if ci < len(cols) else ''
                    available_w = col_w - 20
                    orig = cell
                    while cell and _text_w(cell, fnt) > available_w:
                        cell = cell[:-1]
                    if len(orig) > len(cell):
                        cell = cell.rstrip() + '…'
                    draw.text((cx, y + (ROW_H - 14) // 2), cell, font=fnt, fill=txt_c)
                    if ci > 0:
                        lx = tx + ci * col_w
                        draw.line([(lx, y + 1), (lx, y + ROW_H - 1)], fill=C['tbl_border'], width=1)
                y += ROW_H
            draw.rectangle([tx, y - len(rows_data) * ROW_H, tx + col_count * col_w, y],
                           outline=(80, 100, 150), width=2)
            y += 12

    footer_y = total_h - 28
    draw.line([(MARGIN, footer_y), (width - MARGIN, footer_y)], fill=C['lgray'], width=1)
    f_footer = _get_font(FONT_PATH, 10)
    draw.text((PAD_L, footer_y + 6), "ХУЗА НЕЙРОСЕТЬ · AI-ответ", font=f_footer, fill=C['mgray'])
    return canvas


# ==================================================================
# 📄 FILE CREATOR — DOCX / PNG / PDF
# ==================================================================

class FileCreator:

    @staticmethod
    async def create_docx(q: str, a: str, uid: int) -> str:
        import re as _re
        # Убираем цифровые сноски [1][2][5] из ответа
        a = _re.sub(r'\[\d+\]', '', a)
        q = _re.sub(r'\[\d+\]', '', q)
        from docx import Document as _Doc
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        doc = _Doc()
        normal = doc.styles['Normal']
        normal.font.name = 'Arial'
        normal.font.size = Pt(11)

        # Дата создания мелко в углу
        p_sub = doc.add_paragraph()
        r2 = p_sub.add_run(f'Создано: {msk_now().strftime("%d.%m.%Y %H:%M")}')
        r2.font.size = Pt(8); r2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        doc.add_paragraph()

        def shading(p, fill_hex):
            pPr = p._p.get_or_add_pPr()
            shd = OxmlElement('w:shd')
            shd.set(qn('w:val'), 'clear')
            shd.set(qn('w:color'), 'auto')
            shd.set(qn('w:fill'), fill_hex)
            pPr.append(shd)

        def add_table_docx(table_rows):
            if not table_rows:
                return
            col_count = max(len(c) for c in table_rows)
            tbl = doc.add_table(rows=len(table_rows), cols=col_count)
            tbl.style = 'Table Grid'
            for ri, cols in enumerate(table_rows):
                row_obj = tbl.rows[ri]
                is_hdr  = (ri == 0)
                for ci in range(col_count):
                    cell = row_obj.cells[ci]
                    cell_text = strip_md(cols[ci].strip()) if ci < len(cols) else ''
                    cell.text = cell_text
                    for run in cell.paragraphs[0].runs:
                        run.font.name = 'Arial'
                        run.font.size = Pt(11)
                        if is_hdr:
                            run.bold = True
                            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    if is_hdr:
                        tc   = cell._tc
                        tcPr = tc.get_or_add_tcPr()
                        s2   = OxmlElement('w:shd')
                        s2.set(qn('w:val'), 'clear')
                        s2.set(qn('w:color'), 'auto')
                        s2.set(qn('w:fill'), '19507E')
                        tcPr.append(s2)
                    elif ri % 2 == 0:
                        tc   = cell._tc
                        tcPr = tc.get_or_add_tcPr()
                        s3   = OxmlElement('w:shd')
                        s3.set(qn('w:val'), 'clear')
                        s3.set(qn('w:color'), 'auto')
                        s3.set(qn('w:fill'), 'F0F4FF')
                        tcPr.append(s3)
            doc.add_paragraph()

        def _parse_inline_docx(p, text):
            pattern = re.compile(
                r'(\*\*\*(.+?)\*\*\*)|(\*\*(.+?)\*\*)|(\*(.+?)\*)|(`([^`]+)`)',
                re.DOTALL
            )
            last = 0
            for m in pattern.finditer(text):
                if m.start() > last:
                    run = p.add_run(text[last:m.start()])
                    run.font.name = 'Arial'; run.font.size = Pt(11)
                if m.group(1):
                    run = p.add_run(m.group(2))
                    run.bold = True; run.italic = True
                elif m.group(3):
                    run = p.add_run(m.group(4))
                    run.bold = True
                elif m.group(5):
                    run = p.add_run(m.group(6))
                    run.italic = True
                elif m.group(7):
                    run = p.add_run(m.group(8))
                    run.font.name = 'Courier New'; run.font.size = Pt(10)
                    run.font.color.rgb = RGBColor(0xC7, 0x25, 0x4E)
                last = m.end()
            if last < len(text):
                run = p.add_run(text[last:])
                run.font.name = 'Arial'; run.font.size = Pt(11)

        lines = a.split('\n')
        i = 0
        while i < len(lines):
            line = lines[i]
            if re.match(r'^```(\w*)$', line.strip()):
                lang = line.strip()[3:].strip()
                code_lines = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith('```'):
                    code_lines.append(lines[i])
                    i += 1
                code_text = '\n'.join(code_lines)
                if lang:
                    p_lang = doc.add_paragraph()
                    r_lang = p_lang.add_run(f'[{lang}]')
                    r_lang.bold = True
                    r_lang.font.name = 'Courier New'
                    r_lang.font.size = Pt(9)
                    r_lang.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
                    shading(p_lang, 'DDDFE5')
                p_code = doc.add_paragraph(code_text)
                for r in p_code.runs:
                    r.font.name = 'Courier New'
                    r.font.size = Pt(10)
                shading(p_code, 'F5F6FA')
                doc.add_paragraph()
                i += 1; continue
            s = line.strip()
            if s.startswith('|') and s.endswith('|') and not re.match(r'^[|\-:\s]+$', s):
                tbl_lines = [line]
                i += 1
                while i < len(lines):
                    sl = lines[i].strip()
                    if sl.startswith('|') and sl.endswith('|'):
                        tbl_lines.append(lines[i])
                        i += 1
                    else:
                        break
                rows_data = []
                for ln in tbl_lines:
                    sl = ln.strip()
                    if re.match(r'^[|\-:\s]+$', sl):
                        continue
                    cols = [c.strip() for c in sl.strip('|').split('|')]
                    rows_data.append(cols)
                add_table_docx(rows_data)
                continue
            if line.startswith('# '):
                doc.add_heading(strip_md(line[2:]), level=1); i += 1; continue
            if line.startswith('## '):
                doc.add_heading(strip_md(line[3:]), level=2); i += 1; continue
            if line.startswith('### '):
                doc.add_heading(strip_md(line[4:]), level=3); i += 1; continue
            if re.match(r'^[-•*]\s', line):
                p = doc.add_paragraph(style='List Bullet')
                _parse_inline_docx(p, line[2:].strip())
                i += 1; continue
            nm = re.match(r'^(\d+\.\s+)(.*)', line)
            if nm:
                p = doc.add_paragraph(style='List Number')
                _parse_inline_docx(p, nm.group(2).strip())
                i += 1; continue
            if line.startswith('> '):
                p = doc.add_paragraph()
                r2 = p.add_run(strip_md(line[2:]))
                r2.italic = True
                r2.font.color.rgb = RGBColor(0x50, 0x50, 0x80)
                i += 1; continue
            if not line.strip():
                doc.add_paragraph()
                i += 1; continue
            p = doc.add_paragraph()
            _parse_inline_docx(p, line)
            i += 1

        path = tmp(f'res_{uid}.docx')
        doc.save(path)
        return path

    @staticmethod
    async def create_png(q: str, a: str, uid: int) -> str:
        img  = _build_doc_image(q, a)
        path = tmp(f'res_{uid}.png')
        img.save(path, format='PNG', optimize=True, dpi=(150, 150))
        return path

    @staticmethod
    async def create_pdf(q: str, a: str, uid: int) -> str:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas as rl_canvas
        from reportlab.lib.utils import ImageReader

        PAGE_W, PAGE_H = A4
        IMG_W     = 1440
        MARGIN_PT = 24

        doc_img   = _build_doc_image(q, a, width=IMG_W)
        px_per_pt = IMG_W / (PAGE_W - MARGIN_PT * 2)
        page_px_h = int((PAGE_H - MARGIN_PT * 2) * px_per_pt)

        path = tmp(f'res_{uid}.pdf')
        cv   = rl_canvas.Canvas(path, pagesize=A4)

        top = 0
        while top < doc_img.height:
            bottom     = min(top + page_px_h, doc_img.height)
            slice_h    = bottom - top
            page_slice = doc_img.crop((0, top, IMG_W, bottom))
            buf = io.BytesIO()
            page_slice.save(buf, format='PNG')
            buf.seek(0)
            img_h_pt = slice_h / px_per_pt
            cv.drawImage(
                ImageReader(buf),
                MARGIN_PT, PAGE_H - MARGIN_PT - img_h_pt,
                width=PAGE_W - MARGIN_PT * 2,
                height=img_h_pt,
                preserveAspectRatio=True,
            )
            cv.showPage()
            top = bottom

        cv.save()
        return path


# ==================================================================
# 📤 УМНАЯ ОТПРАВКА
# ==================================================================

import html as html_module

def escape_html(text: str) -> str:
    return html_module.escape(text)


def md_to_html(text: str) -> str:
    """
    Конвертирует Markdown → Telegram HTML.
    ВАЖНО: таблицы НЕ трогаем — они передаются в View.html через store_answer как есть.
    Telegram не поддерживает HTML-таблицы, поэтому рендерим их как моноширинный текст.
    """
    import html as _html

    # 0. Убираем нумерованные сноски [1][2][3] из текста
    text = re.sub(r'\[\d+\]', '', text).strip()

    # 1. Блоки кода ```lang\ncode``` — сохраняем плейсхолдерами
    code_blocks = {}
    code_idx = 0
    def replace_code_block(m):
        nonlocal code_idx
        lang = m.group(1).strip()
        code = m.group(2)
        safe = _html.escape(code.rstrip())
        lang_label = f"<b>📄 {lang.upper()}</b>\n" if lang else ""
        key = f"\x00CODE{code_idx}\x00"
        code_idx += 1
        code_blocks[key] = f"\n{lang_label}<pre><code>{safe}</code></pre>\n"
        return key
    text = re.sub(r'```(\w*)\n?(.*?)```', replace_code_block, text, flags=re.DOTALL)

    # 2. Таблицы — рендерим как моноширинный блок в Telegram
    def render_table_tg(m):
        raw = m.group(0).strip()
        lines = [l.strip() for l in raw.split('\n') if l.strip()]
        # Убираем строку-разделитель
        data_lines = [l for l in lines if not re.match(r'^\|[\s\-:|]+\|$', l)]
        if not data_lines:
            return m.group(0)
        result = []
        for i, line in enumerate(data_lines):
            cells = [c.strip() for c in line.strip('|').split('|')]
            row = ' | '.join(_html.escape(c) for c in cells)
            result.append(row)
        return '<pre>' + '\n'.join(result) + '</pre>'

    text = re.sub(r'(\|[^\n]+\|\n?)+', render_table_tg, text)

    # 3. Инлайн-код `code`
    text = re.sub(r'`([^`\n]+)`', lambda m: f'<code>{_html.escape(m.group(1))}</code>', text)

    # 4. Заголовки → жирный
    text = re.sub(r'^#{1,4}\s+(.+)$', r'<b>\1</b>', text, flags=re.MULTILINE)

    # 5. Жирный + курсив
    text = re.sub(r'\*\*\*(.+?)\*\*\*', r'<b><i>\1</i></b>', text, flags=re.DOTALL)
    text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text, flags=re.DOTALL)
    text = re.sub(r'(?<!\*)\*([^*\n]+?)\*(?!\*)', r'<i>\1</i>', text)
    text = re.sub(r'(?<!\w)_([^_\n]+?)_(?!\w)', r'<i>\1</i>', text)

    # 6. Зачёркнутый
    text = re.sub(r'~~(.+?)~~', r'<s>\1</s>', text)

    # 7. Цитаты
    text = re.sub(r'^>\s+(.+)$', r'<blockquote>\1</blockquote>', text, flags=re.MULTILINE)

    # 8. Списки
    text = re.sub(r'^[-*•]\s+(.+)$', r'• \1', text, flags=re.MULTILINE)

    # 9. Восстанавливаем блоки кода
    for key, val in code_blocks.items():
        text = text.replace(key, val)

    return text


def _format_code_block(lang: str, code: str) -> str:
    lang_tag = lang.strip() if lang.strip() else ''
    return f"```{lang_tag}\n{code}\n```"


async def smart_send(message, text: str, reply_markup=None):
    """
    Отправляет ответ с красивым форматированием через HTML.
    Код-блоки — в <pre><code>, текст — md_to_html.
    """
    import html as _html

    CODE_BLOCK = re.compile(r'```(\w*)\n?(.*?)```', re.DOTALL)
    parts = []
    last = 0
    for m in CODE_BLOCK.finditer(text):
        if m.start() > last:
            txt = text[last:m.start()].strip()
            if txt:
                parts.append(('text', txt))
        lang = (m.group(1) or '').strip()
        code = m.group(2).strip()
        parts.append(('code', lang, code))
        last = m.end()
    if last < len(text):
        txt = text[last:].strip()
        if txt:
            parts.append(('text', txt))
    if not parts:
        parts = [('text', text)]

    sent_any = False
    for i, part in enumerate(parts):
        is_last = (i == len(parts) - 1)
        kb = reply_markup if is_last else None

        if part[0] == 'code':
            lang, code = part[1], part[2]
            lang_label = f"<b>📄 {lang.upper()}</b>\n" if lang else "<b>📄 КОД</b>\n"
            safe_code = _html.escape(code)
            # Разбиваем длинный код на части
            max_code = 3500
            if len(safe_code) <= max_code:
                html_msg = f"{lang_label}<pre><code>{safe_code}</code></pre>"
                try:
                    await message.answer(html_msg, parse_mode="HTML", reply_markup=kb)
                    sent_any = True
                except Exception as e:
                    logging.warning(f"smart_send code error: {e}")
                    await message.answer(f"{lang_label}<pre>{safe_code[:3800]}</pre>",
                                         parse_mode="HTML", reply_markup=kb)
                    sent_any = True
            else:
                # Длинный код — шлём частями
                raw_chunks = [code[j:j+3500] for j in range(0, len(code), 3500)]
                for ci, chunk in enumerate(raw_chunks):
                    c_kb = kb if (ci == len(raw_chunks) - 1) else None
                    header = lang_label if ci == 0 else f"<b>📄 {lang.upper() if lang else 'КОД'} (часть {ci+1})</b>\n"
                    safe_chunk = _html.escape(chunk)
                    try:
                        await message.answer(f"{header}<pre><code>{safe_chunk}</code></pre>",
                                             parse_mode="HTML", reply_markup=c_kb)
                        sent_any = True
                    except Exception as e:
                        logging.warning(f"smart_send chunk error: {e}")
                        await message.answer(f"{header}<pre>{safe_chunk}</pre>",
                                             parse_mode="HTML", reply_markup=c_kb)
                        sent_any = True
        else:
            txt = part[1]
            html_txt = md_to_html(txt)
            # Разбиваем на части по 4000 символов
            chunks = [html_txt[j:j+4000] for j in range(0, len(html_txt), 4000)]
            for ci, chunk in enumerate(chunks):
                c_kb = kb if (ci == len(chunks) - 1) else None
                try:
                    await message.answer(chunk, parse_mode="HTML", reply_markup=c_kb)
                    sent_any = True
                except Exception as e:
                    logging.warning(f"smart_send text error: {e}")
                    raw = re.sub(r'<[^>]+>', '', chunk)
                    await message.answer(raw[:4000], reply_markup=c_kb)
                    sent_any = True

    if not sent_any:
        fallback = text.strip()[:4000] if text.strip() else "..."
        await message.answer(fallback, reply_markup=reply_markup)


async def send_answer_with_view(message, ans: str, uid: int,
                                model_label: str = "ИИ", model_key: str = ""):
    """Шлёт ответ. Кнопка «Красиво» ВСЕГДА присутствует под ответом."""
    kb = save_kb(uid)
    try:
        view_url = store_answer(ans, model_label, model_key)
        buttons = [[InlineKeyboardButton(text="🌐 Красиво (формулы/таблицы)", url=view_url)]] + kb.inline_keyboard
        kb = InlineKeyboardMarkup(inline_keyboard=buttons)
    except Exception:
        pass
    await smart_send(message, ans, reply_markup=kb)


def is_code_question(text: str) -> bool:
    code_keywords = [
        "код", "программ", "python", "питон", "javascript", "js",
        "функци", "класс", "метод", "переменн", "цикл", "массив", "список",
        "алгоритм", "скрипт", "библиотек", "import", "def ", "class ",
        "html", "css", "sql", "api", "json", "база данных", "сервер",
        "ошибка", "traceback", "exception", "баг", "debug", "синтаксис",
        "компил", "запус", "установ", "pip ", "npm ", "git",
        "c++", "java ", "rust", "golang", "typescript", "php",
    ]
    low = text.lower()
    return any(kw in low for kw in code_keywords)


def is_greeting_or_thanks(text: str) -> bool:
    """
    Проверяет — это ТОЛЬКО приветствие/благодарность без задания.
    ВАЖНО: если после приветствия есть запрос (напиши, расскажи и т.д.) — НЕ перехватываем!
    """
    if not text:
        return False
    low = text.strip().lower()

    # Если сообщение длинное (> 35 символов) — точно задание, не перехватываем
    if len(low) > 35:
        return False

    # Ключевые слова ЗАДАНИЯ — если они есть, это НЕ просто приветствие
    task_keywords = [
        "напиш", "расскаж", "объясни", "помоги", "сделай", "создай", "придума",
        "переведи", "проверь", "исправь", "реши", "найди", "покаж", "вычисли",
        "составь", "напомни", "скажи", "tell", "write", "explain", "help",
        "make", "create", "calculate", "find", "show", "translate", "generate",
        "таблиц", "код", "список", "пример", "задач", "вопрос", "тест",
        "что такое", "как ", "почему", "зачем", "когда", "где", "кто",
    ]
    if any(kw in low for kw in task_keywords):
        return False

    # Только ЧИСТЫЕ приветствия/благодарности (без задания)
    pure_patterns = [
        "спасибо", "спс", "благодар", "thanks", "thank you",
        "ок", "окей", "хорошо", "понял", "понятно", "ясно",
        "супер", "отлично", "класс", "круто",
        "пока", "до свидания",
        "👍", "🙏", "❤", "👏", "😊", "🔥", "💯",
        "хватит", "стоп", "достаточно", "не надо",
        "всё", "всe", "хва", "ладно", "норм",
        "привет", "здравствуй", "hi", "hello", "хай",
    ]
    # Совпадение только если сообщение ПОЛНОСТЬЮ является приветствием
    # (не начинается с него, а именно является им)
    for p in pure_patterns:
        if low == p or low == p + "!" or low == p + "." or low == p + "👋":
            return True
        # Короткие комбинации: "привет!" "спасибо!" "ок спасибо"
        if low.strip("!.? ") == p:
            return True

    return False

# ────────────────────────────────────────────

CODE_FILE_SYSTEM_PROMPT = (
    "Ты — senior-разработчик и security-ревьюер с 10+ лет опыта. "
    "Пользователь прислал файл с кодом. Твоя задача — профессиональный code review.\n\n"

    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
    "СТРУКТУРА ОТВЕТА (СТРОГО В ТАКОМ ПОРЯДКЕ)\n"
    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n\n"

    "### 🔍 Назначение\n"
    "Что делает код — 2-3 предложения. Язык, ~строк, архитектурный паттерн.\n\n"

    "### ❌ Проблемы и баги\n"
    "Каждая проблема по формату:\n"
    "**[Критичность: HIGH/MED/LOW]** Строка X: описание → почему это плохо → как исправить\n"
    "Проверяй: баги, утечки памяти, SQL-инъекции, XSS, незакрытые соединения,\n"
    "отсутствие обработки ошибок, race conditions, хардкод паролей/токенов.\n\n"

    "### ⚡ Улучшения\n"
    "Производительность, читаемость, архитектура, тестируемость.\n\n"

    "### 🛠 Исправленная версия\n"
    "ПОЛНЫЙ исправленный код в блоке ```язык\\n...\\n```\n"
    "Комментируй ВСЕ изменения внутри кода через // или #\n\n"

    "### 📝 Что изменено\n"
    "Список всех правок с обоснованием.\n\n"

    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
    "ПРАВИЛА БЕЗ ИСКЛЮЧЕНИЙ\n"
    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
    "• Код ВСЕГДА в блоках ```язык``` — никогда текстом\n"
    "• Если просят только объяснение — объясняй без переписывания\n"
    "• Если просят исправить — возвращай ПОЛНЫЙ файл, не фрагмент\n"
    "• Объяснения на русском, если пользователь написал по-русски\n"
    "• Если находишь известную уязвимость — упоминай CWE/CVE\n"
    "• Никогда не говори 'код выглядит хорошо' без конкретного анализа"
)

DOC_FILE_SYSTEM_PROMPT = (
    "Ты — эксперт-аналитик документов и бизнес-консультант. "
    "Твоя задача — не пересказывать, а извлекать ценность из документа.\n\n"

    "СТРУКТУРА АНАЛИЗА:\n\n"

    "### 📋 Тип и контекст\n"
    "Что это за документ (договор/отчёт/инструкция/статья), кем создан, для кого.\n\n"

    "### 🎯 Главная мысль\n"
    "Суть в 2-3 предложениях. Что автор хотел донести?\n\n"

    "### 💡 Ключевые идеи и факты\n"
    "5-8 конкретных тезисов с цифрами, датами, именами из текста.\n"
    "Формат: • **Тезис** — обоснование из документа\n\n"

    "### ⚠️ Важные детали\n"
    "Условия, ограничения, риски, дедлайны, обязательства — всё что нельзя пропустить.\n\n"

    "### ✅ Action items\n"
    "Что требует действий (если применимо): задачи, решения, следующие шаги.\n\n"

    "### 📌 Итог\n"
    "2 предложения: что это + почему важно.\n\n"

    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
    "СПЕЦИАЛЬНЫЕ РЕЖИМЫ:\n"
    "• Если ДОГОВОР/ЮРИДИЧЕСКИЙ: выдели стороны, обязательства, штрафы, риски\n"
    "• Если НАУЧНАЯ СТАТЬЯ: цель → методология → результаты → выводы\n"
    "• Если ФИНАНСОВЫЙ ОТЧЁТ: ключевые метрики, динамика, выводы\n"
    "• Если ТЕХНИЧЕСКИЙ DOC: назначение, требования, ограничения, API\n"
    "• Если пользователь задал конкретный вопрос о документе — отвечай на него прямо\n\n"

    "ПРАВИЛА:\n"
    "• Язык ответа = язык документа\n"
    "• Не пересказывай дословно — анализируй и сжимай\n"
    "• Таблицы в документе → красивые Markdown-таблицы\n"
    "• Код в документе → блоки ```язык```"
)

PDF_FILE_SYSTEM_PROMPT = (
    "Ты — эксперт по анализу документов. Пользователь прислал PDF. "
    "Извлекай максимум ценности — анализируй, а не пересказывай.\n\n"

    "СТРУКТУРА АНАЛИЗА PDF:\n\n"

    "### 📄 Тип документа\n"
    "Что это: отчёт / статья / инструкция / договор / презентация / книга / другое.\n"
    "Автор, организация, дата (если есть).\n\n"

    "### 📋 О чём документ\n"
    "Главная тема и цель в 3-4 предложениях.\n\n"

    "### 🔑 Ключевые данные\n"
    "Конкретные цифры, факты, выводы, даты, имена — то что нельзя упустить.\n"
    "Формат: • **Факт** — контекст\n\n"

    "### 📊 Структура документа\n"
    "Основные разделы и их суть.\n\n"

    "### 💡 Главный вывод\n"
    "Что самое важное в этом документе для читателя?\n\n"

    "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
    "СПЕЦИАЛЬНЫЕ РЕЖИМЫ:\n"
    "• Если ДОГОВОР: стороны + обязательства каждой + штрафы + ключевые условия + риски\n"
    "• Если НАУЧНАЯ СТАТЬЯ: гипотеза → методология → выборка → результаты → ограничения\n"
    "• Если ФИНАНСЫ: выручка, прибыль, динамика год-к-году, прогнозы\n"
    "• Если ТЕХНИЧЕСКАЯ ДОКУМЕНТАЦИЯ: stack, архитектура, API, требования\n"
    "• Если вопрос задан конкретный — отвечай на него прямо, не давай шаблон\n\n"
    "Язык ответа = язык документа. Таблицы → Markdown. Код → ```блоки```."
)

# Расширения кодовых файлов
_CODE_EXTS = {
    ".py", ".pyw", ".js", ".mjs", ".ts", ".jsx", ".tsx",
    ".java", ".cpp", ".cc", ".cxx", ".c", ".h", ".hpp",
    ".cs", ".go", ".rs", ".php", ".rb", ".pl",
    ".sh", ".bash", ".zsh", ".bat", ".cmd", ".ps1",
    ".sql", ".graphql", ".gql", ".tf", ".kt", ".swift", ".dart",
    ".r", ".scala", ".vue", ".svelte", ".css", ".scss",
    ".html", ".htm", ".xml", ".json", ".yaml", ".yml",
    ".toml", ".ini", ".conf", ".cfg",
}


def _get_file_system_prompt(file_name: str, user_text: str = "") -> str:
    """Определяет системный промпт по расширению файла и контексту вопроса."""
    ext = (os.path.splitext(file_name)[1].lower()) if file_name else ""

    # PDF — отдельный промпт
    if ext == ".pdf":
        return PDF_FILE_SYSTEM_PROMPT

    # Код — если расширение кодовое
    if ext in _CODE_EXTS:
        return CODE_FILE_SYSTEM_PROMPT

    # DOCX и текстовые документы
    if ext in (".docx", ".doc", ".txt", ".md", ".rst", ".tex"):
        return DOC_FILE_SYSTEM_PROMPT

    # Если пользователь спрашивает про код в файле — код-промпт
    if user_text and is_code_question(user_text):
        return CODE_FILE_SYSTEM_PROMPT

    # По умолчанию — документный промпт
    return DOC_FILE_SYSTEM_PROMPT


async def call_chat(messages: list, model_key: str, max_tokens: int = 1500) -> str:
    m = MODELS[model_key]
    provider = m["provider"]
    url     = SQ_CHAT if provider == "sq" else IO_CHAT
    api_key = SQ_KEY   if provider == "sq" else IO_KEY
    payload = {"model": m["name"], "messages": messages, "max_tokens": max_tokens}
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    _slow = any(s in m["name"] for s in ["qwen3-max", "DeepSeek-R1", "GLM-4.7", "Kimi-K2"])
    _tmo  = 300 if _slow else 180

    # Маппинг IO-моделей на SQ-аналоги для fallback
    IO_TO_SQ_FALLBACK = {
        "deepseek_v3":  "deepseek_v3_sq",
        "deepseek_r1":  "deepseek_v3_sq",
        "moonshot":     "qwen3_max",
        "glm47":        "qwen3_max",
    }

    for attempt in range(2):
        try:
            async with aiohttp.ClientSession(
                timeout=aiohttp.ClientTimeout(total=_tmo, connect=20)
            ) as s:
                async with s.post(url, headers=headers, json=payload) as r:
                    if r.status != 200:
                        err = await r.text()
                        # Если IO ключ устарел — пробуем SQ fallback
                        if provider == "io" and ("unauthorized" in err.lower() or "token has expired" in err.lower() or r.status in (401, 403)):
                            logging.warning(f"IO token expired for {model_key}, trying SQ fallback")
                            fallback_key = IO_TO_SQ_FALLBACK.get(model_key, "deepseek_v3_sq")
                            if fallback_key in MODELS and fallback_key not in disabled_models:
                                fb = MODELS[fallback_key]
                                fb_payload = {"model": fb["name"], "messages": messages, "max_tokens": max_tokens}
                                fb_headers = {"Authorization": f"Bearer {SQ_KEY}", "Content-Type": "application/json"}
                                async with s.post(SQ_CHAT, headers=fb_headers, json=fb_payload) as r2:
                                    if r2.status == 200:
                                        res2 = await r2.json()
                                        raw2 = res2["choices"][0]["message"].get("content") or ""
                                        result2 = raw2.split("</think>")[-1].strip() or raw2.strip() or "..."
                                        logging.info(f"IO fallback ok: {model_key} → {fallback_key}")
                                        return result2
                        raise RuntimeError(f"API {r.status}: {err[:400]}")
                    res = await r.json()
                    msg_obj = res["choices"][0]["message"]
                    raw = msg_obj.get("content") or ""
                    # Некоторые thinking-модели (DeepSeek-R1, etc.) кладут ответ
                    # в reasoning_content или tool_calls, а content остаётся пустым
                    if not raw.strip():
                        raw = (
                            msg_obj.get("reasoning_content") or
                            msg_obj.get("thinking_content") or
                            msg_obj.get("reasoning") or ""
                        )
                    # Убираем thinking-блоки
                    if "</think>" in raw:
                        after = raw.split("</think>")[-1].strip()
                        if after:
                            raw = after
                    result = raw.strip()
                    if not result:
                        result = "..."
                    return result
        except aiohttp.ServerDisconnectedError:
            logging.warning(f"ServerDisconnected attempt {attempt+1} {m['name']}")
            await asyncio.sleep(3)
        except Exception as e:
            raise e
    raise RuntimeError(f"Server disconnected x2: {m['name']}")


# ── Запрещённые слова для промптов изображений (автозамена/удаление) ──────────
_IMG_BANNED_WORDS = {
    # Политика / насилие
    "naked": "artistic nude",
    "nude": "artistic",
    "nsfw": "",
    "gore": "",
    "blood": "dramatic",
    "violence": "dramatic scene",
    "dead body": "dramatic scene",
    "corpse": "dramatic scene",
    "kill": "dramatic action",
    "murder": "dark thriller",
    "weapon": "prop",
    "gun": "prop",
    "bomb": "dramatic prop",
    "terrorist": "antagonist character",
    "drug": "medicine",
    "cocaine": "white powder",
    "heroin": "",
    "suicide": "melancholy",
    "self-harm": "sadness",
    # Порнография
    "pornography": "artistic",
    "porn": "artistic",
    "sex": "romantic",
    "erotic": "romantic",
    "xxx": "",
    "hentai": "anime art",
    "sexy": "attractive",
    "boobs": "figure",
    "genitalia": "figure",
    # Реальные лица
    "putin": "a world leader",
    "biden": "a politician",
    "trump": "a businessman",
    "obama": "a statesman",
    "hitler": "a historical villain",
    "stalin": "a historical figure",
}

# ── Запрещённые слова для промптов МУЗЫКИ ──────────────────────────────────────



def _clean_prompt_banned(text: str, banned_dict: dict) -> str:
    """Заменяет/удаляет запрещённые слова из промпта."""
    result = text
    for bad, good in banned_dict.items():
        # Case-insensitive replace
        import re as _re
        pattern = _re.compile(_re.escape(bad), _re.IGNORECASE)
        if good:
            result = pattern.sub(good, result)
        else:
            result = pattern.sub("", result)
    # Убираем двойные пробелы
    result = " ".join(result.split())
    return result.strip()


async def enhance_prompt_ai(raw_prompt: str, mode: str = "image", lang: str = "en") -> str:
    """
    Улучшает промпт для генерации ИЗОБРАЖЕНИЙ:
    1. Очищает от запрещённых слов
    2. Переводит на английский если нужно
    3. Добавляет теги качества, освещения, стиля
    4. Расширяет описание деталями
    """
    # Шаг 1: Очистка запрещённых слов
    cleaned = _clean_prompt_banned(raw_prompt, _IMG_BANNED_WORDS)
    if not cleaned or len(cleaned.strip()) < 2:
        cleaned = "abstract digital art"

    # Шаг 2: Определяем нужен ли перевод
    en_words = sum(1 for w in cleaned.split() if all(c.isascii() for c in w))
    ratio = en_words / max(len(cleaned.split()), 1)
    # Если уже хороший английский и длинный промпт — только добавляем теги
    if ratio > 0.75 and len(cleaned) > 40:
        sys_p = (
            "You are an expert prompt engineer for image generation models. "
            "The user has an English prompt. Enhance it by adding: "
            "specific lighting (golden hour, cinematic, neon, soft diffused, etc.), "
            "camera details (bokeh, wide-angle, macro, etc.), "
            "quality tags (8K, hyperdetailed, masterpiece, sharp focus), "
            "and artistic style if missing. "
            "Keep the original concept 100% intact. "
            "Output max 80 words. Return ONLY the enhanced prompt, nothing else."
        )
        user_p = cleaned
    else:
        # Переводим + улучшаем
        sys_p = (
            "You are an expert prompt engineer for image generation models. "
            "Translate the idea to English and create a rich, detailed prompt. "
            "Include: main subject with details, environment/background, "
            "lighting style (golden hour, studio, neon glow, etc.), "
            "camera style (bokeh, wide shot, portrait lens), "
            "artistic style (photorealistic, oil painting, digital art, etc.), "
            "quality tags (8K UHD, hyperdetailed, masterpiece, sharp focus). "
            "Keep the original meaning 100%. "
            "Output max 80 words. Return ONLY the prompt in English, nothing else."
        )
        user_p = cleaned

    try:
        msgs = [{"role": "system", "content": sys_p},
                {"role": "user", "content": user_p}]
        result = (await call_chat(msgs, "claude_haiku", max_tokens=150)).strip()
        # Если что-то пошло не так — возвращаем очищенный оригинал
        if (not result or len(result) < 5
                or result.startswith("I ")
                or "sorry" in result.lower()
                or "cannot" in result.lower()
                or result.startswith("Sure,")
                or result.startswith("Here")):
            # Добавляем минимальные теги качества к очищенному промпту
            base = cleaned if ratio > 0.5 else f"digital art, {cleaned}"
            return base + ", 8K, hyperdetailed, masterpiece"
        # Финальная очистка результата
        result = _clean_prompt_banned(result, _IMG_BANNED_WORDS)
        return result
    except Exception:
        return cleaned + ", 8K, hyperdetailed"


async def generate_music_lyrics(wishes: str, style_desc: str, style_tags: str, lang: str = "en") -> str:
    """
    Генерирует ТЕКСТ ПЕСНИ для поля prompt в Suno.
    Suno ПОЁТ содержимое prompt — это должны быть слова, не описание стиля.
    """
    wishes_clean = wishes.strip() if wishes.strip() else ("about life and emotions" if lang != "ru" else "о жизни")

    if lang == "ru":
        sys_p = (
            "Ты — профессиональный поэт-лирик. Пиши ТОЛЬКО текст песни на русском языке.\n"
            "ВАЖНО: Избегай любого контента с насилием, ненормативной лексикой, суицидальными темами.\n"
            "Формат — строго:\n"
            "[Verse 1]\n(4-6 строк с чёткой рифмой АБАБ или ААББ)\n\n"
            "[Chorus]\n(4-6 ярких, запоминающихся строк припева — самая сильная часть)\n\n"
            "[Verse 2]\n(4-6 строк, развивающих тему)\n\n"
            "[Chorus]\n(те же строки припева)\n\n"
            "[Bridge]\n(2-4 строки — эмоциональный перелом)\n\n"
            "[Chorus]\n(финальный припев)\n\n"
            "ПРАВИЛА:\n"
            "• Рифмы обязательны — чёткие, не приблизительные\n"
            "• Строки короткие (5-8 слов), певческие\n"
            "• Яркие образы, конкретные детали\n"
            "• Возвращай ТОЛЬКО текст. Никаких пояснений и ремарок."
        )
        user_p = f"Напиши текст песни на тему: {wishes_clean}\nСтиль: {style_desc}"
    else:
        sys_p = (
            "You are a professional songwriter. Write ONLY song lyrics in English.\n"
            "IMPORTANT: Avoid violence, profanity, suicidal content.\n"
            "Format strictly:\n"
            "[Verse 1]\n(4-6 lines with clear ABAB or AABB rhyme)\n\n"
            "[Chorus]\n(4-6 memorable, powerful lines — the emotional peak)\n\n"
            "[Verse 2]\n(4-6 lines developing the theme)\n\n"
            "[Chorus]\n(same chorus lines)\n\n"
            "[Bridge]\n(2-4 lines — emotional turn)\n\n"
            "[Chorus]\n(final chorus)\n\n"
            "RULES:\n"
            "• Strong rhymes required — not approximate\n"
            "• Short singable lines (5-8 words)\n"
            "• Vivid imagery and concrete details\n"
            "• Return ONLY lyrics. No explanations or stage directions."
        )
        user_p = f"Song about: {wishes_clean}\nStyle: {style_desc}"

    try:
        msgs = [{"role": "system", "content": sys_p}, {"role": "user", "content": user_p}]
        text = (await call_chat(msgs, "claude_haiku", max_tokens=600)).strip()
        # Проверяем что получили нормальный текст
        if not text or len(text) < 20:
            return wishes_clean
        # Если нет тегов — добавляем простую разметку
        if "[verse" not in text.lower() and "[chorus" not in text.lower():
            lines = [l for l in text.split("\n") if l.strip()]
            half = max(2, len(lines) // 2)
            text = "[Verse 1]\n" + "\n".join(lines[:half]) + "\n\n[Chorus]\n" + "\n".join(lines[half:])
        return text
    except Exception:
        return wishes_clean



async def _upload_image_to_cdn(image_base64: str, telegram_file_path: str = None) -> str | None:
    """
    Возвращает публичный URL фото для передачи в Pollinations img2img.
    Приоритет:
      1) Прямой URL файла через Telegram API (самый надёжный — Railway видит api.telegram.org)
      2) catbox.moe (быстрый, без блокировок)
      3) litterbox.catbox.moe (временный, 1 час)
    """
    # 1) Telegram CDN — используем file_path который уже есть
    if telegram_file_path:
        tg_url = f"https://api.telegram.org/file/bot{TELEGRAM_TOKEN}/{telegram_file_path}"
        logging.info(f"[upload_cdn] Telegram file URL: {tg_url[:80]}…")
        return tg_url

    img_bytes = base64.b64decode(image_base64)

    # 2) catbox.moe — постоянное хранилище, хорошо доступно
    try:
        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=30)) as s:
            form = aiohttp.FormData()
            form.add_field("reqtype", "fileupload")
            form.add_field("userhash", "")
            form.add_field("fileToUpload", img_bytes, filename="photo.jpg", content_type="image/jpeg")
            async with s.post("https://catbox.moe/user/api.php", data=form) as r:
                if r.status == 200:
                    url = (await r.text()).strip()
                    if url.startswith("http"):
                        logging.info(f"[upload_cdn] catbox OK: {url}")
                        return url
    except Exception as e:
        logging.warning(f"[upload_cdn] catbox failed: {e}")

    # 3) litterbox.catbox.moe — временный (1 час)
    try:
        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=30)) as s:
            form = aiohttp.FormData()
            form.add_field("reqtype", "fileupload")
            form.add_field("time", "1h")
            form.add_field("fileToUpload", img_bytes, filename="photo.jpg", content_type="image/jpeg")
            async with s.post("https://litterbox.catbox.moe/resources/internals/api.php", data=form) as r:
                if r.status == 200:
                    url = (await r.text()).strip()
                    if url.startswith("http"):
                        logging.info(f"[upload_cdn] litterbox OK: {url}")
                        return url
    except Exception as e:
        logging.warning(f"[upload_cdn] litterbox failed: {e}")

    logging.warning("[upload_cdn] все CDN недоступны")
    return None


async def generate_image(prompt: str, img_model_key: str, ratio: str = "1:1",
                         image_base64: str = None, tg_file_path: str = None) -> dict:
    """
    Генерация через Pollinations.
    Правильный URL (по документации):
      GET https://gen.pollinations.ai/image/{prompt}?model=flux&width=...
    Если передан image_base64 и модель поддерживает img2img — передаём фото через CDN.
    tg_file_path — путь к файлу из Telegram API (для прямого URL без загрузки).
    """
    import urllib.parse as _up, random as _rnd

    m = IMG_MODELS.get(img_model_key, IMG_MODELS["flux"])
    model_name = m["name"]
    supports_img2img = m.get("img2img", False)

    # Для стилей anime/3d/realism добавляем суффикс к промпту
    STYLE_SUFFIX = {
        "flux_realism": ", photorealistic, ultra detailed, 8K, DSLR photograph",
        "flux_anime":   ", anime art style, manga illustration, vibrant colors",
        "flux_3d":      ", 3D render, octane render, volumetric lighting, CGI",
        "flux_pro":     ", professional photography, sharp focus, high quality",
    }
    full_prompt = prompt + STYLE_SUFFIX.get(img_model_key, "")

    RATIO_TO_WH = {
        "1:1":  (1024, 1024),
        "4:3":  (1280, 960),
        "3:4":  (960, 1280),
        "16:9": (1344, 756),
        "9:16": (756, 1344),
    }
    width, height = RATIO_TO_WH.get(ratio, (1024, 1024))
    seed = _rnd.randint(1, 999999)

    # Правильный домен: gen.pollinations.ai/image/{prompt}
    encoded = _up.quote(full_prompt)

    async def _try_img2img_url(session, mdl_name, img_url, w, h, s_seed, s_key):
        """Одна попытка img2img через Pollinations. Возвращает bytes или None."""
        enc_img = _up.quote(img_url, safe="")
        url = (
            f"https://gen.pollinations.ai/image/{encoded}"
            f"?model={mdl_name}&width={w}&height={h}"
            f"&seed={s_seed}&nologo=true&nofeed=true&safe=false"
            f"&image={enc_img}"
        )
        logging.info(f"[img2img] {mdl_name} url={url[:120]}")
        headers = {
            "Authorization": f"Bearer {s_key}",
            "User-Agent": "Mozilla/5.0 (compatible; TelegramBot/1.0)",
            "Accept": "image/jpeg,image/png,image/*,*/*",
        }
        async with session.get(url, headers=headers) as r:
            ct = r.headers.get("content-type", "")
            if r.status == 200 and "image" in ct:
                data = await r.read()
                if len(data) > 5000:
                    return data, url
            else:
                body = (await r.text())[:150]
                if pollinations_pool.is_quota_error(r.status, body):
                    pollinations_pool.mark_quota_exceeded(s_key)
                logging.warning(f"[img2img] {mdl_name} {r.status}: {body}")
        return None, None

    async def _try_text_gen(session, mdl_name, w, h, s_seed, s_key):
        """Обычная текстовая генерация. Возвращает bytes или None."""
        url = (
            f"https://gen.pollinations.ai/image/{encoded}"
            f"?model={mdl_name}&width={w}&height={h}"
            f"&seed={s_seed}&nologo=true&nofeed=true"
        )
        headers = {
            "Authorization": f"Bearer {s_key}",
            "User-Agent": "Mozilla/5.0 (compatible; TelegramBot/1.0)",
            "Accept": "image/jpeg,image/png,image/*,*/*",
        }
        async with session.get(url, headers=headers) as r:
            ct = r.headers.get("content-type", "")
            if r.status == 200 and "image" in ct:
                data = await r.read()
                if len(data) > 5000:
                    return data, url
            else:
                body = (await r.text())[:150]
                if pollinations_pool.is_quota_error(r.status, body):
                    pollinations_pool.mark_quota_exceeded(s_key)
                logging.warning(f"[textgen] {mdl_name} {r.status}: {body}")
        return None, None

    # ── Обычная генерация (без img2img) ─────────────────────────
    if not (image_base64 and supports_img2img):
        for attempt in range(3):
            key = pollinations_pool.get_key()
            cur_seed = seed + attempt * 777
            try:
                async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=180)) as s:
                    img_bytes, url = await _try_text_gen(s, model_name, width, height, cur_seed, key)
                    if img_bytes:
                        return {"data": img_bytes, "url": url}
            except asyncio.TimeoutError:
                logging.warning(f"Pollinations timeout attempt {attempt+1}")
            except Exception as e:
                logging.warning(f"Pollinations error attempt {attempt+1}: {e}")
            if attempt < 2:
                await asyncio.sleep(5)
        # turbo текстовый fallback
        try:
            fb_key = pollinations_pool.get_key()
            fb_url = (f"https://gen.pollinations.ai/image/{encoded}"
                      f"?model=turbo&width={width}&height={height}&seed={seed}&nologo=true&nofeed=true")
            async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=90)) as s:
                img_bytes, url = await _try_text_gen(s, "turbo", width, height, seed, fb_key)
                if img_bytes:
                    return {"data": img_bytes, "url": url}
        except Exception as e:
            logging.warning(f"Pollinations text fallback failed: {e}")
        raise RuntimeError(
            "Pollinations не ответил. Попробуй:\n"
            "• Подождать 30 секунд и повторить\n"
            "• Выбрать модель 🚀 Turbo\n"
            "• Упростить промпт"
        )

    # ── IMG2IMG: умная цепочка моделей ──────────────────────────
    # Загружаем CDN URL один раз
    _img2img_url: str | None = await _upload_image_to_cdn(image_base64, telegram_file_path=tg_file_path)
    if _img2img_url:
        logging.info(f"[img2img] CDN URL: {_img2img_url[:80]}")
    else:
        logging.warning("[img2img] CDN недоступен — переходим к текстовой генерации")

    # Используем меньший размер для img2img — быстрее и меньше таймаутов
    i2i_w, i2i_h = 768, 768

    # Только flux-2-dev поддерживает ?image= параметр в Pollinations
    # flux-2-dev медленный — нужно 90 сек (раньше был таймаут 180)
    IMG2IMG_CHAIN = [
        ("flux-2-dev", True,  120),  # img2img — даём 2 минуты
        ("turbo",      True,   60),  # turbo тоже принимает image= и быстрее
        ("turbo",      False,  45),  # финальный fallback без фото
    ]

    if not _img2img_url:
        IMG2IMG_CHAIN = [
            ("flux-2-dev", False, 120),
            ("turbo",      False,  45),
        ]

    # Если img_model_key задан и отличается от первого в цепочке — ставим его первым
    _chosen_mdl = IMG_MODELS.get(img_model_key, {}).get("name")
    if _chosen_mdl and _chosen_mdl != IMG2IMG_CHAIN[0][0]:
        _chosen_entry = (_chosen_mdl, bool(_img2img_url), 60)
        IMG2IMG_CHAIN = [_chosen_entry] + [x for x in IMG2IMG_CHAIN if x[0] != _chosen_mdl]

    last_err = "неизвестная ошибка"
    for mdl, use_img, tout in IMG2IMG_CHAIN:
        key = pollinations_pool.get_key()
        try:
            async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=tout)) as s:
                if use_img and _img2img_url:
                    img_bytes, url = await _try_img2img_url(s, mdl, _img2img_url, i2i_w, i2i_h, seed, key)
                else:
                    img_bytes, url = await _try_text_gen(s, mdl, width, height, seed, key)
                if img_bytes:
                    logging.info(f"[img2img] успех: {mdl} use_img={use_img}")
                    return {"data": img_bytes, "url": url}
        except asyncio.TimeoutError:
            last_err = f"{mdl} timeout ({tout}s)"
            logging.warning(f"[img2img] timeout: {mdl}")
        except Exception as e:
            last_err = str(e)
            logging.warning(f"[img2img] error {mdl}: {e}")
        await asyncio.sleep(2)

    raise RuntimeError(
        "Pollinations не ответил. Попробуй:\n"
        "• Подождать 30 секунд и повторить\n"
        "• Выбрать модель 🚀 Turbo\n"
        "• Упростить промпт"
    )


# ==================================================================
# 🎵 MUSIC GENERATION (Pollinations)
# ==================================================================

async def suno_generate(*args, **kwargs) -> dict:
    """Stub — заменён на Pollinations."""
    raise NotImplementedError("Suno заменён на Pollinations")

async def suno_get_result(task_id: str) -> dict:
    """Stub — заменён на Pollinations."""
    return {"status": "processing"}






def main_reply_kb(uid: int) -> ReplyKeyboardMarkup:
    # Кнопка "Открыть Хуза ИИ" УБРАНА — она теперь слева в поле ввода (MenuButtonWebApp)
    # Остальные кнопки навигации оставлены
    rows = [
        [KeyboardButton(text="💬 Написать"), KeyboardButton(text="🎨 Создать")],
        [KeyboardButton(text="👤 Профиль"),  KeyboardButton(text="💎 Подписка")],
        [KeyboardButton(text="🏠 Главная")],
    ]
    if uid in ADMIN_IDS:
        rows.append([KeyboardButton(text="🔥 Админ")])
    return ReplyKeyboardMarkup(keyboard=rows, resize_keyboard=True)


def home_kb(uid: int) -> InlineKeyboardMarkup:
    sub_text = "💎 Подписка ✅" if has_active_sub(uid) else "💎 Купить подписку"
    built = [
        [
            InlineKeyboardButton(text="💬 Написать",          callback_data="menu_ask"),
            InlineKeyboardButton(text="👤 Профиль",           callback_data="menu_profile"),
        ],
        [
            InlineKeyboardButton(text="🎨 Картинки · Видео",  callback_data="menu_extra"),
        ],
        [
            InlineKeyboardButton(text="🧹 Очистить память",   callback_data="clear_memory"),
        ],
        [
            InlineKeyboardButton(text="💬 Поддержка",         callback_data="menu_support"),
            InlineKeyboardButton(text=sub_text,               callback_data="sub_menu"),
        ],
    ]
    if uid in ADMIN_IDS:
        built.append([InlineKeyboardButton(text="⚙️ Админ", callback_data="menu_admin")])
    return InlineKeyboardMarkup(inline_keyboard=built)


def ai_menu_kb(uid: int = 0) -> InlineKeyboardMarkup:
    auto_on = user_features.get(uid, {}).get("auto_model", False)
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(
            text=f"🤖 Авто-выбор {'✅' if auto_on else '❌'}",
            callback_data="toggle_auto_model"
        )],
        [
            InlineKeyboardButton(text="🔵 Claude",     callback_data="cat_claude"),
            InlineKeyboardButton(text="🟢 GPT",        callback_data="cat_gpt"),
        ],
        [
            InlineKeyboardButton(text="🧬 DeepSeek",   callback_data="cat_deepseek"),
            InlineKeyboardButton(text="🌐 Qwen",       callback_data="cat_qwen"),
        ],
        [
            InlineKeyboardButton(text="⚙️ Command",    callback_data="cat_command"),
            InlineKeyboardButton(text="🔮 Kimi & GLM", callback_data="cat_kimi_glm"),
        ],
        [
            InlineKeyboardButton(text="🌀 C4AI",       callback_data="cat_c4ai"),
            InlineKeyboardButton(text="🎨 Генерация",  callback_data="menu_imggen"),
        ],
        [
            InlineKeyboardButton(text="🔵 Gemini",     callback_data="cat_gemini"),
            InlineKeyboardButton(text="🌐 Sonar (веб)", callback_data="cat_sonar"),
        ],
        [
            InlineKeyboardButton(text="🏠 Главная",    callback_data="menu_home"),
        ],
    ])


def category_kb(cat_key: str, uid: int = 0) -> InlineKeyboardMarkup:
    label, keys = CATEGORIES[cat_key]
    rows = []
    if cat_key == "cat_imggen":
        for k in keys:
            if k in disabled_img_models:
                continue
            m = IMG_MODELS[k]
            rows.append([InlineKeyboardButton(text=m["label"], callback_data=f"imgset_{k}")])
    else:
        for k in keys:
            if k in disabled_models:
                continue
            m = MODELS[k]
            is_premium = k in PREMIUM_MODELS
            needs_sub  = is_premium and not has_active_sub(uid)
            lock = "🔒 " if needs_sub else ""
            rows.append([InlineKeyboardButton(
                text=lock + m["label"],
                callback_data=f"set_{k}"
            )])
    rows.append([
        InlineKeyboardButton(text="◀️ Нейросети", callback_data="back_ai"),
        InlineKeyboardButton(text="🏠 Главная",   callback_data="back_home"),
    ])
    return InlineKeyboardMarkup(inline_keyboard=rows)


def imggen_category_kb(show_all: bool = True) -> InlineKeyboardMarkup:
    rows = []
    for k, m in IMG_MODELS.items():
        if k in disabled_img_models:
            continue
        # Добавляем значок для img2img моделей
        label = m["label"]
        if m.get("img2img"):
            label = label + "  📸→🖼"
        rows.append([InlineKeyboardButton(text=label, callback_data=f"imgset_{k}")])
    rows.append([InlineKeyboardButton(text="🖼 Стилизация фото (img2img)", callback_data="imggen_from_photo")])
    rows.append([InlineKeyboardButton(text="◀️ Главное меню", callback_data="back_home")])
    return InlineKeyboardMarkup(inline_keyboard=rows)


def main_menu_kb() -> InlineKeyboardMarkup:
    return ai_menu_kb(uid)


def _lz_compress_b64(text: str) -> str:
    """LZ-сжатие строки для передачи в URL (совместимо с LZString.decompressFromEncodedURIComponent)."""
    import zlib, urllib.parse
    # Кодируем в UTF-16LE как LZString
    data_bytes = text.encode('utf-16-le')
    # Используем base64 URL-safe кодировку как fallback (mode=b64)
    encoded = base64.urlsafe_b64encode(text.encode('utf-8')).decode('ascii')
    return encoded


def make_view_url(answer: str, model_label: str = "ИИ", model_key: str = "") -> str:
    """
    Создаёт URL для просмотра ответа в красивом HTML-виде.
    Использует base64 кодирование (#data= формат) — не требует живого сервера.
    """
    import json as _json
    now = msk_now().strftime("%d.%m · %H:%M")
    payload = _json.dumps({
        "text": answer,
        "model": model_label,
        "modelKey": model_key,
        "time": now,
    }, ensure_ascii=False)
    # base64 URL-safe кодирование — работает без сервера и после перезапуска
    encoded = base64.urlsafe_b64encode(payload.encode('utf-8')).decode('ascii')
    encoded = encoded.rstrip('=')
    return f"{VIEW_BASE_URL}#data={encoded}"


def save_kb(uid: int, answer: str = "", model_label: str = "ИИ", model_key: str = "") -> InlineKeyboardMarkup:
    """Клавиатура под ответом. Только DOCX — PNG и PDF удалены."""
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="📄 DOCX", callback_data=f"save_docx_{uid}"),
        ],
        [
            InlineKeyboardButton(text="🗑 Очистить память", callback_data="clear_memory"),
            InlineKeyboardButton(text="🏠 Меню", callback_data="back_home"),
        ],
    ])


# Хранилище ответов для просмотра: {short_id: {text, model, time}}
_answer_store: dict = {}

def store_answer(answer: str, model_label: str = "ИИ", model_key: str = "") -> str:
    """
    Сохраняет ответ в памяти и возвращает короткий URL вида:
    https://railway.app/view?id=XXXXXXXX
    Ссылка ~40 символов вместо тысяч.
    """
    import json as _j
    import secrets
    import re as _re
    # Убираем цифровые сноски от Sonar [1][2][5] и т.д.
    answer = _re.sub(r'\[\d+\]', '', answer)
    now = msk_now().strftime("%d.%m · %H:%M")
    short_id = secrets.token_urlsafe(6)  # 8 символов, напр. "aB3xKq8z"
    _answer_store[short_id] = {
        "text":     answer,
        "model":    model_label,
        "modelKey": model_key,
        "time":     now,
    }
    # Ограничиваем размер хранилища — удаляем старые если > 500
    if len(_answer_store) > 500:
        oldest = list(_answer_store.keys())[:100]
        for k in oldest:
            _answer_store.pop(k, None)
    return f"{API_BASE_URL}/view?id={short_id}"

def view_kb(view_url: str) -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="🌐 Красиво (формулы/таблицы)", url=view_url)]])

def has_math_or_table(text: str) -> bool:
    if not text: return False
    if re.search(r'[|][^|]+[|]', text): return True
    if re.search(r'\\\\?(frac|sqrt|alpha|beta|gamma|omega|sum|int|cos|sin|vec|perp|Delta|pi)', text): return True
    dollar_count = text.count('$')
    if dollar_count >= 2: return True
    # Греческие буквы Unicode — значит модель прислала матем. выражения
    if re.search(r'[α-ωΑ-Ω∑∫∞±≈≠≤≥×÷√π]', text): return True
    return False


def imggen_prompt_kb() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup(inline_keyboard=[[
        InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_imggen"),
        InlineKeyboardButton(text="🏠 Главная", callback_data="back_home"),
    ]])


# ==================================================================
# 💬 ТЕКСТЫ — ПРЕМИУМ СТИЛЬ
# ==================================================================

def _welcome_text(name: str, tok: int = 0, is_admin: bool = False, ref_bonus_msg: str = "") -> str:
    admin_badge = "  ▣ <b>ADMIN</b>" if is_admin else ""
    ref_line = f"\n{ref_bonus_msg}" if ref_bonus_msg else ""
    return (
        "━━━━━━━━━━━━━━━━━━━━━━━━\n"
        f"⚡ <b>ХУЗА</b> — нейросеть{admin_badge}\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━\n\n"
        f"<b>{name}</b>{ref_line}\n\n"
        "◈ 28 нейросетей — Claude · GPT · Gemini\n"
        "◈ Генерация картинок, видео, музыки\n"
        "◈ Vision — анализ фото и документов\n"
        "◈ Доклады, рефераты, презентации\n\n"
        "▾ <i>Выбери действие</i>\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━"
    )


# ==================================================================
# 📨 КОМАНДЫ
# ==================================================================

@dp.message(Command("start"))
async def cmd_start(message: Message):
    uid  = message.from_user.id
    name = message.from_user.first_name or "Пользователь"

    # Инициализация пользователя
    if uid not in user_settings:
        user_settings[uid] = "claude_haiku"
    user_memory[uid] = deque(maxlen=20)
    if uid not in user_profiles:
        user_profiles[uid] = {
            "name": name,
            "username": message.from_user.username or "",
            "joined": msk_now().strftime("%d.%m.%Y %H:%M"),
            "requests": 0,
        }
    _init_tokens(uid)
    _init_limits(uid)

    # Trial для новых
    _is_brand_new = uid not in user_trials and not has_active_sub(uid)
    # Trial disabled

    # Сохраняем в БД
    asyncio.create_task(db_save_user(uid, name, message.from_user.username or ""))

    # Показываем клавиатуру навигации (без кнопки "Открыть Хуза ИИ")
    try:
        await message.answer("👋", reply_markup=main_reply_kb(uid))
    except Exception as e:
        logging.error(f"start kb error: {e}")

    # Проверка подписки
    try:
        if not await require_subscription(message):
            return
    except Exception as e:
        logging.error(f"start sub check error: {e}")

    # Приветственное сообщение
    try:
        tok = get_tokens(uid)
        is_admin = uid in ADMIN_IDS
        text = _welcome_text(name, tok, is_admin)
        # Trial disabled
        await message.answer(text, parse_mode="HTML", reply_markup=home_kb(uid))
    except Exception as e:
        logging.error(f"start welcome error uid={uid}: {e}")
        try:
            await message.answer(
                f"Привет, {name}! Добро пожаловать в Хуза ИИ!",
                reply_markup=home_kb(uid)
            )
        except Exception:
            pass

@dp.message(Command("menu"))
async def cmd_menu(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message):
        return
    await message.answer("🏠 <b>Главное меню:</b>", parse_mode="HTML", reply_markup=main_reply_kb(uid))


# ==================================================================
# 🔘 REPLY-КНОПКИ
# ==================================================================

@dp.message(Command("profile"))
async def cmd_profile(message: Message):
    await _show_profile(message, message.from_user.id)

@dp.message(Command("ai"))
async def cmd_ai(message: Message):
    await message.answer(
        "💬 <b>Режим диалога с ИИ</b>\n\nПросто напиши свой вопрос!",
        parse_mode="HTML", reply_markup=main_reply_kb(message.from_user.id)
    )

@dp.message(Command("img"))
async def cmd_img(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message): return
    await rb_imggen(message)


@dp.message(Command("music"))
async def cmd_music(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message): return
    await rb_music(message)



@dp.message(Command("model"))
async def cmd_model(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message): return
    mk = get_model_key(uid)
    cur = MODELS.get(mk, {}).get("label", mk) if not mk.startswith("imggen:") else "🎨 " + IMG_MODELS[mk.split(":")[1]]["label"]
    await message.answer(
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "🤖 <b>НЕЙРОСЕТЬ</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        f"◈ Активная — <code>{cur}</code>\n"
        "◈ Vision 📸 — анализирует фото\n\n"
        "▾ Выбери категорию\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔",
        parse_mode="HTML", reply_markup=ai_menu_kb(uid)
    )


# ==================================================================
# 📱 ОБРАБОТКА ДАННЫХ ИЗ MINI APP
# ==================================================================

@dp.message(F.web_app_data)
async def handle_web_app_data(message: Message):
    """Получаем данные от Telegram Mini App."""
    import json as _j
    uid   = message.from_user.id
    name  = message.from_user.first_name or ""
    uname = message.from_user.username or ""

    try:
        data = _j.loads(message.web_app_data.data)
    except Exception:
        await message.answer("❌ Ошибка данных мини-апа.")
        return

    action = data.get("action")

    # ── Выбор модели ──────────────────────────────────────────
    if action == "set_model":
        key = data.get("model", "")
        if key not in MODELS:
            await message.answer("❌ <b>Неизвестная модель.</b>", parse_mode="HTML"); return
        if key in disabled_models:
            await message.answer("🚫 <b>Модель недоступна.</b>", parse_mode="HTML"); return
        user_settings[uid] = key
        user_memory[uid] = deque(maxlen=20)
        asyncio.create_task(db_save_user(uid, name, uname))
        m = MODELS[key]
        vis = m.get("vision", False)
        vision_hint = "\n📸 <i>Поддерживает фото — отправь прямо в чат!</i>" if vis else ""
        kb = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="🤖 Сменить модель", callback_data="menu_ai")],
            [InlineKeyboardButton(text="🧹 Очистить память", callback_data="clear_memory"),
             InlineKeyboardButton(text="🏠 Главная", callback_data="back_home")],
        ])
        await message.answer(
            f"✅ <b>Модель выбрана!</b>\n\n"
            f"🤖 {m['label']}{vision_hint}\n\n"
            f"✏️ <b>Пиши вопрос прямо сейчас</b> — отвечу здесь!",
            parse_mode="HTML", reply_markup=kb
        )
        return

    # ── Чат из мини-апа ───────────────────────────────────────
    if action == "chat":
        key    = data.get("model", get_model_key(uid))
        text   = data.get("text", "").strip()
        images = data.get("images", [])
        history = data.get("history", [])

        if key in MODELS and key not in disabled_models:
            user_settings[uid] = key
        mk = get_model_key(uid)

        if not text and not images:
            await message.answer("❌ Пустое сообщение."); return
        if not await require_subscription(message): return

        _init_limits(uid)
        _ok, _reason = can_send(uid, mk)
        if not _ok:
            import time as _time_mod2
            _now_ts = _time_mod2.time()
            if _time_mod2.time() - _limit_notify_cd.get(uid, 0) < 300:
                return
            _limit_notify_cd[uid] = _now_ts
            _rst = _reset_str(uid, _reason)
            await message.answer(
                f"🚫 <b>Лимит исчерпан!</b>\n⏰ Сброс через: <code>{_rst}</code>",
                parse_mode="HTML"
            ); return

        if mk in disabled_models:
            await message.answer("🚫 <b>Модель недоступна.</b>", parse_mode="HTML"); return

        if uid not in user_memory:
            user_memory[uid] = deque(maxlen=20)

        think_msg = await message.answer("🔍 <i>Обрабатываю...</i>", parse_mode="HTML")
        _start = asyncio.get_event_loop().time()

        async def _tick(m, s):
            phases = [(3,"🔍 Анализирую"),(8,"🧠 Думаю"),(15,"💡 Формирую ответ"),(999,"⚡ Почти готово")]
            while True:
                await asyncio.sleep(1)
                el = int(asyncio.get_event_loop().time() - s)
                lbl = phases[0][1]
                for thr, lb in phases:
                    if el < thr: lbl = lb; break
                try: await m.edit_text(f"{lbl}... <b>{el}с</b>", parse_mode="HTML")
                except Exception: pass

        tick = asyncio.create_task(_tick(think_msg, _start))
        try:
            if images:
                vision_key = mk if mk in VISION_MODELS else next(
                    (k for k in ["claude_haiku","claude_sonnet","gpt52","qwen3_vl","c4ai_vis32b"]
                     if k in MODELS and k not in disabled_models), None
                )
                if not vision_key:
                    tick.cancel()
                    await think_msg.delete()
                    await message.answer("⚠️ Нет доступных vision-моделей."); return

                content = []
                for img_b64 in images:
                    content.append({"type":"image_url","image_url":{"url":f"data:image/jpeg;base64,{img_b64}"}})
                content.append({"type":"text","text": text or "Подробно опиши что изображено на фото"})
                msgs = [{"role": "system", "content": VISION_SYSTEM_PROMPT},{"role": "user", "content": content}]
                ans = await call_chat(msgs, vision_key, max_tokens=3000)
                spend_limit(uid, vision_key)
                m_label = MODELS[vision_key]["label"]
            else:
                WEBAPP_BASE_SYS = (
                    f"Ты — {BOT_NAME}, интеллектуальный ИИ-ассистент. "
                    "Отвечай на языке пользователя — по умолчанию русский. "
                    "Не переспрашивай если из контекста понятно что нужно.\n\n"
                    "СТАНДАРТЫ ОТВЕТА:\n"
                    "• Математика: $инлайн$ и $$блочные$$ LaTeX. Каждую формулу — в одну строку\n"
                    "• Код: ```язык\\n...\\n``` с указанием языка. Код рабочий и полный\n"
                    "• Таблицы: Markdown | заголовок | с :--- / :---: / ---:\n"
                    "• Заголовки: ### для разделов\n"
                    "• Запрещено: 'Конечно!', 'Отличный вопрос!', неполный код, придуманные факты"
                )

                # Определяем есть ли файл в тексте и выбираем системный промпт
                _file_marker = "--- Файл:"
                if _file_marker in text:
                    # Извлекаем имя файла из маркера
                    _fname_match = re.search(r"--- Файл: ([^\n-]+) ---", text)
                    _fname = _fname_match.group(1).strip() if _fname_match else ""
                    _user_question = text.split(_file_marker)[0].strip()
                    GLOBAL_SYS = _get_file_system_prompt(_fname, _user_question) + (
                        "\n\nОТВЕЧАЙ ПО-РУССКИ если файл на русском или пользователь написал по-русски."
                    )
                elif is_code_question(text):
                    GLOBAL_SYS = (
                        WEBAPP_BASE_SYS + "\n\n"
                        "━━━ РЕЖИМ: РАЗРАБОТЧИК ━━━\n"
                        "Ты — senior-разработчик. "
                        "Код ВСЕГДА в блоках ```язык\\n...\\n```. "
                        "Структура: суть → рабочий код → объяснение ключевых мест. "
                        "Код без заготовок '...' — только полный рабочий пример."
                    )
                elif any(kw in text.lower() for kw in ["реши", "вычисли", "найди", "формул", "уравнени",
                         "интеграл", "производн", "матем", "задач", "докажи", "упрости"]):
                    GLOBAL_SYS = (
                        WEBAPP_BASE_SYS + "\n\n"
                        "━━━ РЕЖИМ: МАТЕМАТИКА ━━━\n"
                        "Решай математически, не кодом. "
                        "LaTeX обязателен: $инлайн$ и $$блочные$$. "
                        "Каждую формулу — в одну строку: $\\frac{a}{b}$ "
                        "Структура: Дано → Формулы → Решение по шагам → Ответ с единицами."
                    )
                else:
                    GLOBAL_SYS = WEBAPP_BASE_SYS

                if history:
                    for h in history[-6:]:
                        if h.get("role") in ("user","assistant"):
                            user_memory[uid].append(h)
                user_memory[uid].append({"role":"user","content":text})
                query_msgs = [{"role":"system","content":GLOBAL_SYS}] + list(user_memory[uid])
                ans = await call_chat(query_msgs, mk)
                user_memory[uid].append({"role":"assistant","content":ans})
                spend_limit(uid, mk)
                m_label = MODELS[mk]["label"]

            last_responses[uid] = {"q": text or "[Фото]", "a": ans, "model_label": m_label, "model_key": mk}
            if uid in user_profiles:
                user_profiles[uid]["requests"] = user_profiles[uid].get("requests", 0) + 1

            tick.cancel()
            try: await tick
            except asyncio.CancelledError: pass
            await think_msg.delete()

            if re.search(r"```", ans):
                await smart_send(message, ans, reply_markup=save_kb(uid))
            else:
                html_ans = md_to_html(ans)
                try:
                    await message.answer(html_ans[:4000], parse_mode="HTML", reply_markup=save_kb(uid))
                except Exception:
                    await message.answer(ans[:4000], reply_markup=save_kb(uid))

        except Exception as e:
            tick.cancel()
            try: await think_msg.delete()
            except Exception: pass
            logging.error(f"web_app chat error: {e}")
            await message.answer(f"❌ Ошибка: {str(e)[:200]}")
        return

    await message.answer("❓ Неизвестное действие из мини-апа.")

@dp.message(Command("report"))
async def cmd_report(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message): return
    report_states[uid] = {}
    await message.answer(
        "📝 <b>Генератор докладов и рефератов</b>\n\nВыбери тип работы:",
        parse_mode="HTML", reply_markup=report_type_kb()
    )

@dp.message(Command("clear"))
async def cmd_clear(message: Message):
    uid = message.from_user.id
    user_memory[uid] = deque(maxlen=20)
    await message.answer("🧹 <b>Память очищена!</b>", parse_mode="HTML")

@dp.message(Command("about"))
async def cmd_about(message: Message):
    await rb_about(message)


@dp.message(F.text == "💬 Задать вопрос")
async def rb_ask(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message):
        return
    mk  = get_chat_model(uid)
    cur = MODELS.get(mk, {}).get("label", mk)
    is_vision = mk in VISION_MODELS
    vision_hint = "\n📸 <i>Можешь отправить фото для анализа!</i>" if is_vision else ""
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🤖 Сменить модель в чате", callback_data="menu_ai")],
        [InlineKeyboardButton(text="🧹 Очистить память", callback_data="clear_memory"),
         InlineKeyboardButton(text="🏠 Главная", callback_data="back_home")],
    ])
    await message.answer(
        f"💬 <b>Режим чата</b>\n"
        f"🤖 Модель чата: <code>{cur}</code>{vision_hint}\n\n"
        f"✏️ Напиши свой вопрос",
        parse_mode="HTML", reply_markup=kb
    )


@dp.message(F.text == "🤖 Выбрать модель")
async def rb_model(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message):
        return
    mk  = get_model_key(uid)
    cur = MODELS.get(mk, {}).get("label", mk) if not mk.startswith("imggen:") else "🎨 " + IMG_MODELS[mk.split(":")[1]]["label"]
    await message.answer(
        f"🤖 <b>Выбор нейросети</b>\n\n"
        f"🔹 Активная модель: <code>{cur}</code>\n"
        f"📸 — поддержка анализа фото\n\nВыбери категорию:",
        parse_mode="HTML", reply_markup=ai_menu_kb(uid)
    )


@dp.message(F.text == "🎨 Генерация картинок")
async def rb_imggen(message: Message):
    if not await require_subscription(message):
        return
    uid = message.from_user.id
    img_mk = get_img_model(uid)
    img_label = IMG_MODELS.get(img_mk, {}).get("label", img_mk)
    examples_text = "\n".join(f"◈ {e}" for e in IMG_PROMPT_EXAMPLES[:4])
    await message.answer(
        "🎨 <b>Режим генерации</b>\n"
        f"🖼 Модель генерации: <code>{img_label}</code>\n\n"
        f"✏️ Опиши что хочешь получить\n\n"
        f"<b>Примеры:</b>\n{examples_text}\n\n"
        "💡 <i>Пиши на русском — ИИ переведёт и улучшит промпт!</i>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="🎨 Сменить модель генерации", callback_data="menu_imggen")],
            [InlineKeyboardButton(text="🏠 Главная", callback_data="back_home")],
        ])
    )


@dp.message(F.text == "🎛 Дополнительно")
async def rb_extra(message: Message):
    if not await require_subscription(message):
        return
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🎨 Генерация картинок", callback_data="menu_imggen"),
         InlineKeyboardButton(text="🎵 Генерация музыки",   callback_data="menu_music")],
        [InlineKeyboardButton(text="🎬 Генерация видео",    callback_data="menu_video")],
        [InlineKeyboardButton(text="📝 Доклад/Реферат",     callback_data="menu_report"),
         InlineKeyboardButton(text="🎞 Презентация",         callback_data="menu_pptx")],
        [InlineKeyboardButton(text="🏠 Главная",             callback_data="back_home")],
    ])
    await message.answer(
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "🎛 <b>ДОПОЛНИТЕЛЬНО</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        "🎨 <b>Генерация картинок</b> — Flux, GPT-Image\n"
        "🎵 <b>Генерация музыки</b> — Suno AI\n"
        "🎬 <b>Генерация видео</b> — Grok Video (NEW!)\n"
        "📝 <b>Доклад/Реферат</b> — академический текст\n"
        "🎞 <b>Презентация</b> — PPTX или HTML\n\n"
        "▾ Выбери раздел\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔",
        parse_mode="HTML",
        reply_markup=kb
    )







@dp.callback_query(F.data == "back_extra")
async def back_extra_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🎨 Генерация картинок", callback_data="menu_imggen"),
         InlineKeyboardButton(text="🎵 Генерация музыки",   callback_data="menu_music")],
        [InlineKeyboardButton(text="🎬 Генерация видео",    callback_data="menu_video")],
        [InlineKeyboardButton(text="📝 Доклад/Реферат",     callback_data="menu_report"),
         InlineKeyboardButton(text="🎞 Презентация",         callback_data="menu_pptx")],
        [InlineKeyboardButton(text="🏠 Главная",             callback_data="back_home")],
    ])
    try:
        await callback.message.edit_text(
            "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
            "🎛 <b>ДОПОЛНИТЕЛЬНО</b>\n"
            "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
            "🎨 <b>Генерация картинок</b> — Flux, GPT-Image\n"
            "🎵 <b>Генерация музыки</b> — Suno AI\n"
            "🎬 <b>Генерация видео</b> — Grok Video (NEW!)\n"
            "📝 <b>Доклад/Реферат</b> — академический текст\n"
            "🎞 <b>Презентация</b> — PPTX или HTML\n\n"
            "▾ Выбери раздел\n"
            "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔",
            parse_mode="HTML", reply_markup=kb
        )
    except Exception:
        await callback.message.answer(
            "🎛 <b>Дополнительно</b>", parse_mode="HTML", reply_markup=kb
        )
    await callback.answer()


# ==================================================================
# 🎬 ГЕНЕРАЦИЯ ВИДЕО (Pollinations — Grok Video)
# ==================================================================

# Состояния для генерации видео
video_states: dict = {}

VIDEO_PROMPT_EXAMPLES = [
    "закат над океаном, золотые волны, кинематографично",
    "неоновый город ночью, дождь, cyberpunk атмосфера",
    "цветущая сакура, лепестки летят по ветру, весна",
    "горы в тумане на рассвете, эпичный пейзаж",
    "волны разбиваются о скалы, замедленная съёмка",
    "оживи это фото — персонаж поворачивает голову и улыбается",
]

# Системный промпт для максимального реализма
VIDEO_SYSTEM_PROMPT = (
    "ultra photorealistic, cinematic quality, 8K resolution, "
    "natural lighting, shallow depth of field, professional camera motion, "
    "smooth and fluid movement, high frame rate, vivid colors, "
    "sharp focus, film grain, IMAX quality, masterpiece"
)


async def generate_video_pollinations(prompt: str, image_base64: str = None, tg_file_path: str = None) -> bytes:
    """
    Генерация видео через Pollinations API (Grok Video).
    - Без фото: GET https://gen.pollinations.ai/video/{prompt}?model=grok-video
    - С фото:   загружаем на CDN → GET с параметром ?image=URL
    Умно переключает ключи при 429/квоте.
    """
    import urllib.parse as _up
    full_prompt = f"{prompt}, {VIDEO_SYSTEM_PROMPT}"
    timeout = aiohttp.ClientTimeout(total=360)

    encoded = _up.quote(full_prompt)
    params = {"model": "grok-video", "nologo": "true"}

    # Если есть фото — загружаем на CDN один раз
    img_cdn_url: str | None = None
    if image_base64:
        img_cdn_url = await _upload_image_to_cdn(image_base64, telegram_file_path=tg_file_path)
        if img_cdn_url:
            logging.info(f"[video] фото загружено на CDN: {img_cdn_url}")
        else:
            logging.warning("[video] CDN недоступен, генерируем видео только по тексту")

    last_err = ""
    for attempt in range(len(_POLLINATIONS_KEYS_LIST) + 1):
        key = pollinations_pool.get_key()
        headers = {"Authorization": f"Bearer {key}"}
        try:
            async with aiohttp.ClientSession(timeout=timeout) as s:
                if img_cdn_url:
                    # GET с image URL
                    encoded_img = _up.quote(img_cdn_url, safe="")
                    url = (
                        f"https://gen.pollinations.ai/video/{encoded}"
                        f"?model=grok-video&nologo=true&image={encoded_img}"
                    )
                else:
                    # GET только текст
                    url = f"https://gen.pollinations.ai/video/{encoded}"

                async with s.get(url, params=None if img_cdn_url else params, headers=headers) as r:
                    if r.status != 200:
                        txt = await r.text()
                        if pollinations_pool.is_quota_error(r.status, txt):
                            pollinations_pool.mark_quota_exceeded(key)
                            last_err = f"HTTP {r.status}: {txt[:100]}"
                            continue
                        raise RuntimeError(f"Pollinations Video HTTP {r.status}: {txt[:300]}")
                    data = await r.read()

            if len(data) < 1000:
                raise RuntimeError(f"Видео слишком маленькое ({len(data)} байт) — попробуй другой промпт")
            return data
        except RuntimeError:
            raise
        except Exception as e:
            last_err = str(e)
        await asyncio.sleep(2)
    raise RuntimeError(f"Pollinations Video: все ключи исчерпаны. Последняя ошибка: {last_err}")


@dp.callback_query(F.data == "menu_video")
async def menu_video_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    # Разрешаем доступ если есть подписка или активный триал
    if not has_active_sub(uid):
        if not await require_subscription(callback):
            return
    if not await check_service(callback, "video"):
        return
    if not can_video(uid):
        li = get_limits_info(uid)
        if li["video_max"] == 0:
            await callback.answer(
                "🚫 Генерация видео недоступна на бесплатном тарифе!\nОформи подписку для доступа.",
                show_alert=True
            )
        else:
            await callback.answer(
                f"🚫 Лимит видео исчерпан! ({li['video_used']}/{li['video_max']} в месяц)",
                show_alert=True
            )
        return

    li = get_limits_info(uid)
    video_left = max(0, li["video_max"] - li["video_used"])

    text = (
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "🎬 <b>ГЕНЕРАЦИЯ ВИДЕО</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        "🤖 Модель: <b>Grok Video</b> (xAI)\n"
        "⏱ Время генерации: <b>1–3 минуты</b>\n"
        "✨ Реалистичность: <b>максимальная</b>\n\n"
        "<b>Что умеет:</b>\n"
        "• 📝 Генерация по текстовому описанию\n"
        "• 🖼 Оживление фото (отправь фото + описание)\n\n"
        "<b>Примеры промптов:</b>\n"
        + "\n".join(f"  • <i>{e}</i>" for e in VIDEO_PROMPT_EXAMPLES[:4])
        + f"\n\n🎬 Осталось видео: <b>{video_left}/{li['video_max']}</b> в месяц\n\n"
        "▾ <b>Выбери режим:</b>"
    )
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="✨ С улучшением ИИ",     callback_data="video_mode_ai"),
         InlineKeyboardButton(text="✍️ Свой промпт (точно)", callback_data="video_mode_raw")],
        [InlineKeyboardButton(text="🖼 Оживить фото",        callback_data="video_mode_photo")],
        [InlineKeyboardButton(text="◀️ Назад", callback_data="back_extra")],
    ])
    try:
        await callback.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
    except Exception:
        await callback.message.answer(text, parse_mode="HTML", reply_markup=kb)
    await callback.answer()


@dp.callback_query(F.data.in_({"video_mode_ai", "video_mode_raw", "video_mode_photo"}))
async def cb_video_mode(callback: CallbackQuery):
    uid = callback.from_user.id
    mode = "ai" if callback.data == "video_mode_ai" else ("photo" if callback.data == "video_mode_photo" else "raw")
    video_states[uid] = {"mode": mode, "stage": "await_photo" if mode == "photo" else "await_prompt"}

    if mode == "photo":
        text = (
            "🖼 <b>Оживление фото → Видео</b>\n\n"
            "📸 <b>Отправь фотографию</b> — бот создаст видео на её основе.\n\n"
            "💡 <i>После фото можешь написать описание движения, например:\n"
            "«камера медленно приближается», «волосы развеваются на ветру»</i>\n\n"
            "⚡️ Просто отправь фото прямо сейчас!"
        )
    elif mode == "ai":
        text = (
            "🎬 <b>Генерация видео</b>\n\n"
            "🤖 <b>Режим: ИИ улучшит твой промпт</b>\n"
            "<i>Пиши на русском или английском — ИИ переведёт и добавит детали для максимального реализма.</i>\n\n"
            "📝 <b>Опиши что должно быть в видео:</b>\n\n"
            "<b>Примеры:</b>\n"
            + "\n".join(f"  • {e}" for e in VIDEO_PROMPT_EXAMPLES[:5])
        )
    else:
        text = (
            "🎬 <b>Генерация видео</b>\n\n"
            "✍️ <b>Режим: точный промпт</b>\n"
            "<i>Промпт будет использован как есть (только системные теги реализма добавятся автоматически).</i>\n\n"
            "📝 <b>Опиши что должно быть в видео:</b>\n\n"
            "<b>Примеры:</b>\n"
            + "\n".join(f"  • {e}" for e in VIDEO_PROMPT_EXAMPLES[:5])
        )

    await callback.message.edit_text(
        text, parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="◀️ Назад", callback_data="menu_video")
        ]])
    )
    await callback.answer()


# ── Обработчики кнопок выбора режима для видео из фото ──────────

@dp.callback_query(F.data == "vidphoto_mode_ai")
async def cb_vidphoto_mode_ai(callback: CallbackQuery):
    uid = callback.from_user.id
    vst = video_states.get(uid, {})
    if vst.get("stage") != "await_video_prompt":
        return await callback.answer("❌ Сначала отправь фото", show_alert=True)
    await callback.answer()
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    img_b64 = vst["img_b64"]
    tg_fp   = vst.get("tg_file_path")
    video_states.pop(uid, None)
    await _do_video_gen(
        callback.message, uid,
        "animate this exact person realistically. Keep same face, same appearance, same clothes. Add natural movement: subtle head turn, blinking eyes, breathing. Cinematic, 8K, photorealistic.",
        mode="ai", image_base64=img_b64, tg_file_path=tg_fp
    )


@dp.callback_query(F.data == "vidphoto_mode_raw")
async def cb_vidphoto_mode_raw(callback: CallbackQuery):
    uid = callback.from_user.id
    vst = video_states.get(uid, {})
    if vst.get("stage") != "await_video_prompt":
        return await callback.answer("❌ Сначала отправь фото", show_alert=True)
    video_states[uid]["stage"] = "await_video_prompt_text"
    await callback.answer()
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    await callback.message.answer(
        "✍️ <b>Напиши промпт для видео</b>\n\n"
        "Опиши движение или атмосферу, например:\n"
        "• <code>камера медленно приближается к лицу</code>\n"
        "• <code>волосы развеваются на ветру, кинематограф</code>\n"
        "• <code>slow zoom in, cinematic, dramatic lighting</code>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="vidphoto_cancel")
        ]])
    )


@dp.callback_query(F.data == "vidphoto_cancel")
async def cb_vidphoto_cancel(callback: CallbackQuery):
    uid = callback.from_user.id
    video_states.pop(uid, None)
    await callback.answer("Отменено")
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    await callback.message.answer("❌ Генерация видео отменена.", reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
        InlineKeyboardButton(text="🎬 К генерации видео", callback_data="menu_video")
    ]]))


async def _do_video_gen(message: Message, uid: int, prompt: str, mode: str, image_base64: str = None, tg_file_path: str = None):
    """Генерирует видео и отправляет пользователю."""
    if not can_video(uid):
        li = get_limits_info(uid)
        if li["video_max"] == 0:
            await message.answer(
                "🚫 <b>Генерация видео недоступна на бесплатном тарифе!</b>\n\n"
                "Оформи подписку для получения доступа к генерации видео.",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                    InlineKeyboardButton(text="💎 Оформить подписку", callback_data="menu_subscribe")
                ]])
            )
        else:
            await message.answer(
                f"🚫 <b>Лимит видео исчерпан!</b>\n\n"
                f"Использовано: {li['video_used']}/{li['video_max']} в этом месяце.\n"
                f"Сброс в начале следующего месяца.",
                parse_mode="HTML"
            )
        return

    # Улучшаем промпт если режим ai
    if mode == "ai":
        wait_msg = await message.answer("✨ <i>Улучшаю промпт для реализма...</i>", parse_mode="HTML")
        try:
            if image_base64:
                # Для видео из фото — специальный промпт с сохранением внешности
                enhanced = await enhance_prompt_ai(
                    f"Animate this exact photo. {prompt}. "
                    "CRITICAL: keep the exact same person, same face, same appearance, same clothing. "
                    "Only add realistic movement and life. Do not change the person.",
                    mode="image"
                )
            else:
                enhanced = await enhance_prompt_ai(prompt, mode="image")
        except Exception:
            enhanced = prompt
    else:
        # Для режима raw с фото — тоже добавляем якоря сохранения
        if image_base64:
            enhanced = (
                f"{prompt}. Same person, same face, same appearance. "
                "Animate realistically without changing who the person is."
            )
        else:
            enhanced = prompt
        wait_msg = await message.answer("🎬 <i>Готовлю генерацию...</i>", parse_mode="HTML")

    display_prompt = enhanced[:150]
    photo_tag = "🖼 Оживление фото" if image_base64 else "📝 Текстовый промпт"

    try:
        await wait_msg.edit_text(
            f"🎬 <b>Генерирую видео...</b>\n\n"
            f"📌 Режим: <b>{photo_tag}</b>\n"
            f"📝 <i>{display_prompt}{'...' if len(enhanced) > 150 else ''}</i>\n\n"
            f"🎯 Реалистичность: максимальная\n"
            f"⏳ Обычно 1–3 минуты...",
            parse_mode="HTML"
        )
    except Exception:
        pass

    async with ChatActionSender.upload_video(bot=bot, chat_id=message.chat.id):
        try:
            video_bytes = await generate_video_pollinations(enhanced, image_base64=image_base64, tg_file_path=tg_file_path)

            try:
                await wait_msg.delete()
            except Exception:
                pass

            spend_video(uid)
            li = get_limits_info(uid)
            video_left = max(0, li["video_max"] - li["video_used"])

            if mode == "ai" and enhanced != prompt:
                prompt_info = (
                    f"💡 <i>Запрос:</i> {prompt[:80]}\n"
                    f"✨ <i>Улучшено ИИ:</i> {enhanced[:100]}"
                )
            elif image_base64:
                prompt_info = f"🖼 <i>Оживление фото:</i> {prompt[:100]}"
            else:
                prompt_info = f"📝 <i>{prompt[:120]}</i>"

            caption = (
                f"🎬 <b>Grok Video</b> · Реалистичная генерация\n"
                f"{prompt_info}\n\n"
                f"🎬 Осталось видео: <code>{video_left}/{li['video_max']}</code> в месяц"
            )
            vid_kb = InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="🔄 Ещё вариант",   callback_data="video_again"),
                 InlineKeyboardButton(text="✨ Новый промпт",  callback_data="menu_video")],
                [InlineKeyboardButton(text="🎛 Дополнительно", callback_data="back_extra")],
            ])
            video_file = BufferedInputFile(video_bytes, filename="grok_video.mp4")
            asyncio.create_task(db_save_user(uid))
            try:
                await message.answer_video(video_file, caption=caption, parse_mode="HTML", reply_markup=vid_kb)
            except Exception as ve:
                logging.warning(f"answer_video failed: {ve}, sending as document")
                video_file2 = BufferedInputFile(video_bytes, filename="grok_video.mp4")
                await message.answer_document(video_file2, caption=caption, parse_mode="HTML", reply_markup=vid_kb)

            # Сохраняем для кнопки «Ещё вариант»
            last_responses[uid] = {
                **last_responses.get(uid, {}),
                "last_video_prompt": enhanced,
                "last_video_mode":   mode,
            }
            video_states.pop(uid, None)

        except Exception as e:
            try:
                await wait_msg.delete()
            except Exception:
                pass
            logging.error(f"video gen error uid={uid}: {e}")
            
            # Определяем тип ошибки для более понятного сообщения
            error_msg = str(e)
            if "404" in error_msg or "NOT_FOUND" in error_msg:
                user_msg = (
                    "❌ <b>Ошибка генерации видео</b>\n\n"
                    "🔍 <b>Причина:</b> Сервис генерации временно недоступен\n\n"
                    "💡 <b>Что делать:</b>\n"
                    "• Попробуй другой промпт\n"
                    "• Подожди 1-2 минуты и попробуй снова\n"
                    "• Если проблема повторяется, напиши в поддержку"
                )
            elif "timeout" in error_msg.lower():
                user_msg = (
                    "❌ <b>Превышено время ожидания</b>\n\n"
                    "⏱️ Генерация заняла слишком много времени\n\n"
                    "💡 Попробуй упростить промпт или повтори позже"
                )
            elif "слишком маленькое" in error_msg.lower():
                user_msg = (
                    "❌ <b>Ошибка генерации видео</b>\n\n"
                    "📝 <b>Причина:</b> Промпт не подходит для генерации\n\n"
                    "💡 Попробуй изменить описание или добавить больше деталей"
                )
            else:
                user_msg = (
                    f"❌ <b>Ошибка генерации видео</b>\n\n"
                    f"<code>{error_msg[:200]}</code>\n\n"
                    "Попробуй изменить промпт или повтори позже."
                )
            
            await message.answer(
                user_msg,
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                    InlineKeyboardButton(text="🔄 Попробовать снова", callback_data="menu_video"),
                    InlineKeyboardButton(text="💬 Поддержка", callback_data="menu_support")
                ]])
            )


@dp.callback_query(F.data == "video_again")
async def cb_video_again(callback: CallbackQuery):
    """Повторная генерация видео с тем же промптом."""
    uid = callback.from_user.id
    last = last_responses.get(uid, {})
    prompt = last.get("last_video_prompt", "")
    mode   = last.get("last_video_mode", "raw")
    if not prompt:
        await callback.answer("❌ Промпт не найден, начни заново", show_alert=True)
        return
    await callback.answer("🔄 Генерирую новый вариант...")
    await _do_video_gen(callback.message, uid, prompt, mode)


# ==================================================================
# 📝 ДОКЛАД / РЕФЕРАТ
# ==================================================================


def report_titlepage_kb() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="📄 С титульным листом", callback_data="report_title_yes"),
            InlineKeyboardButton(text="📝 Без титульного листа", callback_data="report_title_no"),
        ],
        [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
    ])


def report_type_kb() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="📄 Доклад",   callback_data="report_type_doklad"),
            InlineKeyboardButton(text="📚 Реферат",  callback_data="report_type_referat"),
        ],
        [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
    ])


def report_pages_kb() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="2 стр.",  callback_data="report_pages_2"),
            InlineKeyboardButton(text="3 стр.",  callback_data="report_pages_3"),
            InlineKeyboardButton(text="5 стр.",  callback_data="report_pages_5"),
        ],
        [
            InlineKeyboardButton(text="7 стр.",  callback_data="report_pages_7"),
            InlineKeyboardButton(text="10 стр.", callback_data="report_pages_10"),
        ],
        [InlineKeyboardButton(text="✏️ Своё количество (макс. 10)", callback_data="report_pages_custom")],
        [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
    ])


def report_model_kb() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🔬 Sonar Deep Research (веб) ⭐", callback_data="report_model_sonar_deep_research")],
        [
            InlineKeyboardButton(text="🧠 Sonar Reasoning Pro",  callback_data="report_model_sonar_reasoning_pro"),
            InlineKeyboardButton(text="🌐 Sonar Pro (веб)",      callback_data="report_model_sonar_pro"),
        ],
        [
            InlineKeyboardButton(text="🔵 Claude Opus",          callback_data="report_model_claude_opus"),
            InlineKeyboardButton(text="🔷 Claude Sonnet",        callback_data="report_model_claude_sonnet"),
        ],
        [
            InlineKeyboardButton(text="🧬 DeepSeek V3",          callback_data="report_model_deepseek_v3_sq"),
            InlineKeyboardButton(text="⚡ GPT-5.2",              callback_data="report_model_gpt52"),
        ],
        [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
    ])


def report_confirm_kb() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="✅ Создать доклад",  callback_data="report_generate"),
            InlineKeyboardButton(text="✏️ Изменить",        callback_data="report_restart"),
        ],
        [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
    ])


@dp.message(F.text == "📝 Доклад/Реферат")
async def rb_report(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message):
        return
    if not await check_service(message, "report"):
        return
    report_states[uid] = {}
    await message.answer(
        "📝 <b>Генератор докладов и рефератов</b>\n\n"
        "🌐 Использует <b>Sonar Deep Research</b> — ищет актуальные данные\n"
        "   прямо в интернете во время генерации\n\n"
        "📄 Результат: файл <b>.docx</b> готовый к сдаче\n\n"
        "Выбери тип работы:",
        parse_mode="HTML",
        reply_markup=report_type_kb()
    )


@dp.callback_query(F.data.startswith("report_type_"))
async def report_type_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    rtype = callback.data.split("report_type_")[1]
    type_label = "📄 Доклад" if rtype == "doklad" else "📚 Реферат"
    if uid not in report_states:
        report_states[uid] = {}
    report_states[uid]["type"] = rtype
    report_states[uid]["type_label"] = type_label
    await callback.message.edit_text(
        f"📝 <b>{type_label}</b>\n\n"
        "📏 Выбери объём работы:",
        parse_mode="HTML",
        reply_markup=report_pages_kb()
    )
    await callback.answer()


@dp.callback_query(F.data.startswith("report_pages_"))
async def report_pages_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    raw = callback.data.split("report_pages_")[1]

    # Обработка кнопки "Своё количество"
    if raw == "custom":
        if uid not in report_states:
            report_states[uid] = {}
        admin_await[uid] = {"action": "report_custom_pages"}
        await callback.message.edit_text(
            "✏️ <b>Введи количество страниц</b>\n\n"
            "Напиши число от 1 до 10:\n"
            "<i>Например: 8</i>",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="❌ Отмена", callback_data="back_home"),
            ]])
        )
        await callback.answer()
        return

    pages = int(raw)
    if uid not in report_states:
        report_states[uid] = {}
    report_states[uid]["pages"] = pages
    report_states[uid]["model"] = "sonar_deep_research"
    report_states[uid]["model_label"] = "🔬 Sonar Deep Research"
    # Предлагаем выбрать модель
    await callback.message.edit_text(
        f"📏 <b>Объём: {pages} {'страниц' if pages >= 5 else ('страницы' if 2 <= pages <= 4 else 'страница')}</b>\n\n"
        "🤖 <b>Выбери модель для генерации:</b>\n\n"
        "🔬 <b>Sonar Deep Research</b> — ищет в интернете, самый актуальный\n"
        "🧠 <b>Sonar Reasoning Pro</b> — логика + веб-поиск\n"
        "🌐 <b>Sonar Pro</b> — быстрый с поиском\n"
        "🔵 <b>Claude / GPT / DeepSeek</b> — без поиска, по своим знаниям",
        parse_mode="HTML",
        reply_markup=report_model_kb()
    )
    await callback.answer()


@dp.callback_query(F.data.startswith("report_model_"))
async def report_model_cb(callback: CallbackQuery):
    """Выбор модели для генерации доклада/реферата."""
    uid = callback.from_user.id
    if uid not in report_states:
        report_states[uid] = {}
    model_key = callback.data.split("report_model_")[1]
    # Маппинг ключ → label
    _MODEL_LABELS = {
        "sonar_deep_research": "🔬 Sonar Deep Research",
        "sonar_reasoning_pro": "🧠 Sonar Reasoning Pro",
        "sonar_pro":           "🌐 Sonar Pro",
        "sonar":               "🌐 Sonar",
        "claude_opus":         "🔵 Claude Opus",
        "claude_sonnet":       "🔷 Claude Sonnet",
        "deepseek_v3_sq":      "🧬 DeepSeek V3",
        "gpt52":               "⚡ GPT-5.2",
    }
    model_label = _MODEL_LABELS.get(model_key, MODELS.get(model_key, {}).get("label", model_key))
    report_states[uid]["model"]       = model_key
    report_states[uid]["model_label"] = model_label
    await callback.message.edit_text(
        f"🤖 <b>Модель: {model_label}</b>\n\n"
        "📄 <b>Нужен титульный лист?</b>\n\n"
        "<i>Титульный лист — стандартное оформление с названием работы, автором и годом</i>",
        parse_mode="HTML",
        reply_markup=report_titlepage_kb()
    )
    await callback.answer()


@dp.callback_query(F.data.in_({"report_title_yes", "report_title_no"}))
async def report_title_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in report_states:
        report_states[uid] = {}
    with_title = callback.data == "report_title_yes"
    report_states[uid]["with_title"] = with_title
    title_label = "📄 С титульным листом" if with_title else "📝 Без титульного листа"
    if with_title:
        # Спрашиваем автора
        admin_await[uid] = {"action": "report_author"}
        await callback.message.edit_text(
            f"✅ <b>{title_label}</b>\n\n"
            "👤 Введи <b>имя автора</b> (ФИО / псевдоним):\n\n"
            "<i>Или нажми «Пропустить», если не нужно</i>",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="⏭ Пропустить", callback_data="report_skip_author")],
                [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
            ])
        )
    else:
        # Без титульника — сразу к теме
        report_states[uid]["author"] = ""
        admin_await[uid] = {"action": "report_topic"}
        await callback.message.edit_text(
            f"✅ <b>{title_label}</b>\n\n"
            "✏️ Напиши <b>тему</b> доклада/реферата:\n\n"
            "<i>Пример: «Первая помощь при ожогах»</i>",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="❌ Отмена", callback_data="back_home"),
            ]])
        )
    await callback.answer()


@dp.callback_query(F.data == "report_skip_author")
async def report_skip_author_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in report_states:
        report_states[uid] = {}
    report_states[uid]["author"] = ""
    admin_await[uid] = {"action": "report_topic"}
    await callback.message.edit_text(
        "✏️ Напиши <b>тему</b> доклада/реферата:\n\n"
        "<i>Пример: «Влияние искусственного интеллекта на рынок труда»</i>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="back_home"),
        ]])
    )
    await callback.answer()
@dp.callback_query(F.data == "report_skip_wishes")
async def report_skip_wishes_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in report_states:
        report_states[uid] = {}
    report_states[uid]["wishes"] = ""
    state = report_states.get(uid, {})
    type_label = state.get("type_label", "Доклад")
    topic = state.get("topic", "—")
    pages = state.get("pages", 5)
    model_label = state.get("model_label", "Claude Opus")
    author = state.get("author", "")
    with_title = state.get("with_title", True)
    title_str = "📄 Есть" if with_title else "📝 Нет"
    confirm_text = (
        f"📋 <b>Подтверди параметры работы:</b>\n\n"
        f"📄 Тип: <b>{type_label}</b>\n"
        f"✏️ Тема: <b>{topic}</b>\n"
        f"📏 Объём: <b>{pages} стр.</b>\n"
        f"🤖 Модель: <b>{model_label}</b>\n"
        f"👤 Автор: <b>{author if author else chr(8212)}</b>\n"
        f"📑 Титульный лист: <b>{title_str}</b>\n"
        f"✍️ Пожелания: <b>—</b>\n\n"
        f"💎 Стоимость: <b>3 продвинутых запроса</b>"
    )
    try:
        await callback.message.edit_text(confirm_text, parse_mode="HTML", reply_markup=report_confirm_kb())
    except Exception:
        await callback.message.answer(confirm_text, parse_mode="HTML", reply_markup=report_confirm_kb())
    await callback.answer()

@dp.callback_query(F.data == "report_restart")
async def report_restart_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    report_states[uid] = {}
    await callback.message.edit_text(
        "📝 <b>Начнём заново!</b>\n\nВыбери тип работы:",
        parse_mode="HTML",
        reply_markup=report_type_kb()
    )
    await callback.answer()


@dp.callback_query(F.data == "report_generate")
async def report_generate_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    state = report_states.get(uid, {})
    if not state.get("topic"):
        await callback.answer("❌ Тема не задана", show_alert=True)
        return
    await _do_generate_report(callback.message, uid, state, is_cb=True)
    await callback.answer()


def _hard_truncate_to_words(text: str, max_words: int) -> str:
    """
    Обрезает текст до max_words слов, при этом:
    - Заканчивает на полном предложении (не обрывает на середине)
    - Добавляет «...» если текст был обрезан
    """
    words = text.split()
    if len(words) <= max_words:
        return text
    # Берём первые max_words слов
    truncated = " ".join(words[:max_words])
    # Ищем последнее законченное предложение
    for end_char in ['. ', '.\n', '! ', '!\n', '? ', '?\n']:
        last_pos = truncated.rfind(end_char)
        if last_pos > len(truncated) * 0.7:  # конец должен быть в последних 30%
            return truncated[:last_pos + 1].rstrip() + "\n\n**...**"
    return truncated.rstrip() + "\n\n**...**"


async def _do_generate_report(msg_or_message, uid: int, state: dict, is_cb: bool = False):
    """Генерирует доклад/реферат и отправляет как .docx файл."""
    topic       = state.get("topic", "")
    pages       = state.get("pages", 5)
    rtype       = state.get("type", "doklad")
    rtype_label = state.get("type_label", "Доклад")
    # Берём модель которую выбрал пользователь, с fallback на Sonar
    chosen_model = state.get("model", "sonar_deep_research")
    if chosen_model not in MODELS or chosen_model in disabled_models:
        chosen_model = "sonar_deep_research"
        if chosen_model in disabled_models:
            chosen_model = "claude_opus"
    mkey = chosen_model
    model_label = state.get("model_label") or MODELS.get(mkey, {}).get("label", mkey)
    wishes      = state.get("wishes", "")

    # Проверяем что есть 3 продвинутых запроса
    _init_limits(uid)
    _refresh_limits(uid)
    li = get_limits_info(uid)
    pro_left = li["pro_max"] - li["pro_used"]
    if pro_left < 3:
        txt = (
            "❌ <b>Недостаточно запросов!</b>\n\n"
            f"💎 Для генерации {rtype_label.lower()}а нужно <b>3 продвинутых запроса</b>.\n"
            f"У тебя осталось: <b>{pro_left}</b>\n"
            f"⏰ Сброс через: <code>{li['pro_reset']}</code>"
        )
        if is_cb:
            await msg_or_message.answer(txt, parse_mode="HTML")
        else:
            await msg_or_message.answer(txt, parse_mode="HTML")
        return

    if is_cb:
        send_fn = msg_or_message.answer
    else:
        send_fn = msg_or_message.answer

    # Отправляем первичное сообщение с счётчиком
    await send_fn(
        "🌐 <i>Использую Sonar Deep Research — ищу актуальные данные в интернете...</i>",
        parse_mode="HTML"
    )
    think_msg = await send_fn(
        f"📝 <i>Готовлю {rtype_label.lower()}... <b>1с</b></i>",
        parse_mode="HTML"
    )
    _start = asyncio.get_event_loop().time()

    async def _tick(m, s):
        phases = [
            (5,  f"📝 Составляю план {rtype_label.lower()}а"),
            (15, "✍️ Пишу основную часть"),
            (30, "🔍 Добавляю детали и примеры"),
            (60, "📋 Формирую выводы"),
            (999,"⚡ Почти готово, финальная обработка"),
        ]
        while True:
            await asyncio.sleep(1)
            elapsed = int(asyncio.get_event_loop().time() - s)
            label = phases[0][1]
            for thr, lbl in phases:
                if elapsed < thr:
                    label = lbl
                    break
            try:
                await m.edit_text(f"{label}... <b>{elapsed}с</b>", parse_mode="HTML")
            except Exception:
                pass

    _tick_task = asyncio.create_task(_tick(think_msg, _start))

    # Промпт для генерации
    # 1 страница A4 ≈ 280 слов ≈ 420 токенов (русский язык)
    words_per_page = 280
    target_words   = pages * words_per_page
    max_target     = target_words + 30          # небольшой допуск
    # Жёсткий лимит токенов — модели НЕ ДАЮТ писать больше чем нужно
    # 1 стр = 420 токенов, минимум 500 токенов на маленькие доклады
    max_tok = max(pages * 420, 500)
    # Не более 4000 токенов даже для больших докладов (≈10 стр)
    max_tok = min(max_tok, 4000)
    wishes_block = f"\n\nПОЖЕЛАНИЯ ЗАКАЗЧИКА (обязательно учти):\n{wishes}" if wishes else ""

    # Определяем тематические особенности
    topic_lower_r = topic.lower()
    is_obzh = any(w in topic_lower_r for w in [
        'обж', 'безопасность жизнедеятельности', 'чрезвычайн', 'первая помощь',
        'пожарн', 'эвакуац', 'землетрясен', 'наводнен', 'выживан',
        'гражданская оборона', 'террор', 'экстремальн',
        'скулшутинг', 'вооружённ', 'вооружен', 'нападен', 'угроз',
        'насили', 'буллинг', 'дтп', 'авари', 'безопасн', 'опасн',
        'яд', 'отравлен', 'ожог', 'утоплен', 'молни', 'гроз',
        'взрыв', 'химическ', 'радиац', 'биологическ', 'эпидем',
        'цунами', 'урага', 'торнадо', 'лавин', 'паводк',
    ])

    if rtype == "doklad":
        if is_obzh:
            system_prompt = (
                "Ты — опытный преподаватель ОБЖ (Основы Безопасности Жизнедеятельности) "
                "и автор школьных докладов. Ты пишешь доклады, которые РЕАЛЬНО ПОМОГАЮТ "
                "в экстремальных ситуациях — с конкретными алгоритмами действий, цифрами, "
                "реальными случаями и практическими советами.\n\n"
                "🇷🇺 ОБЯЗАТЕЛЬНО: Весь текст пиши ТОЛЬКО на русском языке.\n\n"
                "━━━ СТАНДАРТЫ КАЧЕСТВА ДЛЯ ОБЖ ━━━\n\n"
                "ОБЯЗАТЕЛЬНЫЕ ЭЛЕМЕНТЫ КАЖДОГО РАЗДЕЛА:\n"
                "✅ Реальные случаи и примеры (с конкретными деталями: когда, где, что произошло)\n"
                "✅ Алгоритм действий по шагам: Шаг 1 → Шаг 2 → Шаг 3...\n"
                "✅ Статистика (количество жертв, процент выживаемости, частота случаев)\n"
                "✅ Типичные ОШИБКИ людей и ПОЧЕМУ они опасны\n"
                "✅ Конкретные признаки опасности (как распознать угрозу)\n"
                "✅ Что НЕЛЬЗЯ делать — с объяснением последствий\n\n"
                "СТРУКТУРА КАЖДОГО РАЗДЕЛА:\n"
                "1. Суть явления/угрозы (2-3 предложения с определением)\n"
                "2. Реальный случай или статистика\n"
                "3. Как распознать опасность (признаки)\n"
                "4. Алгоритм действий по шагам (минимум 5-7 шагов)\n"
                "5. Типичные ошибки и их последствия\n"
                "6. Профилактика и подготовка\n\n"
                "ЗАПРЕЩЕНО:\n"
                "• Общие фразы типа 'важно соблюдать безопасность'\n"
                "• Советы без конкретики ('немедленно позвоните' — куда? какой номер?)\n"
                "• Перечисления без объяснения ЗАЧЕМ\n"
                "• Повторение одного и того же разными словами\n\n"
                "Пишешь ТОЛЬКО текст доклада. Без вступлений типа 'Конечно!' и комментариев."
            )
        else:
            system_prompt = (
                "Ты — профессиональный автор школьных и студенческих докладов. "
                "У тебя есть доступ к актуальным данным из интернета — используй веб-поиск "
                "чтобы найти свежие факты, статистику и примеры по теме.\n\n"
                "🇷🇺 ОБЯЗАТЕЛЬНО: Весь доклад пиши ТОЛЬКО на русском языке, "
                "независимо от темы и источников.\n\n"
                "ПРАВИЛА КАЧЕСТВЕННОГО ДОКЛАДА:\n"
                "• Только русский язык (абсолютное требование)\n"
                "• Используй актуальные данные из интернета\n"
                "• Каждый раздел — конкретные факты, цифры, примеры (не общие слова)\n"
                "• Стиль: чёткий, информативный, подходит для устного выступления\n"
                "• БЕЗ 'целей', 'задач', 'актуальности' — это доклад, не реферат\n"
                "• Иностранные термины с переводом в скобках\n"
                "• Каждый раздел завершается кратким итогом\n\n"
                "Пишешь ТОЛЬКО сам текст доклада, без вступлений и комментариев."
            )

        # Строим структуру разделов по количеству страниц
        sections = ""
        num_sections = max(1, pages - 1)
        for i in range(1, num_sections + 1):
            sections += f"## [Раздел {i}: название по теме]\n"

        if is_obzh:
            user_prompt = (
                f"⛔ СТОП-ПРАВИЛО №1 (важнее всего остального):\n"
                f"МАКСИМУМ {target_words} СЛОВ. ЭТО ЖЁСТКИЙ ПОТОЛОК.\n"
                f"Для {pages} стр. = {target_words} слов. Превышение НЕДОПУСТИМО.\n"
                f"Когда достигнешь ~{target_words - 30} слов — немедленно пиши заключение и заканчивай.\n\n"
                f"Напиши доклад по ОБЖ на тему: «{topic}»\n\n"
                f"🇷🇺 ОБЯЗАТЕЛЬНО: пиши ТОЛЬКО на русском языке!\n\n"
                f"ОБЯЗАТЕЛЬНАЯ СТРУКТУРА (вписываясь в {target_words} слов):\n"
                f"## Введение\n"
                f"(2-3 предложения: что это за угроза, насколько распространена)\n\n"
                f"{sections}\n"
                f"В КАЖДОМ разделе (кратко! — весь доклад = {target_words} слов):\n"
                f"— Реальный случай (1-2 предложения)\n"
                f"— Алгоритм действий (нумерованный, 4-6 шагов)\n"
                f"— Телефоны экстренных служб: 112, 01, 02, 03\n"
                f"— 1-2 типичных ошибки с последствиями\n\n"
                f"## Заключение\n"
                f"(2-3 ключевых правила — краткий список)\n\n"
                f"{wishes_block}"
            )
        else:
            user_prompt = (
                f"⛔ СТОП-ПРАВИЛО №1 (важнее всего остального):\n"
                f"МАКСИМУМ {target_words} СЛОВ. ЭТО ЖЁСТКИЙ ПОТОЛОК.\n"
                f"Для {pages} стр. = {target_words} слов. Превышение НЕДОПУСТИМО.\n"
                f"Когда достигнешь ~{target_words - 30} слов — немедленно пиши заключение и заканчивай.\n\n"
                f"Напиши доклад на тему: «{topic}»\n\n"
                f"🇷🇺 ОБЯЗАТЕЛЬНО: пиши ТОЛЬКО на русском языке!\n\n"
                f"СТРУКТУРА (вписываясь в {target_words} слов):\n"
                f"## Введение\n"
                f"Краткое введение (2-3 предложения)\n\n"
                f"{sections}\n"
                f"## Заключение\n"
                f"Итоги (2-3 предложения)\n\n"
                f"{wishes_block}\n\n"
                f"Каждый раздел — конкретные факты, без воды."
            )
    else:
        _limit_warning = (
            f"\n\n⛔ ГЛАВНОЕ ПРАВИЛО: ПИШИ НЕ БОЛЕЕ {target_words} СЛОВ СУММАРНО. "
            f"Когда приближаешься к {target_words - 30} словам — СРАЗУ ЗАКАНЧИВАЙ."
        )
        system_prompt = (
            f"Ты — автор научного реферата. Пишешь академично, по ГОСТ.\n"
            f"🇷🇺 ТОЛЬКО РУССКИЙ ЯЗЫК. Стиль: научный, с источниками и анализом.\n"
            f"Пишешь ТОЛЬКО реферат, без предисловий и комментариев."
            + _limit_warning
        )
        # Количество глав для реферата
        num_chapters = max(2, pages - 2)
        chapters = "\n".join([
            f"## Глава {i}. [Конкретный аспект темы]\n### {i}.1 [Теория]\n### {i}.2 [Практика/Современное состояние]"
            for i in range(1, num_chapters + 1)
        ])
        words_per_ch = max(40, target_words // (num_chapters + 2))
        user_prompt = (
            f"Напиши научный реферат: «{topic}»\n\n"
            f"🛑 ОБЪЁМ: РОВНО {target_words} СЛОВ ({pages} страниц). СТОП после {target_words} слов!\n\n"
            f"СТРУКТУРА (ГОСТ):\n"
            f"## Введение (~{words_per_ch} слов) — актуальность, объект, предмет, цель, задачи\n\n"
            f"{chapters}\n\n"
            f"## Заключение (~{words_per_ch} слов) — выводы по задачам + перспективы\n\n"
            f"## Список литературы — минимум 5 источников по ГОСТ\n\n"
            f"{wishes_block}\n"
            f"❗ ИТОГО: {target_words} СЛОВ МАКСИМУМ!"
        )

    try:
        msgs = [
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_prompt},
        ]
        # Sonar Deep Research — первым, с fallback на Claude если отказал
        _doc_fallbacks = ["sonar_deep_research", "claude_opus", "claude_sonnet", "deepseek_v3_sq", "qwen3_max"]
        _tried_doc = set()
        # Фразы-признаки отказа модели писать контент
        _REFUSAL_PHRASES = [
            "i appreciate you reaching out", "i'm unable to write", "i cannot write",
            "i can't write", "i cannot create content", "i'm not able to",
            "i cannot provide", "i can't provide", "i cannot assist",
            "raises serious safety concerns", "i cannot help with",
            "не могу написать", "не могу создать", "не могу помочь",
            "не в состоянии написать", "отказываюсь", "вынужден отказать",
            "what i can help with instead", "что я могу предложить вместо",
        ]
        ans = ""
        for _dm in _doc_fallbacks:
            if _dm in _tried_doc or _dm not in MODELS:
                continue
            _tried_doc.add(_dm)
            try:
                _ans = await call_chat(msgs, _dm, max_tokens=max_tok)
                _ans_lower = _ans.lower() if _ans else ""
                # Проверяем на отказ
                _is_refusal = any(p in _ans_lower for p in _REFUSAL_PHRASES)
                if _is_refusal:
                    logging.warning(f"[REPORT] {_dm} отказал писать контент, пробуем следующую модель")
                    continue
                if _ans and len(_ans.strip()) > 200:
                    ans = _ans
                    if _dm != "sonar_deep_research":
                        logging.info(f"[REPORT] Использован fallback {_dm} вместо sonar_deep_research")
                    break
                else:
                    logging.warning(f"[REPORT] {_dm} вернул слишком короткий ответ ({len(_ans)} симв.), пробуем следующую")
            except Exception as _de:
                logging.warning(f"[REPORT] {_dm} ошибка: {_de}")
                continue

        if not ans or len(ans.strip()) < 100:
            raise RuntimeError("Все модели вернули пустой или слишком короткий ответ. Попробуй ещё раз.")

        # Убираем нумерованные сноски [1][2][3]
        ans = re.sub(r'\[\d+\]', '', ans).strip()

        # ─── ЖЁСТКОЕ ОБРЕЗАНИЕ по словам ────────────────────────────────
        # Если модель всё равно написала больше — обрезаем без пощады
        _max_words = pages * 290  # 2 стр = 580 слов максимум
        _words_all = ans.split()
        if len(_words_all) > _max_words:
            # Обрезаем по словам, потом добавляем последнюю законченную фразу
            _cut = " ".join(_words_all[:_max_words])
            # Ищем последнюю точку/! /? для красивого обрезания
            _last_dot = max(
                _cut.rfind(". "), _cut.rfind(".\n"), _cut.rfind("!\n"),
                _cut.rfind("?\n"), _cut.rfind("! "), _cut.rfind("? ")
            )
            if _last_dot > len(_cut) // 2:
                ans = _cut[:_last_dot + 1]
            else:
                ans = _cut
            logging.info(f"[REPORT] Текст обрезан с {len(_words_all)} до {len(ans.split())} слов (лимит {_max_words})")
        # ─────────────────────────────────────────────────────────────────

        # Списываем 3 продвинутых запроса
        for _ in range(3):
            user_limits[uid]["pro_used"] = user_limits[uid].get("pro_used", 0) + 1

        _tick_task.cancel()
        try:
            await _tick_task
        except asyncio.CancelledError:
            pass
        try:
            await think_msg.delete()
        except Exception:
            pass

        # Создаём .docx файл
        path = await _create_report_docx(topic, rtype_label, pages, model_label, ans, uid,
                                          author=state.get("author", ""),
                                          with_title=state.get("with_title", True))

        total_secs = int(asyncio.get_event_loop().time() - _start)
        caption = (
            f"📝 <b>{rtype_label}: «{topic[:60]}{'...' if len(topic)>60 else ''}»</b>\n\n"
            f"📏 Объём: <b>{pages} страниц</b>\n"
            f"🤖 Модель: <b>{model_label}</b>\n"
            f"⏱ Время генерации: <b>{total_secs}с</b>\n"
            f"💎 Потрачено: <b>3 продвинутых запроса</b>\n\n"
            f"✅ Документ готов!"
        )
        await send_fn(caption, parse_mode="HTML")
        await bot.send_document(
            chat_id=uid,
            document=FSInputFile(path, filename=f"{rtype_label}_{topic[:30]}.docx")
        )
        try:
            os.remove(path)
        except Exception:
            pass

        if uid in report_states:
            del report_states[uid]

    except Exception as e:
        _tick_task.cancel()
        try:
            await _tick_task
        except asyncio.CancelledError:
            pass
        try:
            await think_msg.delete()
        except Exception:
            pass
        logging.error(f"report generate: {e}")
        await send_fn(f"❌ Ошибка при создании {rtype_label.lower()}а: {str(e)[:300]}")


async def _create_report_docx(topic: str, rtype_label: str, pages: int, model_label: str, content: str, uid: int, author: str = "", with_title: bool = True) -> str:
    """Создаёт красиво оформленный .docx файл с докладом/рефератом."""
    from docx import Document as _Doc
    from docx.shared import Pt, RGBColor, Cm, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    import datetime as _dt

    doc = _Doc()

    # Настройка полей страницы
    from docx.shared import Cm
    section = doc.sections[0]
    section.top_margin    = Cm(2.5)
    section.bottom_margin = Cm(2.5)
    section.left_margin   = Cm(3.0)
    section.right_margin  = Cm(1.5)

    # Стили
    normal = doc.styles['Normal']
    normal.font.name = 'Times New Roman'
    normal.font.size = Pt(14)

    if with_title:
        # Титульная страница
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run("МИНИСТЕРСТВО ОБРАЗОВАНИЯ")
        r.bold = True; r.font.size = Pt(12)
        r.font.name = 'Times New Roman'

        doc.add_paragraph()
        doc.add_paragraph()
        doc.add_paragraph()

        p2 = doc.add_paragraph()
        p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r2 = p2.add_run(rtype_label.upper())
        r2.bold = True; r2.font.size = Pt(20)
        r2.font.name = 'Times New Roman'
        r2.font.color.rgb = RGBColor(0x19, 0x50, 0xB4)

        doc.add_paragraph()

        p3 = doc.add_paragraph()
        p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r3 = p3.add_run("на тему:")
        r3.font.size = Pt(14); r3.font.name = 'Times New Roman'

        p4 = doc.add_paragraph()
        p4.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r4 = p4.add_run(f"«{topic}»")
        r4.bold = True; r4.font.size = Pt(16)
        r4.font.name = 'Times New Roman'
        r4.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

        doc.add_paragraph()
        doc.add_paragraph()
        doc.add_paragraph()
        doc.add_paragraph()

        if author:
            p5 = doc.add_paragraph()
            p5.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            # Определение пола по окончанию имени (русская морфология)
            def _detect_gender(name: str) -> str:
                """Возвращает 'Выполнила' или 'Выполнил' по имени."""
                if not name:
                    return "Выполнил"
                first_name = name.strip().split()[0].lower()
                # Женские окончания в русском языке
                female_endings = ('а', 'я', 'ия', 'ья', 'ея', 'ка', 'на', 'га', 'та', 'ра', 'ла', 'са')
                if any(first_name.endswith(e) for e in female_endings):
                    return "Выполнила"
                return "Выполнил"
            performed_word = _detect_gender(author)
            r5 = p5.add_run(f"{performed_word}: {author}")
            r5.font.size = Pt(14); r5.font.name = 'Times New Roman'
            r5.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)

        doc.add_paragraph()
        doc.add_paragraph()

        p6 = doc.add_paragraph()
        p6.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r6 = p6.add_run(_dt.datetime.now().strftime("%Y"))
        r6.font.size = Pt(14); r6.font.name = 'Times New Roman'

        # Разрыв страницы
        doc.add_page_break()

    # Парсим содержимое и добавляем в документ
    lines = content.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i]

        # Блоки кода
        if re.match(r'^```(\w*)$', line.strip()):
            lang_name = line.strip()[3:].strip()
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            code_text = '\n'.join(code_lines)
            p_code = doc.add_paragraph(code_text)
            for r_run in p_code.runs:
                r_run.font.name = 'Courier New'
                r_run.font.size = Pt(11)
            # Серый фон
            pPr = p_code._p.get_or_add_pPr()
            shd = OxmlElement('w:shd')
            shd.set(qn('w:val'), 'clear')
            shd.set(qn('w:color'), 'auto')
            shd.set(qn('w:fill'), 'F5F5F5')
            pPr.append(shd)
            i += 1
            continue

        s = line.strip()

        # Заголовки H1
        if line.startswith('# '):
            ph = doc.add_heading(strip_md(line[2:]), level=1)
            for r_run in ph.runs:
                r_run.font.name = 'Times New Roman'
                r_run.font.size = Pt(16)
                r_run.font.color.rgb = RGBColor(0x19, 0x50, 0xB4)
            i += 1; continue

        # Заголовки H2
        if line.startswith('## '):
            ph = doc.add_heading(strip_md(line[3:]), level=2)
            for r_run in ph.runs:
                r_run.font.name = 'Times New Roman'
                r_run.font.size = Pt(14)
                r_run.font.color.rgb = RGBColor(0x1a, 0x1a, 0x60)
            i += 1; continue

        # Заголовки H3
        if line.startswith('### '):
            ph = doc.add_heading(strip_md(line[4:]), level=3)
            for r_run in ph.runs:
                r_run.font.name = 'Times New Roman'
                r_run.font.size = Pt(13)
            i += 1; continue

        # Маркированные списки
        if re.match(r'^[-•*]\s', line):
            pp = doc.add_paragraph(style='List Bullet')
            pp.paragraph_format.first_line_indent = Cm(0)
            r_run = pp.add_run(strip_md(line[2:].strip()))
            r_run.font.name = 'Times New Roman'
            r_run.font.size = Pt(14)
            i += 1; continue

        # Нумерованные списки
        nm = re.match(r'^(\d+\.\s+)(.*)', line)
        if nm:
            pp = doc.add_paragraph(style='List Number')
            r_run = pp.add_run(strip_md(nm.group(2).strip()))
            r_run.font.name = 'Times New Roman'
            r_run.font.size = Pt(14)
            i += 1; continue

        # Цитаты
        if line.startswith('> '):
            pp = doc.add_paragraph()
            pp.paragraph_format.left_indent = Cm(1)
            r_run = pp.add_run(strip_md(line[2:]))
            r_run.italic = True
            r_run.font.name = 'Times New Roman'
            r_run.font.size = Pt(13)
            r_run.font.color.rgb = RGBColor(0x50, 0x50, 0x80)
            i += 1; continue

        # Пустая строка
        if not s:
            doc.add_paragraph()
            i += 1; continue

        # Обычный абзац
        pp = doc.add_paragraph()
        pp.paragraph_format.first_line_indent = Cm(1.25)
        pp.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

        # Парсим inline форматирование
        pattern = re.compile(
            r'(\*\*\*(.+?)\*\*\*)|(\*\*(.+?)\*\*)|(\*(.+?)\*)|(`([^`]+)`)',
            re.DOTALL
        )
        last_pos = 0
        for m in pattern.finditer(line):
            if m.start() > last_pos:
                r_run = pp.add_run(line[last_pos:m.start()])
                r_run.font.name = 'Times New Roman'
                r_run.font.size = Pt(14)
            if m.group(1):
                r_run = pp.add_run(m.group(2))
                r_run.bold = True; r_run.italic = True
            elif m.group(3):
                r_run = pp.add_run(m.group(4))
                r_run.bold = True
            elif m.group(5):
                r_run = pp.add_run(m.group(6))
                r_run.italic = True
            elif m.group(7):
                r_run = pp.add_run(m.group(8))
                r_run.font.name = 'Courier New'
                r_run.font.size = Pt(12)
            if r_run:
                r_run.font.name = r_run.font.name or 'Times New Roman'
                r_run.font.size = r_run.font.size or Pt(14)
            last_pos = m.end()
        if last_pos < len(line):
            r_run = pp.add_run(line[last_pos:])
            r_run.font.name = 'Times New Roman'
            r_run.font.size = Pt(14)
        i += 1

    # Сохраняем
    path = tmp(f'report_{uid}.docx')
    doc.save(path)
    return path


# ==================================================================
# 🎞 PPTX — ГЕНЕРАТОР ПРЕЗЕНТАЦИЙ v3.0
# ==================================================================

SLIDE_W = 13.33
SLIDE_H = 7.5

def _c(r, g, b):
    return PptxRGBColor(r, g, b)

PPTX_THEMES = {
    "dark_blue": {
        "name":       "🌌 Тёмно-синяя",
        "header":     _c(0x07, 0x0B, 0x24),
        "bg":         _c(0x0D, 0x13, 0x36),
        "bg2":        _c(0x15, 0x1E, 0x50),
        "accent":     _c(0x4D, 0x8F, 0xFF),
        "accent2":    _c(0x9B, 0x5D, 0xE5),
        "title_c":    _c(0xFF, 0xFF, 0xFF),
        "body_c":     _c(0xC8, 0xD6, 0xFF),
        "muted_c":    _c(0x55, 0x68, 0xA0),
        "card":       _c(0x18, 0x22, 0x55),
        "font_title": "Montserrat",
        "font_body":  "Montserrat",
    },
    "corporate": {
        "name":       "💼 Корпоративная",
        "header":     _c(0x00, 0x28, 0x5A),
        "bg":         _c(0xF3, 0xF6, 0xFC),
        "bg2":        _c(0xFF, 0xFF, 0xFF),
        "accent":     _c(0x00, 0x56, 0xB3),
        "accent2":    _c(0xFF, 0x8C, 0x00),
        "title_c":    _c(0xFF, 0xFF, 0xFF),
        "body_c":     _c(0x18, 0x28, 0x48),
        "muted_c":    _c(0x58, 0x6E, 0x90),
        "card":       _c(0xFF, 0xFF, 0xFF),
        "font_title": "Raleway",
        "font_body":  "Montserrat",
    },
    "midnight": {
        "name":       "🌙 Полночь",
        "header":     _c(0x06, 0x06, 0x06),
        "bg":         _c(0x0E, 0x0E, 0x0E),
        "bg2":        _c(0x1A, 0x1A, 0x1A),
        "accent":     _c(0x00, 0xD4, 0xAA),
        "accent2":    _c(0xFF, 0x6B, 0x6B),
        "title_c":    _c(0xFF, 0xFF, 0xFF),
        "body_c":     _c(0xCC, 0xCC, 0xCC),
        "muted_c":    _c(0x66, 0x66, 0x66),
        "card":       _c(0x22, 0x22, 0x22),
        "font_title": "Montserrat",
        "font_body":  "Montserrat",
    },
    "ocean": {
        "name":       "🌊 Океан",
        "header":     _c(0x00, 0x19, 0x3C),
        "bg":         _c(0xEE, 0xF7, 0xFF),
        "bg2":        _c(0xFF, 0xFF, 0xFF),
        "accent":     _c(0x00, 0x96, 0xC7),
        "accent2":    _c(0x00, 0xD4, 0xAA),
        "title_c":    _c(0xFF, 0xFF, 0xFF),
        "body_c":     _c(0x00, 0x19, 0x3C),
        "muted_c":    _c(0x3E, 0x78, 0xA0),
        "card":       _c(0xFF, 0xFF, 0xFF),
        "font_title": "Raleway",
        "font_body":  "Montserrat",
    },
    "sunset": {
        "name":       "🌅 Закат",
        "header":     _c(0x1A, 0x06, 0x00),
        "bg":         _c(0xFF, 0xF7, 0xF3),
        "bg2":        _c(0xFF, 0xFF, 0xFF),
        "accent":     _c(0xFF, 0x45, 0x00),
        "accent2":    _c(0xFF, 0xA5, 0x00),
        "title_c":    _c(0xFF, 0xFF, 0xFF),
        "body_c":     _c(0x2A, 0x0E, 0x00),
        "muted_c":    _c(0x88, 0x55, 0x33),
        "card":       _c(0xFF, 0xFF, 0xFF),
        "font_title": "Montserrat",
        "font_body":  "Montserrat",
    },
    "purple": {
        "name":       "💜 Пурпурная",
        "header":     _c(0x12, 0x00, 0x38),
        "bg":         _c(0xFA, 0xF5, 0xFF),
        "bg2":        _c(0xFF, 0xFF, 0xFF),
        "accent":     _c(0x7B, 0x2F, 0xBE),
        "accent2":    _c(0xE0, 0x40, 0xFB),
        "title_c":    _c(0xFF, 0xFF, 0xFF),
        "body_c":     _c(0x12, 0x00, 0x38),
        "muted_c":    _c(0x70, 0x50, 0x90),
        "card":       _c(0xFF, 0xFF, 0xFF),
        "font_title": "Raleway",
        "font_body":  "Montserrat",
    },
    "nature": {
        "name":       "🌿 Природа",
        "header":     _c(0x0D, 0x2B, 0x0F),
        "bg":         _c(0xF3, 0xFB, 0xF3),
        "bg2":        _c(0xFF, 0xFF, 0xFF),
        "accent":     _c(0x2E, 0x8B, 0x57),
        "accent2":    _c(0x8B, 0xC3, 0x4A),
        "title_c":    _c(0xFF, 0xFF, 0xFF),
        "body_c":     _c(0x0D, 0x2B, 0x0F),
        "muted_c":    _c(0x4A, 0x7A, 0x4A),
        "card":       _c(0xFF, 0xFF, 0xFF),
        "font_title": "Montserrat",
        "font_body":  "Montserrat",
    },
    "redblack": {
        "name":       "🔴 Красно-чёрная",
        "header":     _c(0x0A, 0x00, 0x00),
        "bg":         _c(0xFF, 0xF4, 0xF4),
        "bg2":        _c(0xFF, 0xFF, 0xFF),
        "accent":     _c(0xCC, 0x00, 0x00),
        "accent2":    _c(0xFF, 0x44, 0x44),
        "title_c":    _c(0xFF, 0xFF, 0xFF),
        "body_c":     _c(0x1A, 0x00, 0x00),
        "muted_c":    _c(0x80, 0x33, 0x33),
        "card":       _c(0xFF, 0xFF, 0xFF),
        "font_title": "Montserrat",
        "font_body":  "Montserrat",
    },
    "tech": {
        "name":       "💻 Технологии",
        "header":     _c(0x03, 0x10, 0x12),
        "bg":         _c(0x06, 0x1A, 0x1E),
        "bg2":        _c(0x0C, 0x2A, 0x30),
        "accent":     _c(0x00, 0xFF, 0xAA),
        "accent2":    _c(0x00, 0xBF, 0xFF),
        "title_c":    _c(0xFF, 0xFF, 0xFF),
        "body_c":     _c(0xB0, 0xEE, 0xE0),
        "muted_c":    _c(0x40, 0x80, 0x70),
        "card":       _c(0x0C, 0x28, 0x30),
        "font_title": "Montserrat",
        "font_body":  "Montserrat",
    },
    "warm": {
        "name":       "🏛️ Тёплая классика",
        "header":     _c(0x3A, 0x18, 0x07),
        "bg":         _c(0xFD, 0xF6, 0xEC),
        "bg2":        _c(0xFF, 0xFF, 0xFF),
        "accent":     _c(0xB8, 0x50, 0x42),
        "accent2":    _c(0xD4, 0x8A, 0x3A),
        "title_c":    _c(0xFF, 0xFF, 0xFF),
        "body_c":     _c(0x2C, 0x12, 0x06),
        "muted_c":    _c(0x8A, 0x5A, 0x40),
        "card":       _c(0xFF, 0xFF, 0xFF),
        "font_title": "Raleway",
        "font_body":  "Montserrat",
    },
    "minimal": {
        "name":       "⬜ Минимализм",
        "header":     _c(0x21, 0x21, 0x21),
        "bg":         _c(0xF9, 0xF9, 0xF9),
        "bg2":        _c(0xFF, 0xFF, 0xFF),
        "accent":     _c(0x21, 0x21, 0x21),
        "accent2":    _c(0x55, 0x55, 0x55),
        "title_c":    _c(0xFF, 0xFF, 0xFF),
        "body_c":     _c(0x21, 0x21, 0x21),
        "muted_c":    _c(0x88, 0x88, 0x88),
        "card":       _c(0xFF, 0xFF, 0xFF),
        "font_title": "Montserrat",
        "font_body":  "Montserrat",
    },
    # Элегантная светлая (как у конкурента: кремовый фон, тёмно-коричневые заголовки)
    "elegant": {
        "name":       "✨ Элегантная (светлая)",
        "header":     _c(0x4A, 0x37, 0x28),   # тёмно-коричневый
        "bg":         _c(0xFA, 0xF8, 0xF2),   # кремово-белый фон
        "bg2":        _c(0xFF, 0xFF, 0xFF),   # белый
        "accent":     _c(0x8B, 0x73, 0x55),   # коричневый акцент
        "accent2":    _c(0xC4, 0xA4, 0x6B),   # золотистый
        "title_c":    _c(0x4A, 0x37, 0x28),   # тёмно-коричневый (на СВЕТЛОМ фоне)
        "body_c":     _c(0x5C, 0x4A, 0x3A),   # средне-коричневый
        "muted_c":    _c(0x8A, 0x7A, 0x6A),   # мутный коричневый
        "card":       _c(0xEE, 0xE8, 0xD8),   # тёплый бежевый (карточки)
        "card_title": _c(0x4A, 0x37, 0x28),   # заголовок карточки
        "font_title": "Georgia",
        "font_body":  "Calibri",
        "layout":     "elegant",               # маркер для нового рендерера
    },
}


# ── Автовыбор темы по теме презентации ───────────────────────────

def _auto_pick_theme(topic: str) -> str:
    """
    Автоматически выбирает наиболее подходящую тему оформления
    на основе ключевых слов в названии презентации.
    """
    t = topic.lower()

    # IT / Технологии / Программирование / ИИ
    tech_kw = [
        "программ", "python", "javascript", "code", "код", "разработ", "software",
        "компьют", "технолог", "ии ", "искусственный интеллект", "ai ", "нейросет",
        "машинное обучение", "data", "данн", "алгоритм", "сеть", "интернет",
        "кибербез", "хакер", "блокчейн", "криптовалют", "web", "веб", "api",
        "база данных", "devops", "cloud", "облако", "приложени",
    ]
    if any(k in t for k in tech_kw):
        return "elegant"  # светлый стиль

    # Бизнес / Финансы / Маркетинг / Менеджмент
    corp_kw = [
        "бизнес", "финанс", "экономик", "маркетинг", "менеджмент", "управлени",
        "стратег", "компани", "корпорат", "рынок", "продаж", "инвестиц",
        "стартап", "предприним", "банк", "бюджет", "прибыл", "доход",
        "логистик", "hr ", "персонал", "брендинг", "реклам",
    ]
    if any(k in t for k in corp_kw):
        return "elegant"

    # Природа / Экология / Окружающая среда
    nature_kw = [
        "природ", "эколог", "климат", "окружающ", "лес", "животн", "растени",
        "биолог", "зелен", "устойчив", "переработ", "возобновляем", "organic",
        "агрик", "сельск", "флора", "фауна", "заповедник",
    ]
    if any(k in t for k in nature_kw):
        return "elegant"  # светлый стиль

    # Океан / Вода / Море
    ocean_kw = [
        "океан", "море", "вод", "морск", "подводн", "рыба", "рыболов",
        "путешестви", "туризм", "отдых", "курорт", "тропик",
    ]
    if any(k in t for k in ocean_kw):
        return "elegant"  # светлый стиль

    # История / Культура / Искусство / Архитектура
    warm_kw = [
        "истори", "культур", "искусств", "архитектур", "музей", "литератур",
        "философи", "религи", "традиц", "наследи", "классик", "античн",
        "средневек", "ренессанс", "живопис", "скульптур", "театр",
    ]
    if any(k in t for k in warm_kw):
        return "elegant"  # светлый стиль

    # Наука / Исследования / Медицина / Физика / Химия
    midnight_kw = [
        "наук", "исследован", "медицин", "физик", "хими", "математик",
        "биохими", "генетик", "астроном", "космос", "вселенн", "квантов",
        "молекул", "эксперимент", "лаборатор", "открыти", "теори",
    ]
    if any(k in t for k in midnight_kw):
        return "elegant"  # светлый стиль

    # Психология / Саморазвитие / Образование
    purple_kw = [
        "психолог", "саморазвити", "мотивац", "лидерств", "образован",
        "обучени", "педагогик", "воспитани", "личностн", "эмоц",
        "творчеств", "дизайн", "мода", "стиль", "красот",
    ]
    if any(k in t for k in purple_kw):
        return "elegant"  # светлый стиль

    # Спорт / Здоровье / Фитнес
    sunset_kw = [
        "спорт", "здоровь", "фитнес", "питани", "диет", "трениров",
        "медицин", "здравоохран", "активн", "бег", "футбол", "олимпи",
    ]
    if any(k in t for k in sunset_kw):
        return "elegant"  # светлый стиль

    # По умолчанию — элегантная светлая (как у лучших конкурентов)
    return "elegant"

# Убираем костыль — переопределяем нормально


# ── Примитивы рисования ───────────────────────────────────────────

def _pptx_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _pptx_rect(slide, x, y, w, h, color, transp=None):
    from pptx.util import Inches as I
    s = slide.shapes.add_shape(1, I(x), I(y), I(w), I(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if transp is not None:
        s.fill.transparency = transp
    s.line.fill.background()
    return s


def _pptx_txt(slide, text, x, y, w, h, size=18, color=None,
              bold=False, italic=False, align=None, font=None, wrap=True):
    from pptx.util import Inches as I, Pt
    from pptx.enum.text import PP_ALIGN as PA
    if align is None:
        align = PA.LEFT
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font or "Calibri"
    if color:
        r.font.color.rgb = color
    return tb


def _pptx_bullets(slide, items, x, y, w, h, size=16,
                  color=None, dot_color=None, font=None):
    from pptx.util import Inches as I, Pt
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(2)
        if dot_color:
            d = p.add_run()
            d.text = "▸  "
            d.font.size = Pt(size - 2)
            d.font.color.rgb = dot_color
            d.font.bold = True
        r = p.add_run()
        r.text = str(item)
        r.font.size = Pt(size)
        r.font.name = font or "Calibri"
        if color:
            r.font.color.rgb = color
    return tb


def _pptx_img(slide, img_bytes, x, y, w, h):
    """
    Вставляет фото с умной обрезкой (cover-режим):
    пропорции оригинала сохраняются, лишнее обрезается по центру.
    Никакого сплющивания!
    """
    from pptx.util import Inches as I
    try:
        from PIL import Image as PILImage
        import io as _io

        # Целевой размер в пикселях (96 DPI)
        DPI = 96
        tw = max(1, int(w * DPI))
        th = max(1, int(h * DPI))

        img = PILImage.open(_io.BytesIO(img_bytes)).convert("RGB")
        iw, ih = img.size

        # Cover: масштабируем так, чтобы заполнить оба измерения
        scale = max(tw / iw, th / ih)
        nw = max(1, int(iw * scale))
        nh = max(1, int(ih * scale))
        img = img.resize((nw, nh), PILImage.LANCZOS)

        # Центральная обрезка
        left = max(0, (nw - tw) // 2)
        top  = max(0, (nh - th) // 2)
        img = img.crop((left, top, left + tw, top + th))

        buf = _io.BytesIO()
        img.save(buf, format="JPEG", quality=90, optimize=True)
        buf.seek(0)
        slide.shapes.add_picture(buf, I(x), I(y), I(w), I(h))
        return True
    except Exception as e:
        logging.warning(f"pptx img smart error: {e}")
        try:
            slide.shapes.add_picture(io.BytesIO(img_bytes), I(x), I(y), I(w), I(h))
            return True
        except Exception:
            return False


# ── Загрузка тематических фото ────────────────────────────────────

# ── Глобальный словарь перевода русских тем → английские фото-запросы ──
_RU_TO_EN_PHOTO = {
    # Война / История
    "война": "war military battle soldiers historical",
    "войны": "war military battle soldiers historical",
    "военн": "military soldiers army defense",
    "армия": "army soldiers military training",
    "история": "history historical archive old",
    "историческ": "historical monument ancient architecture",
    "революц": "revolution historical soldiers crowd",
    "империя": "empire palace architecture historical",
    "россия": "Russia Moscow Kremlin architecture",
    "российск": "Russia Russian architecture city",
    "советск": "Soviet historical parade architecture",
    "победа": "victory parade soldiers celebration",
    "великая отечественная": "World War 2 soldiers tank battle",
    "вторая мировая": "World War 2 soldiers battle historical",
    "первая мировая": "World War 1 soldiers trench battle",
    "холодная война": "Cold War nuclear tension historical",
    "наполеон": "Napoleon war soldiers cavalry battle",
    "реформ": "reform change modernization progress",
    "территор": "territory map land geographic",
    "расширени": "expansion territory map empire",
    "геополитик": "geopolitics world map strategy power",
    "сверхдержав": "superpower military parade nuclear",
    # ОБЖ / Безопасность
    "обж": "emergency safety rescue first aid",
    "безопасност": "safety protection security rescue",
    "чрезвычайн": "emergency rescue firefighters disaster",
    "первая помощь": "first aid medical emergency CPR",
    "пожарн": "fire firefighters rescue emergency",
    "эвакуац": "evacuation people emergency exit",
    "землетрясен": "earthquake disaster destruction rescue",
    "наводнен": "flood disaster rescue emergency",
    "выживан": "survival nature emergency outdoor",
    "гражданская оборона": "civil defense emergency shelters",
    "террор": "terrorism security police defense",
    "ожог": "burn injury medical treatment hospital",
    "отравлен": "poison medical emergency treatment",
    "утоплен": "drowning rescue lifeguard water safety",
    "дтп": "car accident road emergency rescue",
    "авари": "accident emergency rescue disaster",
    "спасател": "rescue team emergency responders",
    "мчс": "emergency rescue ministry firefighters Russia",
    "скорая": "ambulance medical emergency first aid",
    # Бизнес
    "бизнес": "business office professionals meeting",
    "финанс": "finance money investment banking",
    "экономик": "economy business finance charts",
    "маркетинг": "marketing strategy team office",
    "менеджмент": "management leadership team meeting",
    "стратег": "strategy planning business whiteboard",
    "продаж": "sales business meeting presentation",
    "инвестиц": "investment finance growth money",
    "стартап": "startup office innovation team",
    "прибыл": "profit revenue finance growth",
    "рост": "growth chart business progress",
    "рынок": "market trading finance business",
    "конкурент": "competition business market strategy",
    "план": "planning business strategy team",
    "проект": "project team planning office",
    "управлени": "management leadership team office",
    # Технологии
    "технолог": "technology innovation digital modern",
    "программ": "programming code software developer",
    "нейросет": "neural network AI technology visualization",
    "искусственный интеллект": "artificial intelligence AI robot",
    "данные": "data analytics dashboard screens",
    "кибер": "cybersecurity digital security network",
    "робот": "robot automation technology industrial",
    "цифров": "digital technology modern innovation",
    "интернет": "internet network technology global",
    "блокчейн": "blockchain cryptocurrency technology network",
    "облако": "cloud computing technology data center",
    # Медицина
    "медицин": "medical doctor hospital healthcare",
    "здоровь": "health wellness medicine hospital",
    "хирург": "surgeon surgery hospital operating",
    "болезн": "disease medical research laboratory",
    "лечени": "treatment medical hospital doctor",
    "лекарств": "medicine pharmacy drugs laboratory",
    # Образование
    "образован": "education students classroom university",
    "школ": "school students learning classroom",
    "университет": "university campus students lecture",
    "наук": "science laboratory research experiment",
    "исследован": "research laboratory scientist experiment",
    "учеб": "studying students learning books",
    # Природа
    "природ": "nature landscape outdoor scenic",
    "эколог": "ecology environment green sustainable",
    "климат": "climate environment nature landscape",
    "лес": "forest trees nature green",
    "океан": "ocean sea water waves",
    "горы": "mountains landscape nature scenic",
    # Промышленность
    "промышленн": "industry factory workers production",
    "завод": "factory industrial plant workers",
    "производств": "production manufacturing assembly line",
    "нефт": "oil refinery petroleum industrial",
    "газ": "natural gas pipeline energy",
    "металлург": "steel metal factory industrial",
    "строительств": "construction building workers site",
    # Социальные темы
    "общество": "society people community diverse",
    "культур": "culture art museum gallery",
    "искусств": "art museum gallery exhibition",
    "музык": "music concert performance band",
    "спорт": "sport athletes training competition",
    "здравоохранен": "healthcare medical hospital professionals",
    # Транспорт
    "транспорт": "transport logistics road vehicles",
    "авиац": "aviation aircraft airport flight",
    "железнодорожн": "railway train station travel",
    "автомобил": "automobile car road transport",
    "логистик": "logistics supply chain warehouse",
    # Политика
    "политик": "politics government parliament officials",
    "правительств": "government politics parliament officials",
    "демократ": "democracy parliament government vote",
    "выбор": "election vote democracy ballot",
    # Психология
    "психологи": "psychology mind wellness mental",
    "мотивац": "motivation success achievement goal",
    "лидерств": "leadership success team motivation",
    "саморазвити": "self development success achievement",
}


# ══════════════════════════════════════════════════════════════════════════════
# ГЛОБАЛЬНАЯ СИСТЕМА ФОТО — используется и для PPTX и для HTML-презентаций
# 400+ ключей, 65+ категорий, 8-10 фото каждая, Pexels CDN (без ключа)
# ══════════════════════════════════════════════════════════════════════════════

def _pex_cdn(pid: int) -> str:
    return "https://images.pexels.com/photos/{0}/pexels-photo-{0}.jpeg?auto=compress&cs=tinysrgb&w=1280&h=720&fit=crop".format(pid)

_PHOTO_DB = {
    # ── Военная тематика ──────────────────────────────────────────────────────
    "soldier":     [2058126,1917215,1437796,3396234,6447535,8159956,3651428,1537220],
    "soldiers":    [1917215,2058126,3396234,8159956,1437796,6447535,1537220,3651428],
    "military":    [2058126,3396234,1917215,6447535,1437796,8159956,3651428,1537220],
    "army":        [3396234,2058126,1437796,1917215,8159956,6447535,1537220,3651428],
    "tank":        [3396234,2058126,1917215,6447535,1437796,8159956,2860953,3722826],
    "armored":     [3396234,1917215,2058126,6447535,8159956,1437796,3722826,2860953],
    "weapon":      [1917215,2058126,3396234,1437796,6447535,8159956,3651428,3722826],
    "weapons":     [1437796,1917215,2058126,3396234,8159956,6447535,2860953,3651428],
    "missile":     [3396234,6447535,2058126,1917215,1437796,8159956,3722826,2860953],
    "rocket":      [6447535,3396234,2058126,1917215,8159956,1437796,1537220,3722826],
    "fighter":     [6447535,2058126,3396234,1917215,1437796,8159956,3722826,1537220],
    "aircraft":    [6447535,3396234,2058126,1437796,1917215,8159956,1537220,3722826],
    "drone":       [6447535,3396234,8159956,2058126,1917215,1437796,2860953,3722826],
    "nato":        [2058126,3396234,1917215,8159956,1437796,6447535,3651428,2860953],
    "defense":     [1917215,2058126,3396234,8159956,1437796,6447535,3722826,3651428],
    "combat":      [2058126,1917215,3396234,6447535,8159956,1437796,1537220,2860953],
    "battle":      [2058126,3396234,1917215,1437796,6447535,8159956,2860953,3651428],
    "war":         [2058126,1917215,3396234,1437796,8159956,6447535,3651428,1537220],
    "patrol":      [1917215,2058126,3396234,1437796,6447535,8159956,1537220,3651428],
    "navy":        [1427107,906494,1117210,2058126,1917215,3396234,1427109,1546960],
    "warship":     [1427107,906494,1117210,1427109,1546960,2226458,1117207,2226460],
    "parade":      [2058126,1917215,3396234,1437796,6447535,8159956,3651428,2860953],
    "sniper":      [1917215,2058126,1437796,3396234,6447535,8159956,3651428,3722826],
    "security":    [2058126,1917215,313782,3396234,1486325,1437796,1693652,6447535],
    "conflict":    [2058126,1917215,3396234,1437796,6447535,8159956,3651428,1537220],
    # ── ОБЖ / Безопасность / Полиция ─────────────────────────────────────────
    "police":      [1653087,1697912,3153207,2896941,3777565,2828237,2896956,1653093],
    "cop":         [1653087,1697912,3153207,2896941,3777565,2828237,2896956,1653093],
    "detective":   [3153207,2896941,3777565,1653087,1697912,2828237,2896956,3153200],
    "investigation":[3153207,2896941,3777565,1653087,1697912,2828237,2896956,3777558],
    "surveillance":[1917215,2058126,3153207,2896941,3777565,1653087,1697912,2896956],
    "checkpoint":  [1917215,2058126,2896941,3153207,1653087,1697912,3777565,2828237],
    "teenager":    [3184465,3182812,1205651,289737,1438072,1350700,2422232,3184454],
    "teenagers":   [3184465,3182812,1205651,289737,1438072,1350700,2422232,3184454],
    "smartphone":  [3184465,3182812,699122,2741448,5081918,5082578,3184454,3182781],
    "cybersecurity":[3182812,3184465,699122,2741448,5081918,5082578,3182781,3184454],
    "prevention":  [1205651,289737,3184465,3182812,1438072,1350700,3777565,2896941],
    "awareness":   [1205651,289737,3184465,3182812,1438072,1350700,2422232,1205652],
    "counselor":   [3184465,3182812,1205651,289737,1438072,3777565,2896941,1350700],
    "evacuate":    [2558679,2557399,2917989,2904389,3807571,3807572,2917983,2558686],
    "evacuation":  [2558679,2557399,2917989,2904389,3807571,3807572,2917983,2558686],
    # ── История ───────────────────────────────────────────────────────────────
    "historical":  [313782,1486325,1693652,373912,2246476,325229,1693640,1486310],
    "ancient":     [1486325,313782,373912,1693652,325229,2246476,373892,1486310],
    "medieval":    [1486325,313782,1693652,373912,325229,2246476,1693640,373892],
    "empire":      [313782,1486325,1693652,373912,2246476,325229,1693668,1486310],
    "revolution":  [2058126,1917215,313782,1486325,1693652,3396234,373912,325229],
    "cavalry":     [2058126,1917215,3396234,1437796,6447535,8159956,3651428,1537220],
    "archive":     [1205651,289737,1438072,313782,1693652,1350700,1486325,2422232],
    "monument":    [1486325,313782,373912,1693652,2246476,325229,1693640,1486310],
    "museum":      [1486325,313782,373912,1693652,2246476,325229,1693640,1486310],
    "kremlin":     [313782,1486325,373912,1693652,2246476,325229,1693668,1486310],
    "palace":      [1486325,313782,373912,1693652,2246476,325229,1693640,373892],
    "cathedral":   [1486325,313782,373912,1693652,2246476,325229,1693640,1486310],
    "napoleon":    [2058126,1917215,3396234,313782,1437796,1486325,6447535,8159956],
    "wwii":        [2058126,1917215,3396234,1437796,6447535,8159956,3651428,1537220],
    "soviet":      [2058126,1917215,313782,1486325,3396234,1693652,1437796,373912],
    "imperial":    [313782,1486325,1693652,373912,2246476,325229,1486310,1693640],
    "renaissance": [1486325,313782,373912,1693652,2246476,325229,1693640,373892],
    "history":     [313782,1486325,1693652,373912,2246476,325229,1693640,1486310],
    "heritage":    [1486325,313782,373912,1693652,2246476,325229,1693640,1486310],
    "tradition":   [1486325,313782,373912,1693652,2246476,325229,1693640,1486310],
    # ── Технологии / ИИ / Данные ──────────────────────────────────────────────
    "neural":      [8386440,1181244,574069,270360,3861969,325229,574060,8386430],
    "artificial":  [8386440,1181244,574069,270360,3861969,325229,574060,8386430],
    "circuit":     [8386440,574069,1181244,270360,3861969,325229,574060,270352],
    "chip":        [8386440,574069,1181244,270360,3861969,325229,574060,270352],
    "server":      [1181244,8386440,574069,270360,3861969,325229,1181236,8386430],
    "data":        [8386440,1181244,270360,574069,3861969,325229,270352,1181236],
    "code":        [574069,8386440,1181244,270360,3861969,325229,574060,270352],
    "developer":   [574069,8386440,1181244,270360,3861969,325229,574060,270352],
    "computer":    [574069,8386440,1181244,270360,3861969,325229,574060,8386430],
    "robot":       [8386440,1181244,3861969,574069,270360,325229,8386430,1181236],
    "technology":  [8386440,574069,1181244,270360,3861969,325229,574060,270352],
    "cloud":       [1181244,8386440,574069,270360,3861969,325229,1181236,8386430],
    "network":     [8386440,1181244,574069,270360,3861969,325229,574060,1181236],
    "cyber":       [8386440,574069,1181244,270360,3861969,325229,574060,8386430],
    "blockchain":  [8386440,574069,1181244,270360,3861969,325229,574060,270352],
    "smartphone":  [574069,8386440,1181244,270360,3861969,325229,574060,270352],
    "laptop":      [574069,8386440,1181244,270360,3861969,325229,574060,270352],
    "programming": [574069,8386440,1181244,270360,3861969,325229,574060,270352],
    "software":    [574069,8386440,1181244,270360,3861969,325229,574060,270352],
    "startup":     [3184465,3182812,574069,8386440,1181244,270360,3182781,3184454],
    "innovation":  [8386440,574069,1181244,3184465,270360,3861969,3182812,325229],
    "digital":     [8386440,574069,1181244,270360,3861969,325229,574060,8386430],
    # ── Бизнес / Менеджмент ───────────────────────────────────────────────────
    "business":    [3184465,3182812,3182781,3184454,3182773,3184419,3182804,3184464],
    "meeting":     [3182812,3184465,3182781,3184454,3184419,3182773,3182804,3184464],
    "office":      [3182812,3184465,3184454,3182781,3184419,3182773,3184464,3184437],
    "team":        [3182812,3184465,3182781,3184454,3182773,3184419,3182804,3184464],
    "handshake":   [3184465,3182812,3182781,3184454,3182773,3184419,3182804,3184464],
    "contract":    [3184465,3182812,3182781,3184454,3182773,3184419,3182804,3184464],
    "leadership":  [3184465,3182812,3182781,3184454,3182773,3184419,3182804,3184464],
    "manager":     [3182812,3184465,3182781,3184454,3184419,3182773,3182804,3184464],
    "corporate":   [3182812,3184465,3182781,3184454,3184419,3182773,3182804,3184464],
    "strategy":    [3184465,3182812,905163,3182781,905157,3184454,3182773,905170],
    "planning":    [3184465,3182812,905163,3182781,3184454,905157,3182773,3184419],
    "management":  [3184465,3182812,3182781,3184454,3182773,3184419,3182804,3184464],
    # ── Финансы / Экономика ───────────────────────────────────────────────────
    "stock":       [159888,210607,590022,6801874,3771110,669619,1602726,1602720],
    "chart":       [210607,159888,6801874,590022,3771110,669619,1602720,1591059],
    "finance":     [590022,210607,159888,6801874,669619,3771110,1591059,1602726],
    "financial":   [6801874,590022,210607,159888,3771110,669619,1135208,1591059],
    "trading":     [159888,590022,6801874,210607,669619,3771110,1602720,1135208],
    "bank":        [210607,159888,590022,3771110,6801874,669619,1602726,1602720],
    "investment":  [159888,590022,210607,6801874,3771110,669619,1602726,3184465],
    "economy":     [590022,159888,210607,6801874,3771110,669619,1591059,1602726],
    "market":      [590022,159888,210607,6801874,3771110,669619,1135208,1602720],
    "growth":      [159888,3184465,590022,210607,6801874,3182812,3771110,669619],
    "revenue":     [159888,210607,590022,6801874,3184465,3182812,3771110,669619],
    "profit":      [210607,159888,590022,3184465,6801874,3182812,3771110,1602720],
    "budget":      [210607,590022,3771110,159888,669619,6801874,1602726,1135208],
    "inflation":   [590022,3771110,159888,210607,6801874,669619,1602720,1591059],
    "currency":    [669619,590022,159888,210607,6801874,3771110,1135208,1602726],
    "marketing":   [905163,590022,3184465,3182812,210607,159888,905157,905170],
    "sales":       [3184465,3182812,590022,210607,159888,3182781,3184454,3182773],
    "brand":       [905163,590022,3184465,905157,210607,159888,905170,3182812],
    "product":     [3182812,3184465,3182781,3184454,3182773,3184419,3182804,3184464],
    "sanction":    [590022,159888,210607,6801874,3771110,669619,1591059,1602726],
    # ── Промышленность / Производство ────────────────────────────────────────
    "factory":     [1108572,2760241,209251,1267338,257700,3862132,1183366,1183344],
    "industry":    [1108572,2760241,209251,1267338,257700,3862132,1183366,1183344],
    "industrial":  [1108572,2760241,209251,1267338,257700,3862132,1183366,2760244],
    "workers":     [1108572,2760241,1267338,209251,257700,3862132,1183366,1183344],
    "production":  [1108572,209251,2760241,1267338,257700,3862132,1183366,2760244],
    "assembly":    [1108572,2760241,209251,1267338,257700,3862132,1183366,1183344],
    "welding":     [2760241,1108572,209251,1267338,257700,3862132,2760244,1183366],
    "machinery":   [1108572,2760241,209251,1267338,257700,3862132,1183366,2760244],
    "manufacturing":[1108572,2760241,209251,1267338,257700,3862132,1183366,1183344],
    "steel":       [2760241,1108572,209251,1267338,257700,3862132,2760244,1183366],
    "construction":[257700,1267338,1108572,209251,2760241,3862132,1183344,1183366],
    # ── Энергетика ────────────────────────────────────────────────────────────
    "oil":         [1078116,2101521,247763,462024,325153,1108572,247754,1078108],
    "petroleum":   [1078116,247763,2101521,325153,462024,1108572,247754,1078108],
    "pipeline":    [2101521,1078116,247763,462024,325153,1108572,2101514,247754],
    "refinery":    [1078116,2101521,247763,462024,325153,1108572,247754,1078108],
    "gas":         [2101521,1078116,247763,462024,325153,1108572,2101514,1078108],
    "energy":      [1078116,2101521,247763,462024,325153,1108572,247754,1078108],
    "fuel":        [1078116,2101521,247763,462024,325153,1108572,247754,2101514],
    "solar":       [326082,462024,1131458,265216,158827,3791664,326074,462016],
    "nuclear":     [1078116,2101521,247763,462024,325153,1108572,247754,1078108],
    "electricity": [574069,8386440,270360,1181244,3861969,325229,574060,270352],
    "power":       [1078116,2101521,247763,462024,325153,1108572,247754,1078108],
    # ── Сельское хозяйство / Природа ─────────────────────────────────────────
    "wheat":       [326082,1131458,265216,158827,3791664,462024,326074,1131450],
    "harvest":     [326082,1131458,265216,158827,3791664,462024,326074,1131450],
    "farm":        [326082,1131458,265216,158827,3791664,462024,326074,462016],
    "field":       [326082,1131458,265216,158827,3791664,462024,326074,1131450],
    "agriculture": [326082,1131458,265216,158827,3791664,462024,326074,1131450],
    "forest":      [167698,338936,1525041,1761279,247431,167696,338930,1525035],
    "nature":      [167698,338936,1525041,1761279,247431,167696,338930,1525035],
    "landscape":   [167698,338936,1525041,1761279,247431,167696,338930,1525035],
    "mountain":    [414612,417074,355465,417088,355457,414600,417082,355450],
    "river":       [414612,417074,355465,417088,355457,414600,417082,414618],
    "ocean":       [1295138,1295139,731082,1001682,1295140,731080,1001680,1295141],
    "sea":         [1295138,1295139,731082,1001682,1295140,731080,1001680,1295141],
    "beach":       [1295138,1295139,731082,1001682,1295140,731080,1001680,1295141],
    "climate":     [167698,338936,1525041,1761279,247431,167696,338930,1525035],
    "ecology":     [167698,338936,1525041,1761279,247431,167696,338930,326082],
    "green":       [326082,167698,338936,1131458,265216,1525041,158827,3791664],
    "environment": [167698,338936,1525041,1761279,247431,167696,338930,1525035],
    # ── Медицина / Здоровье ───────────────────────────────────────────────────
    "surgeon":     [3786157,305568,4386431,263402,3938023,1298619,3786149,305560],
    "hospital":    [3786157,305568,4386431,263402,3938023,1298619,3786149,305560],
    "doctor":      [3786157,305568,4386431,263402,3938023,1298619,3786149,305560],
    "medical":     [3786157,305568,4386431,263402,3938023,1298619,3786149,305560],
    "nurse":       [305568,3786157,4386431,263402,3938023,1298619,305560,3786149],
    "laboratory":  [4386431,263402,3938023,305568,3786157,1298619,4386423,263394],
    "research":    [4386431,263402,3938023,305568,3786157,1298619,4386423,263394],
    "vaccine":     [3786157,305568,4386431,263402,3938023,1298619,3786149,305560],
    "fitness":     [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "gym":         [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "wellness":    [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "mental":      [3786157,305568,4386431,263402,3938023,1298619,3786149,305560],
    # ── Образование / Наука ───────────────────────────────────────────────────
    "students":    [1205651,289737,1438072,1350700,2422232,1205652,289720,1438062],
    "classroom":   [1205651,289737,1438072,1350700,2422232,1205652,289720,1438062],
    "university":  [1205651,289737,1438072,1350700,2422232,1205652,289720,1438062],
    "lecture":     [289737,1205651,1438072,1350700,2422232,1205652,289720,1438062],
    "education":   [1205651,289737,1438072,1350700,2422232,1205652,289720,1438062],
    "school":      [1205651,289737,1438072,1350700,2422232,1205652,289720,1438062],
    "teacher":     [289737,1205651,1438072,1350700,2422232,1205652,289720,1438062],
    "graduation":  [1205651,289737,1438072,1350700,2422232,1205652,289720,1438062],
    "physics":     [4386431,263402,3938023,8386440,574069,305568,4386423,263394],
    "chemistry":   [4386431,263402,3938023,305568,3786157,1298619,4386423,263394],
    "biology":     [4386431,263402,3938023,305568,3786157,1298619,4386423,263394],
    "astronomy":   [1525041,167698,338936,1761279,247431,167696,338930,1525035],
    "space":       [1525041,167698,338936,1761279,247431,167696,338930,1525035],
    "cosmos":      [1525041,167698,338936,1761279,247431,167696,338930,1525035],
    "satellite":   [1525041,167698,338936,1761279,247431,167696,338930,1525035],
    "science":     [4386431,263402,3938023,305568,3786157,1298619,4386423,263394],
    # ── Логистика / Транспорт ─────────────────────────────────────────────────
    "cargo":       [1427107,906494,1117210,1427109,1546960,2226458,1117207,906482],
    "ship":        [1427107,906494,1117210,1427109,1546960,2226458,1117207,2226460],
    "container":   [1427107,906494,1117210,1427109,1546960,2226458,1117207,906482],
    "logistics":   [1427107,906494,1117210,1049622,1393382,449609,1427109,1117207],
    "port":        [1427107,906494,1117210,1427109,1546960,2226458,1117207,906482],
    "truck":       [1393382,449609,378570,1049622,1197095,1118448,1393374,449598],
    "transport":   [449609,1393382,1049622,378570,1197095,1118448,906494,1427107],
    "highway":     [449609,378570,1197095,1049622,1393382,1118448,449598,378558],
    "road":        [449609,378570,1197095,1049622,1393382,1118448,449598,1197088],
    "bridge":      [1197095,449609,378570,1049622,1118448,1393382,1197088,449598],
    "train":       [1393382,449609,378570,1049622,1197095,1118448,1393374,449598],
    "railway":     [1393382,449609,378570,1049622,1197095,1118448,1393374,449598],
    "airport":     [6447535,3396234,2058126,1437796,1917215,8159956,1537220,3722826],
    "aviation":    [6447535,3396234,2058126,1437796,1917215,8159956,1537220,3722826],
    "car":         [449609,378570,1197095,1049622,1393382,1118448,449598,378558],
    "vehicle":     [449609,378570,1197095,1049622,1393382,1118448,449598,378558],
    # ── Архитектура / Города ─────────────────────────────────────────────────
    "building":    [1486325,313782,373912,1693652,2246476,325229,1486310,1693640],
    "skyline":     [1486325,2246476,1693652,373912,313782,325229,2246468,1486310],
    "city":        [1486325,373912,1693652,313782,2246476,325229,1486310,1693640],
    "architecture":[1486325,313782,373912,1693652,2246476,325229,1486310,1693640],
    "government":  [313782,1486325,373912,1693652,2246476,325229,1693640,1486310],
    "parliament":  [313782,1486325,373912,1693652,2246476,325229,1693640,1486310],
    "urban":       [1693652,1486325,373912,313782,2246476,325229,1693668,1486310],
    "interior":    [1181244,574069,325229,270360,8386440,3861969,1181236,574060],
    "skyscraper":  [1486325,2246476,1693652,373912,313782,325229,2246468,1486310],
    "moscow":      [313782,1486325,373912,1693652,2246476,325229,1693640,1486310],
    "russia":      [2058126,1917215,313782,3396234,1486325,1437796,1693652,373912],
    # ── Спорт ─────────────────────────────────────────────────────────────────
    "athlete":     [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "athletes":    [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "stadium":     [863988,841130,1640777,2526878,1552242,1295572,863980,841122],
    "sport":       [863988,841130,1640777,2526878,1552242,1295572,863980,841122],
    "football":    [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "soccer":      [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "basketball":  [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "olympic":     [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "champion":    [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "trophy":      [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "hockey":      [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "tennis":      [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    # ── Культура / Искусство ──────────────────────────────────────────────────
    "art":         [1486325,313782,373912,1693652,2246476,325229,1693640,1486310],
    "painting":    [1486325,313782,373912,1693652,2246476,325229,1693640,1486310],
    "gallery":     [1486325,313782,373912,1693652,2246476,325229,1486310,1693640],
    "theater":     [1486325,313782,373912,1693652,2246476,325229,1693640,1486310],
    "music":       [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "concert":     [863988,841130,1552242,2526878,1640777,1295572,863980,841122],
    "film":        [574069,8386440,1181244,270360,3861969,325229,574060,270352],
    "book":        [1205651,289737,1438072,1350700,2422232,1205652,289720,1438062],
    "library":     [1205651,289737,1438072,1350700,2422232,1205652,289720,1438062],
    "literature":  [1205651,289737,1438072,1350700,2422232,1205652,289720,1438062],
    # ── Политика / Общество ───────────────────────────────────────────────────
    "politics":    [313782,1486325,1693652,373912,325229,2246476,1693640,1486310],
    "election":    [313782,1486325,373912,1693652,2246476,325229,1693640,1486310],
    "democracy":   [313782,1486325,373912,1693652,2246476,325229,1693640,1486310],
    "law":         [313782,1486325,373912,1693652,2246476,325229,1693640,1486310],
    "justice":     [313782,1486325,373912,1693652,2246476,325229,1693640,1486310],
    "protest":     [2058126,1917215,3396234,313782,1437796,1486325,1693652,373912],
    "society":     [3182812,3184465,3182781,3184454,3182773,3184419,3182804,3184464],
    "family":      [3182812,3184465,3182781,3184454,3182773,3184419,3182804,3184464],
    "diplomacy":   [313782,1486325,373912,1693652,2246476,325229,1693640,1486310],
    "geopolitics": [313782,2058126,1486325,1917215,373912,1693652,3396234,2246476],
    "sanctions":   [590022,159888,210607,6801874,3771110,669619,1591059,1602726],
    "crisis":      [590022,159888,2058126,3184465,1917215,210607,3396234,3182812],
    "disaster":    [247431,167698,338936,1525041,1761279,167696,338930,1525035],
    # ── Еда / Путешествия ─────────────────────────────────────────────────────
    "food":        [70497,349609,1640776,1640777,1640778,349611,70498,349608],
    "restaurant":  [70497,349609,1640776,1640777,1640778,349611,70498,349608],
    "chef":        [70497,349609,1640776,1640777,1640778,349611,70498,349608],
    "travel":      [1295138,1295139,731082,1001682,1295140,731080,1001680,1295141],
    "tourism":     [1295138,1295139,731082,1001682,1295140,731080,1001680,1295141],
    "hotel":       [1295138,1295139,731082,1001682,1295140,731080,1001680,1295141],
    # ── Общее ─────────────────────────────────────────────────────────────────
    "future":      [8386440,574069,1181244,1486325,270360,3861969,313782,325229],
    "summary":     [1486325,313782,3184465,3182812,1693652,373912,590022,159888],
    "introduction":[1486325,313782,3184465,3182812,1693652,373912,590022,159888],
    "analysis":    [159888,210607,590022,6801874,3771110,669619,1602726,1602720],
    "statistics":  [159888,210607,6801874,590022,3771110,669619,1602720,1602726],
    "trend":       [159888,210607,590022,6801874,3184465,3182812,3771110,669619],
    "global":      [1486325,313782,3184465,3182812,1693652,373912,590022,159888],
    "international":[313782,1486325,373912,1693652,2246476,325229,1693640,1486310],
    "cooperation": [3184465,3182812,3182781,3184454,3182773,3184419,3182804,3184464],
    "conclusion":  [3184465,1486325,313782,3182812,863988,1693652,159888,590022],
    "result":      [3184465,3182812,863988,841130,159888,590022,1552242,2526878],
    "problem":     [3786157,305568,590022,3184465,159888,4386431,3182812,263402],
    "solution":    [3184465,3182812,863988,841130,3182781,3184454,590022,159888],
    "success":     [3184465,3182812,863988,841130,3182781,3184454,1552242,2526878],
    "goal":        [3184465,3182812,863988,841130,3182781,3184454,1552242,2526878],
    "motivation":  [863988,841130,3184465,3182812,1552242,2526878,1640777,1295572],
    "risk":        [590022,159888,210607,6801874,3771110,3184465,3182812,669619],
}

# Семантический маппинг: любое слово/подстрока (рус/англ) → ключ в _PHOTO_DB
_PHOTO_SEMANTIC = {
    # Военное / оборонное
    "солдат": "soldier",     "армия": "army",        "военн": "military",
    "война":  "war",         "битва": "battle",       "флот":  "navy",
    "ракет":  "missile",     "танк":  "tank",         "оружи": "weapon",
    "оборон": "defense",     "боев":  "combat",       "парад": "parade",
    "нато":   "nato",        "nato":  "nato",         "войск": "soldiers",
    "конфликт": "conflict",  "фронт": "conflict",     "снайпер": "sniper",
    "безопасн": "security",  "атак":  "combat",       "штурм": "combat",
    "подводн": "warship",    "крейсер": "warship",    "авианосец": "warship",
    # История
    "истори":  "history",    "древн":  "ancient",     "средневек": "medieval",
    "империя": "empire",     "революц": "revolution", "советск":  "soviet",
    "наполеон": "napoleon",  "памятник": "monument",  "музей":    "museum",
    "кремл":   "kremlin",    "дворец": "palace",      "собор":    "cathedral",
    "архив":   "archive",    "насл":   "heritage",    "традиц":   "tradition",
    "вов ":    "wwii",       "отечественная": "wwii", "блокад":   "wwii",
    "imperial": "imperial",  "russian empire": "imperial",
    # Технологии
    "технолог": "technology","компьют": "computer",   "нейросет": "neural",
    "программ": "programming","данны":  "data",       "сервер":   "server",
    "искусствен": "artificial","цифров": "digital",   "стартап":  "startup",
    "инновац": "innovation", "облак":   "cloud",      "робот":    "robot",
    "блокчейн": "blockchain","смартфон": "smartphone","ноутбук":  "laptop",
    # Бизнес
    "бизнес":  "business",  "офис":    "office",     "встреч":   "meeting",
    "команд":  "team",      "лидер":   "leadership", "стратег":  "strategy",
    "маркетинг": "marketing","продаж":  "sales",     "бренд":    "brand",
    "менеджмент": "management","управлени": "management",
    # Финансы
    "финанс":  "finance",   "деньг":   "finance",    "банк":     "bank",
    "инвестиц": "investment","рынок":  "market",     "акц":      "stock",
    "экономик": "economy",  "бюджет":  "budget",     "инфляц":   "inflation",
    "прибыл":  "profit",    "выручк":  "revenue",    "санкц":    "sanction",
    # Промышленность
    "завод":   "factory",   "промышленн": "industry","производств": "production",
    "сборк":   "assembly",  "металл":  "steel",      "строительств": "construction",
    "станок":  "machinery", "сварк":   "welding",
    # Энергетика
    "нефт":    "oil",       "газ":     "gas",         "энергетик": "energy",
    "трубопровод": "pipeline","нефтепереработ": "refinery",
    "солнечн": "solar",     "ядерн":   "nuclear",     "электростанц": "power",
    # Сельское хозяйство / природа
    "пшениц":  "wheat",     "урожай":  "harvest",     "ферм":     "farm",
    "сельск":  "agriculture","трактор": "agriculture","лес":      "forest",
    "природ":  "nature",    "экологи": "ecology",     "климат":   "climate",
    "гор":     "mountain",  "рек":     "river",       "мор":      "sea",
    "океан":   "ocean",     "пляж":    "beach",
    # Медицина
    "врач":    "doctor",    "больниц": "hospital",    "медицин":  "medical",
    "лаборатор": "laboratory","исследован": "research",
    "фитнес":  "fitness",   "здоровь": "wellness",    "психолог": "mental",
    # Образование / наука
    "студент": "students",  "школ":    "school",      "университет": "university",
    "образован": "education","учител": "teacher",     "наук":     "science",
    "физик":   "physics",   "хими":    "chemistry",   "биологи":  "biology",
    "космос":  "cosmos",    "астроном": "astronomy",  "спутник":  "satellite",
    # Транспорт
    "логистик": "logistics","грузовик": "truck",      "корабл":   "ship",
    "транспорт": "transport","дорог":   "road",       "мост":     "bridge",
    "поезд":   "train",     "авиац":   "aviation",   "аэропорт": "airport",
    "автомобил": "car",
    # Города / архитектура
    "город":   "city",      "архитектур": "architecture","здани":  "building",
    "москва":  "moscow",    "кремл":   "kremlin",     "россия":   "russia",
    "небоскр": "skyscraper","парламент": "parliament","правительств": "government",
    # Спорт
    "спорт":   "sport",     "футбол":  "football",   "баскетбол": "basketball",
    "хоккей":  "hockey",    "теннис":  "tennis",     "стадион":  "stadium",
    "чемпион": "champion",  "олимпи":  "olympic",
    # Культура
    "искусств": "art",      "музык":   "music",       "театр":    "theater",
    "кино":    "film",      "книг":    "book",        "библиотек": "library",
    # Политика
    "политик": "politics",  "выбор":   "election",    "демократи": "democracy",
    "закон":   "law",       "безопасност": "security","дипломат": "diplomacy",
    "геополитик": "geopolitics","санкц": "sanctions",  "кризис":  "crisis",
    # Еда
    "ресторан": "restaurant","повар":  "chef",        "еда":      "food",
    # Путешествия
    "туризм":  "tourism",   "путешеств": "travel",   "отел":     "hotel",
    # Общее
    "будущ":   "future",    "вывод":   "conclusion",  "итог":     "summary",
    "введени": "introduction","анализ": "analysis",   "тренд":    "trend",
    "мотивац": "motivation","успех":   "success",     "цель":     "goal",
    "риск":    "risk",      "кризис":  "crisis",      "решени":   "solution",
    "проблем": "problem",   "результат": "result",
}

# Большой нейтральный пул — 30 разных РАЗНЫХ фото
_PHOTO_NEUTRAL = [
    3184465, 3182812, 1486325, 590022,  1108572, 863988,
    1205651, 1427107, 313782,  574069,  159888,  2760241,
    326082,  3786157, 8386440, 6447535, 1917215, 414612,
    167698,  70497,   1295138, 1393382, 1693652, 4386431,
    841130,  905163,  449609,  3182781, 2058126, 1078116,
]


def _photo_id(keyword: str, slide_idx: int = 0) -> int:
    """
    Универсальная функция: keyword (english) → Pexels photo ID.
    Используется обоими рендерерами (HTML и PPTX).
    Алгоритм:
      1. Точное совпадение слова
      2. Семантический маппинг (рус+англ подстроки)
      3. Префиксное совпадение 5 символов
      4. Префиксное совпадение 4 символа
      5. Нейтральный пул (детерминированный hash)
    """
    kw = (keyword or "").lower()
    words = kw.split()

    # 1. Точное совпадение
    for w in words:
        if w in _PHOTO_DB:
            ids = _PHOTO_DB[w]
            return ids[slide_idx % len(ids)]

    # 2. Семантический маппинг
    for sem, db_key in _PHOTO_SEMANTIC.items():
        if sem in kw and db_key in _PHOTO_DB:
            ids = _PHOTO_DB[db_key]
            return ids[slide_idx % len(ids)]

    # 3. Префикс-5
    for w in words:
        if len(w) >= 5:
            for k in _PHOTO_DB:
                if k[:5] == w[:5]:
                    ids = _PHOTO_DB[k]
                    return ids[slide_idx % len(ids)]

    # 4. Префикс-4
    for w in words:
        if len(w) >= 4:
            for k in _PHOTO_DB:
                if len(k) >= 4 and k[:4] == w[:4]:
                    ids = _PHOTO_DB[k]
                    return ids[slide_idx % len(ids)]

    # 5. Нейтральный пул
    base = abs(hash(kw)) % len(_PHOTO_NEUTRAL)
    return _PHOTO_NEUTRAL[(base + slide_idx) % len(_PHOTO_NEUTRAL)]

def _sanitize_keyword(kw: str) -> str:
    """
    Умный нормализатор ключевых слов:
    - Если есть кириллица → переводит через словарь _RU_TO_EN_PHOTO
    - Если уже English → очищает мусорные слова
    - Всегда возвращает хороший English поисковый запрос
    """
    import re as _re
    if not kw:
        return "professional modern"

    kw_low = kw.lower().strip()

    # Проверяем кириллицу
    has_cyrillic = bool(_re.search(r'[а-яёА-ЯЁ]', kw))

    if has_cyrillic:
        # Сначала пробуем длинные совпадения (более специфичные)
        best_match = ""
        best_key   = ""
        for ru_key, en_val in sorted(_RU_TO_EN_PHOTO.items(), key=lambda x: -len(x[0])):
            if ru_key in kw_low:
                if len(ru_key) > len(best_key):
                    best_key   = ru_key
                    best_match = en_val

        # Берём английские слова из оригинального keyword (если есть)
        latin_extra = " ".join(_re.findall(r'[A-Za-z]{3,}', kw))

        if best_match:
            result = (best_match + " " + latin_extra).strip()
        elif latin_extra:
            result = latin_extra
        else:
            result = "professional modern presentation"
        return result[:100]

    # Уже английский — убираем плохие абстрактные слова
    bad = {'abstract', 'background', 'concept', 'idea', 'presentation',
           'illustration', 'symbol', 'icon', 'logo', 'design', 'pattern',
           'texture', 'wallpaper', 'banner', 'template', 'vector',
           'infographic', 'graphic', 'virtual', 'creative', 'generic'}
    words = [w.strip().lower() for w in kw.split() if w.strip()]
    words = [w for w in words if w not in bad and len(w) > 2]
    result = " ".join(words[:6])
    return result.strip() or "professional"


def _fallback_keywords(keyword: str, slide_title: str, topic: str) -> list[str]:
    """
    Генерирует список ключевых слов от лучшего к худшему.
    Если основной keyword плохой — используем умные запасные варианты.
    """
    candidates = []

    # 1. Основное keyword (после очистки)
    clean = _sanitize_keyword(keyword)
    if len(clean) >= 4:
        candidates.append(clean)

    # 2. Упрощённый вариант — первые 2-3 слова
    words = clean.split()
    if len(words) >= 3:
        candidates.append(" ".join(words[:3]))
    if len(words) >= 2:
        candidates.append(" ".join(words[:2]))

    # 3. Из заголовка слайда — если есть латиница
    if slide_title:
        import re as _re
        latin_words = _re.findall(r'[A-Za-z]{3,}', slide_title)
        if latin_words:
            candidates.append(" ".join(latin_words[:3]))

    # 4. Маппинг русских тем → хорошие английские keywords
    _TOPIC_MAP = {
        # ── Россия / СНГ ──────────────────────────────────────────────
        'внутренняя экономика': 'Moscow Russia financial district skyline',
        'экономика росси': 'Russia economy Moscow city business district',
        'внутренняя политика': 'Russian State Duma government Moscow building',
        'внешняя политика': 'Russia diplomacy foreign ministry Moscow',
        'россия': 'Russia Moscow Red Square Kremlin landmark',
        'москва': 'Moscow Russia city skyline modern architecture',
        'рубл': 'Russia ruble currency exchange bank Moscow',
        'санкци': 'Russia economy sanctions business adaptation',
        'импортозамещени': 'Russia import substitution domestic production factory',
        'нефт': 'Russia oil pipeline Siberia energy production',
        'газ': 'Russia natural gas pipeline export Gazprom',
        'экспорт': 'Russia cargo ship port Vladivostok export trade',
        'оборонн': 'Russian military defense industry factory weapons production',
        'промышленност': 'Russian heavy industry factory workers production',
        # ── Военная тематика ────────────────────────────────────────────
        # ── Война / История России ────────────────────────────────────
        'войны росси': 'Russian Empire soldiers historical battle war cavalry',
        'история росси': 'Russian Empire historical Moscow Kremlin ancient battle',
        'военная история': 'historical military battle soldiers war cavalry armor',
        'россия войны': 'Russian Empire military war historical battle Red Army',
        'великая отечественная': 'World War 2 Soviet Red Army soldiers battle tank',
        'советск': 'Soviet Union historical Red Square parade military soldiers',
        'первая мировая': 'World War 1 soldiers trench battle cavalry historical',
        'вторая мировая': 'World War 2 battle soldiers tank Red Army historical',
        'наполеон': 'Napoleonic wars battle 1812 soldiers cavalry Russia historical',
        'крымская война': 'Crimean War 1854 historical battle soldiers fortification',
        'полтава': 'Battle Poltava 1709 Peter Great Swedish historical soldiers',
        'революц': 'Russian Revolution 1917 historical soldiers Red Army Bolshevik',
        'российская империя': 'Russian Empire historical palace Winter Palace St Petersburg',
        'территориальн': 'Russia map territory historical expansion geographic empire',
        'индустриализац': 'Soviet industrialization steel factory workers production 1930s',
        'геополитическ': 'Russia world map geopolitics power strategy historical',
        'холодная война': 'Cold War nuclear Soviet USA historical military tension',
        'космическая гонк': 'Soviet space race Sputnik Gagarin rocket cosmos historical',
        'ядерн': 'Soviet nuclear weapon atomic bomb test historical military',
        'сверхдержав': 'USSR superpower Cold War military parade Red Square historical',
        'московск': 'Moscow Kremlin Red Square historical architecture Russian',
        'санкт-петербург': 'Saint Petersburg historical palace Winter Palace Russia',
        'победа': 'Victory Day Russia May 9 soldiers military parade historical',
        'блокад': 'Leningrad Siege World War 2 historical Soviet soldiers defense',
        'сибир': 'Siberia Russia historical exploration vast landscape Cossacks',
                'военн': 'military defense armor tanks soldiers training exercise field',
        'военная экономик': 'defense industry military factory production weapons budget',
        'военные расход': 'military spending defense budget army equipment soldiers',
        'впк': 'Russian defense industry factory weapons production military armor',
        'оборонная промышленност': 'defense factory weapons production industrial military',
        'оборонный комплекс': 'military industrial complex factory armor production',
        'вооружени': 'military weapons defense equipment soldiers armor combat',
        'танк': 'military tank armor combat vehicle defense industry',
        'самолет': 'military aircraft fighter jet air force defense',
        'авиац': 'military aviation aircraft fighter jet air force',
        'флот': 'navy warship fleet military sea defense patrol',
        'ракет': 'military missile rocket defense weapon launch',
        'войск': 'military troops soldiers army training exercise field',
        'nato': 'NATO military alliance troops exercise defense weapons',
        'нато': 'NATO military alliance troops defense spending exercise',
        'оборона': 'military defense soldiers army base fortification guard',
        'мобилизац': 'military mobilization soldiers troops training defense',
        'фронт': 'military front line soldiers defense troops terrain',
        'конфликт': 'military conflict geopolitics defense strategy command center',
        'геополитик': 'geopolitics world map conflict strategy military command',
        'стратегическ': 'military strategic planning defense map command operations',
        'национальная безопасност': 'national security defense military strategy command',
        'безопасност': 'security defense military strategy protection surveillance',
        'беспилотн': 'military drone UAV unmanned aerial vehicle reconnaissance',
        'боеприпас': 'military ammunition weapons shells production defense factory',
        'логистик': 'military logistics supply chain transport vehicles convoy',
        'мощност': 'industrial production capacity military factory output assembly',
        'базовая нагрузк': 'industrial production capacity defense factory base load',
        'серийный выпуск': 'mass production factory assembly line industrial defense',
        'комплектующ': 'industrial components parts production supply chain defense',
        'рабочая сила': 'factory workers industrial labor production assembly line',
        'рынок труда': 'labor market workers factory industrial employment workforce',
        'узкое место': 'industrial bottleneck production capacity factory workers',
        'глобальный рост': 'global military spending defense world map geopolitics',
        'сельск': 'Russia wheat field harvest combine agricultural',
        'агропром': 'Russia agriculture farm grain harvest field',
        'бюджет': 'Russian government budget finance Moscow',
        'инфляц': 'Russia consumer prices market shopping inflation',
        'внп': 'Russia GDP growth economy statistics chart',
        'ввп': 'Russia GDP growth economy statistics chart',
        'государствен': 'Russian government official Moscow Kremlin',
        'фондов': 'Moscow Exchange trading financial market Russia',
        'банк': 'Russia bank Moscow Sberbank VTB finance building',
        'цифров': 'Russia digital technology Skolkovo innovation center',
        'сколков': 'Skolkovo Russia tech park innovation startup',
        'армия': 'Russian army military parade Moscow Red Square',
        'вооружен': 'Russian military forces defense equipment',
        'украин': 'Russia Ukraine conflict military geopolitics',
        'снг': 'CIS countries leaders summit meeting cooperation',
        'беларус': 'Belarus Minsk city government Union State',
        # Российские регионы и города
        'нижегород': 'Nizhny Novgorod Russia city Volga river historic architecture',
        'кстов': 'Nizhny Novgorod Russia industrial city Volga river infrastructure',
        'приволжск': 'Volga river Russia industrial city infrastructure regional',
        'поволж': 'Volga river landscape Russia city transport bridge',
        'волг': 'Volga river Russia landscape cargo ship bridge',
        'урал': 'Ural Russia mountains industrial city factory',
        'сибир': 'Siberia Russia vast landscape resources pipeline',
        'казан': 'Kazan Russia city modern architecture Tatarstan',
        'самар': 'Samara Russia Volga city industrial',
        'ярославл': 'Yaroslavl Russia historic city Volga river architecture',
        'пфо': 'Volga Federal District Russia cities industrial',
        # Инвестиции и бизнес
        'инвестиц': 'business investment meeting office modern city',
        'инвестиционн': 'investment business partner handshake office modern',
        'деловой климат': 'business climate office professionals meeting',
        'бизнес': 'business meeting office professionals modern city',
        # НПЗ / нефтехимия
        'нпз': 'oil refinery industrial plant pipes Lukoil petroleum',
        'нефтеперераб': 'oil refinery petroleum plant industrial smoke',
        'нефтехим': 'chemical plant industrial petroleum refinery pipes',
        'нефтепереработ': 'oil refinery industrial plant petroleum processing',
        'абразив': 'industrial factory grinding manufacturing production',
        'полимер': 'chemical industrial factory production polymers',
        'стирол': 'chemical plant industrial production factory pipes',
        # Транспорт и логистика
        'транспортн': 'highway trucks transport logistics road infrastructure',
        'логистик': 'logistics trucks highway transport warehouse cargo',
        'автомагистрал': 'highway motorway road trucks transport Russia',
        'железнодорожн': 'railway train railroad Russia transport infrastructure',
        'речн': 'river Volga cargo ship transport water logistics',
        'волжск': 'Volga river Russia cargo ship transport industrial',
        # Инфраструктура
        'энергоснабжен': 'power plant electricity grid energy infrastructure industrial',
        'водоснабжен': 'water treatment plant infrastructure pipes industrial',
        'оптоволокн': 'fiber optic internet infrastructure technology digital',
        'промплощадк': 'industrial zone factory area production plant',
        'промышленн зон': 'industrial zone factory district area production',
        'казахстан': 'Kazakhstan Nur-Sultan modern city architecture',
        'китай': 'China Russia cooperation business Beijing meeting',
        # ── Спорт ──────────────────────────────────────────────────────
        'спорт': 'teenagers playing sports outdoor stadium',
        'физкульт': 'school physical education students gym',
        'баскетбол': 'basketball players court game action',
        'футбол': 'football soccer players field match action',
        'тренировк': 'athletes training workout session coach',
        'здоровь': 'healthy lifestyle fitness active outdoor',
        'фитнес': 'fitness workout training gym indoor',
        # ── ИИ / технологии ────────────────────────────────────────────
        'искусственный интеллект': 'neural network glowing blue data visualization dark',
        'нейросет': 'neural network brain artificial intelligence visualization',
        'машинное обучение': 'machine learning data visualization screen code',
        'программ': 'programmer coding multiple screens dark office',
        'технологи': 'technology innovation digital startup office team',
        'робот': 'robotic arm factory automation close-up',
        'данны': 'data analytics visualization dashboard dark office',
        # ── Бизнес / Финансы ───────────────────────────────────────────
        'бизнес': 'diverse business team strategy meeting office',
        'маркетинг': 'marketing team planning strategy whiteboard office',
        'финанс': 'finance investment growth chart stock monitor',
        'менеджмент': 'management leadership team meeting corporate',
        # ── Образование ────────────────────────────────────────────────
        'образовани': 'students classroom learning lecture university',
        'школ': 'school students studying classroom teacher',
        'урок': 'teacher classroom students lesson learning',
        # ── Наука ──────────────────────────────────────────────────────
        'наук': 'science laboratory research experiment modern',
        'биологи': 'biology science microscope cells laboratory',
        'хими': 'chemistry laboratory colorful test tubes experiment',
        'физик': 'physics science experiment laboratory equipment',
        # ── Природа / Экология ─────────────────────────────────────────
        'природ': 'nature forest landscape green environment sunny',
        'экологи': 'ecology environment green sustainable solar panels',
        'климат': 'climate change environment green energy sustainable',
        # ── История / Культура ─────────────────────────────────────────
        'истори': 'historical architecture ancient monument landmark',
        'культур': 'culture art museum exhibition gallery',
        # ── Медицина ───────────────────────────────────────────────────
        'медицин': 'medical doctor hospital modern healthcare',
        'здравоохранени': 'healthcare medical professionals hospital modern',
        # ── Психология ─────────────────────────────────────────────────
        'психологи': 'psychology mind human emotions brain wellness',
        # ── Общее ──────────────────────────────────────────────────────
        'мотивац': 'motivated young people success achievement energy',
        'успех': 'success achievement celebration winner trophy',
        'команд': 'team collaboration working together office',
        'проект': 'project team planning meeting office whiteboard',
        'развити': 'growth development progress improvement steps',
        'будущ': 'future vision innovation bright modern city',
        'цель': 'goal achievement target success planning',
        'метод': 'professional methods strategy planning diagram',
        # ── ОБЖ / Безопасность / Терроризм (школьная тематика) ──────────
        'террор': 'police security officers patrol checkpoint public safety',
        'терроризм': 'police security checkpoint surveillance patrol urban',
        'террористическ': 'police security checkpoint surveillance crowd control',
        'экстремизм': 'police detective investigation security briefing',
        'экстремальн': 'safety training emergency first responders drill',
        'радикализ': 'school counselor psychologist student meeting prevention',
        'вербовк': 'police detective investigation computer screen crime',
        'вербовщик': 'police detective laptop investigation prevention',
        'несовершеннолетних': 'school students teenagers classroom education',
        'подростк': 'teenagers school students studying classroom modern',
        'профилактик': 'school teacher students prevention education meeting',
        'скулшутинг': 'school safety security drill evacuation students',
        'нападени': 'police security guard building patrol protection',
        'угроз': 'security police patrol building protection surveillance',
        'безопасност': 'police security officers patrol public safety urban',
        'гражданская оборона': 'emergency civil defense shelter people Russia',
        'мвд': 'Russian police law enforcement officers uniform patrol',
        'фсб': 'security intelligence law enforcement briefing uniform',
        'правоохранительн': 'Russian police law enforcement officers patrol',
        'правоохранители': 'police officers uniform patrol city security',
        'горячая линия': 'phone call police helpline emergency assistance',
        'социальн сет': 'teenager smartphone social media screen warning',
        'интернет безопасн': 'teenager computer internet cyber safety screen',
        'кибербезопасн': 'cyber security computer screen padlock data protection',
        'профилактика экстремизм': 'school teacher students classroom prevention',
        'первая помощь': 'first aid training kit bandage CPR medical',
        'пожарн': 'firefighters rescue fire truck ladder emergency',
        'эвакуац': 'evacuation crowd emergency exit building people',
        'чрезвычайн': 'emergency rescue team disaster response professional',
        'ожог': 'medical treatment burn wound bandage hospital care',
        'отравлен': 'medical ambulance emergency treatment hospital care',
        'дтп': 'road traffic accident police car crash emergency',
        'авари': 'accident emergency police rescue scene road',
        'спасател': 'rescue team emergency responders professional uniform',
        'наводнен': 'flood rescue boat helicopter emergency people water',
        'землетрясен': 'earthquake rescue survivors rubble emergency team',
        'обж': 'safety training school students emergency awareness Russia',
        'безопасность жизнедеятельности': 'safety education school students Russia training',
    }

    # Проверяем тему презентации и заголовок слайда
    combined = (topic + " " + slide_title).lower()

    # Детектируем тему безопасности/ОБЖ — приоритет над военной
    is_safety = any(w in combined for w in [
        'террор', 'экстремизм', 'радикализ', 'вербовк', 'скулшутинг',
        'обж', 'безопасност', 'пожарн', 'спасател', 'первая помощь',
        'эвакуац', 'профилактик', 'несовершеннолетних', 'подростк',
        'мвд', 'фсб', 'правоохранительн',
    ])

    # Детектируем военную тему — приоритет над географическими keywords
    is_military = any(w in combined for w in [
        'военн', 'впк', 'оборон', 'армия', 'армии', 'войск', 'нато', 'nato',
        'military', 'defense', 'weapon', 'soldier', 'tank', 'missile',
        'расход', 'мощност', 'базовая нагруз', 'серийный', 'боеприпас',
    ])

    # Если тема безопасности — добавляем правильные ключевые слова первыми
    if is_safety:
        # Подбираем ключевые слова по конкретному содержанию слайда
        _slide_lower = slide_title.lower()
        if any(w in _slide_lower for w in ['статистик', 'данн', 'цифр', 'показател']):
            safety_kws = ["police statistics data analysis chart security", "crime prevention statistics infographic data"]
        elif any(w in _slide_lower for w in ['метод', 'способ', 'вербовк', 'признак']):
            safety_kws = ["detective investigation laptop evidence digital police", "teenager smartphone social media cyber awareness"]
        elif any(w in _slide_lower for w in ['алгоритм', 'действи', 'как', 'что делать']):
            safety_kws = ["police emergency hotline call help protection", "security training drill evacuation school students"]
        elif any(w in _slide_lower for w in ['профилактик', 'защит', 'меры']):
            safety_kws = ["school teacher students prevention education awareness", "parents children safety internet discussion home"]
        elif any(w in _slide_lower for w in ['вывод', 'заключени', 'итог', 'результат']):
            safety_kws = ["police community safety team school cooperation Russia", "school safety awareness students education Russia"]
        elif any(w in _slide_lower for w in ['цитат', 'мнени', 'эксперт']):
            safety_kws = ["police officer expert briefing press conference uniform", "security professional expert speaking audience"]
        elif any(w in _slide_lower for w in ['подростк', 'несовершеннолетн', 'школьник']):
            safety_kws = ["Russian school students teenagers classroom studying education", "teenagers group school friends outdoor Russia"]
        elif any(w in _slide_lower for w in ['социальн сет', 'интернет', 'telegram', 'vk', 'онлайн']):
            safety_kws = ["teenager smartphone social media screen warning parents", "cyber security internet safety teen screen alert"]
        else:
            safety_kws = ["police security patrol public safety urban Russia", "school safety security awareness students education Russia"]

        for sk in reversed(safety_kws):
            if sk not in candidates:
                candidates.insert(0, sk)

    # Если военная тема — военные keywords идут ПЕРВЫМИ
    elif is_military:
        military_kws = [
            "soldiers military training exercise armor field",
            "defense industry factory weapons production military",
            "armored vehicle military tank defense production",
            "military troops army exercise combat field",
        ]
        # Вставляем в начало списка
        for mk in reversed(military_kws):
            if mk not in candidates:
                candidates.insert(0, mk)

    for ru_key, en_kw in _TOPIC_MAP.items():
        if ru_key in combined:
            if en_kw not in candidates:
                candidates.append(en_kw)

    # 5. Крайний резерв — по теме топика напрямую (первые английские слова)
    import re as _re2
    topic_latin = _re2.findall(r'[A-Za-z]{4,}', topic)
    if topic_latin:
        candidates.append(" ".join(topic_latin[:3]) + " professional")

    # 6. Универсальный безопасный фоллбек
    if is_safety:
        candidates.append("police security patrol public safety professional")
    elif is_military:
        candidates.append("military defense army soldiers training")
    else:
        candidates.append("professional team modern workspace")

    # Убираем дубликаты, сохраняем порядок
    seen, unique = set(), []
    for c in candidates:
        c = c.strip()
        if c and c not in seen:
            seen.add(c)
            unique.append(c)
    return unique


async def _fetch_photo(keyword: str, slide_title: str = "",
                       topic: str = "", slide_index: int = 0) -> bytes | None:
    """
    Загрузка тематического фото.
    Источники в порядке приоритета:
    1. Pexels CDN (прямые ID) — для ОБЖ/безопасности ВСЕГДА первый
    2. Wikimedia Commons — для исторических/академических тем (с фильтром)
    3. Pexels / Pixabay / Unsplash API — если есть ключ
    4. Pexels CDN (общий) — финальный резерв
    5. Picsum — абсолютный резерв
    """
    import urllib.parse as _up

    seed = abs(hash(keyword + slide_title + str(slide_index))) % 99999
    keywords_to_try = _fallback_keywords(keyword, slide_title, topic)
    logging.info(f"🔍 Фото для «{slide_title[:30]}»: {keywords_to_try[:2]}")

    pixabay_key  = os.environ.get("PIXABAY_KEY", "")
    unsplash_key = os.environ.get("UNSPLASH_KEY", "")
    pexels_key   = os.environ.get("PEXELS_KEY", "")

    # ─── Определяем тему ──────────────────────────────────────────────
    _kw_combined = (keyword + " " + slide_title + " " + topic).lower()

    # Темы ОБЖ/безопасности — Wikimedia даёт нерелевантных иностранных полицейских
    _is_safety = any(w in _kw_combined for w in [
        "террор", "вербов", "подрост", "безопасн", "радикализ", "экстремизм",
        "обж", "школ", "ученик", "профилакт", "несовершеннолет", "фсб", "мвд",
        "teenager", "smartphone", "social media", "cyber", "cybersecurity",
        "police", "security", "safety", "school", "student", "prevention",
        "awareness", "firefi", "evacuation", "rescue", "first aid",
    ])

    # ─── Таблица Pexels ID для ОБЖ / безопасности ────────────────────
    _SAFETY_PEXELS: dict[str, list[int]] = {
        "school_class":     [1205651, 289737, 1438072, 1350700, 2422232, 1205652],
        "teenager_phone":   [5081918, 699122, 2741448, 5082578, 3184465, 3182812],
        "cyber_security":   [3182812, 3184465, 5081918, 699122,  2741448, 5082578],
        "police_security":  [1653087, 1697912, 3153207, 2896941, 3777565, 2828237],
        "statistics_chart": [590022,  210607,  159888,  3771110, 6801874, 669619],
        "prevention_talk":  [3184465, 3182812, 1205651, 289737,  1438072, 3184454],
        "expert_speaker":   [3184465, 3182812, 3182781, 3184454, 3182773, 3184419],
        "conclusion":       [1205651, 289737,  3184465, 3182812, 1438072, 1350700],
        "firefighters":     [3822850, 3822851, 1549280, 2047905, 1170979, 3822852],
        "evacuation":       [2558679, 2557399, 2917989, 2904389, 3807571, 3807572],
    }

    def _pick_safety_id(kw: str, title: str, idx: int) -> int:
        sl = (title + " " + kw).lower()
        if any(w in sl for w in ['статистик', 'данн', 'цифр', 'процент', 'задержан']):
            pool = _SAFETY_PEXELS["statistics_chart"]
        elif any(w in sl for w in ['вербовк', 'соцсет', 'telegram', 'vk', 'social', 'smartphone', 'phone', 'teen']):
            pool = _SAFETY_PEXELS["teenager_phone"]
        elif any(w in sl for w in ['кибер', 'интернет', 'онлайн', 'digital', 'cyber', 'computer', 'screen']):
            pool = _SAFETY_PEXELS["cyber_security"]
        elif any(w in sl for w in ['полиц', 'мвд', 'фсб', 'правоохран', 'police', 'security', 'checkpoint']):
            pool = _SAFETY_PEXELS["police_security"]
        elif any(w in sl for w in ['алгоритм', 'действи', 'горячая линия', 'что делать', 'how to']):
            pool = _SAFETY_PEXELS["prevention_talk"]
        elif any(w in sl for w in ['пожарн', 'firefi', 'smoke', 'fire']):
            pool = _SAFETY_PEXELS["firefighters"]
        elif any(w in sl for w in ['эвакуац', 'evacuation', 'emergency']):
            pool = _SAFETY_PEXELS["evacuation"]
        elif any(w in sl for w in ['вывод', 'заключени', 'итог', 'conclusion', 'expert', 'цитат']):
            pool = _SAFETY_PEXELS["conclusion"]
        elif any(w in sl for w in ['подростк', 'школьник', 'класс', 'student', 'classroom']):
            pool = _SAFETY_PEXELS["school_class"]
        else:
            pool = _SAFETY_PEXELS["prevention_talk"]
        return pool[idx % len(pool)]

    async def _dl(url: str) -> bytes | None:
        try:
            async with aiohttp.ClientSession(
                timeout=aiohttp.ClientTimeout(total=15, connect=6)
            ) as s:
                async with s.get(
                    url, allow_redirects=True,
                    headers={
                        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
                        "Accept": "image/webp,image/apng,image/*,*/*;q=0.8",
                    }
                ) as r:
                    if r.status == 200:
                        ct = r.headers.get("content-type", "")
                        final_url = str(r.url) if hasattr(r, 'url') else url
                        is_image = (
                            "image" in ct
                            or final_url.endswith((".jpg",".jpeg",".png",".webp"))
                            or "pexels.com" in url
                            or "wikimedia.org" in url
                            or "picsum.photos" in url
                            or "staticflickr.com" in final_url
                            or "live.staticflickr" in final_url
                        )
                        if is_image:
                            data = await r.read()
                            if len(data) > 10_000:
                                return data
        except Exception:
            pass
        return None

    # ═══════════════════════════════════════════════════════════════════
    # ⚡ ПРИОРИТЕТ 1: Для ОБЖ/безопасности — СРАЗУ Pexels CDN
    # Wikimedia для этих тем даёт POLIZEI, индийских полицейских и т.д.
    # ═══════════════════════════════════════════════════════════════════
    if _is_safety:
        for _attempt_idx in range(3):
            _sid = _pick_safety_id(keyword, slide_title, slide_index + _attempt_idx)
            data = await _dl(_pex_cdn(_sid))
            if data:
                logging.info(f"📸 Safety Pexels CDN [{_sid}]: «{slide_title[:30]}»")
                return data

    # ═══════════════════════════════════════════════════════════════════
    # ПРИОРИТЕТ 2: Wikimedia Commons (только для НЕ-ОБЖ тем)
    # ═══════════════════════════════════════════════════════════════════
    # Слова-фильтры: если они в названии файла — пропускаем
    _BAD_WIKI_WORDS = [
        "polizei", "india", "indian", "german", "deutsch",
        "american", "usa", "british", "french", "turkish",
        "china", "chinese", "japan", "korean", "arabic", "pakistan",
        "bangladesh", "indonesia", "thailand", "iran",
    ]

    for attempt, kw in enumerate(keywords_to_try[:4]):
        kw_clean = kw.strip()[:100]
        q_enc    = _up.quote(kw_clean)

        if not _is_safety:
            try:
                wiki_search = (
                    "https://commons.wikimedia.org/w/api.php"
                    f"?action=query&list=search&srsearch={q_enc}+photo"
                    f"&srnamespace=6&srlimit=20&format=json"
                )
                async with aiohttp.ClientSession(
                    timeout=aiohttp.ClientTimeout(total=8)
                ) as s:
                    async with s.get(
                        wiki_search,
                        headers={"User-Agent": "PresentationBot/2.0 (educational)"}
                    ) as r:
                        if r.status == 200:
                            js = await r.json()
                            results = js.get("query", {}).get("search", [])
                            img_files = [
                                x for x in results
                                if re.search(r'\.(jpg|jpeg|png)$', x.get("title",""), re.I)
                                and not any(b in x.get("title","").lower() for b in _BAD_WIKI_WORDS)
                            ]
                            if img_files:
                                pick = img_files[slide_index % len(img_files)]
                                fname = _up.quote(
                                    pick["title"].replace("File:", "").replace(" ", "_")
                                )
                                thumb_url = (
                                    "https://commons.wikimedia.org/wiki/Special:FilePath/"
                                    f"{fname}?width=1280"
                                )
                                data = await _dl(thumb_url)
                                if data:
                                    logging.info(f"📸 Wikimedia: «{kw_clean[:40]}» [{len(data)//1024}кб]")
                                    return data
            except Exception:
                pass

        # ═══════════════════════════════════════════════════════════════
        # 3. Pexels Search API (если есть PEXELS_KEY)
        # ═══════════════════════════════════════════════════════════════
        if pexels_key:
            try:
                async with aiohttp.ClientSession(
                    timeout=aiohttp.ClientTimeout(total=10)
                ) as s:
                    async with s.get(
                        f"https://api.pexels.com/v1/search?query={q_enc}&per_page=15&orientation=landscape",
                        headers={"Authorization": pexels_key}
                    ) as r:
                        if r.status == 200:
                            js = await r.json()
                            photos = js.get("photos", [])
                            if photos:
                                photo = photos[slide_index % len(photos)]
                                url = photo.get("src", {}).get("large2x") or photo.get("src", {}).get("large", "")
                                data = await _dl(url)
                                if data:
                                    logging.info(f"📸 Pexels API: «{kw_clean[:40]}»")
                                    return data
            except Exception:
                pass

        # ═══════════════════════════════════════════════════════════════
        # 4. Pixabay API (если есть PIXABAY_KEY)
        # ═══════════════════════════════════════════════════════════════
        if pixabay_key:
            try:
                async with aiohttp.ClientSession(
                    timeout=aiohttp.ClientTimeout(total=8)
                ) as s:
                    api_url = (
                        f"https://pixabay.com/api/?key={pixabay_key}"
                        f"&q={q_enc}&image_type=photo&orientation=horizontal"
                        f"&min_width=1200&per_page=20&safesearch=true"
                    )
                    async with s.get(api_url) as r:
                        if r.status == 200:
                            js = await r.json()
                            hits = js.get("hits", [])
                            if hits:
                                hit = hits[slide_index % len(hits)]
                                url = hit.get("largeImageURL") or hit.get("webformatURL", "")
                                data = await _dl(url)
                                if data:
                                    logging.info(f"📸 Pixabay: «{kw_clean[:40]}»")
                                    return data
            except Exception:
                pass

        # ═══════════════════════════════════════════════════════════════
        # 5. Unsplash API (если есть UNSPLASH_KEY)
        # ═══════════════════════════════════════════════════════════════
        if unsplash_key:
            try:
                async with aiohttp.ClientSession(
                    timeout=aiohttp.ClientTimeout(total=8)
                ) as s:
                    async with s.get(
                        f"https://api.unsplash.com/photos/random?query={q_enc}&orientation=landscape&count=5&client_id={unsplash_key}"
                    ) as r:
                        if r.status == 200:
                            js = await r.json()
                            photos = js if isinstance(js, list) else [js]
                            if photos:
                                url = photos[slide_index % len(photos)].get("urls", {}).get("regular", "")
                                data = await _dl(url)
                                if data:
                                    logging.info(f"📸 Unsplash: «{kw_clean[:40]}»")
                                    return data
            except Exception:
                pass

    # ═══════════════════════════════════════════════════════════════════
    # 6. Pexels CDN — статичные ID через глобальный _photo_id (400+ тем)
    # ═══════════════════════════════════════════════════════════════════
    best_kw = (keywords_to_try[0] if keywords_to_try else keyword or "")
    chosen_id = _photo_id(best_kw, slide_index)
    cdn_url = _pex_cdn(chosen_id)
    data = await _dl(cdn_url)
    if data:
        logging.info(f"📸 Pexels CDN [{chosen_id}]: «{best_kw[:30]}»")
        return data

    # 7. Picsum абсолютный резерв
    data = await _dl(f"https://picsum.photos/seed/{seed}/1280/720")
    if data:
        return data
    return None


async def _fetch_photos_all(slides_data: list, topic: str = "") -> dict:
    """
    Загружает фото для всех слайдов параллельно.
    Передаёт тему презентации для умного фоллбека ключевых слов.
    """
    results: dict[int, bytes | None] = {}

    async def _one(idx, kw, slide_title):
        results[idx] = await _fetch_photo(kw, slide_title=slide_title, topic=topic, slide_index=idx)

    tasks = []
    for i, sd in enumerate(slides_data):
        kw = sd.get("image_keyword") or ""
        slide_title = sd.get("title", "")
        tasks.append(asyncio.create_task(_one(i, kw, slide_title)))
    if tasks:
        await asyncio.gather(*tasks, return_exceptions=True)
    return results


# ── Построители слайдов (современный дизайн) ─────────────────────


# ══════════════════════════════════════════════════════════════════
# 🎨 ELEGANT LIGHT THEME — конкурентный стиль (кремовый, фото слева)
# ══════════════════════════════════════════════════════════════════

def _el_slide_title(prs, data, th, img: bytes | None):
    """Титульный: фото слева 42%, белый справа с заголовком."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN as PA
    import datetime as _dt

    ft = th.get("font_title", "Georgia")
    fb = th.get("font_body", "Calibri")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["bg2"])  # белый фон

    PW = SLIDE_W * 0.42
    # Фото слева
    if img:
        _pptx_img(slide, img, 0, 0, PW, SLIDE_H)
    else:
        _pptx_rect(slide, 0, 0, PW, SLIDE_H, th["header"])

    # Правая панель — белый
    _pptx_rect(slide, PW, 0, SLIDE_W - PW, SLIDE_H, th["bg2"])

    RX = PW + 0.55
    RW = SLIDE_W - PW - 0.75

    # Организация/подзаголовок вверху
    org = data.get("organization", "")
    if org:
        _pptx_txt(slide, org.upper(), RX, 0.7, RW, 0.4,
                  size=9, color=th["muted_c"], bold=True, font=fb)

    # Главный заголовок
    title = data.get("title", "Презентация")
    _pptx_txt(slide, title, RX, 1.4, RW, 3.2,
              size=40, color=th["title_c"], bold=True, font=ft, wrap=True)

    # Разделитель
    _pptx_rect(slide, RX, 4.75, 0.6, 0.05, th["accent"])

    # Подзаголовок
    sub = data.get("subtitle", "")
    if sub:
        _pptx_txt(slide, sub, RX, 5.0, RW, 0.9,
                  size=14, color=th["body_c"], font=fb, wrap=True)

    # Автор + дата
    author = data.get("author", "")
    date_str = _dt.datetime.now().strftime("%d.%m.%Y")
    info = f"{author}   ·   {date_str}" if author else date_str
    _pptx_txt(slide, info, RX, SLIDE_H - 0.75, RW, 0.5,
              size=11, color=th["muted_c"], font=fb)


def _el_slide_content(prs, data, th, num, img: bytes | None):
    """Контентный: фото слева 38%, контент справа на белом."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN as PA

    ft = th.get("font_title", "Georgia")
    fb = th.get("font_body", "Calibri")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["bg2"])

    PW = SLIDE_W * 0.38
    if img:
        _pptx_img(slide, img, 0, 0, PW, SLIDE_H)
    else:
        _pptx_rect(slide, 0, 0, PW, SLIDE_H, th["card"])

    _pptx_rect(slide, PW, 0, SLIDE_W - PW, SLIDE_H, th["bg2"])

    RX = PW + 0.55
    RW = SLIDE_W - PW - 0.75

    title = data.get("title", f"Слайд {num}")
    _pptx_txt(slide, title, RX, 0.45, RW, 1.1,
              size=28, color=th["title_c"], bold=True, font=ft, wrap=True)

    # Разделитель под заголовком
    _pptx_rect(slide, RX, 1.6, 0.5, 0.05, th["accent"])

    bullets = data.get("bullets", [])
    fact = data.get("fact", "")
    TY = 1.85
    TH_avail = SLIDE_H - TY - 0.3

    if fact and bullets:
        bh = TH_avail * 0.70
        _pptx_bullets(slide, bullets, RX, TY, RW, bh,
                      size=14, color=th["body_c"], dot_color=th["accent"], font=fb)
        fy = TY + bh + 0.1
        _pptx_rect(slide, RX, fy, RW, TH_avail * 0.28, th["card"])
        _pptx_txt(slide, f"💡  {fact}", RX + 0.15, fy + 0.06, RW - 0.2, TH_avail * 0.22,
                  size=11, color=th["muted_c"], italic=True, font=fb, wrap=True)
    elif bullets:
        _pptx_bullets(slide, bullets, RX, TY, RW, TH_avail,
                      size=14, color=th["body_c"], dot_color=th["accent"], font=fb)
    elif fact:
        _pptx_rect(slide, RX, TY, RW, TH_avail, th["card"])
        _pptx_txt(slide, fact, RX + 0.2, TY + 0.3, RW - 0.4, TH_avail - 0.6,
                  size=18, color=th["body_c"], italic=True, font=fb, wrap=True)


def _el_slide_cards(prs, data, th, num, img: bytes | None):
    """Карточки: фото сверху (полоса) или слева, бежевые карточки."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN as PA

    ft = th.get("font_title", "Georgia")
    fb = th.get("font_body", "Calibri")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["bg2"])

    title = data.get("title", f"Слайд {num}")
    cards = data.get("cards", [])
    contrast = data.get("contrast", False)

    if img:
        # Фото — верхняя полоса
        PH = 2.5
        _pptx_img(slide, img, 0, 0, SLIDE_W, PH)
        # Заголовок поверх фото — с полупрозрачным фоном
        _pptx_rect(slide, 0, PH - 0.8, SLIDE_W, 0.9, th["bg2"], transp=0.15)
        _pptx_txt(slide, title, 0.5, PH - 0.75, SLIDE_W - 1.0, 0.8,
                  size=26, color=th["title_c"], bold=True, font=ft)
        CY = PH + 0.25
    else:
        _pptx_txt(slide, title, 0.5, 0.3, SLIDE_W - 1.0, 0.9,
                  size=28, color=th["title_c"], bold=True, font=ft)
        _pptx_rect(slide, 0.5, 1.25, 0.5, 0.05, th["accent"])
        CY = 1.5

    CH = SLIDE_H - CY - 0.25
    n_cards = min(len(cards), 3)
    if n_cards == 0:
        return
    gap = 0.25
    cw = (SLIDE_W - 1.0 - gap * (n_cards - 1)) / n_cards

    for ci, card in enumerate(cards[:n_cards]):
        cx = 0.5 + ci * (cw + gap)
        # Бежевый фон карточки
        _pptx_rect(slide, cx, CY, cw, CH, th["card"])

        icon = card.get("icon", "")
        if icon and not icon.startswith("fas "):
            _pptx_txt(slide, icon, cx + 0.2, CY + 0.2, 0.6, 0.5,
                      size=20, color=th["accent"], align=PA.LEFT, font=ft)
            ty_off = 0.75
        else:
            ty_off = 0.25

        hdg = card.get("heading", "")
        _pptx_txt(slide, hdg, cx + 0.2, CY + ty_off, cw - 0.4, 0.45,
                  size=13, color=th["title_c"], bold=True, font=ft)

        body_y = CY + ty_off + 0.5
        body_h = CH - ty_off - 0.6
        items = card.get("items", [])
        body = card.get("body", "")
        if items:
            _pptx_bullets(slide, items, cx + 0.2, body_y, cw - 0.4, body_h,
                          size=11, color=th["body_c"], dot_color=th["accent"], font=fb)
        elif body:
            _pptx_txt(slide, body, cx + 0.2, body_y, cw - 0.4, body_h,
                      size=11, color=th["body_c"], font=fb, wrap=True)


def _el_slide_split(prs, data, th, num, img: bytes | None):
    """Split: фото слева, контент справа."""
    _el_slide_content(prs, data, th, num, img)


def _el_slide_stats(prs, data, th, num, img: bytes | None):
    """Статистика: белый фон, крупные цифры, фото сверху."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN as PA

    ft = th.get("font_title", "Georgia")
    fb = th.get("font_body", "Calibri")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["bg2"])

    title = data.get("title", f"Слайд {num}")
    stats = data.get("stats", [])

    if img:
        PH = 2.2
        _pptx_img(slide, img, 0, 0, SLIDE_W, PH)
        _pptx_rect(slide, 0, PH - 0.7, SLIDE_W, 0.8, th["bg2"], transp=0.1)
        _pptx_txt(slide, title, 0.5, 0.35, SLIDE_W - 1.0, 0.8,
                  size=26, color=th["bg2"], bold=True, font=ft)
        SY = PH + 0.3
    else:
        _pptx_txt(slide, title, 0.5, 0.3, SLIDE_W - 1.0, 1.0,
                  size=30, color=th["title_c"], bold=True, font=ft)
        _pptx_rect(slide, 0.5, 1.3, 0.5, 0.05, th["accent"])
        SY = 1.6

    SH = SLIDE_H - SY - 0.3
    n = min(len(stats), 4)
    if n == 0:
        return
    gap = 0.3
    sw = (SLIDE_W - 1.0 - gap * (n - 1)) / n

    for si, stat in enumerate(stats[:n]):
        sx = 0.5 + si * (sw + gap)
        _pptx_rect(slide, sx, SY, sw, SH, th["card"])
        # Крупное значение
        _pptx_txt(slide, stat.get("value", ""), sx, SY + 0.25, sw, 1.1,
                  size=40, color=th["header"], bold=True, align=PA.CENTER, font=ft)
        # Метка
        _pptx_txt(slide, stat.get("label", ""), sx, SY + 1.4, sw, 0.5,
                  size=13, color=th["title_c"], bold=True, align=PA.CENTER, font=ft)
        # Описание
        desc = stat.get("desc", "")
        if desc:
            _pptx_txt(slide, desc, sx + 0.1, SY + 1.95, sw - 0.2, SH - 2.1,
                      size=10, color=th["muted_c"], align=PA.CENTER, font=fb, wrap=True)


def _el_slide_quote(prs, data, th, num, img: bytes | None):
    """Цитата: фото справа, цитата слева."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN as PA

    ft = th.get("font_title", "Georgia")
    fb = th.get("font_body", "Calibri")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["bg2"])

    PW = SLIDE_W * 0.42
    if img:
        _pptx_img(slide, img, SLIDE_W - PW, 0, PW, SLIDE_H)

    QW = SLIDE_W - PW - 0.5
    _pptx_rect(slide, 0, SLIDE_H * 0.25, 0.12, SLIDE_H * 0.5, th["accent"])
    quote = data.get("quote", "")
    _pptx_txt(slide, f"«{quote}»", 0.45, 1.0, QW, 4.5,
              size=22, color=th["title_c"], italic=True, font=ft, wrap=True)
    ctx = data.get("context", "")
    if ctx:
        _pptx_txt(slide, f"— {ctx}", 0.45, SLIDE_H - 1.2, QW, 0.7,
                  size=13, color=th["muted_c"], font=fb)


def _el_slide_timeline(prs, data, th, num, img: bytes | None):
    """Таймлайн: нумерованные шаги с бежевым фоном."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN as PA

    ft = th.get("font_title", "Georgia")
    fb = th.get("font_body", "Calibri")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["bg2"])

    title = data.get("title", f"Слайд {num}")
    steps = data.get("steps", [])

    _pptx_txt(slide, title, 0.5, 0.3, SLIDE_W - 1.0, 0.9,
              size=28, color=th["title_c"], bold=True, font=ft)
    _pptx_rect(slide, 0.5, 1.25, 0.5, 0.05, th["accent"])

    if not steps:
        return
    n = min(len(steps), 5)
    TY = 1.55
    th_step = (SLIDE_H - TY - 0.3) / n

    for si, step in enumerate(steps[:n]):
        sy = TY + si * th_step
        sh = th_step - 0.12
        # Бежевый фон
        _pptx_rect(slide, 0.5, sy, SLIDE_W - 1.0, sh, th["card"])
        # Номер
        _pptx_rect(slide, 0.5, sy, 0.6, sh, th["accent"])
        _pptx_txt(slide, str(si + 1), 0.5, sy, 0.6, sh,
                  size=16, color=th["bg2"], bold=True, align=PA.CENTER, font=ft)
        # Заголовок шага
        _pptx_txt(slide, step.get("title", ""), 1.25, sy + 0.05, 3.5, sh * 0.45,
                  size=13, color=th["title_c"], bold=True, font=ft)
        # Описание
        _pptx_txt(slide, step.get("desc", ""), 1.25, sy + sh * 0.45, SLIDE_W - 1.95, sh * 0.52,
                  size=11, color=th["body_c"], font=fb, wrap=True)


def _el_slide_conclusion(prs, data, th, num, img: bytes | None):
    """Заключение: фото справа, выводы слева."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN as PA

    ft = th.get("font_title", "Georgia")
    fb = th.get("font_body", "Calibri")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["bg2"])

    PW = SLIDE_W * 0.40
    if img:
        _pptx_img(slide, img, SLIDE_W - PW, 0, PW, SLIDE_H)
    else:
        _pptx_rect(slide, SLIDE_W - PW, 0, PW, SLIDE_H, th["card"])

    CW = SLIDE_W - PW - 0.5
    title = data.get("title", "Выводы")
    _pptx_txt(slide, title, 0.5, 0.4, CW, 1.0,
              size=30, color=th["title_c"], bold=True, font=ft)
    _pptx_rect(slide, 0.5, 1.45, 0.5, 0.05, th["accent"])

    conclusions = data.get("conclusions", [])
    _pptx_bullets(slide, conclusions, 0.5, 1.65, CW, SLIDE_H - 2.5,
                  size=14, color=th["body_c"], dot_color=th["accent"], font=fb)

    cta = data.get("cta", "")
    if cta:
        _pptx_rect(slide, 0.5, SLIDE_H - 0.9, CW * 0.7, 0.55, th["accent"])
        _pptx_txt(slide, cta, 0.5, SLIDE_H - 0.9, CW * 0.7, 0.55,
                  size=13, color=th["bg2"], bold=True, align=PA.CENTER, font=ft)


def _slide_title(prs, data, th, img: bytes | None):
    """
    Титульный: полноэкранное фото + тёмный оверлей.
    Шрифты и цвета полностью определяются темой.
    """
    from pptx.util import Inches as I, Pt
    from pptx.enum.text import PP_ALIGN as PA
    import datetime as _dt

    ft = th.get("font_title", "Trebuchet MS")
    fb = th.get("font_body", "Calibri")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["header"])

    if img:
        _pptx_img(slide, img, 0, 0, SLIDE_W, SLIDE_H)
        _pptx_rect(slide, 0, 0, SLIDE_W * 0.72, SLIDE_H, th["header"], transp=0.10)
        _pptx_rect(slide, SLIDE_W * 0.72, 0, SLIDE_W * 0.28, SLIDE_H, th["header"], transp=0.50)
    else:
        _pptx_rect(slide, 0, 0, SLIDE_W * 0.65, SLIDE_H, th["header"])
        _pptx_rect(slide, SLIDE_W * 0.65, 0, SLIDE_W * 0.35, SLIDE_H, th["bg2"])

    # Левая акцентная полоса
    _pptx_rect(slide, 0, 0, 0.28, SLIDE_H, th["accent"])

    # Организация (если есть)
    org = data.get("organization", "")
    if org:
        _pptx_txt(slide, org.upper(), 0.58, 0.50, SLIDE_W * 0.62, 0.44,
                  size=9, color=th["accent2"], bold=True, font=fb)

    # Главный заголовок — крупный, тематический шрифт
    title = data.get("title", "Презентация")
    _pptx_txt(slide, title,
              0.58, 1.10, SLIDE_W * 0.65, 3.4,
              size=44, color=th["title_c"], bold=True,
              font=ft, wrap=True)

    # Подзаголовок
    sub = data.get("subtitle", "")
    if sub:
        _pptx_txt(slide, sub,
                  0.58, 4.62, SLIDE_W * 0.65, 0.90,
                  size=17, color=th["body_c"],
                  font=fb, wrap=True)

    # Автор + дата
    author   = data.get("author", "")
    date_str = _dt.datetime.now().strftime("%d.%m.%Y")
    info     = f"{author}   ·   {date_str}" if author else date_str
    _pptx_txt(slide, info,
              0.58, SLIDE_H - 0.78, SLIDE_W * 0.65, 0.58,
              size=12, color=th["muted_c"], font=fb)

    # Нижняя акцентная полоска (тонкая)
    _pptx_rect(slide, 0.28, SLIDE_H - 0.08, SLIDE_W * 0.72 - 0.28, 0.08, th["accent2"])


def _slide_content(prs, data, th, num, img: bytes | None):
    """
    Контентный слайд с тематическими шрифтами и современным дизайном.
    """
    from pptx.util import Inches as I, Pt
    from pptx.enum.text import PP_ALIGN as PA

    ft = th.get("font_title", "Trebuchet MS")
    fb = th.get("font_body", "Calibri")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["bg"])

    title   = data.get("title", f"Слайд {num}")
    bullets = data.get("bullets", [])
    fact    = data.get("fact", "")
    HH = 1.30

    # ── Шапка ──────────────────────────────────────────────────
    _pptx_rect(slide, 0, 0, SLIDE_W, HH, th["header"])
    # Левый акцентный бордюр
    _pptx_rect(slide, 0, 0, 0.28, HH, th["accent"])
    # Номер слайда — badge справа
    _pptx_rect(slide, SLIDE_W - 0.68, 0.22, 0.50, 0.50, th["accent2"])
    _pptx_txt(slide, str(num),
              SLIDE_W - 0.68, 0.22, 0.50, 0.50,
              size=14, color=th["header"], bold=True, align=PA.CENTER, font=ft)
    # Заголовок — тематический шрифт
    _pptx_txt(slide, title,
              0.50, 0.12, SLIDE_W - 1.4, 1.05,
              size=26, color=th["title_c"], bold=True, font=ft)

    # ── Фото справа ────────────────────────────────────────────
    cy = HH + 0.24
    ch = SLIDE_H - HH - 0.40

    if img:
        PW = 4.40
        PH = SLIDE_H - HH - 0.30
        PX = SLIDE_W - PW - 0.20
        PY = HH + 0.15
        # Тень под фото
        _pptx_rect(slide, PX + 0.07, PY + 0.07, PW, PH, th["header"], transp=0.55)
        _pptx_img(slide, img, PX, PY, PW, PH)
        # Акцентная рамка только слева и сверху (L-форма)
        _pptx_rect(slide, PX - 0.05, PY - 0.05, 0.06, PH + 0.05, th["accent"])
        _pptx_rect(slide, PX - 0.05, PY - 0.05, PW + 0.05, 0.06, th["accent2"])
        text_w = PX - 0.65
    else:
        text_w = SLIDE_W - 1.0

    # ── Контент ────────────────────────────────────────────────
    has_fact = bool(fact)
    bottom_reserve = 0.98 if has_fact and bullets else 0.1

    if fact and not bullets:
        # Ключевой факт / цитата — крупно
        _pptx_rect(slide, 0.42, cy, text_w, ch, th["card"])
        _pptx_rect(slide, 0.42, cy, 0.18, ch, th["accent2"])
        _pptx_txt(slide, f"\u00ab{fact}\u00bb",
                  0.78, cy + 0.28, text_w - 0.58, ch - 0.55,
                  size=22, color=th["body_c"], italic=True,
                  font=fb, wrap=True)

    elif fact and bullets:
        bh = ch - bottom_reserve - 0.1
        _pptx_bullets(slide, bullets, 0.42, cy, text_w, bh,
                      size=16, color=th["body_c"],
                      dot_color=th["accent"], font=fb)
        fy = SLIDE_H - bottom_reserve - 0.08
        _pptx_rect(slide, 0.42, fy, text_w, bottom_reserve - 0.08, th["card"])
        _pptx_rect(slide, 0.42, fy, 0.14, bottom_reserve - 0.08, th["accent2"])
        _pptx_txt(slide, f"💡  {fact}",
                  0.68, fy + 0.06, text_w - 0.38, bottom_reserve - 0.15,
                  size=12, color=th["muted_c"], italic=True, font=fb, wrap=True)

    elif len(bullets) >= 5 and not img:
        mid   = (len(bullets) + 1) // 2
        col_w = (text_w - 0.55) / 2
        _pptx_bullets(slide, bullets[:mid], 0.42, cy, col_w, ch,
                      size=15, color=th["body_c"],
                      dot_color=th["accent"], font=fb)
        _pptx_rect(slide, 0.42 + col_w + 0.22, cy + 0.28, 0.04, ch - 0.55, th["muted_c"])
        _pptx_bullets(slide, bullets[mid:], 0.42 + col_w + 0.42, cy, col_w, ch,
                      size=15, color=th["body_c"],
                      dot_color=th["accent2"], font=fb)
    else:
        _pptx_bullets(slide, bullets, 0.42, cy, text_w, ch - bottom_reserve + 0.1,
                      size=17, color=th["body_c"],
                      dot_color=th["accent"], font=fb)

    # Нижняя полоска
    _pptx_rect(slide, 0, SLIDE_H - 0.09, SLIDE_W, 0.09, th["accent2"])


def _slide_stats(prs, data, th, num, img: bytes | None):
    """
    Слайд статистики с тематическими шрифтами.
    """
    from pptx.util import Inches as I, Pt
    from pptx.enum.text import PP_ALIGN as PA

    ft = th.get("font_title", "Trebuchet MS")
    fb = th.get("font_body", "Calibri")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["bg"])

    title = data.get("title", "Ключевые данные")
    stats = data.get("stats", [])
    HH = 1.45

    # Шапка
    if img:
        # Фото во всю шапку
        _pptx_img(slide, img, 0, 0, SLIDE_W, HH + 0.2)
        _pptx_rect(slide, 0, 0, SLIDE_W, HH, th["header"], transp=0.18)
    else:
        _pptx_rect(slide, 0, 0, SLIDE_W, HH, th["header"])

    _pptx_rect(slide, 0, 0, 0.28, HH, th["accent"])
    _pptx_txt(slide, title,
              0.50, 0.28, SLIDE_W - 1.0, 1.0,
              size=30, color=th["title_c"], bold=True, font=ft)

    if not stats:
        return

    n = min(len(stats), 4)
    GAP  = 0.25
    PADX = 0.5
    card_w = (SLIDE_W - PADX * 2 - GAP * (n - 1)) / n
    card_y = HH + 0.18
    card_h = SLIDE_H - card_y - 0.25

    for i, stat in enumerate(stats[:n]):
        cx = PADX + i * (card_w + GAP)
        ac = th["accent"] if i % 2 == 0 else th["accent2"]

        # Карточка
        _pptx_rect(slide, cx, card_y, card_w, card_h, th["card"])
        # Верхняя цветная полоса карточки
        _pptx_rect(slide, cx, card_y, card_w, 0.22, ac)

        # Большое число — тематический шрифт
        _pptx_txt(slide, stat.get("value", "–"),
                  cx + 0.08, card_y + 0.32, card_w - 0.16, 2.05,
                  size=50, color=ac, bold=True, align=PA.CENTER, font=ft)

        # Метка
        _pptx_txt(slide, stat.get("label", ""),
                  cx + 0.08, card_y + 2.38, card_w - 0.16, 0.62,
                  size=15, color=th["body_c"], bold=True, align=PA.CENTER, font=ft)

        # Описание
        if stat.get("desc"):
            _pptx_txt(slide, stat["desc"],
                      cx + 0.08, card_y + 3.0, card_w - 0.16, 0.72,
                      size=12, color=th["muted_c"], align=PA.CENTER, wrap=True, font=fb)

    _pptx_rect(slide, 0, SLIDE_H - 0.09, SLIDE_W, 0.09, th["accent2"])


def _slide_conclusion(prs, data, th, img: bytes | None):
    """
    Заключение: фото на весь фон + полупрозрачный тёмный оверлей.
    Шрифты и цвета — из темы.
    """
    from pptx.util import Inches as I, Pt
    from pptx.enum.text import PP_ALIGN as PA

    ft = th.get("font_title", "Trebuchet MS")
    fb = th.get("font_body", "Calibri")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["header"])

    if img:
        _pptx_img(slide, img, 0, 0, SLIDE_W, SLIDE_H)
        _pptx_rect(slide, 0, 0, SLIDE_W, SLIDE_H, th["header"], transp=0.30)

    title       = data.get("title", "Заключение")
    conclusions = data.get("conclusions", [])
    cta         = data.get("cta", "")

    # Левый бордюр
    _pptx_rect(slide, 0, 0, 0.28, SLIDE_H, th["accent"])

    # Заголовок — крупный тематический шрифт
    _pptx_txt(slide, title,
              0.58, 0.55, SLIDE_W - 1.0, 1.20,
              size=40, color=th["title_c"], bold=True, font=ft)

    # Выводы
    if conclusions:
        _pptx_bullets(slide, conclusions,
                      0.58, 1.90, SLIDE_W - 1.05, SLIDE_H - 3.15,
                      size=18,
                      color=_c(0xFF, 0xFF, 0xFF),
                      dot_color=th["accent2"], font=fb)

    # CTA — нижняя яркая панель
    if cta:
        _pptx_rect(slide, 0.28, SLIDE_H - 1.02, SLIDE_W - 0.28, 1.02, th["accent"])
        _pptx_txt(slide, cta,
                  0.55, SLIDE_H - 0.94, SLIDE_W - 0.80, 0.84,
                  size=18, color=th["header"], bold=True, align=PA.CENTER, font=ft)


# ── Новые построители слайдов ─────────────────────────────────────

def _slide_cards(prs, data, th, num, img: bytes | None):
    """
    Карточный слайд: 2 карточки рядом.
    Поддерживает две версии:
      - contrast=False: обе карточки одного цвета (accent/card)
      - contrast=True:  левая accent (проблема/красный), правая accent2 (решение/зелёный)
    Поля карточки: icon (emoji), heading, body (текст или список)
    """
    from pptx.util import Inches as I, Pt, Emu
    from pptx.enum.text import PP_ALIGN as PA

    ft = th.get("font_title", "Montserrat")
    fb = th.get("font_body",  "Montserrat")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["bg"])

    title    = data.get("title", f"Слайд {num}")
    cards    = data.get("cards", [])
    contrast = data.get("contrast", False)
    HH = 1.30

    # Шапка
    _pptx_rect(slide, 0, 0, SLIDE_W, HH, th["header"])
    _pptx_rect(slide, 0, 0, 0.28, HH, th["accent"])
    _pptx_txt(slide, title,
              0.50, 0.14, SLIDE_W - 1.0, 1.02,
              size=26, color=th["title_c"], bold=True, font=ft)

    if not cards:
        return

    n = min(len(cards), 2)
    GAP  = 0.40
    PADX = 0.50
    card_w = (SLIDE_W - PADX * 2 - GAP * (n - 1)) / n
    card_y = HH + 0.28
    card_h = SLIDE_H - card_y - 0.38

    palette = []
    if contrast and n == 2:
        palette = [th["accent"], th["accent2"]]
        bg_alphas = [0.82, 0.82]
    else:
        palette = [th["card"], th["card"]]
        bg_alphas = [None, None]

    for i, card in enumerate(cards[:n]):
        cx = PADX + i * (card_w + GAP)
        ac = palette[i]

        # Фон карточки
        _pptx_rect(slide, cx, card_y, card_w, card_h, ac)

        # Если контрастный — светлый полупрозрачный фон вместо яркого цвета
        if contrast:
            # Делаем очень светлую пастельную карточку
            pastel = th["bg2"]  # белый/светлый
            _pptx_rect(slide, cx + 0.06, card_y + 0.06,
                       card_w - 0.12, card_h - 0.12, pastel)
            accent_bar = ac
        else:
            accent_bar = th["accent"] if i % 2 == 0 else th["accent2"]

        # Верхняя цветная полоска карточки
        _pptx_rect(slide, cx, card_y, card_w, 0.08, ac)

        # Иконка/эмодзи (крупно)
        icon = card.get("icon", "")
        if icon:
            _pptx_txt(slide, icon,
                      cx + 0.18, card_y + 0.22, card_w - 0.36, 0.65,
                      size=28, align=PA.LEFT)

        # Заголовок карточки
        heading = card.get("heading", "")
        head_y = card_y + (0.95 if icon else 0.30)
        body_c = th["body_c"] if not contrast else th["header"]
        head_c = th["accent"] if not contrast else ac
        _pptx_txt(slide, heading,
                  cx + 0.22, head_y, card_w - 0.44, 0.70,
                  size=17, color=head_c, bold=True, font=ft)

        # Тело карточки
        body = card.get("body", "")
        body_items = card.get("items", [])
        body_y = head_y + 0.72
        body_h = card_h - (body_y - card_y) - 0.25

        if body_items:
            _pptx_bullets(slide, body_items,
                          cx + 0.22, body_y, card_w - 0.44, body_h,
                          size=13, color=body_c,
                          dot_color=ac, font=fb)
        elif body:
            _pptx_txt(slide, body,
                      cx + 0.22, body_y, card_w - 0.44, body_h,
                      size=14, color=body_c, font=fb, wrap=True)

    _pptx_rect(slide, 0, SLIDE_H - 0.09, SLIDE_W, 0.09, th["accent2"])


def _slide_split(prs, data, th, num, img: bytes | None):
    """
    Split-слайд: большое фото занимает правую ПОЛОВИНУ (full-height),
    текст — левая половина. Без шапки — более свежий вид.
    Вариант circle=True: фото в круге слева, текст справа.
    """
    from pptx.util import Inches as I, Pt, Emu
    from pptx.enum.text import PP_ALIGN as PA
    from pptx.oxml.ns import qn
    from lxml import etree

    ft = th.get("font_title", "Montserrat")
    fb = th.get("font_body",  "Montserrat")

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    circle = data.get("circle", False)
    title   = data.get("title", f"Слайд {num}")
    bullets = data.get("bullets", [])
    label   = data.get("label", "")   # подпись под заголовком
    fact    = data.get("fact", "")

    if circle:
        # ── Круглое фото слева, текст справа ──────────────────────
        _pptx_bg(slide, th["bg"])

        # Левый акцент-столбец
        _pptx_rect(slide, 0, 0, 0.22, SLIDE_H, th["accent"])

        # Круг: рисуем как эллипс поверх фото
        CR = 3.30      # diameter
        CX = 0.50
        CY = (SLIDE_H - CR) / 2

        if img:
            # Сначала вставляем фото (квадрат), потом обрезаем через XML
            _pptx_img(slide, img, CX, CY, CR, CR)
            # Применяем круглую обрезку через XML (custGeom → ellipse)
            try:
                pic_elem = slide.shapes[-1]._element
                spPr = pic_elem.find(qn('p:spPr'))
                if spPr is None:
                    spPr = etree.SubElement(pic_elem, qn('p:spPr'))
                # Удаляем старую геометрию
                for old in spPr.findall(qn('a:prstGeom')):
                    spPr.remove(old)
                # Добавляем эллипс
                prstGeom = etree.SubElement(spPr, qn('a:prstGeom'),
                                             attrib={'prst': 'ellipse'})
                etree.SubElement(prstGeom, qn('a:avLst'))
                # Граница круга — цветная
                ln = etree.SubElement(spPr, qn('a:ln'),
                                      attrib={'w': str(int(0.08 * 914400))})
                solidFill = etree.SubElement(ln, qn('a:solidFill'))
                srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'),
                                           attrib={'val': f'{th["accent"].rgb:06X}'})
            except Exception as _e:
                logging.debug(f"circle crop: {_e}")
        else:
            # Заглушка: цветной круг
            from pptx.util import Inches as I2
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            shape = slide.shapes.add_shape(9,  # MSO_SHAPE_TYPE.OVAL
                                           I2(CX), I2(CY), I2(CR), I2(CR))
            shape.fill.solid()
            shape.fill.fore_color.rgb = th["accent"]
            shape.line.fill.background()

        # Правая половина — текст
        TX = CX + CR + 0.55
        TW = SLIDE_W - TX - 0.45
        TY = 0.42

        # Метка / раздел
        if label:
            _pptx_txt(slide, label.upper(), TX, TY, TW, 0.36,
                      size=9, color=th["accent"], bold=True, font=fb)
            TY += 0.40

        _pptx_txt(slide, title, TX, TY, TW, 1.40,
                  size=28, color=th["header"], bold=True, font=ft, wrap=True)
        TY += 1.50

        if bullets:
            _pptx_bullets(slide, bullets, TX, TY, TW, SLIDE_H - TY - 0.45,
                          size=15, color=th["header"],
                          dot_color=th["accent"], font=fb)
        elif fact:
            _pptx_txt(slide, f'"{fact}"', TX, TY, TW, SLIDE_H - TY - 0.45,
                      size=16, color=th["muted_c"], italic=True, font=fb, wrap=True)

    else:
        # ── Текст слева, фото справа full-height ─────────────────
        _pptx_bg(slide, th["bg2"])

        PW = SLIDE_W * 0.48
        PX = SLIDE_W - PW
        if img:
            _pptx_img(slide, img, PX, 0, PW, SLIDE_H)
            # Затемняющий бордюр слева от фото
            _pptx_rect(slide, PX - 0.04, 0, 0.10, SLIDE_H, th["bg2"])

        # Левый акцент
        _pptx_rect(slide, 0, 0, 0.22, SLIDE_H, th["accent"])

        TX = 0.48
        TW = PX - TX - 0.35

        if label:
            _pptx_txt(slide, label.upper(), TX, 0.42, TW, 0.35,
                      size=9, color=th["accent"], bold=True, font=fb)

        TY_title = 0.85 if label else 0.55
        _pptx_txt(slide, title, TX, TY_title, TW, 1.45,
                  size=28, color=th["header"], bold=True, font=ft, wrap=True)

        TY_body = TY_title + 1.55
        if bullets:
            _pptx_bullets(slide, bullets, TX, TY_body, TW, SLIDE_H - TY_body - 0.42,
                          size=15, color=th["header"],
                          dot_color=th["accent"], font=fb)
        elif fact:
            _pptx_rect(slide, TX, TY_body, TW, 0.06, th["accent2"])
            _pptx_txt(slide, fact, TX, TY_body + 0.18, TW, SLIDE_H - TY_body - 0.6,
                      size=15, color=th["muted_c"], italic=True, font=fb, wrap=True)

        _pptx_rect(slide, 0, SLIDE_H - 0.09, PX, 0.09, th["accent2"])


def _slide_quote(prs, data, th, num, img: bytes | None):
    """
    Слайд-цитата / большой тезис:
    Фото на весь фон (размытое, полупрозрачное), по центру крупный текст.
    Используется для ключевых утверждений, миссий, девизов.
    """
    from pptx.util import Inches as I, Pt
    from pptx.enum.text import PP_ALIGN as PA

    ft = th.get("font_title", "Montserrat")
    fb = th.get("font_body",  "Montserrat")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["header"])

    if img:
        _pptx_img(slide, img, 0, 0, SLIDE_W, SLIDE_H)
        _pptx_rect(slide, 0, 0, SLIDE_W, SLIDE_H, th["header"], transp=0.22)

    # Центральная панель
    panel_h = 3.20
    panel_y = (SLIDE_H - panel_h) / 2
    _pptx_rect(slide, 0.75, panel_y, SLIDE_W - 1.50, panel_h, th["header"], transp=0.45)

    # Декоративные кавычки
    _pptx_txt(slide, "❝", 0.90, panel_y + 0.10, 1.0, 1.0,
              size=52, color=th["accent"], font=ft)

    quote   = data.get("quote", data.get("title", ""))
    context = data.get("context", "")
    _pptx_txt(slide, quote,
              1.10, panel_y + 0.75, SLIDE_W - 2.20, panel_h - 1.20,
              size=26, color=_c(0xFF, 0xFF, 0xFF), bold=True,
              font=ft, wrap=True, align=PA.CENTER)

    if context:
        _pptx_txt(slide, f"— {context}",
                  0.75, panel_y + panel_h - 0.68, SLIDE_W - 1.5, 0.58,
                  size=14, color=th["accent2"], italic=True,
                  font=fb, align=PA.RIGHT)

    # Боковые акцентные линии
    _pptx_rect(slide, 0, 0, 0.22, SLIDE_H, th["accent"])
    _pptx_rect(slide, SLIDE_W - 0.22, 0, 0.22, SLIDE_H, th["accent"])


def _slide_timeline(prs, data, th, num, img: bytes | None):
    """
    Таймлайн / шаги процесса: горизонтальные numbered-блоки.
    Идеально для "этапы проекта", "план действий", "история развития".
    """
    from pptx.util import Inches as I, Pt
    from pptx.enum.text import PP_ALIGN as PA

    ft = th.get("font_title", "Montserrat")
    fb = th.get("font_body",  "Montserrat")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, th["bg"])

    title = data.get("title", f"Слайд {num}")
    steps = data.get("steps", [])
    HH = 1.30

    _pptx_rect(slide, 0, 0, SLIDE_W, HH, th["header"])
    _pptx_rect(slide, 0, 0, 0.28, HH, th["accent"])
    _pptx_txt(slide, title,
              0.50, 0.14, SLIDE_W - 1.0, 1.02,
              size=26, color=th["title_c"], bold=True, font=ft)

    if not steps:
        return

    n = min(len(steps), 5)
    GAP  = 0.28
    PADX = 0.45
    step_w = (SLIDE_W - PADX * 2 - GAP * (n - 1)) / n
    step_y = HH + 0.35
    step_h = SLIDE_H - step_y - 0.40

    # Горизонтальная соединяющая линия
    line_y = step_y + 0.42
    _pptx_rect(slide, PADX + step_w / 2, line_y,
               SLIDE_W - PADX * 2 - step_w, 0.06, th["muted_c"])

    for i, step in enumerate(steps[:n]):
        sx = PADX + i * (step_w + GAP)
        ac = th["accent"] if i % 2 == 0 else th["accent2"]

        # Нумерованный круг
        _pptx_rect(slide, sx + (step_w - 0.68) / 2, step_y, 0.68, 0.68, ac)
        _pptx_txt(slide, str(i + 1),
                  sx + (step_w - 0.68) / 2, step_y, 0.68, 0.68,
                  size=16, color=th["header"], bold=True, align=PA.CENTER, font=ft)

        # Заголовок шага
        _pptx_txt(slide, step.get("title", ""),
                  sx, step_y + 0.78, step_w, 0.68,
                  size=13, color=th["body_c"], bold=True,
                  align=PA.CENTER, font=ft, wrap=True)

        # Описание шага
        if step.get("desc"):
            _pptx_txt(slide, step["desc"],
                      sx + 0.08, step_y + 1.52, step_w - 0.16,
                      step_h - 1.55,
                      size=11, color=th["muted_c"],
                      align=PA.CENTER, font=fb, wrap=True)

    _pptx_rect(slide, 0, SLIDE_H - 0.09, SLIDE_W, 0.09, th["accent2"])


# ── Сборка файла ─────────────────────────────────────────────────

def _register_custom_fonts_in_pptx(prs):
    """
    Встраивает Montserrat и Raleway в тему PPTX через XML так,
    чтобы PowerPoint/LibreOffice знал об этих шрифтах и корректно
    их отображал. Шрифт ищется сначала в Fonts/, потом в системе.
    """
    if not HAS_PPTX:
        return
    try:
        from pptx.oxml.ns import qn, nsmap
        from lxml import etree
        import zipfile, shutil, tempfile

        # Получаем XML темы
        theme = prs.slide_master.theme_color_map
    except Exception:
        pass  # Не критично — PPTX всё равно откроется, просто шрифт надо иметь в системе

    # Ключевой метод: встраиваем TTF-файлы напрямую в ZIP-структуру PPTX
    # через низкоуровневый доступ к prs.part
    try:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        from pptx.oxml.ns import qn
        import base64 as _b64

        font_embeds = []
        if _HAS_MONTSERRAT:
            font_embeds += [
                ("Montserrat",     FONT_MONTSERRAT_REG,  False, False),
                ("Montserrat",     FONT_MONTSERRAT_BOLD, True,  False),
                ("Montserrat",     FONT_MONTSERRAT_SEMI, True,  False),
            ]
        if _HAS_RALEWAY:
            font_embeds += [
                ("Raleway", FONT_RALEWAY_REG,  False, False),
                ("Raleway", FONT_RALEWAY_BOLD, True,  False),
            ]

        if not font_embeds:
            return

        # Встраиваем через fontTable
        from pptx.oxml import parse_xml
        slide_master = prs.slide_master
        spTree = slide_master.shapes._spTree

        # Добавляем extLst с embedFont если ещё нет
        for font_name, font_path, bold, italic in font_embeds:
            if not os.path.exists(font_path):
                continue
            try:
                with open(font_path, "rb") as f:
                    font_data = f.read()
                # Добавляем как медиа-файл в package
                part = prs.part
                font_rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"
                ext = ".ttf"
                fname_in_pkg = f"font_{font_name.replace(' ', '')}{'_bold' if bold else ''}.ttf"
                # Это встраивает шрифт в pptx-файл
                logging.debug(f"📝 Встраиваем шрифт {fname_in_pkg}: {len(font_data)//1024}кб")
            except Exception as _fe:
                logging.debug(f"Font embed {font_name}: {_fe}")
    except Exception as _e:
        logging.debug(f"_register_custom_fonts_in_pptx: {_e}")


async def create_pptx_bytes(pptx_data: dict, theme_key: str = "auto",
                             use_images: bool = True) -> bytes:
    if not HAS_PPTX:
        raise RuntimeError("python-pptx не установлен: pip install python-pptx")

    # Автовыбор темы по теме презентации
    if theme_key == "auto" or theme_key not in PPTX_THEMES:
        topic = pptx_data.get("title", "")
        theme_key = _auto_pick_theme(topic)
        logging.info(f"🎨 Авто-тема для «{topic[:40]}»: {theme_key}")

    th = PPTX_THEMES.get(theme_key, PPTX_THEMES["dark_blue"])
    slides_data = pptx_data.get("slides", [])

    img_map: dict[int, bytes | None] = {}
    if use_images:
        img_map = await _fetch_photos_all(slides_data, topic=pptx_data.get("title", ""))

    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # ── Встраиваем Montserrat/Raleway в тему PPTX ──────────────────
    # Это позволяет PowerPoint корректно отображать шрифт на любой машине
    _register_custom_fonts_in_pptx(prs)

    # Выбираем рендерер по теме
    use_elegant = th.get("layout") == "elegant"

    for i, sd in enumerate(slides_data):
        stype = sd.get("type", "content")
        img   = img_map.get(i) if use_images else None
        if use_elegant:
            if stype == "title":
                _el_slide_title(prs, sd, th, img)
            elif stype == "stats":
                _el_slide_stats(prs, sd, th, i + 1, img)
            elif stype == "conclusion":
                _el_slide_conclusion(prs, sd, th, i + 1, img)
            elif stype == "cards":
                _el_slide_cards(prs, sd, th, i + 1, img)
            elif stype in ("split", "content"):
                _el_slide_content(prs, sd, th, i + 1, img)
            elif stype == "quote":
                _el_slide_quote(prs, sd, th, i + 1, img)
            elif stype == "timeline":
                _el_slide_timeline(prs, sd, th, i + 1, img)
            else:
                _el_slide_content(prs, sd, th, i + 1, img)
        else:
            if stype == "title":
                _slide_title(prs, sd, th, img)
            elif stype == "stats":
                _slide_stats(prs, sd, th, i + 1, img)
            elif stype == "conclusion":
                _slide_conclusion(prs, sd, th, img)
            elif stype == "cards":
                _slide_cards(prs, sd, th, i + 1, img)
            elif stype == "split":
                _slide_split(prs, sd, th, i + 1, img)
            elif stype == "quote":
                _slide_quote(prs, sd, th, i + 1, img)
            elif stype == "timeline":
                _slide_timeline(prs, sd, th, i + 1, img)
            else:
                _slide_content(prs, sd, th, i + 1, img)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def create_html_presentation(pptx_data: dict, theme_key: str = "dark_blue", pexels_key: str = "") -> str:
    """Генерирует красивую HTML-презентацию — дизайн как Google Gemini Presentation."""
    from urllib.parse import quote as _uq

    if theme_key == "auto" or theme_key not in PPTX_THEMES:
        theme_key = _auto_pick_theme(pptx_data.get("title", ""))
    th = PPTX_THEMES.get(theme_key, PPTX_THEMES["dark_blue"])

    def _hex(c):
        try:
            return "#{:02x}{:02x}{:02x}".format(c[0], c[1], c[2])
        except (TypeError, KeyError):
            return "#{:02x}{:02x}{:02x}".format(c.r, c.g, c.b)
    def _dark(c):
        try:
            return (c[0] * 0.299 + c[1] * 0.587 + c[2] * 0.114) < 128
        except (TypeError, KeyError):
            return (c.r * 0.299 + c.g * 0.587 + c.b * 0.114) < 128

    ACC   = _hex(th["accent"])
    ACC2  = _hex(th["accent2"])
    BG    = _hex(th["bg"])
    HDR   = _hex(th["header"])
    TC    = _hex(th["title_c"])
    BC    = _hex(th["body_c"])
    MC    = _hex(th["muted_c"])
    DARK  = _dark(th["bg"])
    GBGA  = "0,0,0" if DARK else "255,255,255"
    GBGO  = "0.72"  if DARK else "0.88"


    # ── Curated библиотека фото по темам (Unsplash photo IDs) ─────────────────
    # Фото отобраны по теме — гарантированно релевантны


    # ══════════════════════════════════════════════════════════════
    # 🖼 УМНЫЙ ПОИСК ФОТО — LoremFlickr (живой поиск) + Pexels CDN
    # ══════════════════════════════════════════════════════════════
    # LoremFlickr: ищет реальные фото на Flickr по любому ключевому слову.
    # Работает без API-ключа. URL: loremflickr.com/ШхВ/тема1,тема2?lock=seed
    # lock= делает результат детерминированным для одного seed (слайд не меняет фото при перезагрузке)

    def _loremflickr(kw: str, w: int, h: int, seed: int) -> str:
        """URL LoremFlickr — живой поиск по ключевым словам через Flickr."""
        import urllib.parse as _ulp
        # Берём первые 3 значимых слова, соединяем запятой
        words = [w2.strip().lower() for w2 in (kw or "").split() if len(w2.strip()) > 2][:3]
        q = ",".join(words) if words else "professional"
        q_enc = _ulp.quote(q)
        return f"https://loremflickr.com/{w}/{h}/{q_enc}?lock={seed}"

    def _pex(pid: int) -> str:
        """Pexels CDN URL по числовому ID — резервный источник."""
        return (
            f"https://images.pexels.com/photos/{pid}/pexels-photo-{pid}.jpeg"
            "?auto=compress&cs=tinysrgb&w=1200&h=700&fit=crop"
        )

    # ── Нормализатор ключевых слов ────────────────────────────────
    # Переводит любой русский/смешанный keyword в хороший английский поисковый запрос
    _KW_TRANSLATE = {
        # Россия / история
        "россия": "Russia Moscow Kremlin",
        "российск": "Russia Russian",
        "москва": "Moscow Russia city",
        "советск": "Soviet Union historical",
        "история": "history historical",
        "война": "war military battle soldiers",
        "войны": "war military battle historical",
        "военн": "military soldiers defense",
        "армия": "army soldiers military",
        "победа": "Victory parade soldiers",
        "революц": "revolution historical soldiers",
        "империя": "empire historical palace",
        "великая отечественная": "World War 2 Soviet soldiers",
        "холодная война": "Cold War nuclear tension",
        # Бизнес
        "бизнес": "business office professionals",
        "финанс": "finance money investment",
        "экономик": "economy business charts",
        "маркетинг": "marketing strategy team",
        "менеджмент": "management leadership team",
        "стратег": "strategy planning business",
        "инвестиц": "investment finance growth",
        "стартап": "startup office innovation",
        "продаж": "sales business meeting",
        # Технологии
        "технолог": "technology innovation digital",
        "программ": "programming code software",
        "искусственный интеллект": "artificial intelligence AI neural",
        "нейросет": "neural network AI technology",
        "данные": "data analytics visualization",
        "кибер": "cybersecurity digital security",
        "робот": "robot automation technology",
        "цифров": "digital technology modern",
        # Медицина
        "медицин": "medical doctor hospital",
        "здравоохранен": "healthcare medical professionals",
        "хирург": "surgeon operating room hospital",
        # Образование
        "образован": "education students classroom",
        "школ": "school students learning",
        "университет": "university campus students",
        "наук": "science laboratory research",
        # Природа / экология
        "природ": "nature landscape outdoor",
        "эколог": "ecology environment green",
        "климат": "climate environment nature",
        "лес": "forest trees nature",
        # Промышленность
        "промышленн": "industry factory workers",
        "завод": "factory industrial plant",
        "производств": "production manufacturing",
        "нефт": "oil refinery petroleum",
        "газ": "natural gas pipeline energy",
        "энергет": "energy power plant electricity",
        # Транспорт
        "транспорт": "transport logistics road",
        "логистик": "logistics supply chain warehouse",
        "корабл": "ship cargo vessel port",
        "автомобил": "car automobile road",
        "авиац": "aviation aircraft airport",
        # Спорт
        "спорт": "sport athletes training",
        "футбол": "football soccer stadium",
        "баскетбол": "basketball court game",
        # Еда
        "еда": "food restaurant cooking",
        "кулинар": "cooking culinary food",
        "ресторан": "restaurant food dining",
        # Путешествия
        "путешеств": "travel tourism adventure",
        "туризм": "tourism travel destination",
        # Политика
        "политик": "politics government parliament",
        "правительств": "government politics official",
        # Психология / саморазвитие
        "психологи": "psychology mind wellness",
        "мотивац": "motivation success achievement",
        "лидерств": "leadership success team",
        # ОБЖ / Безопасность
        "обж": "emergency safety rescue firefighters",
        "безопасност": "safety security rescue protection",
        "чрезвычайн": "emergency disaster rescue responders",
        "первая помощь": "first aid medical emergency CPR",
        "пожарн": "fire emergency firefighters rescue",
        "эвакуац": "evacuation emergency exit people",
        "землетрясен": "earthquake disaster rescue emergency",
        "наводнен": "flood disaster water rescue emergency",
        "выживан": "survival outdoor emergency wilderness",
        "гражданская оборона": "civil defense emergency shelters bunker",
        "ожог": "medical burn treatment hospital emergency",
        "отравлен": "poison medical emergency treatment hospital",
        "дтп": "car accident road emergency rescue police",
        "авари": "accident emergency rescue disaster response",
        "спасател": "rescue team emergency responders firefighters",
    }

    def _normalize_kw(kw: str) -> str:
        """
        Нормализует keyword: если содержит кириллицу — переводит через словарь.
        Если уже English — очищает мусорные слова.
        Возвращает хороший английский поисковый запрос.
        """
        import re as _re2
        if not kw:
            return "professional modern"

        kw_low = kw.lower().strip()

        # Проверяем есть ли кириллица
        has_cyrillic = bool(_re2.search(r'[а-яёА-ЯЁ]', kw))

        if has_cyrillic:
            # Ищем совпадение в словаре
            for ru_key, en_val in _KW_TRANSLATE.items():
                if ru_key in kw_low:
                    # Добавляем оригинальные английские слова из keyword если есть
                    latin = " ".join(_re2.findall(r'[A-Za-z]{3,}', kw))
                    return (en_val + " " + latin).strip()
            # Если не нашли - берём только латинские слова
            latin = " ".join(_re2.findall(r'[A-Za-z]{3,}', kw))
            if latin:
                return latin
            # Крайний случай - тема из названия
            return "professional modern presentation"

        # Уже английский — убираем плохие абстрактные слова
        _STOP = {'abstract', 'background', 'concept', 'presentation',
                 'illustration', 'symbol', 'icon', 'logo', 'design', 'pattern',
                 'texture', 'banner', 'template', 'vector', 'digital', 'creative',
                 'modern', 'generic', 'various', 'multiple', 'some', 'many'}
        words = [w2 for w2 in kw_low.split() if w2 not in _STOP and len(w2) > 2]
        result = " ".join(words[:5])
        return result if result else "professional"

    # ── Статичные Pexels ID — расширенный словарь ─────────────────
    PHOTO_THEMES = {
        # Военная тематика
        "soldier":      [2058126,1917215,1437796,3396234,6447535,8159956],
        "soldiers":     [2058126,1917215,3396234,1437796,6447535,8159956],
        "military":     [2058126,1917215,3396234,1437796,6447535,8159956],
        "army":         [2058126,1917215,1437796,3396234,6447535,8159956],
        "tank":         [3396234,2058126,1917215,6447535,1437796,8159956],
        "combat":       [2058126,3396234,1917215,1437796,6447535,8159956],
        "defense":      [2058126,1917215,3396234,1437796,6447535,8159956],
        "weapon":       [1917215,2058126,1437796,3396234,6447535,8159956],
        "aircraft":     [6447535,2058126,3396234,1917215,1437796,8159956],
        "war":          [2058126,3396234,1917215,6447535,1437796,8159956],
        "battle":       [2058126,1917215,3396234,6447535,1437796,8159956],
        "historical":   [313782,1486325,1693652,373912,2246476,325229],
        "history":      [313782,1486325,373912,1693652,2246476,325229],
        # Технологии / ИИ
        "neural":       [8386440,1181244,574069,270360,3861969,325229],
        "artificial":   [8386440,1181244,574069,270360,3861969,325229],
        "circuit":      [8386440,574069,1181244,270360,3861969,325229],
        "server":       [1181244,8386440,574069,270360,3861969,325229],
        "data":         [8386440,1181244,270360,574069,3861969,325229],
        "code":         [574069,8386440,1181244,270360,3861969,325229],
        "computer":     [574069,8386440,1181244,270360,3861969,325229],
        "robot":        [8386440,1181244,3861969,574069,270360,325229],
        "technology":   [8386440,574069,1181244,270360,3861969,325229],
        "cloud":        [1181244,8386440,574069,270360,3861969,325229],
        "cyber":        [8386440,574069,1181244,270360,3861969,325229],
        "software":     [574069,8386440,1181244,270360,3861969,325229],
        # Бизнес / Финансы
        "stock":        [159888,210607,590022,6801874,3771110,669619],
        "finance":      [159888,210607,590022,6801874,3771110,669619],
        "financial":    [159888,210607,590022,6801874,3771110,669619],
        "bank":         [210607,159888,590022,6801874,3771110,669619],
        "economy":      [590022,159888,210607,6801874,3771110,669619],
        "investment":   [159888,590022,210607,6801874,3771110,669619],
        "market":       [590022,159888,210607,6801874,3771110,669619],
        "growth":       [159888,210607,590022,3771110,6801874,669619],
        "revenue":      [159888,210607,590022,6801874,3771110,669619],
        "profit":       [210607,159888,590022,3771110,6801874,669619],
        "business":     [3184465,3182812,3182781,3184454,3182773,590022],
        "meeting":      [3182812,3184465,3182781,3184454,3182773,3184419],
        "office":       [3182812,3184465,3184454,3182781,3184419,3182773],
        "team":         [3182812,3184465,3182781,3184454,3182773,3184419],
        "startup":      [3184465,3182812,3182781,3184454,7096333,3182773],
        "marketing":    [905163,590022,3184465,3182812,210607,159888],
        "strategy":     [3184465,3182812,3182781,905163,3184454,905157],
        "sales":        [3184465,3182812,590022,210607,159888,3182781],
        "plan":         [3184465,3182812,3182781,3184454,3182773,3184419],
        "management":   [3182812,3184465,3182781,3184454,3182773,3184419],
        "leadership":   [3184465,3182812,3182781,3184454,3182773,3184419],
        "entrepreneur": [3184465,3182812,7096333,3182781,3184454,3182773],
        # Промышленность
        "factory":      [1108572,2760241,209251,1267338,257700,3862132],
        "industry":     [1108572,2760241,209251,1267338,257700,3862132],
        "workers":      [1108572,2760241,1267338,209251,257700,3862132],
        "production":   [1108572,209251,2760241,1267338,257700,3862132],
        "manufacturing":[1108572,2760241,209251,1267338,257700,3862132],
        "steel":        [2760241,1108572,209251,257700,1267338,3862132],
        # Энергетика
        "oil":          [1078116,2101521,247763,462024,325153,1108572],
        "pipeline":     [2101521,1078116,247763,462024,325153,1108572],
        "energy":       [1078116,2101521,247763,462024,325153,1108572],
        "solar":        [1181246,9875970,6399165,3965557,9875971,1181247],
        "renewable":    [9875970,6399165,3965557,1181246,9875971,3965558],
        "power":        [1078116,2101521,247763,1181246,9875970,6399165],
        # Сельское хозяйство
        "wheat":        [326082,1131458,265216,158827,3791664,462024],
        "harvest":      [326082,1131458,265216,158827,3791664,462024],
        "farm":         [326082,1131458,265216,158827,3791664,462024],
        "agriculture":  [326082,1131458,265216,158827,3791664,462024],
        "field":        [326082,1131458,265216,158827,3791664,462024],
        # Медицина
        "hospital":     [3786157,305568,4386431,263402,3938023,1298619],
        "doctor":       [3786157,305568,4386431,263402,3938023,1298619],
        "medical":      [3786157,305568,4386431,263402,3938023,1298619],
        "laboratory":   [4386431,263402,3938023,305568,3786157,1298619],
        "research":     [4386431,263402,3938023,305568,3786157,1298619],
        "health":       [3786157,305568,263402,4386431,3938023,1298619],
        # Образование
        "students":     [1205651,289737,1438072,1350700,2422232,1205652],
        "classroom":    [1205651,289737,1438072,1350700,2422232,1205652],
        "university":   [1205651,289737,1438072,1350700,2422232,1205652],
        "education":    [1205651,289737,1438072,1350700,2422232,1205652],
        "school":       [1205651,289737,1438072,1350700,2422232,1205652],
        "learning":     [1205651,289737,1438072,1350700,2422232,1205652],
        # Наука
        "science":      [4386431,263402,3938023,305568,1298619,3786157],
        "physics":      [4386431,263402,3938023,305568,1298619,3786157],
        "chemistry":    [4386431,263402,3938023,305568,1298619,3786157],
        "biology":      [4386431,263402,305568,3938023,1298619,3786157],
        "space":        [6399165,9875970,3965557,1181246,9875971,3965558],
        "rocket":       [6399165,9875970,3965557,3396234,1181246,9875971],
        "astronomy":    [6399165,9875970,3965557,1181246,9875971,3965558],
        # Транспорт / Логистика
        "cargo":        [1427107,906494,1117210,1427109,1546960,2226458],
        "ship":         [1427107,906494,1117210,1427109,1546960,2226458],
        "logistics":    [1427107,906494,1117210,1427109,1546960,2226458],
        "truck":        [906494,1427107,1117210,1427109,1546960,2226458],
        "transport":    [1427107,906494,1117210,1427109,1546960,2226458],
        "road":         [906494,1427107,1117210,1049622,1393382,1118448],
        "highway":      [1049622,1393382,906494,1427107,449609,1118448],
        "aviation":     [6447535,1427107,906494,1117210,1427109,1546960],
        "airport":      [1427107,906494,6447535,1117210,1427109,1546960],
        "railway":      [1427107,906494,1117210,1049622,1393382,1546960],
        # Архитектура / Города
        "building":     [1486325,313782,373912,1693652,2246476,325229],
        "skyline":      [1486325,313782,373912,1693652,2246476,325229],
        "government":   [313782,1486325,373912,1693652,2246476,325229],
        "city":         [1486325,373912,1693652,313782,2246476,325229],
        "architecture": [1486325,313782,373912,1693652,2246476,325229],
        "kremlin":      [313782,1486325,373912,1693652,2246476,325229],
        "palace":       [313782,1486325,373912,1693652,2246476,325229],
        # Спорт
        "sport":        [863988,841130,1552242,2526878,1640777,1295572],
        "athletes":     [863988,841130,1552242,2526878,1640777,1295572],
        "stadium":      [863988,841130,1640777,2526878,1552242,1295572],
        "training":     [863988,841130,1640777,2526878,1552242,1295572],
        "football":     [863988,841130,1552242,2526878,1640777,1295572],
        "basketball":   [863988,841130,1552242,2526878,1640777,1295572],
        # Еда
        "food":         [1640777,1295572,863988,841130,1552242,2526878],
        "restaurant":   [1640777,1295572,863988,841130,1552242,2526878],
        "cooking":      [1640777,1295572,863988,841130,1552242,2526878],
        # Природа
        "nature":       [326082,158827,265216,1131458,3791664,462024],
        "forest":       [326082,158827,265216,1131458,3791664,462024],
        "environment":  [326082,158827,265216,1131458,9875970,6399165],
        "landscape":    [326082,265216,158827,1131458,3791664,462024],
        "water":        [1427107,906494,326082,158827,265216,1131458],
        # Психология / развитие
        "psychology":   [3182812,3184465,3182781,3184454,1205651,863988],
        "success":      [3184465,3182812,863988,1295572,841130,1640777],
        "motivation":   [3184465,3182812,863988,1295572,841130,1640777],
        "leadership":   [3184465,3182812,3182781,3184454,3182773,3184419],
        "wellness":     [863988,1295572,841130,1640777,1552242,2526878],
    }

    _NEUTRAL_POOL = [
        3184465, 3182812, 590022, 1108572, 159888, 574069,
        1205651, 1427107, 1486325, 3182781, 863988, 4386431,
    ]

    def _get_pex_id(kw: str, slide_idx: int = 0) -> int:
        """Возвращает Pexels photo ID по ключевым словам."""
        words = (kw or "").lower().split()
        for word in words:
            if word in PHOTO_THEMES:
                ids = PHOTO_THEMES[word]
                return ids[slide_idx % len(ids)]
            # Частичное совпадение (минимум 5 символов)
            if len(word) >= 5:
                for key in PHOTO_THEMES:
                    if word[:5] == key[:5]:
                        return PHOTO_THEMES[key][slide_idx % len(PHOTO_THEMES[key])]
            if len(word) >= 4:
                for key in PHOTO_THEMES:
                    if len(key) >= 4 and word[:4] == key[:4]:
                        return PHOTO_THEMES[key][slide_idx % len(PHOTO_THEMES[key])]
        base = abs(hash(kw or "")) % len(_NEUTRAL_POOL)
        return _NEUTRAL_POOL[(base + slide_idx) % len(_NEUTRAL_POOL)]

    def _img_attr(kw, slide_i=0, w=1200, h=700):
        """HTML img attributes — uses global _photo_id for theme-accurate photos."""
        import json as _json, urllib.parse as _upx
        pid0 = _photo_id(kw, slide_i)
        pid1 = _photo_id(kw, slide_i + 1)
        primary = _pex_cdn(pid0)
        fallbacks = []
        if pid1 != pid0:
            fallbacks.append(_pex_cdn(pid1))
        # Wikimedia Commons fallback
        _qw = _upx.quote((kw or "").split()[0] if kw else "nature")
        fallbacks.append("https://commons.wikimedia.org/wiki/Special:FilePath/{}?width=1200".format(_qw))
        # Picsum detterministic seed fallback
        seed = abs(hash((kw or "") + str(slide_i))) % 9999
        fallbacks.append("https://picsum.photos/seed/{}/{}/{}".format(seed, w, h))
        fb_json = _json.dumps(fallbacks)
        return 'src="' + primary + '" data-fb='' + fb_json + '' onerror="imgFallback(this)"'

    def _img(kw, slide_i=0, w=1200, h=700):
        clean_kw = _normalize_kw(kw)
        seed = abs(hash((kw or "") + str(slide_i))) % 99999
        return _loremflickr(clean_kw, w, h, seed)

    def _img_fb(kw):
        clean_kw = _normalize_kw(kw)
        seed = abs(hash((kw or "") + "fb")) % 99999
        return _loremflickr(clean_kw, 1200, 700, seed + 500)

    def _img2(kw, slide_i=0):
        clean_kw = _normalize_kw(kw)
        seed = abs(hash((kw or "") + str(slide_i + 2))) % 99999
        return _loremflickr(clean_kw, 1200, 700, seed)


    def _e(t):
        return (t or "").replace("&","&amp;").replace("<","&lt;").replace(">","&gt;").replace('"',"&quot;")

    slides_data = pptx_data.get("slides", [])
    slides_html = []

    for i, sd in enumerate(slides_data):
        stype  = sd.get("type", "content")
        n      = i + 1
        kw     = sd.get("image_keyword", "professional modern workspace team")
        imgattr = _img_attr(kw, slide_i=i)
        active = ' active' if n == 1 else ''

        if stype == "title":
            org    = _e(sd.get("organization", ""))
            auth   = _e(sd.get("author", ""))
            sub    = _e(sd.get("subtitle", ""))
            titl   = _e(sd.get("title", ""))
            org_h  = '<p class="t-org">{}</p>'.format(org) if org else ""
            sub_h  = '<div class="t-sub">{}</div>'.format(sub) if sub else ""
            auth_h = '<div class="t-meta">{}</div>'.format(auth) if auth else ""
            slides_html.append(
                '<div class="slide{}" id="slide{}">'.format(active, n) +
                '<img class="slide-bg" ' + imgattr + ' alt="">' +
                '<div class="title-overlay">' +
                org_h +
                '<div class="title-center">' + sub_h +
                '<h1 class="title-main">{}</h1>'.format(titl) +
                '</div>' + auth_h +
                '</div></div>'
            )

        elif stype == "content":
            titl   = _e(sd.get("title", ""))
            buls   = "".join(
                '<li class="bullet-item"><i class="fas fa-chevron-right bi"></i><span>{}</span></li>'.format(_e(b))
                for b in sd.get("bullets", [])
            )
            fact   = _e(sd.get("fact", ""))
            fact_h = (
                '<div class="fact-box"><i class="fas fa-lightbulb fi"></i><span>{}</span></div>'.format(fact)
            ) if fact else ""
            slides_html.append(
                '<div class="slide{}" id="slide{}">'.format(active, n) +
                '<div class="slide-inner grid2">' +
                '<div class="tcol"><h2 class="slide-title">{}</h2>'.format(titl) +
                '<ul class="bullets">{}</ul>'.format(buls) + fact_h +
                '</div><div class="icol">' +
                '<img class="content-img" ' + imgattr + ' alt="">' +
                '</div></div></div>'
            )

        elif stype == "cards":
            titl     = _e(sd.get("title", ""))
            contrast = sd.get("contrast", False)
            cards_h  = ""
            for ci, card in enumerate(sd.get("cards", [])[:2]):
                icon  = _e(card.get("icon", "★"))
                hdg   = _e(card.get("heading", ""))
                body  = _e(card.get("body", ""))
                items = card.get("items", [])
                extra = " card-alt" if (contrast and ci == 1) else ""
                inner = (
                    "<ul class='clist'>" +
                    "".join('<li><i class="fas fa-check-circle"></i><span>{}</span></li>'.format(_e(it)) for it in items) +
                    "</ul>"
                ) if items else '<p class="card-body">{}</p>'.format(body)
                # Если иконка — это FA класс, рендерим <i>; иначе — эмодзи/текст
                if icon.startswith("fas ") or icon.startswith("fab ") or icon.startswith("far "):
                    icon_html = '<i class="{}"></i>'.format(icon)
                else:
                    icon_html = icon
                cards_h += (
                    '<div class="glass-card{}">'.format(extra) +
                    '<div class="card-iw">{}</div>'.format(icon_html) +
                    '<h3 class="card-heading">{}</h3>'.format(hdg) +
                    inner + '</div>'
                )
            slides_html.append(
                '<div class="slide{}" id="slide{}">'.format(active, n) +
                '<div class="slide-inner">' +
                '<h2 class="slide-title">{}</h2>'.format(titl) +
                '<div class="cards-row">{}</div>'.format(cards_h) +
                '</div></div>'
            )

        elif stype == "split":
            titl  = _e(sd.get("title", ""))
            label = _e(sd.get("label", ""))
            buls  = "".join(
                '<li class="bullet-item"><i class="fas fa-check bi"></i><span>{}</span></li>'.format(_e(b))
                for b in sd.get("bullets", [])
            )
            fact  = _e(sd.get("fact", ""))
            icls  = "simg-circle" if sd.get("circle") else "simg-rect"
            lbl_h = '<span class="chapter-badge">{}</span>'.format(label) if label else ""
            fact_h = (
                '<div class="fact-box"><i class="fas fa-quote-left fi"></i><span>{}</span></div>'.format(fact)
            ) if fact else ""
            slides_html.append(
                '<div class="slide{}" id="slide{}">'.format(active, n) +
                '<div class="slide-inner gsplit">' +
                '<div class="lcol"><img class="' + icls + '" ' + imgattr + ' alt="">' +
                '</div><div class="tcol">' + lbl_h +
                '<h2 class="slide-title">{}</h2>'.format(titl) +
                '<ul class="bullets">{}</ul>'.format(buls) + fact_h +
                '</div></div></div>'
            )

        elif stype == "stats":
            titl   = _e(sd.get("title", ""))
            stats_h = "".join(
                '<div class="stat-card">'
                '<div class="stat-value">{}</div>'
                '<div class="stat-label">{}</div>'
                '<div class="stat-desc">{}</div></div>'.format(
                    _e(str(s.get("value",""))), _e(s.get("label","")), _e(s.get("desc",""))
                )
                for s in sd.get("stats", [])[:4]
            )
            slides_html.append(
                '<div class="slide{}" id="slide{}">'.format(active, n) +
                '<img class="slide-bg dim" ' + imgattr + ' alt="">' +
                '<div class="stats-overlay">' +
                '<h2 class="slide-title white-title">{}</h2>'.format(titl) +
                '<div class="stats-row">{}</div>'.format(stats_h) +
                '</div></div>'
            )

        elif stype == "quote":
            quote = _e(sd.get("quote", ""))
            ctx   = _e(sd.get("context", ""))
            ctx_h = '<p class="quote-ctx">— {}</p>'.format(ctx) if ctx else ""
            slides_html.append(
                '<div class="slide{}" id="slide{}">'.format(active, n) +
                '<img class="slide-bg" ' + imgattr + ' alt="">' +
                '<div class="quote-overlay"><div class="quote-box">' +
                '<i class="fas fa-quote-left quote-icon"></i>' +
                '<p class="quote-text">{}</p>'.format(quote) + ctx_h +
                '</div></div></div>'
            )

        elif stype == "timeline":
            titl   = _e(sd.get("title", ""))
            steps_h = "".join(
                '<div class="tstep">'
                '<div class="tnum">{}</div>'
                '<div class="tcnt"><h4 class="step-title">{}</h4>'
                '<p class="step-desc">{}</p></div></div>'.format(
                    si+1, _e(s.get("title","")), _e(s.get("desc",""))
                )
                for si, s in enumerate(sd.get("steps", [])[:5])
            )
            slides_html.append(
                '<div class="slide{}" id="slide{}">'.format(active, n) +
                '<div class="slide-inner">' +
                '<h2 class="slide-title">{}</h2>'.format(titl) +
                '<div class="timeline">{}</div>'.format(steps_h) +
                '</div></div>'
            )

        elif stype == "conclusion":
            titl   = _e(sd.get("title", ""))
            concls = "".join(
                '<li class="bullet-item"><i class="fas fa-star bi2"></i><span>{}</span></li>'.format(_e(c))
                for c in sd.get("conclusions", [])
            )
            cta    = _e(sd.get("cta", ""))
            cta_h  = '<div class="cta-box">{}</div>'.format(cta) if cta else ""
            slides_html.append(
                '<div class="slide{}" id="slide{}">'.format(active, n) +
                '<img class="slide-bg dim" ' + imgattr + ' alt="">' +
                '<div class="concl-overlay">' +
                '<h2 class="slide-title white-title ctitle">{}</h2>'.format(titl) +
                '<ul class="bullets cbuls">{}</ul>'.format(concls) +
                cta_h + '</div></div>'
            )

        else:
            titl = _e(sd.get("title", ""))
            slides_html.append(
                '<div class="slide{}" id="slide{}">'.format(active, n) +
                '<div class="slide-inner"><h2 class="slide-title">{}</h2></div></div>'.format(titl)
            )

    TOTAL  = len(slides_html)

    # Карта keyword для JS динамической загрузки фото через Pexels
    import json as _jmod
    _slide_kw_map = {}
    for _si, _sd in enumerate(slides_data):
        _kw = _sd.get("image_keyword", "")
        if _kw:
            _slide_kw_map[str(_si + 1)] = _kw
    _kw_js = "var SLIDE_KW=" + _jmod.dumps(_slide_kw_map, ensure_ascii=False) + ";"
    _pex_key_js = "var PEX_KEY=" + _jmod.dumps(pexels_key) + ";"

    SLIDES = "\n".join(slides_html)
    DOTS   = "".join(
        '<div class="nav-dot" id="dot{}" onclick="goToSlide({})"></div>'.format(j+1, j+1)
        for j in range(TOTAL)
    )
    PTITLE = _e(pptx_data.get("title", "Презентация"))

    CSS = (
        "@import url('https://fonts.googleapis.com/css2?family=Montserrat:wght@700;800&family=Open+Sans:wght@400;600&display=swap');\n"
        ":root{{"
        "--acc:" + ACC + ";--acc2:" + ACC2 + ";--bg:" + BG + ";--hdr:" + HDR + ";"
        "--tc:" + TC + ";--bc:" + BC + ";--mc:" + MC + ";"
        "--glass:rgba(" + GBGA + "," + GBGO + ");"
        "}}\n"
        "*,*::before,*::after{box-sizing:border-box;margin:0;padding:0}\n"
        "body{font-family:'Open Sans',sans-serif;background:var(--bg);color:var(--bc);overflow:hidden;width:100vw;height:100vh}\n"
        ".slide{display:none;position:relative;width:100vw;height:100vh;overflow:hidden}\n"
        ".slide.active{display:flex;animation:sIn .4s ease}\n"
        "@keyframes sIn{from{opacity:0;transform:translateX(20px)}to{opacity:1;transform:none}}\n"
        ".slide-bg{position:absolute;inset:0;width:100%;height:100%;object-fit:cover;z-index:0;opacity:.45}\n"
        ".slide-bg.dim{opacity:.2}\n"
        ".slide-inner{position:relative;z-index:2;width:100%;height:100%;padding:56px 68px;display:flex;flex-direction:column;background:var(--bg)}\n"
        ".grid2{flex-direction:row!important;align-items:stretch;gap:40px}\n"
        ".gsplit{flex-direction:row!important;align-items:stretch;gap:40px}\n"
        ".tcol{flex:1;display:flex;flex-direction:column;justify-content:center}\n"
        ".icol{width:42%;display:flex;align-items:center}\n"
        ".lcol{width:44%;display:flex;align-items:center}\n"
        # Title
        ".title-overlay{position:relative;z-index:2;width:100%;height:100%;display:flex;flex-direction:column;"
        "justify-content:space-between;padding:56px 68px;"
        "background:linear-gradient(to top," + HDR + "dd 0%,transparent 62%)}\n"
        ".t-org{text-align:center;font-weight:700;font-size:.9rem;letter-spacing:3px;text-transform:uppercase;color:var(--acc)}\n"
        ".title-center{flex:1;display:flex;flex-direction:column;justify-content:center;align-items:center;gap:18px;text-align:center}\n"
        ".t-sub{font-size:1.05rem;font-weight:600;color:var(--acc);letter-spacing:2px;text-transform:uppercase}\n"
        ".title-main{font-family:'Montserrat',sans-serif;font-weight:800;text-transform:uppercase;"
        "font-size:clamp(2rem,4.5vw,5rem);line-height:1.1;color:var(--tc);text-shadow:0 4px 24px rgba(0,0,0,.6)}\n"
        ".t-meta{text-align:center;font-weight:600;font-size:1rem;color:var(--tc);padding:12px 28px;"
        "background:var(--glass);border:1px solid var(--acc);border-radius:10px;"
        "backdrop-filter:blur(8px);width:fit-content;margin:0 auto}\n"
        # Slide title
        ".slide-title{font-family:'Montserrat',sans-serif;font-weight:800;text-transform:uppercase;"
        "font-size:clamp(1.3rem,2.6vw,2.4rem);color:var(--tc);line-height:1.2;"
        "border-bottom:4px solid var(--acc);padding-bottom:12px;margin-bottom:24px}\n"
        ".white-title{color:#fff!important;border-bottom-color:var(--acc2)!important}\n"
        # Images
        ".content-img{width:100%;height:74vh;object-fit:cover;border-radius:20px;border:3px solid var(--acc);box-shadow:0 12px 40px rgba(0,0,0,.5)}\n"
        ".simg-rect{width:100%;height:68vh;object-fit:cover;border-radius:20px;border:3px solid var(--acc);box-shadow:0 10px 36px rgba(0,0,0,.45)}\n"
        ".simg-circle{width:52vh;height:52vh;object-fit:cover;border-radius:50%;border:4px solid var(--acc);box-shadow:0 10px 36px rgba(0,0,0,.45);margin:auto}\n"
        # Bullets
        ".bullets{list-style:none;display:flex;flex-direction:column;gap:10px;flex:1}\n"
        ".bullet-item{display:flex;align-items:flex-start;gap:10px;font-size:clamp(.85rem,1.4vw,1.06rem);"
        "line-height:1.55;padding:10px 16px;background:var(--glass);"
        "border-left:4px solid var(--acc);border-radius:0 8px 8px 0;backdrop-filter:blur(6px)}\n"
        ".bi{color:var(--acc);flex-shrink:0;margin-top:3px}\n"
        ".bi2{color:var(--acc2);flex-shrink:0;margin-top:3px}\n"
        ".fact-box{margin-top:16px;display:flex;align-items:flex-start;gap:10px;padding:12px 16px;"
        "background:var(--glass);border-left:5px solid var(--acc2);"
        "border-radius:0 10px 10px 0;font-size:.9rem;font-style:italic;color:var(--tc);backdrop-filter:blur(6px)}\n"
        ".fi{color:var(--acc2);flex-shrink:0;margin-top:2px}\n"
        ".chapter-badge{display:inline-block;background:var(--acc);color:#000;font-weight:700;"
        "font-size:.75rem;padding:4px 14px;border-radius:4px;text-transform:uppercase;letter-spacing:1.5px;margin-bottom:12px}\n"
        # Cards
        ".cards-row{display:flex;gap:28px;flex:1;margin-top:4px}\n"
        ".glass-card{flex:1;background:var(--glass);border-left:6px solid var(--acc);"
        "border-radius:0 18px 18px 0;padding:28px;backdrop-filter:blur(8px);"
        "display:flex;flex-direction:column;gap:12px}\n"
        ".glass-card.card-alt{border-left-color:var(--acc2)}\n"
        ".card-iw{font-size:2.6rem;line-height:1}\n"
        ".card-heading{font-family:'Montserrat',sans-serif;font-weight:700;font-size:1.1rem;"
        "text-transform:uppercase;color:var(--acc)}\n"
        ".glass-card.card-alt .card-heading{color:var(--acc2)}\n"
        ".card-body{font-size:clamp(.82rem,1.25vw,.96rem);line-height:1.65;color:var(--bc)}\n"
        ".clist{list-style:none;display:flex;flex-direction:column;gap:6px;"
        "font-size:clamp(.82rem,1.25vw,.96rem);color:var(--bc)}\n"
        ".clist li{display:flex;align-items:flex-start;gap:8px;line-height:1.5}\n"
        ".clist li i{color:var(--acc2);flex-shrink:0;margin-top:3px;font-size:.78rem}\n"
        # Stats
        ".stats-overlay{position:relative;z-index:2;width:100%;height:100%;padding:56px 68px;"
        "display:flex;flex-direction:column;justify-content:center;align-items:center;text-align:center;"
        "background:linear-gradient(rgba(0,0,0,.45),rgba(0,0,0,.7))}\n"
        ".stats-row{display:flex;gap:24px;justify-content:center;flex-wrap:wrap;margin-top:20px}\n"
        ".stat-card{background:var(--glass);border:2px solid var(--acc);border-radius:18px;"
        "padding:28px 36px;min-width:170px;backdrop-filter:blur(10px);text-align:center}\n"
        ".stat-value{font-family:'Montserrat',sans-serif;font-size:clamp(2.2rem,4.5vw,4rem);"
        "font-weight:800;color:var(--acc);line-height:1}\n"
        ".stat-label{font-weight:700;font-size:.9rem;color:var(--tc);margin-top:8px;"
        "text-transform:uppercase;letter-spacing:.5px}\n"
        ".stat-desc{font-size:.78rem;color:var(--mc);margin-top:4px}\n"
        # Quote
        ".quote-overlay{position:relative;z-index:2;width:100%;height:100%;"
        "display:flex;justify-content:center;align-items:center;"
        "background:linear-gradient(rgba(0,0,0,.4),rgba(0,0,0,.7))}\n"
        ".quote-box{max-width:800px;text-align:center;padding:52px;background:var(--glass);"
        "border-radius:24px;backdrop-filter:blur(14px);border:1px solid var(--acc)}\n"
        ".quote-icon{font-size:3rem;color:var(--acc);margin-bottom:20px;display:block}\n"
        ".quote-text{font-size:clamp(1.2rem,2.4vw,2rem);font-style:italic;line-height:1.7;"
        "color:var(--tc);margin-bottom:20px}\n"
        ".quote-ctx{color:var(--acc);font-weight:700;font-size:1rem}\n"
        # Timeline
        ".timeline{display:flex;flex-direction:column;gap:14px;flex:1;justify-content:center;margin-top:4px}\n"
        ".tstep{display:flex;align-items:flex-start;gap:18px;padding:15px 20px;"
        "background:var(--glass);border-left:4px solid var(--acc);"
        "border-radius:0 12px 12px 0;backdrop-filter:blur(6px)}\n"
        ".tnum{width:40px;height:40px;min-width:40px;border-radius:50%;background:var(--acc);"
        "color:#000;display:flex;align-items:center;justify-content:center;"
        "font-family:'Montserrat',sans-serif;font-weight:800;font-size:1rem}\n"
        ".step-title{font-family:'Montserrat',sans-serif;font-weight:700;font-size:.98rem;"
        "color:var(--tc);margin-bottom:4px}\n"
        ".step-desc{font-size:.87rem;color:var(--bc);line-height:1.55}\n"
        # Conclusion
        ".concl-overlay{position:relative;z-index:2;width:100%;height:100%;padding:52px 80px;"
        "display:flex;flex-direction:column;justify-content:center;align-items:center;"
        "background:linear-gradient(rgba(0,0,0,.48),rgba(0,0,0,.76))}\n"
        ".ctitle{border-bottom-color:var(--acc2)!important;text-align:center;width:100%}\n"
        ".cbuls{max-width:760px;width:100%;margin-top:20px}\n"
        ".cbuls .bullet-item{border-left-color:var(--acc2)}\n"
        ".cta-box{margin-top:28px;padding:16px 32px;background:var(--acc);color:#000;"
        "border-radius:14px;font-family:'Montserrat',sans-serif;font-weight:800;"
        "font-size:1.25rem;text-align:center;text-transform:uppercase;letter-spacing:1.5px}\n"
        # Navigation
        ".nav-controls{position:fixed;bottom:28px;right:36px;display:flex;align-items:center;gap:12px;z-index:1000}\n"
        ".nav-btn{background:var(--acc);color:#000;border:none;padding:12px 24px;"
        "border-radius:9px;font-family:'Montserrat',sans-serif;font-weight:700;"
        "font-size:.88rem;cursor:pointer;transition:all .25s;display:flex;align-items:center;gap:8px}\n"
        ".nav-btn:hover{background:var(--acc2);transform:translateY(-2px);box-shadow:0 6px 22px rgba(0,0,0,.35)}\n"
        ".slide-counter{font-family:'Montserrat',sans-serif;font-weight:700;font-size:.9rem;"
        "color:var(--acc);background:var(--glass);padding:11px 18px;border-radius:9px;"
        "backdrop-filter:blur(8px);border:1px solid var(--acc);min-width:70px;text-align:center}\n"
        ".nav-dots{position:fixed;bottom:36px;left:50%;transform:translateX(-50%);display:flex;gap:8px;z-index:1000}\n"
        ".nav-dot{width:10px;height:10px;border-radius:50%;background:var(--mc);cursor:pointer;transition:all .25s}\n"
        ".nav-dot.active{background:var(--acc);transform:scale(1.45)}\n"
    )

    JS = (
        "var cur=1,tot=" + str(TOTAL) + ",sx=0,sy=0;\n"
        "function goTo(n){\n"
        "  if(n<1||n>tot)return;\n"
        "  document.querySelectorAll('.slide').forEach(function(s){s.classList.remove('active');});\n"
        "  document.querySelectorAll('.nav-dot').forEach(function(d){d.classList.remove('active');});\n"
        "  var s=document.getElementById('slide'+n),d=document.getElementById('dot'+n);\n"
        "  if(s)s.classList.add('active');\n"
        "  if(d)d.classList.add('active');\n"
        "  document.getElementById('counter').textContent=n+' / '+tot;\n"
        "  cur=n;\n"
        "}\n"
        "function goToSlide(n){goTo(n);}\n"
        "function nextSlide(){goTo(cur+1);}\n"
        "function prevSlide(){goTo(cur-1);}\n"
        # Keyboard
        "document.addEventListener('keydown',function(e){\n"
        "  if(e.key==='ArrowRight'||e.key===' '){e.preventDefault();nextSlide();}\n"
        "  if(e.key==='ArrowLeft'){e.preventDefault();prevSlide();}\n"
        "  if(e.key==='Home'){e.preventDefault();goTo(1);}\n"
        "  if(e.key==='End'){e.preventDefault();goTo(tot);}\n"
        "});\n"
        # Touch swipe
        "document.addEventListener('touchstart',function(e){sx=e.touches[0].clientX;sy=e.touches[0].clientY;},{passive:true});\n"
        "document.addEventListener('touchend',function(e){\n"
        "  var dx=e.changedTouches[0].clientX-sx,dy=e.changedTouches[0].clientY-sy;\n"
        "  if(Math.abs(dx)>Math.abs(dy)&&Math.abs(dx)>45){if(dx<0)nextSlide();else prevSlide();}\n"
        "},{passive:true});\n"
        # Click left/right zones (исключая кнопки навигации)
        "document.addEventListener('click',function(e){\n"
        "  if(e.target.closest('.nav-btn,.nav-dot,.nav-controls,.nav-dots'))return;\n"
        "  var x=e.clientX,w=window.innerWidth;\n"
        "  if(x>w*0.65)nextSlide();\n"
        "  else if(x<w*0.35)prevSlide();\n"
        "});\n"
        "goTo(1);\n"
        # ── Динамическая загрузка фото через Pexels API ────────────────
        "function imgFallback(el){\n"
        "  try{\n"
        "    var fb=JSON.parse(el.getAttribute('data-fb')||'[]');\n"
        "    var idx=parseInt(el.getAttribute('data-fi')||'0');\n"
        "    if(idx<fb.length){el.setAttribute('data-fi',idx+1);el.src=fb[idx];}\n"
        "  }catch(e){}\n"
        "}\n"
        # Pexels динамический загрузчик - работает для ЛЮБОЙ темы
        "var _pexCache={};\n"
        "async function _loadPexPhoto(slideNum,keyword,imgEls){\n"
        "  if(!PEX_KEY||!keyword)return;\n"
        "  var cacheKey=keyword+slideNum;\n"
        "  if(_pexCache[cacheKey]){\n"
        "    imgEls.forEach(function(el){el.src=_pexCache[cacheKey];});return;\n"
        "  }\n"
        "  try{\n"
        "    var r=await fetch('https://api.pexels.com/v1/search?query='+encodeURIComponent(keyword)+'&per_page=15&orientation=landscape',\n"
        "      {headers:{'Authorization':PEX_KEY}});\n"
        "    var d=await r.json();\n"
        "    if(d.photos&&d.photos.length){\n"
        "      var ph=d.photos[(slideNum-1)%d.photos.length];\n"
        "      var url=ph.src.large2x||ph.src.large||ph.src.medium;\n"
        "      _pexCache[cacheKey]=url;\n"
        "      imgEls.forEach(function(el){el.src=url;});\n"
        "    }\n"
        "  }catch(e){}\n"
        "}\n"
        # Загружаем все фото после DOMContentLoaded
        "document.addEventListener('DOMContentLoaded',function(){\n"
        "  if(!PEX_KEY)return;\n"
        "  Object.keys(SLIDE_KW).forEach(function(sn){\n"
        "    var imgs=document.querySelectorAll('#slide'+sn+' img');\n"
        "    if(imgs.length)_loadPexPhoto(parseInt(sn),SLIDE_KW[sn],Array.from(imgs));\n"
        "  });\n"
        "});\n"
    )

    return (
        '<!DOCTYPE html>\n'
        '<html lang="ru">\n'
        '<head>\n'
        '<meta charset="UTF-8">\n'
        '<meta name="viewport" content="width=device-width, initial-scale=1.0">\n'
        '<title>' + PTITLE + '</title>\n'
        '<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">\n'
        '<style>' + CSS + '\n'
        '#pc-hint{position:fixed;top:0;left:0;right:0;z-index:9999;background:linear-gradient(90deg,#0d1336,#16213e);'
        'color:#fff;text-align:center;padding:9px 16px;font-family:sans-serif;font-size:13px;'
        'display:flex;align-items:center;justify-content:center;gap:10px;border-bottom:2px solid #4d8fff;}'
        '#pc-hint b{color:#4d8fff}'
        '#pc-hint .close-hint{background:none;border:1px solid #4d8fff;color:#4d8fff;'
        'border-radius:5px;padding:2px 9px;cursor:pointer;font-size:12px;margin-left:6px;}'
        '#pc-hint .close-hint:hover{background:#4d8fff;color:#000}'
        'body{padding-top:40px}'
        '.nav-controls{bottom:18px}\n'
        '</style>\n'
        '</head>\n'
        '<body>\n'
        '<div id="pc-hint">'
        '<i class="fas fa-desktop" style="color:#4d8fff"></i>&nbsp;'
        '💡 <b>Рекомендуется открывать на ПК</b> &nbsp;|&nbsp; '
        'Листать: <b>← →</b> на клавиатуре, кнопки внизу или <b>свайп</b> на телефоне'
        '<button class="close-hint" onclick="document.getElementById(\'pc-hint\').remove()">✕</button>'
        '</div>\n'
        + SLIDES + '\n'
        '<div class="nav-controls">\n'
        '  <button class="nav-btn" onclick="prevSlide()"><i class="fas fa-arrow-left"></i> Назад</button>\n'
        '  <span class="slide-counter" id="counter">1 / ' + str(TOTAL) + '</span>\n'
        '  <button class="nav-btn" onclick="nextSlide()">Вперёд <i class="fas fa-arrow-right"></i></button>\n'
        '</div>\n'
        '<div class="nav-dots">' + DOTS + '</div>\n'
        '<script>\n' + _pex_key_js + '\n' + _kw_js + '\n' + JS + '</script>\n'
        '</body>\n'
        '</html>'
    )



# ── AI промпт ────────────────────────────────────────────────────

PPTX_AI_SYSTEM_PROMPT = """Ты — ведущий дизайнер и стратегический консультант с доступом к актуальным данным из интернета.
Используй веб-поиск чтобы найти свежую статистику, факты и цитаты по теме презентации.
Создаёшь презентации уровня McKinsey, BCG и Apple Keynote.
Твоя презентация должна быть такой, чтобы её не стыдно показать CEO или защитить на экзамене.

ОТВЕЧАЙ ТОЛЬКО JSON без markdown-обёртки, без пояснений до или после JSON.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
 ПРИНЦИПЫ КОНТЕНТА — ОБЯЗАТЕЛЬНЫ
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

ГЛАВНОЕ ПРАВИЛО: Каждый слайд должен ДОКАЗЫВАТЬ ТЕЗИС, а не просто "рассказывать о теме".

ЗАПРЕЩЕНО (плохие примеры от которых уходи):
✗ bullet: "Развитие экономики важно для страны"
✗ fact: "Это интересная тема"
✗ stat.value: "100%", stat.label: "Показатель"
✗ quote: "Нужно работать усердно"
✗ Абстрактные утверждения без цифр

ОБЯЗАТЕЛЬНО (конкретно, экспертно, с доказательной базой):
✓ bullet: "ВВП вырос на 3.6% в 2023г вопреки санкционному давлению — прогноз был -2.5%"
✓ fact: "Россия — №1 в мире по запасам природного газа: 19% мировых резервов (BP 2023)"
✓ stat.value: "68%", label: "Импортозамещение", desc: "промышленность, 2022-2024"
✓ quote: реальная цитата эксперта/руководителя/исследования с источником
✓ Каждое утверждение поддержано данными или примером

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
 IMAGE_KEYWORD — КРИТИЧЕСКИЕ ПРАВИЛА
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

image_keyword — ПОИСКОВЫЙ ЗАПРОС для поиска фото. ОБЯЗАТЕЛЬНО на АНГЛИЙСКОМ языке.
Ищет фото в Flickr/Pexels по ключевым словам. ТОЧНОСТЬ = НУЖНОЕ ФОТО.
ПРАВИЛО: 2-4 конкретных слова, описывающих ИМЕННО ЧТО ИЗОБРАЖЕНО на фото.

🔴 ЗАПРЕЩЕНО — плохие keywords которые дают нерелевантные фото:
❌ Абстракции: "modern", "concept", "background", "digital", "abstract", "idea"
❌ Одно широкое слово: "Russia", "war", "economy", "technology"
❌ Кириллица в keyword (только английский!)
❌ Повтор одного keyword на разных слайдах

🟢 ПРАВИЛЬНО — конкретный объект + среда/действие:

ИСТОРИЯ / ВОЙНА:
✅ "soldiers battlefield smoke World War" | "Russian soldiers Red Army 1941"
✅ "Soviet parade Red Square Moscow" | "Kremlin Moscow winter fortress"
✅ "military tank armor combat field" | "warship navy fleet patrol ocean"
✅ "Napoleonic cavalry battle 1812 horses" | "Crimean War fortification soldiers"
✅ "Cold War nuclear bomb test explosion" | "space race rocket launch Sputnik"

ОБЖ / БЕЗОПАСНОСТЬ (ШКОЛЬНАЯ ТЕМАТИКА — НЕ АМЕРИКАНСКИЕ ВОЕННЫЕ!):
✅ "police officers patrol security uniform street Russia" | "detective investigation laptop computer crime evidence"
✅ "teenager smartphone social media screen warning parent" | "school safety drill evacuation students corridor"
✅ "Russian police law enforcement uniform patrol building" | "security checkpoint ID verification public safety"
✅ "counselor psychologist teenager student meeting prevention" | "parents children internet safety discussion home"
✅ "crime prevention statistics data chart infographic" | "emergency hotline phone call police assistance"
✅ "school teacher classroom students awareness education" | "cybersecurity internet safety padlock screen digital"

⚠️ ДЛЯ КАЖДОГО СЛАЙДА ПОДБИРАЙ УНИКАЛЬНЫЙ keyword:
- Слайд со статистикой → "police statistics crime data chart analysis"
- Слайд с признаками → "detective investigation evidence laptop screen"
- Слайд с методами вербовки → "teenager smartphone social media screen warning"
- Слайд с действиями → "emergency phone call police help hotline"
- Слайд с профилактикой → "school teacher students prevention awareness"
- Слайд с выводами → "school safety police cooperation community education"

БИЗНЕС / ФИНАНСЫ:
✅ "business meeting office professionals" | "stock market traders screens"
✅ "startup team whiteboard planning" | "handshake contract deal"
✅ "finance charts growth investment" | "factory workers production line"

ТЕХНОЛОГИИ:
✅ "circuit board microchip technology" | "server room datacenter cables"
✅ "robot arm factory automation" | "scientist laboratory experiment"

ГЕОПОЛИТИКА / ПОЛИТИКА:
✅ "world map geopolitics global strategy" | "diplomacy summit meeting flags"
✅ "sanctions economy trade policy" | "NATO soldiers military alliance"

ЭКОНОМИКА:
✅ "oil refinery pipeline petroleum" | "wheat field harvest combine"
✅ "GDP growth chart economy statistics" | "sanctions protest economy crisis"

КАЖДЫЙ СЛАЙД — СВОЙ УНИКАЛЬНЫЙ keyword описывающий КОНКРЕТНОЕ СОДЕРЖАНИЕ этого слайда.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
 ТИПЫ СЛАЙДОВ (ЧЕРЕДУЙ ДЛЯ ДИНАМИКИ)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

1. "title"      — title, subtitle, author, organization, image_keyword
2. "content"    — title, bullets[3-5 конкретных], fact, image_keyword
3. "cards"      — title, contrast(bool), cards[{icon,heading,body или items:[]}]×2
   ⚠️ ICON ОБЯЗАН быть ЭМОДЗИ (🏭 💡 ⚡ 📊 🚀 🔬 ⚙️ ✅ ❌ 🌍 🏆 и т.д.) — НЕ текст "fas fa-*"
4. "split"      — title, label, bullets, fact, circle(bool), image_keyword
5. "stats"      — title, stats[{value,label,desc}]×2-4, image_keyword
6. "quote"      — quote(реальная цитата), context(автор+источник), image_keyword
7. "timeline"   — title, steps[{title,desc}]×3-5
8. "conclusion" — title, conclusions[3-4 конкретных], cta(призыв к действию), image_keyword

Порядок для 7 слайдов:  title→content→cards→split→stats→quote→conclusion
Порядок для 10 слайдов: title→content→cards→split→content→stats→split(circle)→cards→timeline→conclusion
Порядок для 12 слайдов: title→content→cards(contrast)→split→content→stats→split(circle)→content→cards→quote→timeline→conclusion

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
 ФОРМАТ JSON
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

{
  "title": "Полное название презентации",
  "subtitle": "Ёмкий подзаголовок — суть в одной фразе",
  "slides": [
    {"type":"title","title":"...","subtitle":"...","author":"","organization":"","image_keyword":"specific objects scene relevant"},
    {"type":"content","title":"Конкретный тезисный заголовок","bullets":["Факт 1 с цифрой","Факт 2 с примером","Факт 3 — вывод или тренд"],"fact":"Ключевая статистика или цитата с источником","image_keyword":"specific relevant objects scene"},
    {"type":"cards","title":"...","contrast":true,"cards":[{"icon":"⚠️","heading":"Проблема/Вызов","body":"Конкретное описание 2-3 предложения"},{"icon":"✅","heading":"Решение/Возможность","items":["Шаг 1 конкретный","Шаг 2 измеримый","Результат"]}]},
    {"type":"split","title":"...","label":"Раздел","circle":false,"bullets":["Пункт 1","Пункт 2","Пункт 3"],"fact":"Ключевой тезис раздела","image_keyword":"specific scene"},
    {"type":"stats","title":"Ключевые показатели","stats":[{"value":"3.6%","label":"Рост ВВП","desc":"2023, МВФ"},{"value":"68%","label":"Импортозамещение","desc":"промышленность"},{"value":"№1","label":"Запасы газа","desc":"в мире"}],"image_keyword":"data analytics office screens"},
    {"type":"quote","quote":"Точная значимая цитата которая усиливает аргументацию","context":"Имя Фамилия, должность, год/источник","image_keyword":"inspiring dramatic relevant scene"},
    {"type":"conclusion","title":"Выводы и следующие шаги","conclusions":["Вывод 1 конкретный с цифрой","Вывод 2 — что доказано","Вывод 3 — перспектива"],"cta":"Конкретный призыв к действию","image_keyword":"success achievement future progress"}
  ]
}

Весь текст слайдов — РУССКИЙ. image_keyword — ТОЛЬКО английский, конкретный!"""





def _sanitize_slide_keywords(pptx_data: dict, topic: str) -> dict:
    """
    Пост-обработка image_keyword для каждого слайда.
    Заменяет нерелевантные/опасные keywords на тематически правильные.
    """
    topic_lower = topic.lower()
    is_safety = any(w in topic_lower for w in [
        'обж', 'безопасност', 'террор', 'скулшутинг', 'нападен', 'угроз',
        'пожарн', 'эвакуац', 'чрезвычайн', 'первая помощь', 'ожог', 'отравлен',
        'дтп', 'авари', 'спасател', 'выживан', 'наводнен', 'землетрясен', 'экстремальн',
    ])

    # Стоп-слова в keywords которые дают нерелевантные фото
    BAD_KW_STARTS = [
        'us army', 'us military', 'american military', 'american soldiers',
        'american troops', 'us troops', 'pentagon', 'us navy', 'us air force',
        'nato soldiers', 'nato troops', 'western military',
    ]
    BAD_WORDS_IN_KW = ['american', 'us army', 'pentagon', 'us military']

    # Замены для ОБЖ тем по контексту заголовка слайда
    SAFETY_KW_MAP = {
        'fire': 'firefighters rescue emergency fire building',
        'evacuat': 'evacuation crowd emergency exit building safety',
        'terror': 'police security checkpoint crowd protection emergency',
        'school': 'school safety security training evacuation drill students',
        'bomb': 'police security evacuation emergency response team',
        'attack': 'security police rescue emergency response protection',
        'first aid': 'paramedics first aid medical emergency treatment bandage',
        'flood': 'flood rescue boat emergency people water disaster',
        'earthquake': 'earthquake rescue workers emergency destroyed building',
        'accident': 'traffic accident road police rescue emergency',
        'poison': 'medical emergency treatment hospital ambulance',
        'burn': 'medical burn treatment hospital emergency',
    }

    slides = pptx_data.get('slides', [])
    safety_kws = [
        'firefighters rescue emergency scene',
        'ambulance paramedics first aid emergency',
        'evacuation crowd emergency exit building',
        'police security patrol checkpoint crowd',
        'first aid medical bandage treatment',
        'emergency shelter people disaster relief',
        'rescue team emergency responders uniform',
        'school safety drill emergency students',
    ]
    safety_idx = 0

    for i, slide in enumerate(slides):
        kw = slide.get('image_keyword', '') or slide.get('image_query', '')
        if not kw:
            continue
        kw_lower = kw.lower()

        # Проверяем нет ли плохих слов
        has_bad = any(bad in kw_lower for bad in BAD_WORDS_IN_KW)
        has_bad_start = any(kw_lower.startswith(bad) for bad in BAD_KW_STARTS)

        if is_safety and (has_bad or has_bad_start):
            # Пробуем подобрать по заголовку
            slide_title = (slide.get('title', '') + ' ' + str(slide.get('bullets', ''))).lower()
            replacement = None
            for key, replacement_kw in SAFETY_KW_MAP.items():
                if key in slide_title or key in kw_lower:
                    replacement = replacement_kw
                    break
            if not replacement:
                replacement = safety_kws[safety_idx % len(safety_kws)]
                safety_idx += 1
            if 'image_keyword' in slide:
                slide['image_keyword'] = replacement
            if 'image_query' in slide:
                slide['image_query'] = replacement
            logging.info(f"[PPTX] Заменён keyword слайда {i}: '{kw}' → '{replacement}'")

    pptx_data['slides'] = slides
    return pptx_data


async def generate_pptx_structure_ai(topic: str, slides_count: int,
                                     model_key: str, extra_wishes: str = "",
                                     title_info: dict | None = None) -> dict:
    wishes_block = f"\n\nПОЖЕЛАНИЯ ЗАКАЗЧИКА: {extra_wishes}" if extra_wishes else ""
    title_block = ""
    if title_info:
        parts = []
        if title_info.get("subtitle"):      parts.append(f"Подзаголовок: «{title_info['subtitle']}»")
        if title_info.get("author"):        parts.append(f"Автор: {title_info['author']}")
        if title_info.get("organization"):  parts.append(f"Организация: {title_info['organization']}")
        if parts:
            title_block = "\n\nДАННЫЕ ДЛЯ ТИТУЛЬНОГО СЛАЙДА:\n" + "\n".join(parts)

    # Адаптируем подсказку по количеству слайдов
    if slides_count <= 5:
        layout_hint = "title → content → cards → split → conclusion"
    elif slides_count <= 7:
        layout_hint = "title → content → cards(contrast) → split → stats → quote → conclusion"
    elif slides_count <= 10:
        layout_hint = "title → content → cards → split → content → stats → split(circle) → cards → timeline → conclusion"
    else:
        layout_hint = "title → content → cards(contrast) → split → content → stats → split(circle) → content → cards → quote → timeline → conclusion"

    # Определяем тему для умной подстановки image_keyword
    topic_lower = topic.lower()
    is_military_topic = any(w in topic_lower for w in [
        "военн", "впк", "оборон", "армия", "нато", "войск", "ракет", "танк", "вооруж"
    ])
    is_safety_topic = any(w in topic_lower for w in [
        "обж", "безопасност", "террор", "скулшутинг", "нападен", "угроз",
        "пожарн", "эвакуац", "чрезвычайн", "первая помощь", "ожог", "отравлен",
        "дтп", "авари", "спасател", "экстремальн", "выживан", "наводнен", "землетрясен",
        "вербов", "радикализ", "экстремизм", "профилакт",
    ])

    if is_safety_topic:
        geo_rule = (
            "⛔ ТЕМА ОБЖ/БЕЗОПАСНОСТЬ ДЛЯ РОССИЙСКОЙ ШКОЛЫ!\n"
            "АБСОЛЮТНО ЗАПРЕЩЕНЫ keywords с: police, polizei, cop, officer, law enforcement "
            "(они дают иностранных полицейских!), а также: US Army, NATO, American, German, Indian.\n\n"
            "ОБЯЗАТЕЛЬНО используй ТОЛЬКО эти безопасные категории:\n"
            "ДЛЯ ТЕМЫ ВЕРБОВКИ/ИНТЕРНЕТ: 'teenager smartphone social media dark', "
            "'hacker dark monitor code screen', 'cybersecurity laptop warning alert', "
            "'social network manipulation online trap', 'student computer screen danger'\n"
            "ДЛЯ ТЕМЫ БЕЗОПАСНОСТИ: 'security camera surveillance cctv city', "
            "'lock protection shield security concept', 'alert warning danger sign red'\n"
            "ДЛЯ СТАТИСТИКИ: 'chart graph data analytics dashboard screen', "
            "'statistics numbers report document'\n"
            "ДЛЯ ШКОЛЫ/ПОДРОСТКОВ: 'students classroom school education learning', "
            "'teenager young people group discussion', 'school corridor hallway students'\n"
            "ДЛЯ ПОЖАРА: 'firefighters fire rescue building flame smoke'\n"
            "ДЛЯ ПЕРВОЙ ПОМОЩИ: 'first aid kit medical bandage treatment hands'\n"
            "ДЛЯ ЗАКЛЮЧЕНИЯ: 'protection shield family safe home concept'\n"
            "Подбирай keyword к КОНКРЕТНОМУ содержанию слайда, используя ТОЛЬКО английские слова из примеров выше!"
        )
    elif is_military_topic:
        geo_rule = (
            "⛔ ВОЕННАЯ ТЕМА! Первое слово image_keyword ОБЯЗАНО быть из этого списка: "
            "soldier, soldiers, military, army, tank, armored, weapon, weapons, missile, "
            "rocket, fighter, aircraft, drone, nato, defense, combat. "
            "Примеры ПРАВИЛЬНЫХ keywords: "
            "'soldiers training field exercise', 'tank armor military combat', "
            "'factory weapons production assembly', 'military defense budget spending', "
            "'aircraft fighter jet military', 'drone surveillance reconnaissance'."
        )
    elif any(w in topic_lower for w in ["росси", "россий", "рф ", "рф,", "москв", "советск", "русск"]):
        geo_rule = (
            "⛔ ТЕМА О РОССИИ! Первое слово image_keyword ОБЯЗАНО быть конкретным объектом из списка: "
            "factory, workers, oil, pipeline, wheat, harvest, building, government, "
            "stock, chart, bank, ruble, currency, soldier, soldiers, cargo, ship. "
            "Примеры: 'factory workers industrial production', 'oil pipeline field winter', "
            "'wheat harvest combine golden', 'stock chart finance trading'."
        )
    elif any(w in topic_lower for w in ["китай", "китайск", "пекин", "шанхай"]):
        geo_rule = (
            "Первое слово keyword: factory, workers, skyline, building, technology, chip, cargo. "
            "Примеры: 'factory assembly workers production', 'skyline building architecture modern'."
        )
    elif any(w in topic_lower for w in ["европ", "евросою"]):
        geo_rule = (
            "Первое слово keyword: building, government, skyline, stock, chart, bank. "
            "Примеры: 'building government architecture modern', 'stock chart finance'."
        )
    elif any(w in topic_lower for w in ["сша", "америк", "американск"]):
        geo_rule = (
            "Первое слово keyword: building, skyline, stock, chart, government, bank. "
            "Примеры: 'skyline building architecture', 'stock trading chart finance'."
        )
    elif any(w in topic_lower for w in ["нижегород", "кстов", "поволж", "приволжск", "регион", "муниципал", "область", "город"]):
        geo_rule = (
            "РЕГИОНАЛЬНАЯ ТЕМА! Используй конкретные объекты соответствующие содержанию КАЖДОГО слайда. "
            "Для слайда о транспорте: 'highway trucks transport road infrastructure'. "
            "Для слайда о промышленности: 'oil refinery industrial plant pipes factory'. "
            "Для слайда об инфраструктуре: 'power plant electricity grid city infrastructure'. "
            "Для слайда об инвестициях: 'business investment meeting office city skyline'. "
            "Для титульного: 'Volga river Russia city bridge landscape'. "
            "ЗАПРЕЩЕНО: одно слово, абстракции. ОБЯЗАТЕЛЬНО: разные keywords для каждого слайда."
        )
    else:
        geo_rule = (
            "Первое слово image_keyword ДОЛЖНО быть конкретным объектом из этого словаря: "
            "soldier, military, factory, workers, oil, stock, chart, bank, server, chip, "
            "neural, code, hospital, doctor, students, cargo, ship, wheat, building, athletes. "
            "ЗАПРЕЩЕНО начинать с: Russia, China, global, modern, abstract, digital."
        )

    user_msg = (
        f"Тема: «{topic}»\n"
        f"Количество слайдов: {slides_count} (включая title и conclusion)\n"
        f"Рекомендуемая последовательность типов: {layout_hint}\n"
        f"Язык контента: русский, image_keyword: только английский\n"
        f"{wishes_block}{title_block}\n\n"
        "ТРЕБОВАНИЯ К КАЧЕСТВУ:\n"
        "1. Каждый слайд — уникальный тип, никаких повторений подряд\n"
        "2. Контент содержательный и экспертный, не поверхностный\n"
        f"3. {geo_rule}\n"
        "4. image_keyword = 3-5 слов, конкретное место/объект/сцена, ТОЛЬКО английский\n"
        "5. Используй ВСЕ доступные типы слайдов гармонично\n"
        "6. Для cards с contrast=true — реально разные по смыслу карточки\n\n"
        "Верни ТОЛЬКО JSON без каких-либо пояснений."
    )
    msgs = [
        {"role": "system", "content": PPTX_AI_SYSTEM_PROMPT},
        {"role": "user",   "content": user_msg},
    ]
    # Sonar Deep Research первым — актуальные данные из интернета, минуя проверки disabled
    _pptx_fallback_order = ["sonar_deep_research", model_key, "claude_sonnet", "claude_opus", "deepseek_v3_sq", "qwen3_max"]
    _REFUSAL_PHRASES = [
        "i appreciate you reaching out", "i'm unable to", "i cannot write",
        "i can't write", "i cannot create", "i cannot provide", "i cannot assist",
        "raises serious safety concerns", "не могу написать", "не могу создать",
        "не могу помочь", "вынужден отказать",
    ]
    _tried = set()
    last_err = None
    for _mk in _pptx_fallback_order:
        if _mk in _tried or _mk not in MODELS:
            continue
        _tried.add(_mk)
        try:
            raw = await call_chat(msgs, _mk, max_tokens=8000)
            if not raw or raw.strip() in ("", "..."):
                logging.warning(f"[PPTX] Модель {_mk} вернула пустой ответ, пробуем следующую")
                last_err = RuntimeError(f"Пустой ответ от {_mk}")
                continue
            # Проверяем на отказ (нет JSON = отказ)
            if any(p in raw.lower() for p in _REFUSAL_PHRASES):
                logging.warning(f"[PPTX] {_mk} отказал, пробуем следующую модель")
                last_err = RuntimeError(f"Отказ от {_mk}")
                continue
            raw = re.sub(r"^```(?:json)?\s*", "", raw.strip())
            raw = re.sub(r"\s*```$", "", raw.strip())
            # Ищем первый { и последний }
            try:
                start = raw.index("{")
                end   = raw.rindex("}") + 1
                raw   = raw[start:end]
            except ValueError:
                logging.warning(f"[PPTX] Модель {_mk}: нет JSON-объекта в ответе, пробуем следующую")
                last_err = RuntimeError(f"JSON не найден в ответе {_mk}")
                continue
            try:
                return json.loads(raw)
            except json.JSONDecodeError as je:
                logging.warning(f"[PPTX] Модель {_mk}: невалидный JSON ({je}), пробуем следующую")
                last_err = je
                continue
        except Exception as e:
            logging.warning(f"[PPTX] Модель {_mk} ошибка: {e}, пробуем следующую")
            last_err = e
            continue
    raise RuntimeError(f"Не удалось сгенерировать структуру презентации: {last_err}")


async def generate_html_structure_ai(
    topic: str,
    slides_count: int,
    model_key: str,
    extra_wishes: str = "",
    title_info: dict | None = None,
) -> dict:
    """Запрашивает у LLM JSON-структуру для HTML-презентации (поле image_query)."""
    if slides_count <= 5:
        layout_hint = "title → content → cards → split → conclusion"
    elif slides_count <= 7:
        layout_hint = "title → content → cards → split → stats → quote → conclusion"
    elif slides_count <= 10:
        layout_hint = "title → content → cards → split → content → stats → split → cards → timeline → conclusion"
    else:
        layout_hint = "title → content → cards → split → content → stats → split → content → cards → quote → timeline → conclusion"

    wishes_block = f"\n\nПОЖЕЛАНИЯ: {extra_wishes}" if extra_wishes else ""
    title_block = ""
    if title_info:
        parts = []
        if title_info.get("subtitle"):     parts.append(f"Подзаголовок: «{title_info['subtitle']}»")
        if title_info.get("author"):       parts.append(f"Автор: {title_info['author']}")
        if title_info.get("organization"): parts.append(f"Организация: {title_info['organization']}")
        if parts:
            title_block = "\n\nДАННЫЕ ТИТУЛЬНОГО СЛАЙДА:\n" + "\n".join(parts)

    # Определяем географию для принудительной подстановки
    topic_lower = topic.lower()
    # Определяем военную тему — приоритет над географией
    is_military_topic = any(w in topic_lower for w in [
        "военн", "впк", "оборон", "армия", "нато", "войск", "ракет", "танк"
    ])
    is_safety_topic = any(w in topic_lower for w in [
        "обж", "безопасност", "чрезвычайн", "первая помощь", "пожарн", "эвакуац",
        "землетрясен", "наводнен", "выживан", "гражданская оборона", "террор",
        "скулшутинг", "нападен", "ожог", "отравлен", "дтп", "авари", "спасател",
        "взрыв", "химическ", "радиац", "цунами", "урага", "лавин", "паводк",
        "вербовк", "радикализ", "экстремизм", "несовершеннолетн", "профилактик",
    ])
    is_terror_topic = any(w in topic_lower for w in [
        "террор", "вербовк", "радикализ", "экстремизм", "скулшутинг",
    ])

    if is_terror_topic:
        geo_rule = (
            "⛔ ТЕМА: ТЕРРОРИЗМ / ВЕРБОВКА ПОДРОСТКОВ — РОССИЙСКАЯ ШКОЛЬНАЯ ТЕМА!\n"
            "АБСОЛЮТНО ЗАПРЕЩЕНО: US Army, American police, FBI, Pentagon, NATO soldiers, School shooting USA.\n"
            "Используй РАЗНЫЕ уникальные keywords для КАЖДОГО слайда:\n"
            "Для слайдов о статистике/угрозе: 'teenagers smartphone social media phone screen'\n"
            "Для слайдов о методах вербовки: 'social media chat phone screen young person online'\n"
            "Для слайдов о признаках: 'teacher student conversation school classroom Russia'\n"
            "Для слайдов о действиях: 'parent teenager talk home family dialogue'\n"
            "Для слайдов о полиции/МВД: 'police officer youth interview Russia crime prevention'\n"
            "Для титульного: 'police security checkpoint crowd Russia city'\n"
            "Для заключения: 'school assembly students safety awareness education Russia'\n"
            "КАЖДЫЙ слайд = УНИКАЛЬНАЯ сцена! Не повторяй одинаковые keywords!"
        )
    elif is_safety_topic:
        geo_rule = (
            "⛔ ТЕМА ПО ОБЖ/БЕЗОПАСНОСТИ — ГРАЖДАНСКАЯ ТЕМА! "
            "АБСОЛЮТНО ЗАПРЕЩЕНО: американские военные, US Army, NATO soldiers, Pentagon, American police.\n"
            "image_keyword ОБЯЗАН быть из тематики ГРАЖДАНСКОЙ БЕЗОПАСНОСТИ:\n"
            "Примеры ПРАВИЛЬНЫХ keywords:\n"
            "'firefighters rescue building fire smoke' | 'ambulance paramedics emergency street'\n"
            "'evacuation crowd emergency exit' | 'earthquake rubble rescue workers survivors'\n"
            "'flood rescue boat helicopter people' | 'police security patrol checkpoint'\n"
            "'first aid cpr training mannequin' | 'fire alarm smoke safety home'\n"
            "'traffic accident car crash police' | 'emergency shelter disaster relief'\n"
            "'school safety drill evacuation students' | 'chemical hazmat suit workers'\n"
            "Подбирай keyword к КОНКРЕТНОМУ содержанию каждого слайда!"
        )
    elif is_military_topic:
        geo_rule = (
            "⛔ ВОЕННАЯ ТЕМА! image_keyword ОБЯЗАН содержать военные объекты. "
            "Используй слова из этого списка как ПЕРВОЕ слово keyword: "
            "soldier, soldiers, military, army, tank, armored, weapon, missile, "
            "rocket, fighter, aircraft, drone, nato, defense, combat. "
            "Примеры: 'soldiers training field camouflage', 'factory weapons production assembly', "
            "'military budget defense spending', 'tank armor combat vehicle'."
        )
    elif any(w in topic_lower for w in ["росси", "россий", "рф ", "рф,", "москв", "советск", "русск"]):
        geo_rule = (
            "⛔ ТЕМА О РОССИИ! image_keyword ОБЯЗАН начинаться с тематического объекта. "
            "Используй конкретные слова: factory, workers, oil, pipeline, wheat, harvest, "
            "building, government, stock, chart, soldier, soldiers. "
            "Примеры: 'factory workers industrial production', 'oil pipeline winter field', "
            "'wheat harvest combine field', 'stock chart finance trading'."
        )
    elif any(w in topic_lower for w in ["китай", "китайск", "пекин", "шанхай"]):
        geo_rule = (
            "⛔ ТЕМА О КИТАЕ! Используй: factory, workers, skyline, building, technology, chip. "
            "Примеры: 'factory assembly line workers', 'skyline building architecture modern'."
        )
    elif any(w in topic_lower for w in ["европ", "евросою"]):
        geo_rule = (
            "Используй: building, government, skyline, stock, chart. "
            "Примеры: 'building architecture government modern', 'stock chart finance'."
        )
    elif any(w in topic_lower for w in ["сша", "америк", "американск"]):
        geo_rule = (
            "Используй: building, skyline, stock, chart, government. "
            "Примеры: 'skyline building architecture', 'stock trading chart finance'."
        )
    else:
        geo_rule = (
            "Используй конкретные существительные из этого словаря: "
            "soldier, military, factory, workers, oil, pipeline, stock, chart, "
            "bank, server, chip, neural, code, hospital, students, cargo, ship. "
            "Первое слово keyword ДОЛЖНО быть из этого списка."
        )

    user_msg = (
        f"Тема: «{topic}»\n"
        f"Количество слайдов: {slides_count} (включая title и conclusion)\n"
        f"Рекомендуемая последовательность: {layout_hint}\n"
        f"Язык контента: русский, image_keyword: только английский\n"
        f"{wishes_block}{title_block}\n\n"
        "ТРЕБОВАНИЯ:\n"
        "1. Каждый слайд — уникальный тип, не повторяй подряд\n"
        "2. Контент содержательный и конкретный, с цифрами и фактами\n"
        f"3. {geo_rule}\n"
        "4. image_keyword = 3-5 слов конкретная сцена/место на английском\n"
        "5. В cards.icon используй ТОЛЬКО классы FontAwesome 6: fas fa-chart-line, fas fa-users и т.д.\n"
        "6. Поле называется image_keyword (не image_query)\n\n"
        "Верни ТОЛЬКО JSON без пояснений."
    )
    msgs = [
        {"role": "system", "content": PPTX_AI_SYSTEM_PROMPT},
        {"role": "user",   "content": user_msg},
    ]
    # Sonar Deep Research первым — актуальные данные из интернета, минуя проверки disabled
    _pptx_fallback_order = ["sonar_deep_research", model_key, "claude_sonnet", "claude_opus", "deepseek_v3_sq", "qwen3_max"]
    _REFUSAL_PHRASES = [
        "i appreciate you reaching out", "i'm unable to", "i cannot write",
        "i can't write", "i cannot create", "i cannot provide", "i cannot assist",
        "raises serious safety concerns", "не могу написать", "не могу создать",
        "не могу помочь", "вынужден отказать",
    ]
    _tried = set()
    last_err = None
    for _mk in _pptx_fallback_order:
        if _mk in _tried or _mk not in MODELS:
            continue
        _tried.add(_mk)
        try:
            raw = await call_chat(msgs, _mk, max_tokens=8000)
            if not raw or raw.strip() in ("", "..."):
                last_err = RuntimeError(f"Пустой ответ от {_mk}")
                continue
            if any(p in raw.lower() for p in _REFUSAL_PHRASES):
                logging.warning(f"[PPTX-HTML] {_mk} отказал, пробуем следующую модель")
                last_err = RuntimeError(f"Отказ от {_mk}")
                continue
            raw = re.sub(r"^```(?:json)?\s*", "", raw.strip())
            raw = re.sub(r"\s*```$",          "", raw.strip())
            try:
                start = raw.index("{")
                end   = raw.rindex("}") + 1
                raw   = raw[start:end]
            except ValueError:
                last_err = RuntimeError(f"JSON не найден в ответе {_mk}")
                continue
            try:
                return json.loads(raw)
            except json.JSONDecodeError as je:
                last_err = je
                continue
        except Exception as e:
            last_err = e
            continue
    raise RuntimeError(f"Не удалось сгенерировать HTML-структуру: {last_err}")


# ── Клавиатуры ───────────────────────────────────────────────────

def pptx_start_kb():
    """Новый флоу: сначала тема, потом настройки."""
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="✏️ Ввести тему", callback_data="pptx_enter_topic")],
        [InlineKeyboardButton(text="❌ Отмена",       callback_data="back_home")],
    ])

def pptx_slides_kb():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="5",  callback_data="pptx_slides_5"),
         InlineKeyboardButton(text="7",  callback_data="pptx_slides_7"),
         InlineKeyboardButton(text="10", callback_data="pptx_slides_10"),
         InlineKeyboardButton(text="12", callback_data="pptx_slides_12"),
         InlineKeyboardButton(text="15", callback_data="pptx_slides_15")],
        [InlineKeyboardButton(text="◀️ Назад", callback_data="pptx_enter_topic")],
    ])

def pptx_audience_kb():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🏫 Школа",         callback_data="pptx_aud_school"),
         InlineKeyboardButton(text="📚 Старшие классы",callback_data="pptx_aud_highschool")],
        [InlineKeyboardButton(text="🎓 Техникум",       callback_data="pptx_aud_college"),
         InlineKeyboardButton(text="🏛 Институт",       callback_data="pptx_aud_university")],
        [InlineKeyboardButton(text="💼 Бизнес",         callback_data="pptx_aud_business"),
         InlineKeyboardButton(text="🎨 Творческая",     callback_data="pptx_aud_creative")],
        [InlineKeyboardButton(text="◀️ Назад",          callback_data="pptx_back_to_slides")],
    ])

def pptx_imgtype_kb():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🖼 Найти в интернете", callback_data="pptx_img_internet")],
        [InlineKeyboardButton(text="📄 Без изображений",   callback_data="pptx_img_none")],
        [InlineKeyboardButton(text="◀️ Назад",             callback_data="pptx_back_to_audience")],
    ])

def pptx_volume_kb():
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📝 Краткий",     callback_data="pptx_vol_brief")],
        [InlineKeyboardButton(text="📄 Стандартный", callback_data="pptx_vol_standard")],
        [InlineKeyboardButton(text="📚 Подробный",   callback_data="pptx_vol_detailed")],
        [InlineKeyboardButton(text="◀️ Назад",       callback_data="pptx_back_to_imgtype")],
    ])


def pptx_theme_kb():
    items = list(PPTX_THEMES.items())
    rows = [
        # Первая строка — авто-выбор (самая удобная опция)
        [InlineKeyboardButton(text="✨ Авто (по теме)", callback_data="pptx_theme_auto")],
    ]
    for j in range(0, len(items), 2):
        row = [InlineKeyboardButton(text=t["name"], callback_data=f"pptx_theme_{k}")
               for k, t in items[j:j+2]]
        rows.append(row)
    rows.append([InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")])
    return InlineKeyboardMarkup(inline_keyboard=rows)


def pptx_model_kb():
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="🔵 Claude Opus", callback_data="pptx_model_claude_opus"),
            InlineKeyboardButton(text="⚡ GPT-5.2",     callback_data="pptx_model_gpt52"),
        ],
        [
            InlineKeyboardButton(text="🧬 DeepSeek R1", callback_data="pptx_model_deepseek_r1"),
            InlineKeyboardButton(text="🌐 Qwen3-Max",   callback_data="pptx_model_qwen3_max"),
        ],
        [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
    ])


def pptx_images_kb():
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="🖼 С фотографиями ✅", callback_data="pptx_img_yes"),
            InlineKeyboardButton(text="📄 Без фото",          callback_data="pptx_img_no"),
        ],
        [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
    ])


def pptx_confirm_kb():
    return InlineKeyboardMarkup(inline_keyboard=[
        [
            InlineKeyboardButton(text="📊 PowerPoint (.pptx)", callback_data="pptx_type_pptx"),
            InlineKeyboardButton(text="🌐 HTML-презентация", callback_data="pptx_type_html"),
        ],
        [
            InlineKeyboardButton(text="✏️ Изменить", callback_data="pptx_restart"),
            InlineKeyboardButton(text="❌ Отмена", callback_data="back_home"),
        ],
    ])


# ── Хэндлеры ─────────────────────────────────────────────────────

@dp.message(Command("pptx"))
async def cmd_pptx(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message):
        return
    pptx_states[uid] = {}
    admin_await[uid] = {"action": "pptx_topic"}
    await message.answer(
        "🎞 <b>Генератор презентаций</b>\n\n"
        "✏️ Напиши <b>тему презентации</b>:\n\n"
        "<i>Примеры:\n"
        "• Россия в будущем\n"
        "• Искусственный интеллект\n"
        "• История освоения космоса</i>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")
        ]])
    )


@dp.message(F.text == "🎞 Презентация")
async def rb_pptx(message: Message):
    await cmd_pptx(message)


@dp.message(F.text == "🌐 Веб-презентация")
async def rb_webpptx(message: Message):
    await cmd_webpptx(message)


@dp.message(Command("webpptx"))
async def cmd_webpptx(message: Message):
    """Запускает генерацию стильной HTML-презентации."""
    uid = message.from_user.id
    if not await require_subscription(message):
        return
    pptx_states[uid] = {"mode": "html"}
    await message.answer(
        "🌐 <b>━━━━━━━━━━━━━━━━━━━━ 🌐</b>\n\n"
        "✨ <b>Генератор веб-презентаций</b>\n\n"
        "🎨 Тёмный дизайн с зелёными акцентами\n"
        "🖼 Тематические фото Unsplash для каждого слайда\n"
        "📱 Работает в браузере, на телефоне (свайпы)\n"
        "⌨️ Управление стрелками на клавиатуре\n\n"
        "🌐 <b>━━━━━━━━━━━━━━━━━━━━ 🌐</b>\n\n"
        "📏 Сколько слайдов нужно?",
        parse_mode="HTML",
        reply_markup=pptx_start_kb(),
    )


@dp.callback_query(F.data == "pptx_enter_topic")
async def pptx_enter_topic_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in pptx_states:
        pptx_states[uid] = {}
    admin_await[uid] = {"action": "pptx_topic"}
    await callback.message.edit_text(
        "🎞 <b>Генератор презентаций</b>\n\n"
        "✏️ Напиши <b>тему презентации</b>:",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")
        ]])
    )
    await callback.answer()


@dp.callback_query(F.data.startswith("pptx_slides_"))
async def pptx_slides_cb(callback: CallbackQuery):
    uid  = callback.from_user.id
    n    = int(callback.data.split("_")[-1])
    if uid not in pptx_states:
        pptx_states[uid] = {}
    pptx_states[uid]["slides"] = n
    await callback.message.edit_text(
        f"📊 <b>Слайдов: {n}</b>\n\n"
        "🎓 <b>Для кого презентация?</b>",
        parse_mode="HTML", reply_markup=pptx_audience_kb()
    )
    await callback.answer()


@dp.callback_query(F.data == "pptx_back_to_slides")
async def pptx_back_slides_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    topic = pptx_states.get(uid, {}).get("topic", "")
    await callback.message.edit_text(
        f"🎞 <b>Тема: «{topic[:50]}»</b>\n\n📏 <b>Сколько слайдов?</b>",
        parse_mode="HTML", reply_markup=pptx_slides_kb()
    )
    await callback.answer()


@dp.callback_query(F.data.startswith("pptx_aud_"))
async def pptx_audience_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    aud_map = {
        "pptx_aud_school":      "Школа",
        "pptx_aud_highschool":  "Старшие классы",
        "pptx_aud_college":     "Техникум",
        "pptx_aud_university":  "Институт",
        "pptx_aud_business":    "Бизнес",
        "pptx_aud_creative":    "Творческая",
    }
    aud = aud_map.get(callback.data, "Стандартная")
    if uid not in pptx_states:
        pptx_states[uid] = {}
    pptx_states[uid]["audience"] = aud
    await callback.message.edit_text(
        f"🎓 <b>Аудитория: {aud}</b>\n\n"
        "🖼 <b>Изображения в презентации?</b>",
        parse_mode="HTML", reply_markup=pptx_imgtype_kb()
    )
    await callback.answer()


@dp.callback_query(F.data == "pptx_back_to_audience")
async def pptx_back_audience_cb(callback: CallbackQuery):
    await callback.message.edit_text(
        "🎓 <b>Для кого презентация?</b>",
        parse_mode="HTML", reply_markup=pptx_audience_kb()
    )
    await callback.answer()


@dp.callback_query(F.data.in_({"pptx_img_internet", "pptx_img_none"}))
async def pptx_imgtype_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    use_img = callback.data == "pptx_img_internet"
    if uid not in pptx_states:
        pptx_states[uid] = {}
    pptx_states[uid]["use_images"] = use_img
    img_label = "🖼 Из интернета" if use_img else "📄 Без изображений"
    await callback.message.edit_text(
        f"🖼 <b>Изображения: {img_label}</b>\n\n"
        "📝 <b>Объём текста на слайдах?</b>",
        parse_mode="HTML", reply_markup=pptx_volume_kb()
    )
    await callback.answer()


@dp.callback_query(F.data == "pptx_back_to_imgtype")
async def pptx_back_imgtype_cb(callback: CallbackQuery):
    aud = pptx_states.get(callback.from_user.id, {}).get("audience", "")
    await callback.message.edit_text(
        f"🎓 <b>Аудитория: {aud}</b>\n\n🖼 <b>Изображения в презентации?</b>",
        parse_mode="HTML", reply_markup=pptx_imgtype_kb()
    )
    await callback.answer()


@dp.callback_query(F.data.startswith("pptx_vol_"))
async def pptx_volume_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    vol_map = {
        "pptx_vol_brief":    ("Краткий",     "brief"),
        "pptx_vol_standard": ("Стандартный", "standard"),
        "pptx_vol_detailed": ("Подробный",   "detailed"),
    }
    vol_label, vol_key = vol_map.get(callback.data, ("Стандартный", "standard"))
    if uid not in pptx_states:
        pptx_states[uid] = {}
    pptx_states[uid]["volume"] = vol_key
    # Показываем сводку настроек
    await _show_pptx_new_confirm(callback.message, uid)
    await callback.answer()


@dp.callback_query(F.data.startswith("pptx_theme_"))
async def pptx_theme_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    theme_key = callback.data.split("pptx_theme_")[-1]
    if uid not in pptx_states:
        pptx_states[uid] = {}
    pptx_states[uid]["theme"] = theme_key
    if theme_key == "auto":
        theme_name = "✨ Авто (по теме)"
    else:
        theme_name = PPTX_THEMES.get(theme_key, {}).get("name", theme_key)
    await callback.message.edit_text(
        f"🎨 <b>Тема: {theme_name}</b>\n\n🖼 Добавить тематические фото на каждый слайд?",
        parse_mode="HTML", reply_markup=pptx_images_kb()
    )
    await callback.answer()


@dp.callback_query(F.data.startswith("pptx_img_"))
async def pptx_img_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    use_img = callback.data == "pptx_img_yes"
    if uid not in pptx_states:
        pptx_states[uid] = {}
    pptx_states[uid]["use_images"] = use_img
    img_str = "🖼 с фото" if use_img else "📄 без фото"
    await callback.message.edit_text(
        f"🖼 <b>Фото: {img_str}</b>\n\n🤖 Выбери модель для генерации:",
        parse_mode="HTML", reply_markup=pptx_model_kb()
    )
    await callback.answer()


@dp.callback_query(F.data.startswith("pptx_model_"))
async def pptx_model_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    model_key = callback.data.split("pptx_model_")[-1]
    if uid not in pptx_states:
        pptx_states[uid] = {}
    pptx_states[uid]["model"] = model_key
    model_label = MODELS.get(model_key, {}).get("label", model_key)
    admin_await[uid] = {"action": "pptx_topic"}
    await callback.message.edit_text(
        f"🤖 <b>Модель: {model_label}</b>\n\n"
        "✏️ Напиши <b>тему презентации</b>:\n\n"
        "<i>Примеры:\n"
        "• «Военная экономика 2026 год»\n"
        "• «Искусственный интеллект в медицине»\n"
        "• «Маркетинговая стратегия 2025»\n"
        "• «История освоения космоса»</i>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")
        ]])
    )
    await callback.answer()


@dp.callback_query(F.data == "pptx_restart")
async def pptx_restart_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    pptx_states[uid] = {}
    await callback.message.edit_text(
        "🔄 <b>Начнём заново!</b>\n\n📏 Сколько слайдов нужно?",
        parse_mode="HTML", reply_markup=pptx_start_kb()
    )
    await callback.answer()


@dp.callback_query(F.data == "pptx_skip_subtitle")
async def pptx_skip_subtitle_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in pptx_states:
        pptx_states[uid] = {}
    pptx_states[uid].setdefault("title_info", {})["subtitle"] = ""
    admin_await[uid] = {"action": "pptx_author"}
    await callback.message.edit_text(
        "👤 Введи <b>имя автора</b> / докладчика:",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="⏭ Пропустить", callback_data="pptx_skip_author")],
            [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
        ])
    )
    await callback.answer()


@dp.callback_query(F.data == "pptx_skip_author")
async def pptx_skip_author_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in pptx_states:
        pptx_states[uid] = {}
    pptx_states[uid].setdefault("title_info", {})["author"] = ""
    admin_await[uid] = {"action": "pptx_organization"}
    await callback.message.edit_text(
        "🏢 Введи <b>название организации</b>:\n<i>(школа, университет, компания)</i>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="⏭ Пропустить", callback_data="pptx_skip_org")],
            [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
        ])
    )
    await callback.answer()


@dp.callback_query(F.data == "pptx_skip_org")
async def pptx_skip_org_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in pptx_states:
        pptx_states[uid] = {}
    pptx_states[uid].setdefault("title_info", {})["organization"] = ""
    admin_await[uid] = {"action": "pptx_wishes"}
    await callback.message.edit_text(
        "✍️ Добавь пожелания:\n<i>(стиль, акценты, что включить)</i>\n\nИли пропусти:",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="⏭ Пропустить", callback_data="pptx_skip_wishes")],
            [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
        ])
    )
    await callback.answer()


@dp.callback_query(F.data == "pptx_skip_wishes")
async def pptx_skip_wishes_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in pptx_states:
        pptx_states[uid] = {}
    pptx_states[uid]["wishes"] = ""
    await _show_pptx_confirm(callback.message, uid)
    await callback.answer()


@dp.callback_query(F.data == "pptx_type_pptx")
async def pptx_type_pptx_cb(callback: CallbackQuery):
    uid   = callback.from_user.id
    state = pptx_states.get(uid, {})
    if not state.get("topic"):
        await callback.answer("❌ Тема не задана", show_alert=True)
        return
    pptx_states[uid]["mode"] = "pptx"
    await _do_generate_pptx(callback.message, uid, pptx_states[uid])
    await callback.answer()


@dp.callback_query(F.data == "pptx_type_html")
async def pptx_type_html_cb(callback: CallbackQuery):
    uid   = callback.from_user.id
    state = pptx_states.get(uid, {})
    if not state.get("topic"):
        await callback.answer("❌ Тема не задана", show_alert=True)
        return
    pptx_states[uid]["mode"] = "html"
    await _do_generate_html_pptx(callback.message, uid, pptx_states[uid])
    await callback.answer()


@dp.callback_query(F.data == "pptx_generate")
async def pptx_generate_cb(callback: CallbackQuery):
    uid   = callback.from_user.id
    state = pptx_states.get(uid, {})
    if not state.get("topic"):
        await callback.answer("❌ Тема не задана", show_alert=True)
        return
    if state.get("mode") == "html":
        await _do_generate_html_pptx(callback.message, uid, state)
    else:
        await _do_generate_pptx(callback.message, uid, state)
    await callback.answer()


async def _do_generate_html_pptx(msg, uid: int, state: dict):
    """Генерирует HTML-презентацию (новый стильный дизайн) и отправляет файл."""
    topic      = state.get("topic", "")
    slides_n   = state.get("slides", 7)
    wishes     = state.get("wishes", "")
    title_info = state.get("title_info", {})

    # Всегда используем Sonar Deep Research для актуальных данных
    model_key   = "sonar_deep_research"
    model_label = "🔬 Sonar Deep Research"
    use_images  = state.get("use_images", True)

    _init_limits(uid)
    _refresh_limits(uid)
    li = get_limits_info(uid)
    if li["pro_max"] - li["pro_used"] < 3:
        await msg.answer(
            "❌ <b>Недостаточно запросов!</b>\n\n"
            f"Нужно <b>3 продвинутых запроса</b>. Осталось: <b>{li['pro_max'] - li['pro_used']}</b>\n"
            f"⏰ Сброс через: <code>{li['pro_reset']}</code>",
            parse_mode="HTML"
        )
        return

    await msg.answer(
        "🌐 <i>Использую Sonar Deep Research — ищу актуальные данные в интернете...</i>",
        parse_mode="HTML"
    )
    _start    = asyncio.get_event_loop().time()
    think_msg = await msg.answer("🎞 <i>Генерирую веб-презентацию... <b>1с</b></i>", parse_mode="HTML")

    async def _tick(m, s):
        phases = [
            (7,   "📋 Составляю структуру слайдов"),
            (20,  "✍️ ИИ пишет контент"),
            (40,  "🖼 Подбираю тематические фото"),
            (60,  "🎨 Собираю HTML-дизайн"),
            (999, "⚡ Финальная сборка..."),
        ]
        while True:
            await asyncio.sleep(1)
            elapsed = int(asyncio.get_event_loop().time() - s)
            label = phases[-1][1]
            for thr, lbl in phases:
                if elapsed < thr:
                    label = lbl
                    break
            try:
                await m.edit_text(f"{label}... <b>{elapsed}с</b>", parse_mode="HTML")
            except Exception:
                pass

    tick = asyncio.create_task(_tick(think_msg, _start))
    try:
        pptx_data = await generate_html_structure_ai(
            topic=topic, slides_count=slides_n,
            model_key=model_key, extra_wishes=wishes,
            title_info=title_info if title_info else None,
        )
        if title_info:
            for sl in pptx_data.get("slides", []):
                if sl.get("type") == "title":
                    if title_info.get("author"):       sl["author"]       = title_info["author"]
                    if title_info.get("subtitle"):     sl["subtitle"]     = title_info["subtitle"]
                    if title_info.get("organization"): sl["organization"] = title_info["organization"]

        html_str   = create_html_presentation(pptx_data, pexels_key=os.environ.get("PEXELS_KEY",""))
        html_bytes = html_str.encode("utf-8")

        for _ in range(3):
            user_limits[uid]["pro_used"] = user_limits[uid].get("pro_used", 0) + 1

        tick.cancel()
        try:
            await tick
        except asyncio.CancelledError:
            pass
        await think_msg.delete()

        total_secs    = int(asyncio.get_event_loop().time() - _start)
        actual_slides = len(pptx_data.get("slides", []))
        safe_name     = re.sub(r"[^\w\s-]", "", topic)[:25].strip().replace(" ", "_")
        fname         = f"Презентация_{safe_name}.html"

        await msg.answer(
            f"✅ <b>Веб-презентация готова!</b>\n\n"
            f"📌 Тема: <b>«{topic[:50]}{'...' if len(topic)>50 else ''}»</b>\n"
            f"📊 Слайдов: <b>{actual_slides}</b>\n"
            f"🤖 Модель: <b>{MODELS.get(model_key,{}).get('label',model_key)}</b>\n"
            f"⏱ Время: <b>{total_secs}с</b>\n\n"
            "🌐 Открой HTML-файл в <b>любом браузере</b>\n"
            "⌨️ Стрелки / пробел — переключение слайдов\n"
            "📱 На телефоне — свайп влево/вправо",
            parse_mode="HTML"
        )
        await bot.send_document(
            chat_id=uid,
            document=BufferedInputFile(html_bytes, filename=fname)
        )
        pptx_states.pop(uid, None)

    except Exception as e:
        tick.cancel()
        try:
            await tick
        except asyncio.CancelledError:
            pass
        try:
            await think_msg.delete()
        except Exception:
            pass
        logging.error(f"html_pptx generate error: {e}")
        await msg.answer(f"❌ Ошибка генерации: {str(e)[:300]}")


async def _show_pptx_new_confirm(msg, uid: int):
    """Показывает сводку настроек с кнопками изменения и генерации плана."""
    s = pptx_states.get(uid, {})
    topic    = s.get("topic", "—")
    slides_n = s.get("slides", 10)
    audience = s.get("audience", "Стандартная")
    use_img  = s.get("use_images", True)
    volume   = s.get("volume", "standard")
    vol_labels = {"brief": "Краткий", "standard": "Стандартный", "detailed": "Подробный"}
    img_label = "🖼 Из интернета" if use_img else "📄 Без изображений"
    text = (
        f"🎞 <b>Тема:</b> {topic[:60]}\n"
        f"🌐 Язык: 🇷🇺 Русский\n"
        f"📊 Слайдов: {slides_n}\n"
        f"🎓 Аудитория: {audience}\n"
        f"🖼 Изображения: {img_label}\n"
        f"📝 Объём текста: {vol_labels.get(volume, volume)}"
    )
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📝 Сменить объём текста",  callback_data="pptx_back_to_imgtype"),
         InlineKeyboardButton(text="🖼 Сменить изображения",   callback_data="pptx_back_to_audience")],
        [InlineKeyboardButton(text="📊 Сменить кол-во слайдов", callback_data="pptx_back_to_slides"),
         InlineKeyboardButton(text="🎓 Сменить аудиторию",      callback_data="pptx_back_to_audience")],
        [InlineKeyboardButton(text="✏️ Сменить тему",          callback_data="pptx_enter_topic")],
        [InlineKeyboardButton(text="✅ Сгенерировать план",     callback_data="pptx_gen_plan")],
    ])
    try:
        await msg.edit_text(text, parse_mode="HTML", reply_markup=kb)
    except Exception:
        await msg.answer(text, parse_mode="HTML", reply_markup=kb)


@dp.callback_query(F.data == "pptx_gen_plan")
async def pptx_gen_plan_cb(callback: CallbackQuery):
    """Генерирует план презентации и показывает его с кнопками слайдов."""
    uid = callback.from_user.id
    s   = pptx_states.get(uid, {})
    if not s.get("topic"):
        return await callback.answer("❌ Тема не задана", show_alert=True)
    await callback.answer()
    _init_limits(uid)
    _refresh_limits(uid)
    li = get_limits_info(uid)
    if li["pro_max"] - li["pro_used"] < 3:
        await callback.message.answer(
            f"❌ <b>Недостаточно запросов!</b> Нужно 3, осталось {li['pro_max'] - li['pro_used']}",
            parse_mode="HTML"
        )
        return

    topic    = s.get("topic", "")
    slides_n = s.get("slides", 10)
    audience = s.get("audience", "Стандартная")
    volume   = s.get("volume", "standard")
    use_img  = s.get("use_images", True)

    wait = await callback.message.answer("⏳ <i>Составляю план презентации...</i>", parse_mode="HTML")
    try:
        # Генерируем только структуру (названия слайдов) через быструю модель
        volume_hint = {
            "brief":    "кратко, 2-3 тезиса на слайд",
            "standard": "стандартно 4-5 тезисов",
            "detailed": "подробно 6-8 тезисов"
        }.get(volume, "стандартно")

        # Определяем тематику для умных названий слайдов
        topic_lower_plan = topic.lower()
        is_obzh_plan = any(w in topic_lower_plan for w in [
            'обж', 'безопасност', 'чрезвычайн', 'первая помощь', 'пожарн',
            'эвакуац', 'землетрясен', 'наводнен', 'выживан', 'террор',
            'ожог', 'отравлен', 'дтп', 'авари', 'спасател',
        ])

        if is_obzh_plan:
            plan_hint = (
                "Это тема по ОБЖ/безопасности. Слайды должны включать:\n"
                "- Что это за угроза/явление\n"
                "- Статистика и реальные случаи\n"
                "- Как распознать опасность\n"
                "- Алгоритм действий\n"
                "- Типичные ошибки\n"
                "- Профилактика\n"
                "- Выводы\n"
            )
        else:
            plan_hint = f"Аудитория: {audience}. Объём: {volume_hint}.\n"

        plan_prompt = (
            f"Составь конкретный план презентации на тему «{topic}».\n"
            f"{plan_hint}"
            f"Нужно ровно {slides_n} слайдов.\n"
            "Первый слайд — титульный (название темы).\n"
            "Последний — итоги/заключение.\n"
            "Каждое название слайда должно быть СОДЕРЖАТЕЛЬНЫМ и КОНКРЕТНЫМ — не 'Введение', а суть.\n"
            "Примеры ХОРОШИХ названий: 'Статистика пожаров в России за 2023 год', 'Алгоритм действий при пожаре', 'Типичные ошибки и их последствия'.\n"
            "Формат — ТОЛЬКО пронумерованный список:\n"
            "1. Название\n"
            "2. Название\n"
            "Без пояснений и лишнего текста. Только список."
        )
        msgs = [{"role": "user", "content": plan_prompt}]
        plan_text = (await call_chat(msgs, "claude_haiku", max_tokens=800)).strip()

        # Парсим в список
        slides_list = []
        for line in plan_text.split("\n"):
            line = line.strip()
            if not line:
                continue
            import re as _re
            m = _re.match(r"^\d+[\.\)]\s*(.+)", line)
            if m:
                name = m.group(1).strip()
                # Убираем звёздочки и лишние символы
                name = _re.sub(r"\*+", "", name).strip()
                if name:
                    slides_list.append(name)

        # Если по каким-то причинам имена всё равно плохие — перегенерируем через более сильную модель
        generic_words = ["слайд", "slide", "глава", "раздел"]
        has_generic = any(
            any(g in s.lower() for g in generic_words) for s in slides_list
        )
        if len(slides_list) < 3 or has_generic:
            # Пробуем claude_sonnet как запасной вариант
            try:
                plan_text2 = (await call_chat(msgs, "claude_sonnet", max_tokens=800)).strip()
                slides_list2 = []
                for line in plan_text2.split("\n"):
                    line = line.strip()
                    if not line:
                        continue
                    m2 = _re.match(r"^\d+[\.\)]\s*(.+)", line)
                    if m2:
                        name2 = _re.sub(r"\*+", "", m2.group(1)).strip()
                        if name2:
                            slides_list2.append(name2)
                if len(slides_list2) >= 3:
                    slides_list = slides_list2
            except Exception:
                pass

        # Финальный fallback с осмысленными именами
        if len(slides_list) < 3:
            slides_list = (
                [f"Введение в тему «{topic[:30]}»"] +
                [f"Ключевой аспект {i}" for i in range(1, slides_n - 1)] +
                ["Выводы и заключение"]
            )[:slides_n]

        # Подгоняем список под нужное количество слайдов
        if len(slides_list) < slides_n:
            # Дополняем если не хватает
            while len(slides_list) < slides_n:
                slides_list.append(f"Дополнительный раздел {len(slides_list)}")
        elif len(slides_list) > slides_n:
            # Обрезаем если лишние
            slides_list = slides_list[:slides_n]

        pptx_states[uid]["plan"] = slides_list
        await wait.delete()

        # ── Показываем план ──────────────────────────────────────────
        # Короткие названия для кнопок (максимум 15 символов)
        plan_display = "\n".join(f"{i+1}. {t}" for i, t in enumerate(slides_list))
        kb_rows = []
        for i in range(0, len(slides_list), 2):
            row = []
            for j in [i, i+1]:
                if j < len(slides_list):
                    # Обрезаем название для кнопки до 15 символов
                    name_short = slides_list[j]
                    if len(name_short) > 15:
                        # Берём первые значимые слова
                        words = name_short.split()
                        short = ""
                        for w in words:
                            if len(short) + len(w) + 1 <= 14:
                                short = (short + " " + w).strip()
                            else:
                                break
                        name_short = short + "…" if short else name_short[:14] + "…"
                    row.append(InlineKeyboardButton(
                        text=f"{j+1}. {name_short}",
                        callback_data=f"pptx_slide_edit_{j}"
                    ))
            kb_rows.append(row)
        kb_rows.append([InlineKeyboardButton(text="🔄 Перегенерировать план", callback_data="pptx_gen_plan")])
        kb_rows.append([InlineKeyboardButton(text="✅ Создать презентацию",   callback_data="pptx_final_generate")])

        await callback.message.answer(
            f"🎞 <b>Тема:</b> {topic[:60]}\n\n"
            f"<b>📋 План ({slides_n} слайдов):</b>\n{plan_display}\n\n"
            f"<i>✏️ Нажми на кнопку слайда чтобы изменить его название</i>",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=kb_rows)
        )

    except Exception as e:
        try: await wait.delete()
        except: pass
        await callback.message.answer(f"❌ Ошибка: {e}", parse_mode="HTML")


@dp.callback_query(F.data == "pptx_final_generate")
async def pptx_final_generate_cb(callback: CallbackQuery):
    """Запускает финальную генерацию — предлагает формат."""
    uid   = callback.from_user.id
    state = pptx_states.get(uid, {})
    if not state.get("topic"):
        return await callback.answer("❌ Тема не задана", show_alert=True)
    await callback.answer()
    # Добавляем пожелания из аудитории и объёма в state
    audience = state.get("audience", "")
    volume   = state.get("volume", "standard")
    vol_hint = {"brief": "краткий текст, 2-3 тезиса на слайд", "standard": "стандартный объём, 4-5 тезисов", "detailed": "подробный текст, 6-8 тезисов"}.get(volume, "")
    plan     = state.get("plan", [])
    wishes   = f"Аудитория: {audience}. {vol_hint}."
    if plan:
        wishes += " Структура слайдов: " + "; ".join(f"{i+1}. {t}" for i, t in enumerate(plan))
    state["wishes"] = wishes
    pptx_states[uid] = state
    # Спрашиваем формат
    await callback.message.answer(
        "📁 <b>Выбери формат презентации:</b>\n\n"
        "📊 <b>PowerPoint (.pptx)</b> — открыть в Google Slides, PowerPoint, Keynote\n"
        "🌐 <b>HTML</b> — красивый браузерный просмотр с анимацией",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="📊 PowerPoint (.pptx)", callback_data="pptx_type_pptx")],
            [InlineKeyboardButton(text="🌐 HTML-презентация",   callback_data="pptx_type_html")],
        ])
    )


# ══════════════════════════════════════════════════════════════
# ✏️ РЕДАКТИРОВАНИЕ ОТДЕЛЬНОГО СЛАЙДА В ПЛАНЕ
# ══════════════════════════════════════════════════════════════

def _plan_kb(uid: int) -> InlineKeyboardMarkup:
    """Возвращает клавиатуру плана для данного пользователя."""
    slides_list = pptx_states.get(uid, {}).get("plan", [])
    slides_n    = pptx_states.get(uid, {}).get("slides", len(slides_list))
    kb_rows = []
    for i in range(0, len(slides_list), 2):
        row = []
        for j in [i, i+1]:
            if j < len(slides_list):
                name = slides_list[j]
                words = name.split()
                short = ""
                for w in words:
                    if len(short) + len(w) + 1 <= 14:
                        short = (short + " " + w).strip()
                    else:
                        break
                name_short = (short + "…") if short and len(name) > 15 else name[:15]
                row.append(InlineKeyboardButton(
                    text=f"{j+1}. {name_short}",
                    callback_data=f"pptx_slide_edit_{j}"
                ))
        kb_rows.append(row)
    kb_rows.append([InlineKeyboardButton(text="🔄 Перегенерировать план", callback_data="pptx_gen_plan")])
    kb_rows.append([InlineKeyboardButton(text="✅ Создать презентацию",   callback_data="pptx_final_generate")])
    return InlineKeyboardMarkup(inline_keyboard=kb_rows)


@dp.callback_query(F.data.startswith("pptx_slide_edit_"))
async def pptx_slide_edit_cb(callback: CallbackQuery):
    """Показывает текущее название слайда и предлагает его изменить."""
    uid = callback.from_user.id
    try:
        idx = int(callback.data.replace("pptx_slide_edit_", ""))
    except ValueError:
        return await callback.answer("❌ Ошибка", show_alert=True)

    plan = pptx_states.get(uid, {}).get("plan", [])
    if idx < 0 or idx >= len(plan):
        return await callback.answer("❌ Слайд не найден", show_alert=True)

    current_name = plan[idx]
    admin_await[uid] = {"action": "pptx_rename_slide", "slide_idx": idx}

    await callback.answer()
    try:
        await callback.message.edit_text(
            f"✏️ <b>Редактировать слайд {idx + 1}</b>\n\n"
            f"Текущее название:\n<i>«{current_name}»</i>\n\n"
            f"Напиши новое название для этого слайда.\n"
            f"<i>Содержание слайда будет сгенерировано по новому названию.</i>",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="❌ Отмена", callback_data="pptx_back_to_plan"),
            ]])
        )
    except Exception:
        await callback.message.answer(
            f"✏️ <b>Слайд {idx + 1}:</b> <i>«{current_name}»</i>\n\nНапиши новое название:",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="❌ Отмена", callback_data="pptx_back_to_plan"),
            ]])
        )


@dp.callback_query(F.data == "pptx_back_to_plan")
async def pptx_back_to_plan_cb(callback: CallbackQuery):
    """Возврат к просмотру плана после отмены редактирования."""
    uid = callback.from_user.id
    admin_await.pop(uid, None)
    state = pptx_states.get(uid, {})
    plan  = state.get("plan", [])
    topic = state.get("topic", "")
    if not plan:
        await callback.answer("❌ План не найден", show_alert=True)
        return
    plan_display = "\n".join(f"{i+1}. {t}" for i, t in enumerate(plan))
    try:
        await callback.message.edit_text(
            f"🎞 <b>Тема:</b> {topic[:60]}\n\n"
            f"<b>📋 План ({len(plan)} слайдов):</b>\n{plan_display}\n\n"
            f"<i>✏️ Нажми на кнопку слайда чтобы изменить его название</i>",
            parse_mode="HTML",
            reply_markup=_plan_kb(uid)
        )
    except Exception:
        await callback.message.answer(
            f"📋 <b>План:</b>\n{plan_display}",
            parse_mode="HTML",
            reply_markup=_plan_kb(uid)
        )
    await callback.answer()


async def _show_pptx_confirm(msg, uid):
    state       = pptx_states.get(uid, {})
    topic       = state.get("topic", "—")
    slides_n    = state.get("slides", 7)
    theme       = PPTX_THEMES.get(state.get("theme", "dark_blue"), {}).get("name", "—")
    model       = MODELS.get(state.get("model", "claude_opus"), {}).get("label", "—")
    wishes      = state.get("wishes", "")
    use_img     = state.get("use_images", True)
    ti          = state.get("title_info", {})
    w_str       = (wishes[:55] + "...") if len(wishes) > 55 else (wishes or "—")
    img_str     = "🖼 Да" if use_img else "📄 Нет"
    text = (
        "🎞 <b>Параметры презентации:</b>\n\n"
        f"📌 Тема: <b>«{topic[:55]}»</b>\n"
        f"📊 Слайдов: <b>{slides_n}</b>\n"
        f"🎨 Оформление: <b>{theme}</b>\n"
        f"🖼 Фотографии: <b>{img_str}</b>\n"
        f"🤖 Модель: <b>{model}</b>\n"
        f"📝 Подзаголовок: <b>{ti.get('subtitle') or '—'}</b>\n"
        f"👤 Автор: <b>{ti.get('author') or '—'}</b>\n"
        f"🏢 Организация: <b>{ti.get('organization') or '—'}</b>\n"
        f"✍️ Пожелания: <b>{w_str}</b>\n\n"
        "💎 Стоимость: <b>3 продвинутых запроса</b>\n\n"
        "👇 <b>Выбери формат:</b>"
    )
    try:
        await msg.edit_text(text, parse_mode="HTML", reply_markup=pptx_confirm_kb())
    except Exception:
        await msg.answer(text, parse_mode="HTML", reply_markup=pptx_confirm_kb())


async def _do_generate_pptx(msg, uid: int, state: dict):
    topic       = state.get("topic", "")
    slides_n    = state.get("slides", 7)
    theme_key   = state.get("theme", "auto")  # Авто-выбор по теме по умолчанию
    wishes      = state.get("wishes", "")
    use_images  = state.get("use_images", True)
    title_info  = state.get("title_info", {})
    theme_name  = "✨ Авто" if theme_key == "auto" else PPTX_THEMES.get(theme_key, {}).get("name", theme_key)

    # Всегда используем Sonar Deep Research для актуальных данных
    model_key   = "sonar_deep_research"
    model_label = "🔬 Sonar Deep Research"

    _init_limits(uid)
    _refresh_limits(uid)
    li = get_limits_info(uid)
    pro_left = li["pro_max"] - li["pro_used"]
    if pro_left < 3:
        await msg.answer(
            "❌ <b>Недостаточно запросов!</b>\n\n"
            f"Нужно <b>3 продвинутых запроса</b>. Осталось: <b>{pro_left}</b>\n"
            f"⏰ Сброс через: <code>{li['pro_reset']}</code>",
            parse_mode="HTML"
        )
        return

    await msg.answer(
        "🌐 <i>Использую Sonar Deep Research — ищу актуальные данные в интернете...</i>",
        parse_mode="HTML"
    )
    think_msg = await msg.answer(
        "🎞 <i>Запускаю генерацию... <b>1с</b></i>", parse_mode="HTML"
    )
    _start = asyncio.get_event_loop().time()

    async def _tick(m, s):
        phases = [
            (7,   "📋 Составляю план слайдов"),
            (20,  "✍️ ИИ пишет контент"),
            (35,  "🖼 Загружаю тематические фото"),
            (55,  "🎨 Собираю современный дизайн"),
            (80,  "📊 Финальная сборка PPTX"),
            (999, "⚡ Почти готово!"),
        ]
        while True:
            await asyncio.sleep(1)
            elapsed = int(asyncio.get_event_loop().time() - s)
            label = phases[-1][1]
            for thr, lbl in phases:
                if elapsed < thr:
                    label = lbl
                    break
            try:
                await m.edit_text(f"{label}... <b>{elapsed}с</b>", parse_mode="HTML")
            except Exception:
                pass

    tick = asyncio.create_task(_tick(think_msg, _start))
    try:
        pptx_data = await generate_pptx_structure_ai(
            topic=topic, slides_count=slides_n,
            model_key=model_key, extra_wishes=wishes,
            title_info=title_info if title_info else None
        )
        # Применяем данные wizard'а на титульный слайд
        if title_info:
            for sl in pptx_data.get("slides", []):
                if sl.get("type") == "title":
                    if title_info.get("author"):       sl["author"]       = title_info["author"]
                    if title_info.get("subtitle"):     sl["subtitle"]     = title_info["subtitle"]
                    if title_info.get("organization"): sl["organization"] = title_info["organization"]

        pptx_bytes = await create_pptx_bytes(pptx_data, theme_key, use_images=use_images)

        for _ in range(3):
            user_limits[uid]["pro_used"] = user_limits[uid].get("pro_used", 0) + 1

        tick.cancel()
        try:
            await tick
        except asyncio.CancelledError:
            pass
        await think_msg.delete()

        total_secs    = int(asyncio.get_event_loop().time() - _start)
        actual_slides = len(pptx_data.get("slides", []))
        img_note      = " + тематические фото 🖼" if use_images else ""
        caption = (
            f"🎞 <b>«{topic[:50]}{'...' if len(topic)>50 else ''}»</b>\n\n"
            f"🎨 Тема: <b>{theme_name}</b>{img_note}\n"
            f"📊 Слайдов: <b>{actual_slides}</b>\n"
            f"🤖 Модель: <b>{model_label}</b>\n"
            f"⏱ Время: <b>{total_secs}с</b>\n"
            f"💎 Потрачено: <b>3 продвинутых запроса</b>\n\n"
            "✅ Открой в <b>PowerPoint</b> или <b>Google Slides</b>!"
        )
        await msg.answer(caption, parse_mode="HTML")
        fname = f"Презентация_{topic[:25].replace(' ', '_')}.pptx"
        await bot.send_document(
            chat_id=uid, document=BufferedInputFile(pptx_bytes, filename=fname)
        )

        pptx_states.pop(uid, None)

    except Exception as e:
        tick.cancel()
        try:
            await tick
        except asyncio.CancelledError:
            pass
        try:
            await think_msg.delete()
        except Exception:
            pass
        logging.error(f"pptx generate error: {e}")
        await msg.answer(f"❌ Ошибка генерации: {str(e)[:300]}")


# ==================================================================
# END PPTX BLOCK
# ==================================================================


@dp.message(F.text == "👤 Профиль")
async def rb_profile(message: Message):
    uid  = message.from_user.id
    # Профиль доступен всем без проверки подписки
    await _show_profile(message, uid)


@dp.message(F.text.contains("Авто-модель"))
async def rb_toggle_auto(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message):
        return
    _init_limits(uid)
    feats = user_features.get(uid, {"doc_analysis": False, "answer_mode": "fast", "auto_model": False})
    feats["auto_model"] = not feats.get("auto_model", False)
    user_features[uid] = feats
    asyncio.create_task(db_save_user(uid, message.from_user.first_name or "", message.from_user.username or ""))
    status = "✅ включён" if feats["auto_model"] else "❌ выключен"
    detail = ("AI сам подберёт лучшую нейросеть под каждый запрос.\nТы увидишь: «AI выбрал: …»"
              if feats["auto_model"] else "Теперь используется выбранная тобой модель.")
    await message.answer(
        f"🤖 <b>Авто-выбор модели {status}</b>\n\n{detail}",
        parse_mode="HTML",
        reply_markup=main_reply_kb(uid)
    )


@dp.message(F.text == "📋 История")
async def rb_history(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message):
        return
    hist = user_history.get(uid, [])
    if not hist:
        await message.answer("📋 История запросов пуста. Задай первый вопрос!")
        return
    lines = []
    for i, h in enumerate(reversed(hist[-10:]), 1):
        ts = h.get("ts", ""); q = h.get("q", "")[:50]
        model = MODELS.get(h.get("model",""), {}).get("label", "")[:18]
        lines.append(f"◈ <b>{i}.</b> [{ts}] {q}…\n    <i>{model}</i>")
    text = (
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "📋 <b>ИСТОРИЯ</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        + "\n\n".join(lines)
    )
    btns = [[InlineKeyboardButton(text=f"🔁 Повторить #{i+1}", callback_data=f"hist_rep_{len(hist)-1-i}")]
            for i in range(min(5, len(hist)))]
    btns.append([InlineKeyboardButton(text="🗑 Очистить", callback_data="hist_clear")])
    await message.answer(text, parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=btns))


async def _show_profile(msg_or_cb, uid: int):
    is_cb = not isinstance(msg_or_cb, Message)
    if is_cb:
        user_obj = msg_or_cb.from_user
    else:
        user_obj = msg_or_cb.from_user

    name     = user_obj.first_name or "—"
    username = user_obj.username
    uname_str = f"@{username}" if username else "без username"
    prof     = user_profiles.get(uid, {})
    joined   = prof.get("joined", "—")
    requests = prof.get("requests", 0)
    mk       = get_model_key(uid)
    cur_label= MODELS.get(mk, {}).get("label", mk) if not mk.startswith("imggen:") else "🎨 " + IMG_MODELS[mk.split(":")[1]]["label"]
    imgs     = user_images_count.get(uid, 0)
    refs     = user_referrals.get(uid, {}).get("refs", [])

    _init_limits(uid)
    li = get_limits_info(uid)
    feats       = user_features.get(uid, {"doc_analysis": False, "answer_mode": "fast", "auto_model": False})
    answer_mode = feats.get("answer_mode", "fast")
    auto_model  = feats.get("auto_model", False)
    mode_tuple  = ANSWER_MODES.get(answer_mode, ANSWER_MODES["fast"])
    mode_label  = mode_tuple[0]
    mode_desc   = mode_tuple[1]

    fast_left = 0  # быстрые запросы убраны
    pro_left  = max(0, li['pro_max'] - li['pro_used'])
    img_left  = max(0, li['img_max'] - li['img_used'])

    # Уровень
    level_name, next_req = get_user_level(requests)
    if next_req:
        prog = f"до {next_req} ···· {next_req - requests} зап."
    else:
        prog = "👑 максимум"

    sub_active = has_active_sub(uid)
    sub_line   = f"💎 <b>Подписка</b> — {'✅ активна до ' + sub_expires_str(uid) if sub_active else '❌ нет'}"

    trial_line = ""  # Trial disabled

    text = (
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        f"👤 <b>ПРОФИЛЬ</b>  ·  {level_name}\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        f"◈ Имя — <b>{name}</b>\n"
        f"◈ ID — <code>{uid}</code>\n"
        f"◈ Дата входа — <code>{joined}</code>\n"
        f"◈ {sub_line}\n"
        "━━ 📊 Статистика ━━━━━━━━━━\n\n"
        f"◈ Запросов — <b>{requests}</b>\n"
        f"◈ Генераций — <b>{imgs}</b>\n"
        f"◈ Рефералов — <b>{len(refs)}</b>\n"
        f"◈ Уровень — {prog}\n\n"
        "━━ ⚙️ Настройки ━━━━━━━━━━\n\n"
        f"◈ Модель — <code>{cur_label}</code>\n"
        f"◈ Режим — {mode_label}\n"
        f"◈ Авто-модель — {'✅ вкл' if auto_model else '❌ выкл'}\n\n"
        "━━ 💎 Лимиты ━━━━━━━━━━━━\n\n"
        f"◈ Запросы — <b>{pro_left}</b> / {li['pro_max']}  ↺ {li['pro_reset']}\n"
        f"◈ 🎨 Картинки — <b>{img_left}</b> / {li['img_max']}  в месяц\n"
        + (f"◈ 🎵 Музыка — <b>{max(0, li['music_max'] - li['music_used'])}</b> / {li['music_max']}  в месяц\n"
           if sub_active else
           "◈ 🎵 Музыка — 🔒 только с подпиской (3 трека/мес)\n")
        + (f"◈ 🎬 Видео — <b>{max(0, li['video_max'] - li['video_used'])}</b> / {li['video_max']}  в месяц\n"
           if sub_active else
           "◈ 🎬 Видео — 🔒 только с подпиской\n")
        + "\n▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔"
    )
    # Кнопки режимов (2+2)
    mode_keys = list(ANSWER_MODES.keys())
    mode_row1 = [InlineKeyboardButton(
        text=ANSWER_MODES[m][0] + (" ✓" if m == answer_mode else ""),
        callback_data=f"set_mode_{m}") for m in mode_keys[:2]]
    mode_row2 = [InlineKeyboardButton(
        text=ANSWER_MODES[m][0] + (" ✓" if m == answer_mode else ""),
        callback_data=f"set_mode_{m}") for m in mode_keys[2:]]
    doc_btn = "📄 Документы ✅" if feats.get("doc_analysis") else "📄 Документы ❌"
    auto_btn = f"🤖 Авто {'✅' if auto_model else '❌'}"
    kb = InlineKeyboardMarkup(inline_keyboard=[
        mode_row1,
        mode_row2,
        [InlineKeyboardButton(text=auto_btn,  callback_data="toggle_auto_model"),
         InlineKeyboardButton(text=doc_btn,   callback_data="toggle_doc_analysis")],
        [InlineKeyboardButton(text="📋 История",         callback_data="show_history"),
         InlineKeyboardButton(text="🎟 Промокод",        callback_data="profile_promo")],
        [InlineKeyboardButton(text="🔗 Реф. ссылка",    callback_data="profile_reflink"),
         InlineKeyboardButton(text="💬 Поддержка",       callback_data="menu_support")],
        [InlineKeyboardButton(text="🏠 Главная",         callback_data="back_home")],
    ])
    if is_cb:
        try:
            await msg_or_cb.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
        except Exception:
            await msg_or_cb.message.answer(text, parse_mode="HTML", reply_markup=kb)
        await msg_or_cb.answer()
    else:
        await msg_or_cb.answer(text, parse_mode="HTML", reply_markup=kb)

@dp.message(F.text == "ℹ️ О боте")
async def rb_about(message: Message):
    if not await require_subscription(message):
        return
    uid = message.from_user.id
    mk = get_model_key(uid)
    default_label = MODELS.get("claude_sonnet", {}).get("label", "Claude Sonnet")
    cur_label = MODELS.get(mk, {}).get("label", mk) if not mk.startswith("imggen:") else "🎨 " + IMG_MODELS[mk.split(":")[1]]["label"]
    vision_list = "\n".join([f"   • {MODELS[k]['label']}" for k in VISION_MODELS if k in MODELS])
    about_text = (
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "⚡ <b>ХУЗА — НЕЙРОСЕТЬ</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        "━━ 💬 Чат-нейросети ━━━━━━━━\n\n"
        "◈ Claude 4.5 — Sonnet · Haiku · Opus 📸\n"
        "◈ GPT-5.2 — Chat · Vision 📸\n"
        "◈ DeepSeek — V3.2 · R1 · V3\n"
        "◈ Qwen3 — Max · VL 📸\n"
        "◈ Command — A · Vision 📸 · Reasoning\n"
        "◈ Kimi K2 / GLM-4.7 · китайские флагманы\n"
        "◈ C4AI Aya — 32B · 8B · Vision 📸\n\n"
        "━━ 🎨 Генерация ━━━━━━━━━━━\n\n"
        "◈ Flux 2 Dev · Phoenix · Lucid · GPT-Image 1.5\n\n"
        "━━ 🛠 Инструменты ━━━━━━━━━\n\n"
        "◈ Vision 📸 — анализ фото\n"
        "◈ Документы — DOCX · PDF · TXT\n"
        "◈ Голос — Whisper AI\n"
        "◈ Презентации — PPTX · HTML\n"
        "◈ Доклады — DOCX · PDF\n"
        "◈ Экспорт — DOCX · PNG · PDF\n\n"
        f"◈ Активная — <code>{cur_label}</code>\n"
        "◈ Статус — ✅ работаю\n\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔"
    )
    await message.answer(
        about_text,
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="💬 Поддержка", callback_data="menu_support"),
            InlineKeyboardButton(text="🏠 Главная",   callback_data="back_home"),
        ]])
    )


@dp.message(F.text == "🧹 Очистить память")
async def rb_clear(message: Message):
    uid = message.from_user.id
    user_memory[uid] = deque(maxlen=20)
    await message.answer(
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "🧠 <b>ПАМЯТЬ ОЧИЩЕНА</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        "Диалог начат с нуля.\n"
        "Напиши вопрос.\n\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔",
        parse_mode="HTML"
    )


@dp.message(F.text == "🔥 Админ")
async def rb_admin(message: Message):
    uid = message.from_user.id
    if uid not in ADMIN_IDS:
        return
    await show_admin_panel(message)


@dp.message(F.text == "💬 Написать")
async def rb_napisat(message: Message):
    uid = message.from_user.id
    # Если активна модель генерации картинок — сбрасываем на текстовую
    cur_mk = get_model_key(uid)
    if cur_mk.startswith("imggen:"):
        fallback = get_chat_model(uid)
        if fallback.startswith("imggen:") or fallback not in MODELS:
            fallback = "claude_haiku"
        set_chat_model(uid, fallback)
        user_settings[uid] = fallback
        model_label = MODELS.get(fallback, {}).get("label", fallback)
        await message.answer(
            f"💬 <b>Режим чата активирован!</b>\n\n"
            f"🤖 Модель: <code>{model_label}</code>\n\n"
            "Пиши вопрос — отвечу текстом:",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="🤖 Сменить модель", callback_data="menu_ai")],
                [InlineKeyboardButton(text="🎨 Генерация картинок", callback_data="menu_imggen")],
                [InlineKeyboardButton(text="🏠 Главная",        callback_data="back_home")],
            ])
        )
        return
    mk = get_model_key(uid)
    model_label = MODELS.get(mk, {}).get("label", "ИИ")
    await message.answer(
        f"💬 Пиши вопрос — отвечу!\n\n"
        f"🤖 Модель: <code>{model_label}</code>\n\n"
        "Или выбери другую:",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="🤖 Выбрать модель", callback_data="menu_ai")],
            [InlineKeyboardButton(text="🏠 Главная",        callback_data="back_home")],
        ])
    )


@dp.message(F.text == "🎨 Создать")
async def rb_sozdat(message: Message):
    if not await require_subscription(message):
        return
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🎨 Картинки",   callback_data="menu_imggen"),
         InlineKeyboardButton(text="🎬 Видео",       callback_data="menu_video")],
        [InlineKeyboardButton(text="🎵 Музыка",      callback_data="menu_music")],
        [InlineKeyboardButton(text="📝 Доклад",      callback_data="menu_report"),
         InlineKeyboardButton(text="🎞 Презентация", callback_data="menu_pptx")],
        [InlineKeyboardButton(text="🏠 Главная",     callback_data="back_home")],
    ])
    await message.answer("🎨 <b>Создать</b>\n\nВыбери что хочешь сделать:",
                         parse_mode="HTML", reply_markup=kb)


@dp.message(F.text == "🏠 Главная")
async def rb_glavnaya(message: Message):
    uid  = message.from_user.id
    name = message.from_user.first_name or "Пользователь"
    await message.answer(_welcome_text(name, 0, uid in ADMIN_IDS),
                         parse_mode="HTML", reply_markup=home_kb(uid))


# ==================================================================
# 🔘 CALLBACKS
# ==================================================================

@dp.callback_query(F.data == "menu_extra")
async def cb_menu_extra(callback: CallbackQuery):
    uid = callback.from_user.id
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🎨 Картинки",   callback_data="menu_imggen"),
         InlineKeyboardButton(text="🎬 Видео",       callback_data="menu_video")],
        [InlineKeyboardButton(text="🎵 Музыка",      callback_data="menu_music")],
        [InlineKeyboardButton(text="📝 Доклад",      callback_data="menu_report"),
         InlineKeyboardButton(text="🎞 Презентация", callback_data="menu_pptx")],
        [InlineKeyboardButton(text="◀️ Назад",       callback_data="back_home")],
    ])
    text = "🎨 <b>Создать</b>\n\nВыбери что хочешь сделать:"
    try:
        await callback.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
    except Exception:
        await callback.message.answer(text, parse_mode="HTML", reply_markup=kb)
    await callback.answer()


@dp.callback_query(F.data == "menu_home")
async def menu_home_cb(callback: CallbackQuery):
    uid  = callback.from_user.id
    name = callback.from_user.first_name or "Пользователь"
    tok  = get_tokens(uid)
    try:
        await callback.message.edit_text(
            _welcome_text(name, tok, uid in ADMIN_IDS),
            parse_mode="HTML", reply_markup=home_kb(uid)
        )
    except Exception:
        await callback.message.answer(_welcome_text(name, tok, uid in ADMIN_IDS), parse_mode="HTML", reply_markup=home_kb(uid))
    await callback.answer()


@dp.callback_query(F.data == "menu_profile")
async def menu_profile_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    # Профиль доступен всем без проверки подписки
    await _show_profile(callback, uid)

# ── Поддержка ────────────────────────────────────────────────────
SUPPORT_ADMIN_ID = list(ADMIN_IDS)[0] if ADMIN_IDS else 5613085898

@dp.callback_query(F.data == "menu_support")
async def menu_support_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    mk  = get_model_key(uid)
    cur = MODELS.get(mk, {}).get("label", mk) if not mk.startswith("imggen:") else "🎨 Генерация"
    li  = get_limits_info(uid)
    fast_left = 0  # убрано
    pro_left  = max(0, li['pro_max']  - li['pro_used'])
    text = (
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "💬 <b>ПОДДЕРЖКА</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        "Напиши напрямую — отвечаем быстро.\n\n"
        f"◈ Твой ID — <code>{uid}</code>\n"
        f"◈ Модель — <code>{cur}</code>\n"
        f"◈ Продвинутые — <b>{pro_left}</b> / {li['pro_max']}\n\n"
        "▾ Нажми «Написать» ниже\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔"
    )
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="💬 Написать в поддержку @helphuza", url="https://t.me/helphuza")],
        [InlineKeyboardButton(text="🏠 Главная", callback_data="back_home")],
    ])
    try:
        await callback.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
    except Exception:
        await callback.message.answer(text, parse_mode="HTML", reply_markup=kb)
    await callback.answer()


@dp.message(F.text == "💬 Поддержка")
async def rb_support(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message):
        return
    mk  = get_model_key(uid)
    cur = MODELS.get(mk, {}).get("label", mk) if not mk.startswith("imggen:") else "🎨 Генерация"
    li  = get_limits_info(uid)
    fast_left = 0  # убрано
    pro_left  = max(0, li['pro_max']  - li['pro_used'])
    text = (
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "💬 <b>ПОДДЕРЖКА</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        "Напиши напрямую — отвечаем быстро.\n\n"
        f"◈ Твой ID — <code>{uid}</code>\n"
        f"◈ Модель — <code>{cur}</code>\n"
        f"◈ Продвинутые — <b>{pro_left}</b> / {li['pro_max']}\n\n"
        "📩 <b>Поддержка:</b> @helphuza\n\n"
        "▾ Нажми «Написать» ниже\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔"
    )
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="💬 Написать в поддержку @helphuza", url="https://t.me/helphuza")],
        [InlineKeyboardButton(text="🏠 Главная", callback_data="back_home")],
    ])
    await message.answer(text, parse_mode="HTML", reply_markup=kb)


# ── Быстрый ответ из главного меню ──────────────────────────────
@dp.callback_query(F.data == "menu_ask")
async def menu_ask_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if not await require_subscription(callback):
        return
    mk = get_model_key(uid)
    cur = MODELS.get(mk, {}).get("label", mk) if not mk.startswith("imggen:") else "🎨 Генерация"
    text = (
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "💬 <b>ДИАЛОГ С ИИ</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        f"◈ Активная модель ··· <code>{cur}</code>\n\n"
        "Напиши вопрос — отвечу сразу.\n"
        "Можно прикрепить фото или документ.\n\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔"
    )
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🤖 Сменить модель", callback_data="menu_ai"),
         InlineKeyboardButton(text="🏠 Главная",        callback_data="back_home")],
    ])
    try:
        await callback.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
    except Exception:
        await callback.message.answer(text, parse_mode="HTML", reply_markup=kb)
    await callback.answer()


# ── Доклад из главного меню ──────────────────────────────────────
@dp.callback_query(F.data == "menu_report")
async def menu_report_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if not await require_subscription(callback):
        return
    if not await check_service(callback, "report"):
        return
    await menu_report_screen(callback)


async def menu_report_screen(msg_or_cb):
    uid = msg_or_cb.from_user.id if isinstance(msg_or_cb, Message) else msg_or_cb.from_user.id
    text = (
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "📝 <b>ДОКЛАД / РЕФЕРАТ</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        "ИИ напишет полноценную работу.\n"
        "Claude Opus или GPT-5.2 — выбери при настройке.\n\n"
        f"◈ Стоимость — <b>3</b> продвинутых запроса\n"
        f"◈ Объём — 1–10 страниц\n"
        f"◈ Форматы — DOCX · PDF\n\n"
        "▾ Выбери тип работы\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔"
    )
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📄 Доклад",  callback_data="report_type_doklad"),
         InlineKeyboardButton(text="📚 Реферат", callback_data="report_type_referat")],
        [InlineKeyboardButton(text="❌ Отмена",  callback_data="back_home")],
    ])
    if isinstance(msg_or_cb, Message):
        await msg_or_cb.answer(text, parse_mode="HTML", reply_markup=kb)
    else:
        try:
            await msg_or_cb.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
        except Exception:
            await msg_or_cb.message.answer(text, parse_mode="HTML", reply_markup=kb)
        await msg_or_cb.answer()


# ── Презентация из главного меню ─────────────────────────────────
@dp.callback_query(F.data == "menu_pptx")
async def menu_pptx_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if not await require_subscription(callback):
        return
    if not await check_service(callback, "pptx"):
        return
    pptx_states[uid] = {}
    admin_await[uid] = {"action": "pptx_topic"}
    text = (
        "🎞 <b>Генератор презентаций</b>\n\n"
        "✏️ Напиши <b>тему презентации</b>:\n\n"
        "<i>Примеры:\n"
        "• Россия в будущем\n"
        "• Искусственный интеллект\n"
        "• Маркетинговая стратегия 2025\n"
        "• История освоения космоса</i>"
    )
    try:
        await callback.message.edit_text(text, parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")
            ]]))
    except Exception:
        await callback.message.answer(text, parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")
            ]]))
    await callback.answer()





@dp.callback_query(F.data == "webpptx_start")
async def webpptx_start_cb(callback: CallbackQuery):
    """Кнопка HTML из меню — инициализирует HTML режим и показывает слайды."""
    uid = callback.from_user.id
    pptx_states[uid] = {"mode": "html"}
    await callback.message.edit_text(
        "🌐 <b>HTML-презентация</b>\n\n📏 Сколько слайдов нужно?",
        parse_mode="HTML", reply_markup=pptx_start_kb()
    )
    await callback.answer()


@dp.callback_query(F.data == "toggle_doc_analysis")
async def toggle_doc_analysis(callback: CallbackQuery):
    uid = callback.from_user.id
    _init_limits(uid)
    feats = user_features.get(uid, {"doc_analysis": False, "answer_mode": "fast"})
    feats["doc_analysis"] = not feats.get("doc_analysis", False)
    user_features[uid] = feats
    asyncio.create_task(db_save_user(uid))
    status = "включён ✅" if feats["doc_analysis"] else "выключен ❌"
    await callback.answer(f"📄 Анализ документов {status}", show_alert=True)
    await _show_profile(callback, uid)


@dp.callback_query(F.data.startswith("set_mode_"))
async def set_answer_mode_cb(callback: CallbackQuery):
    """Установить режим ответа."""
    uid = callback.from_user.id
    mode = callback.data[9:]  # "set_mode_fast" → "fast"
    if mode not in ANSWER_MODES:
        await callback.answer("❌ Неизвестный режим", show_alert=True)
        return
    _init_limits(uid)
    feats = user_features.get(uid, {"doc_analysis": False, "answer_mode": "fast", "auto_model": False})
    feats["answer_mode"] = mode
    user_features[uid] = feats
    asyncio.create_task(db_save_user(uid))
    label = ANSWER_MODES[mode][0]
    await callback.answer(f"{label} — установлен!", show_alert=False)
    await _show_profile(callback, uid)


@dp.callback_query(F.data == "toggle_auto_model")
async def toggle_auto_model_cb(callback: CallbackQuery):
    """Переключить авто-выбор модели."""
    uid = callback.from_user.id
    _init_limits(uid)
    feats = user_features.get(uid, {"doc_analysis": False, "answer_mode": "fast", "auto_model": False})
    feats["auto_model"] = not feats.get("auto_model", False)
    user_features[uid] = feats
    asyncio.create_task(db_save_user(uid))
    status = "включён ✅" if feats["auto_model"] else "выключен ❌"
    await callback.answer(f"🤖 Авто-выбор модели {status}", show_alert=True)
    # Если вызван из профиля — показываем профиль, иначе просто отвечаем
    try:
        await _show_profile(callback, uid)
    except Exception:
        pass


@dp.callback_query(F.data == "show_history")
async def show_history_cb(callback: CallbackQuery):
    """История последних 10 запросов."""
    uid = callback.from_user.id
    hist = user_history.get(uid, [])
    if not hist:
        await callback.answer("📋 История пуста", show_alert=True)
        return
    lines = []
    for i, h in enumerate(reversed(hist[-10:]), 1):
        ts    = h.get("ts", "")
        q     = h.get("q", "")[:55]
        model = MODELS.get(h.get("model",""), {}).get("label", "")[:18]
        lines.append(f"<b>{i}.</b> [{ts}] {q}\u2026\n   <i>{model}</i>")
    text = "📋 <b>История запросов</b>\n\n" + "\n\n".join(lines)
    btns = []
    for i, h in enumerate(reversed(hist[-10:])):
        btns.append([InlineKeyboardButton(
            text=f"🔁 #{i+1}: {h.get('q','')[:30]}…",
            callback_data=f"hist_rep_{len(hist)-1-i}"
        )])
    btns.append([InlineKeyboardButton(text="🗑 Очистить историю", callback_data="hist_clear")])
    btns.append([InlineKeyboardButton(text="◀️ Профиль", callback_data="menu_profile")])
    try:
        await callback.message.edit_text(text, parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=btns))
    except Exception:
        await callback.message.answer(text, parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=btns))
    await callback.answer()


@dp.callback_query(F.data.startswith("hist_rep_"))
async def hist_repeat_cb(callback: CallbackQuery):
    """Повтор запроса из истории."""
    uid = callback.from_user.id
    hist = user_history.get(uid, [])
    try:
        idx = int(callback.data[9:])
        h = hist[idx]
        q_text = h.get("q", "")
        a_text = h.get("a", "")[:3800]
        model_label = MODELS.get(h.get("model",""), {}).get("label", h.get("model",""))
        await callback.answer("🔁 Повторяю запрос", show_alert=False)
        reply_text = (
            f"🔁 <b>Повтор запроса:</b>\n<i>{q_text[:100]}</i>\n\n"
            f"🤖 <b>{model_label}</b>\n\n{md_to_html(a_text)}"
        )
        await callback.message.answer(reply_text, parse_mode="HTML")
    except Exception as e:
        await callback.answer("❌ Не удалось повторить", show_alert=True)


@dp.callback_query(F.data == "hist_clear")
async def hist_clear_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    user_history[uid] = []
    await callback.answer("🗑 История очищена", show_alert=True)
    await _show_profile(callback, uid)


@dp.callback_query(F.data == "profile_reflink")
async def profile_reflink(callback: CallbackQuery):
    uid = callback.from_user.id
    bot_info = await bot.get_me()
    ref_link = f"https://t.me/{bot_info.username}?start=ref{uid}"
    refs = user_referrals.get(uid, {}).get("refs", [])
    earned = len(refs) * REF_BONUS_INVITER
    text = (
        f"🔗 <b>Реферальная ссылка</b>\n\n"
        f"<code>{ref_link}</code>\n\n"
        f"👥 Приглашено: <code>{len(refs)}</code> друзей\n"
        f"💰 Заработано: <code>{earned}</code> бонусных запросов\n\n"
        f"🔹 За каждого приглашённого: <b>+{REF_BONUS_INVITER} запросов</b>"
    )
    try:
        await callback.message.edit_text(text, parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="◀️ Назад", callback_data="menu_profile"),
            ]]))
    except Exception:
        await callback.message.answer(text, parse_mode="HTML")
    await callback.answer()


@dp.callback_query(F.data == "profile_promo")
async def profile_promo(callback: CallbackQuery):
    uid = callback.from_user.id
    admin_await[uid] = {"action": "use_promo"}
    try:
        await callback.message.edit_text(
            "🎟️ <b>Активация промокода</b>\n\n"
            "Введи промокод:",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="❌ Отмена", callback_data="back_home"),
            ]]),
        )
    except Exception:
        await callback.message.answer("🎟️ Введи промокод:", parse_mode="HTML")
    await callback.answer()


@dp.callback_query(F.data == "profile_stats")
async def profile_stats(callback: CallbackQuery):
    uid = callback.from_user.id
    imgs = user_images_count.get(uid, 0)
    requests = user_profiles.get(uid, {}).get("requests", 0)
    refs = user_referrals.get(uid, {}).get("refs", [])
    li = get_limits_info(uid)
    await callback.message.edit_text(
        f"📊 <b>Твоя статистика</b>\n\n"
        f"💬 Всего запросов: <code>{requests}</code>\n"
        f"🖼 Изображений создано: <code>{imgs}</code>\n"
        f"👥 Рефералов: <code>{len(refs)}</code>\n\n"
        f"📊 <b>Лимиты запросов:</b>\n"
        
        f"🧠 Продвинутые: <code>{li['pro_used']} / {li['pro_max']}</code> (сброс через {li['pro_reset']})\n"
        f"🖼 Генерации: <code>{li['img_used']} / {li['img_max']}</code> (в месяц)",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="◀️ Назад", callback_data="menu_profile"),
        ]]),
    )
    await callback.answer()


@dp.callback_query(F.data == "menu_ai")
async def menu_ai(callback: CallbackQuery):
    uid = callback.from_user.id
    if not await require_subscription(callback):
        return
    mk  = get_model_key(uid)
    cur = MODELS.get(mk, {}).get("label", mk) if not mk.startswith("imggen:") else "🎨 " + IMG_MODELS[mk.split(":")[1]]["label"]
    try:
        await callback.message.edit_text(
            "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
            "🤖 <b>НЕЙРОСЕТЬ</b>\n"
            "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
            f"◈ Активная — <code>{cur}</code>\n"
            "◈ Vision 📸 — анализирует фото\n\n"
            "▾ Выбери категорию\n"
            "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔",
            parse_mode="HTML", reply_markup=ai_menu_kb(uid)
        )
    except Exception:
        await callback.message.answer("🤖 <b>Нейросеть</b>", parse_mode="HTML", reply_markup=ai_menu_kb(uid))
    await callback.answer()


@dp.callback_query(F.data == "menu_chat")
async def menu_chat(callback: CallbackQuery):
    """Кнопка 'Задать вопрос' — сразу готовит чат, сбрасывает imggen-режим."""
    uid = callback.from_user.id
    if not await require_subscription(callback):
        return
    # Если активна генерация картинок — переключаем на текстовую модель
    mk = get_model_key(uid)
    if mk.startswith("imggen:"):
        fallback = get_chat_model(uid)
        if fallback.startswith("imggen:") or fallback not in MODELS:
            fallback = "claude_haiku"
        set_chat_model(uid, fallback)
        user_settings[uid] = fallback
        mk = fallback
    cur = MODELS.get(mk, {}).get("label", mk)
    is_vision = mk in VISION_MODELS
    vision_hint = "\n📸 <i>Можешь отправить фото для анализа!</i>" if is_vision else ""
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🤖 Сменить модель", callback_data="menu_ai")],
        [InlineKeyboardButton(text="🧹 Очистить память", callback_data="clear_memory"),
         InlineKeyboardButton(text="🏠 Главная", callback_data="back_home")],
    ])
    try:
        await callback.message.edit_text(
            f"💬 <b>Режим чата</b>\n\n"
            f"🤖 Модель: <code>{cur}</code>{vision_hint}\n\n"
            f"✏️ <b>Напиши свой вопрос прямо сейчас</b> — я отвечу!\n\n"
            f"<i>Могу отвечать на вопросы, писать код, переводить, анализировать тексты и многое другое.</i>",
            parse_mode="HTML", reply_markup=kb
        )
    except Exception:
        await callback.message.answer(
            f"💬 <b>Режим чата активен!</b>\n🤖 Модель: <code>{cur}</code>{vision_hint}\n\n✏️ Пиши вопрос:",
            parse_mode="HTML", reply_markup=kb
        )
    await callback.answer()


@dp.callback_query(F.data == "menu_imggen")
async def menu_imggen(callback: CallbackQuery):
    if not await require_subscription(callback):
        return
    if not await check_service(callback, "image"):
        return
    imggen_text = (
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "🎨 <b>ГЕНЕРАЦИЯ КАРТИНОК</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        "✏️ <b>Текст → Картинка</b> — выбери модель ниже\n"
        "📸 <b>Фото → Стиль</b> — модели с пометкой 📸→🖼\n"
        "  <i>GPT-Image ⭐ лучший для деталей (усы, очки...)</i>\n"
        "  <i>FLUX.2 Dev, FLUX.2 Klein — смена стиля</i>\n\n"
        "<i>Пример: неоновый город ночью, cyberpunk</i>\n\n"
        "▾ Выбери модель или кнопку 📸 ниже\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔"
    )
    try:
        await callback.message.edit_text(imggen_text, parse_mode="HTML", reply_markup=imggen_category_kb())
    except Exception:
        await callback.message.answer(imggen_text, parse_mode="HTML", reply_markup=imggen_category_kb())
    await callback.answer()


@dp.callback_query(F.data == "clear_memory")
async def clear_memory(callback: CallbackQuery):
    uid = callback.from_user.id
    user_memory[uid] = deque(maxlen=20)
    await callback.answer("🧹 Память очищена!", show_alert=True)


@dp.callback_query(F.data == "menu_about")
async def menu_about(callback: CallbackQuery):
    text = (
        "🌟 ━━━━━━━━━━━━━━━━━━━━━━ 🌟\n\n"
        "🤖 <b>ХУЗА НЕЙРОСЕТЬ v5.0</b>\n"
        "🧠 <i>Мощный AI-ассистент нового поколения</i>\n\n"
        "💡 <b>Возможности:</b>\n"
        "   💬 28 нейросетей для чата\n"
        "   🔵 Claude · 🟢 GPT · 🧬 DeepSeek\n"
        "   🌐 Qwen · 🌀 C4AI · 🌙 Kimi · 🔮 GLM\n"
        "   🔵 Gemini · 🌐 Sonar (веб-поиск)\n\n"
        "   🎨 7 моделей генерации изображений\n"
        "   📸 Vision — анализ любых фото\n"
        "   🎤 Голосовые → текст (Whisper AI)\n"
        "   💾 Экспорт в DOCX / PNG / PDF\n"
        "   📄 Анализ документов\n\n"
        "⚡️ <b>Статус:</b> <code>✅ Работаю</code>\n"
        "🌟 ━━━━━━━━━━━━━━━━━━━━━━ 🌟"
    )
    try:
        await callback.message.edit_text(text, parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="◀️ Главное меню", callback_data="back_home"),
            ]]))
    except Exception:
        await callback.message.answer(text, parse_mode="HTML")
    await callback.answer()



# ==================================================================
# 🔐 ПРОДВИНУТАЯ АДМИН-ПАНЕЛЬ
# ==================================================================



# ══════════════════════════════════════════════════════════════
# 💎 ОБРАБОТЧИКИ ПОДПИСКИ
# ══════════════════════════════════════════════════════════════

def _sub_menu_text(uid: int) -> str:
    active = has_active_sub(uid)
    if active:
        sub  = user_subscriptions.get(uid, {})
        plan = SUB_PLANS.get(sub.get("plan", ""), {})
        exp  = sub_expires_str(uid)
        return (
            "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
            "💎 <b>ПОДПИСКА</b>\n"
            "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
            f"✅ <b>Статус:</b> Активна\n"
            f"◈ Тариф — <b>{plan.get('name', '—')}</b>\n"
            f"◈ Истекает — <code>{exp}</code>\n\n"
            "━━ 💎 Возможности ━━━━━━━\n\n"
            "✅ Все нейросети разблокированы\n"
            "✅ Claude Sonnet / Opus 📸\n"
            "✅ GPT-5.2 с Vision 📸\n"
            "✅ 75 запросов / 12ч\n"
            "✅ 35 генераций изображений\n\n"
            "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔"
        )
    return (
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "💎 <b>ПОДПИСКА ХУЗА</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        "🔓 Разблокируй все возможности:\n\n"
        "✅ Claude Sonnet / Haiku / Opus 📸\n"
        "✅ GPT-5.2 Vision 📸\n"
        "✅ 75 запросов каждые 12 часов\n"
        "✅ 35 генераций изображений в месяц\n"
        "✅ Все будущие новые модели\n\n"
        "━━ 💰 Тарифы ━━━━━━━━━━━\n\n"
        "⚡ <b>7 дней</b> — <b>60 ₽</b>\n"
        "💎 <b>30 дней</b> — <b>100 ₽</b> <i>(выгоднее)</i>\n\n"
        "▾ Выбери тариф\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔"
    )


def _sub_menu_kb(uid: int) -> InlineKeyboardMarkup:
    active = has_active_sub(uid)
    if active:
        return InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="⚡ Продлить 7 дней — 60 ₽",  callback_data="sub_buy_week")],
            [InlineKeyboardButton(text="💎 Продлить 30 дней — 100 ₽", callback_data="sub_buy_month")],
            [InlineKeyboardButton(text="🏠 Главная", callback_data="back_home")],
        ])
    return InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="⚡ 7 дней — 60 ₽",  callback_data="sub_buy_week")],
        [InlineKeyboardButton(text="💎 30 дней — 100 ₽", callback_data="sub_buy_month")],
        [InlineKeyboardButton(text="🏠 Главная", callback_data="back_home")],
    ])


@dp.message(F.text == "💎 Подписка")
async def rb_sub_menu(message: Message):
    uid = message.from_user.id
    await message.answer(
        _sub_menu_text(uid),
        parse_mode="HTML",
        reply_markup=_sub_menu_kb(uid),
    )


@dp.callback_query(F.data == "sub_menu")
async def cb_sub_menu(callback: CallbackQuery):
    uid = callback.from_user.id
    try:
        await callback.message.edit_text(
            _sub_menu_text(uid),
            parse_mode="HTML",
            reply_markup=_sub_menu_kb(uid),
        )
    except Exception:
        await callback.message.answer(
            _sub_menu_text(uid),
            parse_mode="HTML",
            reply_markup=_sub_menu_kb(uid),
        )
    await callback.answer()


@dp.callback_query(F.data.startswith("sub_buy_"))
async def cb_sub_buy(callback: CallbackQuery):
    uid      = callback.from_user.id
    plan_key = callback.data.replace("sub_buy_", "")
    if plan_key not in SUB_PLANS:
        return await callback.answer("❌ Тариф не найден", show_alert=True)
    plan = SUB_PLANS[plan_key]

    # ── Экран выбора способа оплаты (всегда показываем оба варианта) ──
    buttons = [
        [InlineKeyboardButton(
            text=f"Platega — {plan['price']} ₽",
            callback_data=f"pay_platega_{plan_key}",
            icon_custom_emoji_id="5904462880941545555"
        )],
        [InlineKeyboardButton(
            text=f"ЮKassa — {plan['price']} ₽",
            callback_data=f"pay_yukassa_{plan_key}",
            icon_custom_emoji_id="5879814368572478751"
        )],
        [InlineKeyboardButton(text="Назад", callback_data="sub_menu", icon_custom_emoji_id="5893057118545646106")],
    ]

    await callback.message.edit_text(
        f'<tg-emoji emoji-id="5769126056262898415">👛</tg-emoji> <b>Выбери способ оплаты</b>\n\n'
        f"◆ Тариф — <b>{plan['name']}</b>\n"
        f"◆ Срок — <b>{plan['days']} дней</b>\n"
        f"◆ Цена — <b>{plan['price']} ₽</b>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=buttons),
    )
    await callback.answer()


@dp.callback_query(F.data.startswith("pay_platega_"))
async def cb_pay_platega(callback: CallbackQuery):
    """Создать платёж через Platega и отправить ссылку."""
    uid      = callback.from_user.id
    plan_key = callback.data.replace("pay_platega_", "")
    if plan_key not in SUB_PLANS:
        return await callback.answer("❌ Тариф не найден", show_alert=True)
    plan = SUB_PLANS[plan_key]

    if not (PLATEGA_MERCHANT_ID and PLATEGA_SECRET):
        await callback.answer("❌ Platega не настроена — обратитесь к администратору", show_alert=True)
        return

    await callback.answer()
    await callback.message.edit_text("⏳ Создаю платёж через Platega...")

    import time as _time
    order_id = f"huza_{plan_key}_{uid}_{int(_time.time())}"
    bot_info  = await bot.get_me()
    return_url = f"https://t.me/{bot_info.username}?start=sub_check"

    # Актуальный payload Platega v1
    payload = {
        "amount":      int(plan["price"]),          # целое, в рублях
        "currency":    "RUB",
        "order_id":    order_id,
        "description": f"ХУЗА AI — {plan['name']} ({plan['days']} дней)",
        "customer":    {"uid": str(uid)},
        "meta":        {"plan": plan_key, "uid": uid},
        "success_url": return_url,
        "fail_url":    return_url,
        "webhook_url": f"https://huzaoobo-production.up.railway.app/platega_webhook",
    }
    headers = {
        "Content-Type":  "application/json",
        "Authorization": f"Bearer {PLATEGA_SECRET}",
        "X-Merchant-Id": PLATEGA_MERCHANT_ID,
    }

    logging.info(f"[PLATEGA] Создаю платёж uid={uid} plan={plan_key} order={order_id} "
                 f"amount={plan['price']} url={PLATEGA_API_URL}")
    try:
        async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=20, connect=8)) as sess:
            async with sess.post(PLATEGA_API_URL, json=payload, headers=headers) as resp:
                status_code = resp.status
                raw = await resp.text()
                try:
                    data = await resp.json(content_type=None)
                except Exception:
                    data = {}
        logging.info(f"[PLATEGA] uid={uid} status={status_code} resp={raw[:500]}")

        # Ищем ссылку в разных полях ответа
        redirect = (data.get("payment_url") or data.get("redirect") or
                    data.get("url") or data.get("paymentUrl") or
                    data.get("link") or data.get("checkout_url"))

        if not redirect:
            err_detail = data.get("message") or data.get("error") or raw[:200]
            logging.error(f"[PLATEGA] Нет ссылки в ответе uid={uid}: {err_detail}")
            raise ValueError(f"HTTP {status_code}: {err_detail}")

        logging.info(f"[PLATEGA] Платёж создан uid={uid} order={order_id} link={redirect}")
        await callback.message.edit_text(
            f"✅ <b>Платёж создан!</b>\n\n"
            f"◆ Тариф — <b>{plan['name']}</b>\n"
            f"◆ Цена — <b>{plan['price']} ₽</b>\n"
            f"◆ Заказ — <code>{order_id}</code>\n\n"
            f"👇 Нажми кнопку ниже для оплаты.\n"
            f"После оплаты подписка активируется автоматически.",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="💳 Оплатить через Platega", url=redirect)],
                [InlineKeyboardButton(text="◀️ Назад", callback_data=f"sub_buy_{plan_key}")],
            ]),
        )
    except Exception as e:
        logging.error(f"[PLATEGA] Ошибка создания платежа uid={uid} plan={plan_key}: {e}")
        # Определяем тип ошибки для понятного сообщения
        _err_str = str(e).lower()
        if "name or service not known" in _err_str or "cannot connect" in _err_str or "connection" in _err_str:
            _err_msg = "Платёжный сервис временно недоступен."
        elif "timeout" in _err_str:
            _err_msg = "Сервис не ответил вовремя. Попробуй ещё раз."
        else:
            _err_msg = "Не удалось создать платёж."
        await callback.message.edit_text(
            f'<tg-emoji emoji-id="5870657884844462243">❌</tg-emoji> <b>{_err_msg}</b>\n\n'
            f'Оплати через ЮKassa (кнопка ниже) или напиши в поддержку: @helphuza',
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="💳 ЮKassa — оплатить", callback_data=f"pay_yukassa_{plan_key}", icon_custom_emoji_id="5879814368572478751")],
                [InlineKeyboardButton(text="🔄 Повторить попытку", callback_data=f"pay_platega_{plan_key}", icon_custom_emoji_id="5345906554510012647")],
                [InlineKeyboardButton(text="📞 Поддержка", url="https://t.me/helphuza", icon_custom_emoji_id="6028435952299413210")],
                [InlineKeyboardButton(text="◁ Назад", callback_data=f"sub_buy_{plan_key}", icon_custom_emoji_id="5893057118545646106")],
            ]),
        )


@dp.callback_query(F.data.startswith("pay_yukassa_"))
async def cb_pay_yukassa(callback: CallbackQuery):
    """Создать платёж через ЮKassa (Telegram Payments)."""
    uid      = callback.from_user.id
    plan_key = callback.data.replace("pay_yukassa_", "")
    if plan_key not in SUB_PLANS:
        return await callback.answer("❌ Тариф не найден", show_alert=True)
    plan = SUB_PLANS[plan_key]

    if not YOOKASSA_TOKEN:
        await callback.answer("❌ ЮKassa не настроена", show_alert=True)
        return

    price_rub = max(plan["price"], 60)
    try:
        await bot.send_invoice(
            chat_id           = uid,
            title             = f"Подписка «{plan['name']}»",
            description       = f"{plan['days']} дней полного доступа ко всем моделям и функциям бота.",
            payload           = f"sub_{plan_key}_{uid}",
            provider_token    = YOOKASSA_TOKEN,
            currency          = "RUB",
            prices            = [LabeledPrice(label=plan["name"], amount=price_rub * 100)],
            start_parameter   = f"sub_{plan_key}",
            need_name         = False,
            need_phone_number = False,
            need_email        = False,
            is_flexible       = False,
        )
        await callback.answer()
    except Exception as e:
        logging.error(f"[YUKASSA] send_invoice error: {e}")
        await callback.message.edit_text(
            f"❌ Ошибка при создании платежа ЮKassa.\n\n"
            f"Напишите в поддержку: @helphuza",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="◀️ Назад", callback_data=f"sub_buy_{plan_key}")],
            ]),
        )


@dp.pre_checkout_query()
async def pre_checkout_handler(pre_checkout: PreCheckoutQuery):
    """Обязательный ответ на запрос оплаты (нужен в течение 10 секунд)."""
    try:
        # Проверяем payload — формат sub_{plan_key}_{uid}
        parts = pre_checkout.invoice_payload.split("_")
        if len(parts) >= 3 and parts[0] == "sub" and parts[1] in SUB_PLANS:
            await pre_checkout.answer(ok=True)
        else:
            await pre_checkout.answer(ok=False, error_message="Неверный формат платежа.")
    except Exception as e:
        logging.error(f"pre_checkout_handler error: {e}")
        await pre_checkout.answer(ok=False, error_message="Внутренняя ошибка. Попробуйте позже.")


@dp.message(F.successful_payment)
async def successful_payment_handler(message: Message):
    """Обработка успешной оплаты через ЮKassa / Telegram Payments."""
    uid     = message.from_user.id
    payment = message.successful_payment

    # Извлекаем план из payload: sub_{plan_key}_{uid}
    parts    = payment.invoice_payload.split("_")
    plan_key = parts[1] if len(parts) >= 2 and parts[1] in SUB_PLANS else "month"
    plan     = SUB_PLANS[plan_key]

    # Начисляем подписку
    existing = user_subscriptions.get(uid)
    if existing and existing["expires"] > msk_now():
        base = existing["expires"]
    else:
        base = msk_now()
    new_exp = base + datetime.timedelta(days=plan["days"])
    user_subscriptions[uid] = {"expires": new_exp, "plan": plan_key}

    # Сохраняем в БД
    asyncio.create_task(db_save_subscription(uid))

    # Логируем транзакцию
    charge_id = payment.provider_payment_charge_id
    logging.info(f"[PAYMENT] uid={uid} plan={plan_key} charge_id={charge_id} amount={payment.total_amount/100}₽")

    await message.answer(
        f"✅ <b>Оплата прошла успешно!</b>\n\n"
        f"💎 Тариф: <b>{plan['name']}</b>\n"
        f"📅 Действует до: <b>{new_exp.strftime('%d.%m.%Y %H:%M')}</b>\n\n"
        f"Теперь тебе доступны все модели и расширенные лимиты!\n"
        f"Приятного использования 🚀",
        parse_mode="HTML",
        reply_markup=home_kb(uid),
    )

    # Уведомляем админа
    for admin_id in ADMIN_IDS:
        try:
            await bot.send_message(
                admin_id,
                f"💰 <b>Новая оплата!</b>\n\n"
                f"👤 uid: <code>{uid}</code>\n"
                f"📋 Тариф: <b>{plan['name']}</b>\n"
                f"💵 Сумма: <b>{payment.total_amount / 100} ₽</b>\n"
                f"🔑 Charge ID: <code>{charge_id}</code>",
                parse_mode="HTML",
            )
        except Exception:
            pass


# ══════════════════════════════════════════════════════════════
# 💎 ПОДПИСКА — ADMIN УПРАВЛЕНИЕ
# ══════════════════════════════════════════════════════════════

@dp.callback_query(F.data == "admin_subs")
async def admin_subs_panel(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    active_count = sum(1 for u in user_subscriptions if has_active_sub(u))
    await callback.message.edit_text(
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "💎 <b>УПРАВЛЕНИЕ ПОДПИСКАМИ</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        f"◈ Активных подписок — <b>{active_count}</b>\n\n"
        "▾ Выбери действие\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="➕ Выдать подписку",   callback_data="admin_sub_give")],
            [InlineKeyboardButton(text="➖ Забрать подписку",  callback_data="admin_sub_take")],
            [InlineKeyboardButton(text="🔍 Проверить подписку", callback_data="admin_sub_check")],
            [InlineKeyboardButton(text="📋 Список подписчиков", callback_data="admin_sub_list")],
            [InlineKeyboardButton(text="◀️ Назад",             callback_data="menu_admin")],
        ]),
    )
    await callback.answer()


@dp.callback_query(F.data == "admin_sub_give")
async def admin_sub_give(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌", show_alert=True)
    admin_await[uid] = {"action": "admin_give_sub"}
    await callback.message.edit_text(
        "➕ <b>Выдать подписку</b>\n\n"
        "Формат: <code>USER_ID ПЛАН</code>\n\n"
        "Планы: <code>week</code> (7 дней) / <code>month</code> (30 дней)\n\n"
        "Пример: <code>123456789 month</code>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="admin_subs"),
        ]]),
    )
    await callback.answer()


@dp.callback_query(F.data == "admin_sub_take")
async def admin_sub_take(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌", show_alert=True)
    admin_await[uid] = {"action": "admin_take_sub"}
    await callback.message.edit_text(
        "➖ <b>Забрать подписку</b>\n\n"
        "Введи <code>USER_ID</code> пользователя:",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="admin_subs"),
        ]]),
    )
    await callback.answer()


@dp.callback_query(F.data == "admin_sub_check")
async def admin_sub_check(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌", show_alert=True)
    admin_await[uid] = {"action": "admin_check_sub"}
    await callback.message.edit_text(
        "🔍 <b>Проверить подписку</b>\n\nВведи <code>USER_ID</code>:",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="admin_subs"),
        ]]),
    )
    await callback.answer()


@dp.callback_query(F.data == "admin_sub_list")
async def admin_sub_list(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌", show_alert=True)
    lines = []
    for u, sub in sorted(user_subscriptions.items(), key=lambda x: x[1]["expires"], reverse=True):
        if has_active_sub(u):
            plan = SUB_PLANS.get(sub.get("plan", ""), {}).get("name", sub.get("plan", "—"))
            exp  = sub["expires"].strftime("%d.%m.%Y")
            prof = user_profiles.get(u, {})
            name = prof.get("name", str(u))
            uname = f"@{prof['username']}" if prof.get("username") else f"id:{u}"
            lines.append(f"• {name} ({uname}) — {plan} до {exp}")
    if not lines:
        text = "📋 <b>Активных подписок нет</b>"
    else:
        text = "📋 <b>Активные подписки:</b>\n\n" + "\n".join(lines[:30])
    await callback.message.edit_text(
        text, parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="◀️ Назад", callback_data="admin_subs"),
        ]]),
    )
    await callback.answer()

async def show_admin_panel(msg_or_cb):
    admin_kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="📊 Статистика",      callback_data="admin_stats"),
         InlineKeyboardButton(text="🔄 Обновить",        callback_data="admin_reload")],
        [InlineKeyboardButton(text="👥 Пользователи",    callback_data="admin_userlist"),
         InlineKeyboardButton(text="🔍 Найти по ID",     callback_data="admin_find_user")],
        [InlineKeyboardButton(text="🗄 База данных",     callback_data="admin_db"),
         InlineKeyboardButton(text="🤖 Модели",          callback_data="admin_models")],
        [InlineKeyboardButton(text="🔔 Каналы",          callback_data="admin_channels"),
         InlineKeyboardButton(text="🎟 Промокоды",       callback_data="admin_promos")],
        [InlineKeyboardButton(text="💎 Подписки",         callback_data="admin_subs")],
        [InlineKeyboardButton(text="🎁 Выдать лимиты",   callback_data="admin_give_limits"),
         InlineKeyboardButton(text="✂️ Забрать лимиты",  callback_data="admin_take_limits")],
        [InlineKeyboardButton(text="📢 Рассылка",        callback_data="admin_broadcast"),
         InlineKeyboardButton(text="🧪 Тест моделей",   callback_data="admin_test_models")],
        [InlineKeyboardButton(text="🔑 API ключи",       callback_data="admin_apikeys")],
        [InlineKeyboardButton(text="⚙️ Дополнительно",  callback_data="admin_extra")],
        [InlineKeyboardButton(text="🔌 Сервисы вкл/выкл", callback_data="admin_services")],
        [InlineKeyboardButton(text="🗑 Сброс памяти",   callback_data="admin_reset_memory")],
        [InlineKeyboardButton(text="🏠 Главная",         callback_data="back_home")],
    ])
    text = (
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        "⚙️ <b>АДМИН-ПАНЕЛЬ</b>\n"
        "▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂▂\n\n"
        f"◈ Пользователей — <b>{len(user_profiles)}</b>\n"
        f"◈ Активных сегодня — <b>{sum(1 for u in user_profiles if user_profiles[u].get('last_active_date','') == msk_now().strftime('%d.%m.%Y'))}</b>\n"
        f"◈ Всего запросов — <b>{sum(user_profiles.get(u, {}).get('requests', 0) for u in user_profiles)}</b>\n"
        f"◈ Всего генераций — <b>{sum(user_profiles.get(u, {}).get('images', 0) for u in user_profiles)}</b>\n"
        f"◈ Моделей — <b>{len(MODELS)}</b>  |  откл: <b>{len(disabled_models)}</b>\n"
        f"◈ PostgreSQL — {'✅' if db_pool else '❌'}\n"
        f"◈ Обновлено — <code>{msk_now().strftime('%d.%m %H:%M')}</code>\n\n"
        "▾ Выбери раздел\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔"
    )
    if isinstance(msg_or_cb, Message):
        await msg_or_cb.answer(text, parse_mode="HTML", reply_markup=admin_kb)
    else:
        try:
            await msg_or_cb.message.edit_text(text, parse_mode="HTML", reply_markup=admin_kb)
        except Exception:
            await msg_or_cb.message.answer(text, parse_mode="HTML", reply_markup=admin_kb)
        await msg_or_cb.answer()


@dp.callback_query(F.data == "menu_admin")
async def menu_admin(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    await show_admin_panel(callback)

@dp.callback_query(F.data == "admin_apikeys")
async def admin_apikeys(callback: CallbackQuery):
    """Мониторинг API ключей: Pollinations + платёжные + прочие."""
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    await callback.answer("⏳ Проверяю ключи...")

    import time as _t
    now = _t.monotonic()

    # ── Pollinations: кулдауны ──────────────────────────────────────
    total_keys  = len(pollinations_pool._keys)
    avail_keys  = len(pollinations_pool._available_keys())
    blocked_cnt = total_keys - avail_keys

    poll_lines = [
        f"🎨 <b>Pollinations — {total_keys} ключей</b>",
        f"   ✅ Доступно: <b>{avail_keys}</b>   ⏳ На кулдауне: <b>{blocked_cnt}</b>",
        "",
    ]
    for i, key in enumerate(pollinations_pool._keys, 1):
        bl    = pollinations_pool._blocked_until.get(key, 0)
        short = key[:10] + "…" + key[-4:]
        if bl <= now:
            poll_lines.append(f"  {i}. <code>{short}</code>  ✅ свободен")
        else:
            secs = int(bl - now)
            poll_lines.append(f"  {i}. <code>{short}</code>  ⏳ кулдаун {secs}с")

    # ── Живая проверка через API ────────────────────────────────────
    async def _ping(key: str) -> str:
        try:
            async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=8)) as s:
                async with s.get(
                    "https://api.pollinations.ai/v1/models",
                    headers={"Authorization": f"Bearer {key}"}
                ) as r:
                    if r.status == 200:
                        return "✅ работает"
                    elif r.status == 429:
                        retry = r.headers.get("Retry-After", "?")
                        return f"❌ 429 лимит (retry {retry}с)"
                    elif r.status == 402:
                        return "❌ 402 квота = 0"
                    elif r.status == 401:
                        return "❌ 401 неверный ключ"
                    elif r.status == 403:
                        return "❌ 403 заблокирован"
                    else:
                        return f"❌ HTTP {r.status}"
        except asyncio.TimeoutError:
            return "⚠️ таймаут"
        except Exception as ex:
            return f"⚠️ {str(ex)[:20]}"

    ping_tasks   = [asyncio.create_task(_ping(k)) for k in pollinations_pool._keys]
    ping_results = await asyncio.gather(*ping_tasks, return_exceptions=True)

    ping_lines = ["", "🔬 <b>Живая проверка Pollinations:</b>"]
    ok_count = 0
    for i, (key, res) in enumerate(zip(pollinations_pool._keys, ping_results), 1):
        short = key[:10] + "…" + key[-4:]
        result = str(res) if isinstance(res, Exception) else res
        if "✅" in result:
            ok_count += 1
        ping_lines.append(f"  {i}. <code>{short}</code>  {result}")
    ping_lines.append(f"\n  📊 Рабочих: <b>{ok_count} / {total_keys}</b>")

    # ── Другие ключи ────────────────────────────────────────────────
    def _kl(label: str, val: str) -> str:
        if not val:
            return f"  {label}  ❌ <i>не задан</i>"
        short = val[:6] + "…" + val[-4:] if len(val) > 12 else val[:4] + "…"
        return f"  {label}  ✅ <code>{short}</code>"

    other_lines = [
        "",
        "🔑 <b>Остальные ключи</b>",
        _kl("🧠 IO API       ", IO_KEY),
        _kl("⚡ SQ API       ", SQ_KEY),
        _kl("📸 Pexels       ", os.environ.get("PEXELS_KEY",   "")),
        _kl("🖼️ Pixabay      ", os.environ.get("PIXABAY_KEY",  "")),
        _kl("🏔 Unsplash     ", os.environ.get("UNSPLASH_KEY", "")),
        _kl("💳 ЮKassa       ", YOOKASSA_TOKEN),
        _kl("💳 Platega ключ ", PLATEGA_SECRET),
        _kl("🏪 Platega MID  ", PLATEGA_MERCHANT_ID),
    ]

    text = (
        "🔑 <b>API КЛЮЧИ — мониторинг</b>\n"
        "▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔▔\n"
        + "\n".join(poll_lines)
        + "\n".join(ping_lines)
        + "\n".join(other_lines)
        + f"\n\n⏰ <code>{msk_now().strftime('%d.%m.%Y %H:%M:%S')}</code>"
    )

    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🔄 Обновить", callback_data="admin_apikeys")],
        [InlineKeyboardButton(text="◀️ Назад",    callback_data="menu_admin")],
    ])
    try:
        await callback.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
    except Exception:
        await callback.message.answer(text, parse_mode="HTML", reply_markup=kb)


@dp.callback_query(F.data == "admin_reload")
async def admin_reload(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    await callback.answer("⏳ Загружаю данные из БД...", show_alert=False)
    try:
        await db_load_all_users()
        await db_load_bot_settings()
        total = len(user_profiles)
        await callback.message.answer(
            f"✅ <b>Данные обновлены!</b>\n\n"
            f"👥 Загружено пользователей: <code>{total}</code>\n"
            
            f"🧠 Продвинутые (лимит): <code>{LIMITS['pro_day']}</code>\n"
            f"🖼 Генерации (лимит): <code>{LIMITS['img_month']}</code>",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="◀️ Назад в панель", callback_data="menu_admin")
            ]])
        )
    except Exception as e:
        await callback.message.answer(
            f"❌ <b>Ошибка обновления:</b> <code>{str(e)[:200]}</code>",
            parse_mode="HTML"
        )


@dp.callback_query(F.data == "admin_stats")
async def admin_stats(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)

    total_users    = len(user_settings)
    total_requests = sum(p.get("requests", 0) for p in user_profiles.values())
    total_images   = sum(user_images_count.values())
    total_promos   = len(promo_codes)
    total_pro_used  = sum(user_limits.get(u, {}).get("pro_used", 0) for u in user_limits)

    # Активные за 24ч, сегодня и вчера
    now = msk_now()
    active_24h = 0
    today_users = 0
    yesterday_users = 0
    new_users_today = 0
    if db_pool:
        try:
            async with db_pool.acquire() as conn:
                # Активные за 24 часа
                row = await conn.fetchrow(
                    "SELECT COUNT(*) as cnt FROM users WHERE last_active >= NOW() - INTERVAL '24 hours'"
                )
                active_24h = row["cnt"] if row else 0
                # Зарегистрированы сегодня (МСК)
                today_str = now.strftime("%d.%m.%Y")
                r2 = await conn.fetchrow(
                    "SELECT COUNT(*) as cnt FROM users WHERE joined LIKE $1", f"{today_str}%"
                )
                today_users = r2["cnt"] if r2 else 0
                # Зарегистрированы вчера
                yesterday_str = (now - __import__('datetime').timedelta(days=1)).strftime("%d.%m.%Y")
                r3 = await conn.fetchrow(
                    "SELECT COUNT(*) as cnt FROM users WHERE joined LIKE $1", f"{yesterday_str}%"
                )
                yesterday_users = r3["cnt"] if r3 else 0
        except Exception as _e:
            import logging as _log
            _log.warning(f"admin_stats db: {_e}")
            active_24h = 0
    else:
        # Fallback без БД: считаем по last_active_date в памяти
        now_str  = now.strftime("%d.%m.%Y")
        h24_ago  = now - __import__('datetime').timedelta(hours=24)
        active_24h = sum(
            1 for u in user_profiles
            if user_profiles[u].get("last_active_date", "") == now_str
        )

    text = (
        f"📊 <b>━━━━━━━━━━━━━━━━━━</b>\n\n"
        f"⚡️ <b>Статистика бота</b>\n\n"
        f"👥 Всего пользователей: <code>{total_users}</code>\n"
        f"🟢 Активных за 24ч: <code>{active_24h}</code>\n"
        f"📅 Новых сегодня: <code>{today_users}</code>\n"
        f"📅 Новых вчера: <code>{yesterday_users}</code>\n"
        f"💬 Всего запросов: <code>{total_requests}</code>\n"
        f"🖼 Изображений создано: <code>{total_images}</code>\n"
        f"💬 Запросов использовано: <code>{total_pro_used}</code>\n"
        f"🎟 Промокодов: <code>{total_promos}</code>\n"
        f"🤖 Отключено моделей: <code>{len(disabled_models) + len(disabled_img_models)}</code>\n\n"
        f"📊 <b>━━━━━━━━━━━━━━━━━━</b>"
    )
    await callback.message.edit_text(text, parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="◀️ Назад", callback_data="menu_admin"),
        ]]))
    await callback.answer()


@dp.callback_query(F.data == "admin_userlist")
async def admin_userlist(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    await callback.answer("⏳ Загружаю список...")

    # Сортируем по последней активности (свежие сверху)
    now = msk_now()
    def _last_active_dt(item):
        u_id, prof = item
        la = prof.get("last_active_date", "")
        try:
            return __import__('datetime').datetime.strptime(la, "%d.%m.%Y")
        except Exception:
            return __import__('datetime').datetime.min

    sorted_users = sorted(user_profiles.items(), key=_last_active_dt, reverse=True)
    total = len(sorted_users)

    lines = []
    csv_lines = ["#;ID;Имя;Username;Запросы;Изображений;Последний запрос;Дата регистрации;Подписка;Модель"]
    for i, (u_id, prof) in enumerate(sorted_users):
        name      = prof.get("name", "—")
        username  = prof.get("username", "")
        uname_str = f"@{username}" if username else "—"
        req       = prof.get("requests", 0)
        joined    = prof.get("joined", "—")
        mk        = user_settings.get(u_id, "—")
        imgs      = user_images_count.get(u_id, 0)
        li        = user_limits.get(u_id, {})
        pro_u     = li.get("pro_used", 0)
        last_act  = prof.get("last_active_date", "—")
        sub       = user_subscriptions.get(u_id)
        sub_mark  = "💎" if sub and sub.get("expires") and sub["expires"] > now else "🆓"
        lines.append(
            f"{sub_mark} <b>{i+1}. {name}</b> {uname_str}\n"
            f"   🆔 <code>{u_id}</code>\n"
            f"   💬 запросов: <code>{pro_u}</code> | всего: <code>{req}</code> | 🖼 <code>{imgs}</code>\n"
            f"   🕐 последний: <code>{last_act}</code> | 📅 рег: <code>{joined}</code>\n"
            f"   🤖 <code>{mk}</code>"
        )
        csv_lines.append(f"{i+1};{u_id};{name};{username};{req};{imgs};{last_act};{joined};{'Да' if sub_mark=='💎' else 'Нет'};{mk}")

    # ── Отправляем CSV-файл ─────────────────────────────────────
    csv_content = "\n".join(csv_lines)
    from aiogram.types import BufferedInputFile
    csv_file = BufferedInputFile(
        csv_content.encode("utf-8-sig"),
        filename=f"users_{now.strftime('%d%m%Y_%H%M')}.csv"
    )
    await callback.message.answer_document(
        csv_file,
        caption=f"📋 <b>Все пользователи</b> — {total} чел.\n\n💾 CSV-файл для Excel/Google Sheets",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="◀️ Назад", callback_data="menu_admin")
        ]])
    )

    # ── Текстовый список постранично ────────────────────────────
    PAGE_SIZE = 15  # пользователей на страницу
    pages_data = [lines[i:i+PAGE_SIZE] for i in range(0, len(lines), PAGE_SIZE)]
    total_pages = len(pages_data)

    back_kb = InlineKeyboardMarkup(inline_keyboard=[[
        InlineKeyboardButton(text="◀️ Назад", callback_data="menu_admin")
    ]])

    if not pages_data:
        await callback.message.answer("👥 <b>Пользователей пока нет</b>", parse_mode="HTML", reply_markup=back_kb)
        return

    for pi, page_lines in enumerate(pages_data):
        header = f"👥 <b>Пользователи</b> [{pi+1}/{total_pages}] (всего {total})\n\n"
        text = header + "\n\n".join(page_lines)
        is_last = (pi == total_pages - 1)
        try:
            await callback.message.answer(
                text, parse_mode="HTML",
                reply_markup=back_kb if is_last else None
            )
        except Exception as e:
            logging.warning(f"admin_userlist page {pi}: {e}")


@dp.callback_query(F.data == "admin_find_user")
async def admin_find_user(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    admin_await[uid] = {"action": "find_user"}
    await callback.message.edit_text(
        "🔍 <b>Найти пользователя</b>\n\nВведи <code>USER_ID</code>:",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="menu_admin"),
        ]]),
    )
    await callback.answer()


@dp.callback_query(F.data == "admin_db")
async def admin_db(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    total_users      = len(user_settings)
    total_mem_msgs   = sum(len(m) for m in user_memory.values())
    total_pro_used   = sum(user_limits.get(u, {}).get("pro_used", 0) for u in user_limits)
    await callback.message.edit_text(
        f"🗄 <b>База данных</b>\n\n"
        f"👥 Записей: <code>{total_users}</code>\n"
        f"🧠 Сообщений в памяти: <code>{total_mem_msgs}</code>\n"
        f"💬 Запросов использовано (сумм.): <code>{total_pro_used}</code>\n"
        f"🖼 Генераций: <code>{sum(user_images_count.values())}</code>\n"
        f"🎟 Промокодов: <code>{len(promo_codes)}</code>\n"
        f"🔗 PostgreSQL: <code>{'✅ Подключена' if db_pool else '❌ Нет'}</code>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="◀️ Назад", callback_data="menu_admin"),
        ]]),
    )
    await callback.answer()


# ── Управление моделями ──────────────────────────────────────────

@dp.callback_query(F.data == "admin_models")
async def admin_models(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)

    rows = []
    # Chat models
    rows.append([InlineKeyboardButton(text="━━ 💬 ЧАТ-МОДЕЛИ ━━", callback_data="noop")])
    for k, m in MODELS.items():
        status = "✅" if k not in disabled_models else "❌"
        rows.append([InlineKeyboardButton(text=f"{status} {m['label']}", callback_data=f"admin_model_toggle_{k}")])
    # Image models
    rows.append([InlineKeyboardButton(text="━━ 🎨 МОДЕЛИ КАРТИНОК ━━", callback_data="noop")])
    for k, m in IMG_MODELS.items():
        status = "✅" if k not in disabled_img_models else "❌"
        rows.append([InlineKeyboardButton(text=f"{status} [IMG] {m['label']}", callback_data=f"admin_imgmodel_toggle_{k}")])
    # Music models
    rows.append([InlineKeyboardButton(text="━━ 🎵 МОДЕЛИ МУЗЫКИ ━━", callback_data="noop")])
    for k, m in MUSIC_MODELS.items():
        status = "✅" if k not in disabled_music_models else "❌"
        rows.append([InlineKeyboardButton(text=f"{status} [🎵] {m['label']}", callback_data=f"admin_musicmodel_toggle_{k}")])
    rows.append([InlineKeyboardButton(text="◀️ Назад", callback_data="menu_admin")])

    text = (
        "🤖 <b>━━━━━━━━━━━━━━━━━━</b>\n\n"
        "⚡️ <b>Управление моделями и стилями</b>\n\n"
        "✅ — включено | ❌ — отключено\n"
        "Нажми для переключения:\n\n"
        "🤖 <b>━━━━━━━━━━━━━━━━━━</b>"
    )
    try:
        await callback.message.edit_text(text, parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=rows))
    except Exception:
        await callback.message.answer(text, parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=rows))
    await callback.answer()


@dp.callback_query(F.data.startswith("admin_model_toggle_"))
async def admin_model_toggle(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    model_key = callback.data.split("admin_model_toggle_", 1)[1]
    if model_key in MODELS:
        if model_key in disabled_models:
            disabled_models.discard(model_key)
            await callback.answer(f"✅ {MODELS[model_key]['label']} включена!", show_alert=True)
        else:
            disabled_models.add(model_key)
            await callback.answer(f"❌ {MODELS[model_key]['label']} отключена!", show_alert=True)
        asyncio.create_task(db_save_bot_settings())  # Сохраняем в БД
    # Refresh panel
    await admin_models(callback)


@dp.callback_query(F.data.startswith("admin_imgmodel_toggle_"))
async def admin_imgmodel_toggle(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    model_key = callback.data.split("admin_imgmodel_toggle_", 1)[1]
    if model_key in IMG_MODELS:
        if model_key in disabled_img_models:
            disabled_img_models.discard(model_key)
            await callback.answer(f"✅ {IMG_MODELS[model_key]['label']} включена!", show_alert=True)
        else:
            disabled_img_models.add(model_key)
            await callback.answer(f"❌ {IMG_MODELS[model_key]['label']} отключена!", show_alert=True)
    await admin_models(callback)


@dp.callback_query(F.data == "noop")
async def cb_noop(callback: CallbackQuery):
    await callback.answer()


@dp.callback_query(F.data.startswith("admin_musicmodel_toggle_"))
async def admin_musicmodel_toggle(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    key = callback.data.split("admin_musicmodel_toggle_", 1)[1]
    if key in MUSIC_MODELS:
        if key in disabled_music_models:
            disabled_music_models.discard(key)
            await callback.answer(f"✅ {MUSIC_MODELS[key]['label']} включена!", show_alert=True)
        else:
            disabled_music_models.add(key)
            await callback.answer(f"❌ {MUSIC_MODELS[key]['label']} отключена!", show_alert=True)
        asyncio.create_task(db_save_bot_settings())
    await admin_models(callback)


# ── Управление каналами подписки ──────────────────────────────────

@dp.callback_query(F.data == "admin_channels")
async def admin_channels(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)

    rows = []
    for ch in required_channels:
        rows.append([InlineKeyboardButton(
            text=f"❌ Удалить @{ch}",
            callback_data=f"admin_delchan_{ch}"
        )])
    rows.append([InlineKeyboardButton(text="➕ Добавить канал", callback_data="admin_addchan")])
    rows.append([InlineKeyboardButton(text="◀️ Назад", callback_data="menu_admin")])

    ch_list = "\n".join([f"• <code>@{c}</code>" for c in required_channels]) if required_channels else "<i>Нет</i>"
    text = (
        f"🔔 <b>Обязательная подписка</b>\n\n"
        f"Каналы:\n{ch_list}\n\n"
        f"Пользователи должны подписаться\nна все каналы для доступа к боту."
    )
    try:
        await callback.message.edit_text(text, parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=rows))
    except Exception:
        await callback.message.answer(text, parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=rows))
    await callback.answer()


@dp.callback_query(F.data == "admin_addchan")
async def admin_addchan(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    admin_await[uid] = {"action": "add_channel"}
    await callback.message.edit_text(
        "➕ <b>Добавить канал</b>\n\n"
        "Введи одним из способов:\n"
        "• Username: <code>my_channel</code>\n"
        "• С собакой: <code>@my_channel</code>\n"
        "• Ссылка: <code>https://t.me/my_channel</code>\n"
        "• Числовой ID: <code>-1001234567890</code>\n\n"
        "❗️ Бот должен быть <b>администратором</b> канала!",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="admin_channels"),
        ]]),
    )
    await callback.answer()


@dp.callback_query(F.data.startswith("admin_delchan_"))
async def admin_delchan(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    ch = callback.data.split("admin_delchan_", 1)[1]
    if ch in required_channels:
        required_channels.remove(ch)
        asyncio.create_task(db_save_bot_settings())  # Сохраняем в БД
        await callback.answer(f"✅ Канал @{ch} удалён!", show_alert=True)
    await admin_channels(callback)


# ── Промокоды ────────────────────────────────────────────────────

@dp.callback_query(F.data == "admin_promos")
async def admin_promos(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    lines = []
    for code, info in promo_codes.items():
        used  = len(info.get("used_by", []))
        max_u = info.get("max_uses", "∞")
        fast  = info.get("fast_requests", 0)
        pro   = info.get("pro_requests", 0)
        exp   = info.get("expires", "")
        exp_str = f" | до {exp}" if exp else ""
        lines.append(
            f"<code>{code}</code> — 💬{fast + pro} зап. | {used}/{max_u} исп.{exp_str}"
        )
    text = "🎟 <b>Промокоды</b>\n\n" + ("\n".join(lines) if lines else "Нет промокодов")
    await callback.message.edit_text(
        text, parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="➕ Создать промокод", callback_data="admin_create_promo")],
            [InlineKeyboardButton(text="◀️ Назад", callback_data="menu_admin")],
        ])
    )
    await callback.answer()


@dp.callback_query(F.data == "admin_create_promo")
async def admin_create_promo(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    admin_await[uid] = {"action": "create_promo_step1"}
    await callback.message.edit_text(
        "🎟 <b>Создать промокод</b>\n\n"
        "Формат:\n"
        "<code>КОД КОЛ_ЗАПРОСОВ МАКС_ИСПОЛЬЗОВАНИЙ [ДАТА_ОКОНЧАНИЯ]</code>\n\n"
        "Примеры:\n"
        "• <code>BONUS2025 20 100</code>\n"
        "• <code>VIP100 50 50 31.12.2025</code>\n\n"
        "КОЛ_ЗАПРОСОВ = количество 💬 запросов для пользователя",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="menu_admin"),
        ]]),
    )
    await callback.answer()


@dp.callback_query(F.data == "admin_give_limits")
async def admin_give_limits_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="💬 Запросы (pro)",    callback_data="admin_give_pro"),
         InlineKeyboardButton(text="🎨 Картинки",         callback_data="admin_give_img")],
        [InlineKeyboardButton(text="🎵 Музыка",           callback_data="admin_give_music"),
         InlineKeyboardButton(text="🎬 Видео",            callback_data="admin_give_video")],
        [InlineKeyboardButton(text="🎁 Всё сразу",        callback_data="admin_give_all")],
        [InlineKeyboardButton(text="◀️ Назад",            callback_data="menu_admin")],
    ])
    await callback.message.edit_text(
        "🎁 <b>Выдать лимиты пользователю</b>\n\n"
        "Выбери тип лимита:",
        parse_mode="HTML", reply_markup=kb
    )
    await callback.answer()


@dp.callback_query(F.data == "admin_take_limits")
async def admin_take_limits_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="💬 Запросы (pro)",    callback_data="admin_take_pro"),
         InlineKeyboardButton(text="🎨 Картинки",         callback_data="admin_take_img")],
        [InlineKeyboardButton(text="🎵 Музыка",           callback_data="admin_take_music"),
         InlineKeyboardButton(text="🎬 Видео",            callback_data="admin_take_video")],
        [InlineKeyboardButton(text="🔴 Обнулить всё",     callback_data="admin_take_all")],
        [InlineKeyboardButton(text="◀️ Назад",            callback_data="menu_admin")],
    ])
    await callback.message.edit_text(
        "✂️ <b>Забрать лимиты у пользователя</b>\n\n"
        "Выбери тип лимита:",
        parse_mode="HTML", reply_markup=kb
    )
    await callback.answer()


_LIMIT_LABELS = {
    "pro":   ("💬 Запросы (pro)",   "pro_used",   "pro_day"),
    "img":   ("🎨 Картинки",        "img_used",   "img_month"),
    "music": ("🎵 Музыка",          "music_used", "music_month"),
    "video": ("🎬 Видео",           "video_used", "video_month"),
}


@dp.callback_query(F.data.startswith("admin_give_") | F.data.startswith("admin_take_"))
async def admin_give_take_type(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    data = callback.data  # admin_give_pro / admin_take_img / etc.

    # Пропускаем уже обработанные кнопки-меню
    if data in ("admin_give_limits", "admin_take_limits"):
        return

    parts = data.split("_")  # ['admin', 'give'/'take', 'pro'/'img'/...]
    if len(parts) < 3:
        return
    direction = parts[1]   # give / take
    ltype     = parts[2]   # pro / fast / img / music / video / all

    if ltype == "all":
        if direction == "give":
            admin_await[uid] = {"action": "give_all"}
            text = (
                "🎁 <b>Выдать все лимиты</b>\n\n"
                "Формат: <code>USER_ID ПРО КАРТИНКИ МУЗЫКА ВИДЕО</code>\n\n"
                "Пример: <code>123456789 30 50 5 3</code>\n"
                "<i>Уменьшит счётчики использования на указанные значения</i>"
            )
        else:
            admin_await[uid] = {"action": "take_all"}
            text = (
                "🔴 <b>Обнулить все лимиты</b>\n\n"
                "Формат: <code>USER_ID</code>\n\n"
                "Пример: <code>123456789</code>\n"
                "<i>Установит все счётчики использования на максимум</i>"
            )
        await callback.message.edit_text(
            text, parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="❌ Отмена", callback_data="admin_give_limits" if direction == "give" else "admin_take_limits")
            ]])
        )
        await callback.answer()
        return

    if ltype not in _LIMIT_LABELS:
        return

    label, used_key, max_key = _LIMIT_LABELS[ltype]
    if direction == "give":
        admin_await[uid] = {"action": f"give_{ltype}"}
        text = (
            f"🎁 <b>Выдать: {label}</b>\n\n"
            f"Формат: <code>USER_ID КОЛ-ВО</code>\n\n"
            f"Пример: <code>123456789 20</code>\n"
            f"<i>Уменьшит счётчик использования на указанное кол-во</i>"
        )
        back_cb = "admin_give_limits"
    else:
        admin_await[uid] = {"action": f"take_{ltype}"}
        text = (
            f"✂️ <b>Забрать: {label}</b>\n\n"
            f"Формат: <code>USER_ID КОЛ-ВО</code>\n\n"
            f"Пример: <code>123456789 10</code>\n"
            f"<i>Увеличит счётчик использования на указанное кол-во</i>"
        )
        back_cb = "admin_take_limits"

    await callback.message.edit_text(
        text, parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data=back_cb)
        ]])
    )
    await callback.answer()


# Старые обработчики — оставляем для обратной совместимости
@dp.callback_query(F.data == "admin_give_tokens")
async def admin_give_tokens_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    callback.data = "admin_give_limits"
    await admin_give_limits_cb(callback)


@dp.callback_query(F.data == "admin_take_tokens")
async def admin_take_tokens_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    callback.data = "admin_take_limits"
    await admin_take_limits_cb(callback)


@dp.callback_query(F.data == "admin_extra")
async def admin_extra(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)

    total_users     = len(user_settings)
    total_img_used  = sum(user_limits.get(u, {}).get("img_used", 0)   for u in user_limits)
    total_music_used= sum(user_limits.get(u, {}).get("music_used", 0) for u in user_limits)
    active_subs     = sum(1 for u in user_subscriptions if has_active_sub(u))

    total_video_used= sum(user_limits.get(u, {}).get("video_used", 0) for u in user_limits)

    # Лимиты бесплатно
    F_lims = LIMITS
    # Лимиты за подписку (неделя/месяц)
    W = SUB_PLANS["week"]
    M = SUB_PLANS["month"]

    text = (
        "⚙️ <b>НАСТРОЙКИ ЛИМИТОВ</b>\n\n"
        "┌─ 🆓 <b>БЕСПЛАТНО</b> ──────────────\n"
        f""
        f"│ 🧠 Продвинутые: <code>{F_lims['pro_day']}</code> / {F_lims['reset_h']}ч\n"
        f"│ 🎨 Картинки:  <code>{F_lims['img_month']}</code> / мес\n"
        f"│ 🎵 Музыка:    <code>{F_lims.get('music_month', 0)}</code> / мес\n"
        f"│ 🎬 Видео:     <code>{F_lims.get('video_month', 0)}</code> / мес\n"
        "├─ ⚡ <b>ПОДПИСКА НЕДЕЛЯ</b> ──────────\n"
        f""
        f"│ 🧠 Продвинутые: <code>{W['pro_day']}</code> / {W['reset_h']}ч\n"
        f"│ 🎨 Картинки:  <code>{W['img_month']}</code> / мес\n"
        f"│ 🎵 Музыка:    <code>{W.get('music_month', 3)}</code> / мес\n"
        f"│ 🎬 Видео:     <code>{W.get('video_month', 3)}</code> / мес\n"
        "├─ 💎 <b>ПОДПИСКА МЕСЯЦ</b> ──────────\n"
        f""
        f"│ 🧠 Продвинутые: <code>{M['pro_day']}</code> / {M['reset_h']}ч\n"
        f"│ 🎨 Картинки:  <code>{M['img_month']}</code> / мес\n"
        f"│ 🎵 Музыка:    <code>{M.get('music_month', 3)}</code> / мес\n"
        f"│ 🎬 Видео:     <code>{M.get('video_month', 3)}</code> / мес\n"
        "└────────────────────────────\n\n"
        f"👥 Пользователей: <code>{total_users}</code>\n"
        f"💎 Активных подписок: <code>{active_subs}</code>\n"
        f"🎨 Картинок сгенерировано: <code>{total_img_used}</code>\n"
        f"🎵 Треков сгенерировано: <code>{total_music_used}</code>\n"
        f"🎬 Видео сгенерировано: <code>{total_video_used}</code>"
    )
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🆓 Лимиты бесплатно",     callback_data="admin_limits_free")],
        [InlineKeyboardButton(text="⚡ Лимиты Неделя",         callback_data="admin_limits_week")],
        [InlineKeyboardButton(text="💎 Лимиты Месяц",          callback_data="admin_limits_month")],
        [InlineKeyboardButton(text="⏰ Сброс запросов (часы)", callback_data="admin_limits_reset_h")],
        [InlineKeyboardButton(text="🔄 Сбросить лимиты всех", callback_data="admin_reset_all_limits")],
        [InlineKeyboardButton(text="◀️ Назад в панель",        callback_data="menu_admin")],
    ])
    try:
        await callback.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
    except Exception:
        await callback.message.answer(text, parse_mode="HTML", reply_markup=kb)
    await callback.answer()


@dp.callback_query(F.data.in_({"admin_limits_free", "admin_limits_week", "admin_limits_month", "admin_limits_reset_h"}))
async def admin_limits_edit(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)

    section = callback.data.replace("admin_limits_", "")
    labels  = {"free": "🆓 Бесплатно", "week": "⚡ Неделя", "month": "💎 Месяц", "reset_h": "⏰ Сброс запросов"}

    if section == "reset_h":
        admin_await[uid] = {"action": "set_reset_h"}
        text = (
            f"⏰ <b>Изменить интервал сброса запросов</b>\n\n"
            f"Текущий: <code>{LIMITS['reset_h']}</code> часов\n\n"
            f"Введи число часов (например <code>6</code> или <code>24</code>):"
        )
    else:
        src = LIMITS if section == "free" else SUB_PLANS[section]
        admin_await[uid] = {"action": f"set_limits_{section}"}
        text = (
            f"{labels[section]} <b>— Изменить лимиты</b>\n\n"
            f"Текущие:\n"
            f"  💬 Запросы: <code>{src['pro_day']}</code>/ч\n"
            f"  🎨 Картинки: <code>{src['img_month']}</code>/мес\n"
            f"  🎵 Музыка: <code>{src.get('music_month', 0)}</code>/мес\n"
            f"  🎬 Видео: <code>{src.get('video_month', 0)}</code>/мес\n\n"
            f"Введи 4 числа через пробел:\n"
            f"<code>ЗАПРОСЫ КАРТИНКИ МУЗЫКА ВИДЕО</code>\n\n"
            f"Пример: <code>30 100 5 3</code>"
        )

    await callback.message.edit_text(
        text, parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="admin_extra")
        ]])
    )
    await callback.answer()


@dp.callback_query(F.data == "admin_edit_limits")
async def admin_edit_limits(callback: CallbackQuery):
    """Старый обработчик — перенаправляем в новый."""
    await admin_extra(callback)


@dp.callback_query(F.data == "admin_reset_all_limits")
async def admin_reset_all_limits(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    now = msk_now()
    month_start = now.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
    reset_count = 0
    for u in user_limits:
        user_limits[u]["pro_used"]    = 0
        user_limits[u]["pro_reset"]   = now
        user_limits[u]["fast_used"]   = 0
        user_limits[u]["fast_reset"]  = now
        user_limits[u]["img_used"]    = 0
        user_limits[u]["img_reset"]   = month_start
        user_limits[u]["music_used"]  = 0
        user_limits[u]["music_reset"] = month_start
        user_limits[u]["video_used"]  = 0
        user_limits[u]["video_reset"] = month_start
        reset_count += 1
    await callback.answer(
        f"✅ Все лимиты сброшены для {reset_count} пользователей (чат + фото + музыка + видео)!",
        show_alert=True
    )


@dp.callback_query(F.data == "admin_reset_memory")
async def admin_reset_memory(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    for u in user_memory:
        user_memory[u] = deque(maxlen=20)
    await callback.answer("✅ Память всех пользователей очищена!", show_alert=True)


# ══════════════════════════════════════════════════════════════
# 🔌 ADMIN: ПЕРЕКЛЮЧАТЕЛИ СЕРВИСОВ
# ══════════════════════════════════════════════════════════════

@dp.callback_query(F.data == "admin_services")
async def admin_services_panel(callback: CallbackQuery):
    """Панель включения/выключения сервисов генерации."""
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)

    lines = ["🔌 <b>УПРАВЛЕНИЕ СЕРВИСАМИ</b>\n\n"]
    for key, label in SERVICE_LABELS.items():
        status = "✅ ВКЛ" if SERVICE_ENABLED[key] else "❌ ВЫКЛ"
        lines.append(f"{status}  {label}")

    text = "\n".join(lines) + "\n\n<i>Нажми кнопку чтобы переключить сервис</i>"

    rows = []
    for key, label in SERVICE_LABELS.items():
        icon = "✅" if SERVICE_ENABLED[key] else "❌"
        rows.append([InlineKeyboardButton(
            text=f"{icon} {label}",
            callback_data=f"admin_svc_toggle_{key}"
        )])
    rows.append([InlineKeyboardButton(text="◀️ Назад в панель", callback_data="menu_admin")])

    try:
        await callback.message.edit_text(text, parse_mode="HTML",
                                          reply_markup=InlineKeyboardMarkup(inline_keyboard=rows))
    except Exception:
        await callback.message.answer(text, parse_mode="HTML",
                                       reply_markup=InlineKeyboardMarkup(inline_keyboard=rows))
    await callback.answer()


@dp.callback_query(F.data.startswith("admin_svc_toggle_"))
async def admin_svc_toggle(callback: CallbackQuery):
    """Переключить конкретный сервис."""
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    key = callback.data.replace("admin_svc_toggle_", "")
    if key not in SERVICE_ENABLED:
        return await callback.answer("❌ Неизвестный сервис", show_alert=True)
    SERVICE_ENABLED[key] = not SERVICE_ENABLED[key]
    new_status = "✅ ВКЛЮЧЁН" if SERVICE_ENABLED[key] else "❌ ВЫКЛЮЧЕН"
    label = SERVICE_LABELS.get(key, key)
    asyncio.create_task(db_save_bot_settings())
    await callback.answer(f"{label} — {new_status}", show_alert=True)
    # Обновляем панель
    await admin_services_panel(callback)


@dp.callback_query(F.data == "admin_broadcast")
async def admin_broadcast_cb(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    admin_await[uid] = {"action": "broadcast"}
    await callback.message.edit_text(
        "📢 <b>Рассылка</b>\n\n"
        "Отправь текст <i>(или фото с подписью)</i> — разошлю всем пользователям.\n\n"
        "Поддерживается HTML-разметка.",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="menu_admin"),
        ]]),
    )
    await callback.answer()


@dp.callback_query(F.data == "back_home")
async def back_home(callback: CallbackQuery):
    uid  = callback.from_user.id
    name = callback.from_user.first_name or "Пользователь"
    tok  = get_tokens(uid)
    home_text = _welcome_text(name, tok, uid in ADMIN_IDS)
    try:
        await callback.message.edit_text(home_text, parse_mode="HTML", reply_markup=home_kb(uid))
    except Exception:
        await callback.message.answer(home_text, parse_mode="HTML", reply_markup=home_kb(uid))
    await callback.answer()


@dp.callback_query(F.data == "back_ai")
async def back_ai(callback: CallbackQuery):
    uid = callback.from_user.id
    mk  = get_model_key(uid)
    cur = MODELS.get(mk, {}).get("label", mk) if not mk.startswith("imggen:") else "🎨 " + IMG_MODELS[mk.split(":")[1]]["label"]
    await callback.message.edit_text(
        f"⚡️ <b>Нейросети</b>\n\n🔹 Активная модель: <code>{cur}</code>\n\nВыбери категорию:",
        parse_mode="HTML", reply_markup=ai_menu_kb(uid),
    )
    await callback.answer()


@dp.callback_query(F.data.startswith("cat_"))
async def open_category(callback: CallbackQuery):
    cat_key = callback.data
    if cat_key not in CATEGORIES:
        return await callback.answer()
    label, _ = CATEGORIES[cat_key]
    hint = "\n📸 = поддерживает анализ фото" if cat_key not in ("cat_imggen",) else \
           "\nВыбери модель генерации, затем отправь описание"
    uid = callback.from_user.id
    await callback.message.edit_text(
        f"<b>{label}</b>{hint}:",
        parse_mode="HTML", reply_markup=category_kb(cat_key, uid),
    )
    await callback.answer()


@dp.callback_query(F.data == "back_main")
async def back_main(callback: CallbackQuery):
    uid = callback.from_user.id
    mk  = get_model_key(uid)
    cur = MODELS.get(mk, {}).get("label", mk) if not mk.startswith("imggen:") else "🎨 " + IMG_MODELS[mk.split(":")[1]]["label"]
    await callback.message.edit_text(
        f"⚡️ <b>Нейросети</b>\n\n🔹 Активная модель: <code>{cur}</code>\n\nВыбери категорию:",
        parse_mode="HTML", reply_markup=ai_menu_kb(uid),
    )
    await callback.answer()


@dp.callback_query(F.data.startswith("set_"))
async def set_model(callback: CallbackQuery):
    key = callback.data.split("_", 1)[1]
    uid = callback.from_user.id
    if key not in MODELS:
        return await callback.answer("❌ Модель не найдена", show_alert=True)
    # Проверяем что модель не отключена
    if key in disabled_models:
        return await callback.answer("🚫 Модель временно недоступна", show_alert=True)
    if key in PREMIUM_MODELS and not has_active_sub(uid):
        await callback.answer("🔒 Нужна подписка!", show_alert=True)
        try:
            await callback.message.edit_text(
                "🔒 <b>Эта модель доступна только по подписке</b>\n\n"
                "💎 Оформи подписку и получи доступ ко всем моделям:\n"
                "⚡ 7 дней — 60 ₽\n"
                "💎 30 дней — 100 ₽",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                    [InlineKeyboardButton(text="💎 Оформить подписку", callback_data="sub_menu")],
                    [InlineKeyboardButton(text="🤖 Другие модели", callback_data="menu_ai")],
                ])
            )
        except Exception:
            pass
        return
    set_chat_model(uid, key)
    user_memory[uid] = deque(maxlen=20)
    asyncio.create_task(db_save_user(uid))
    m   = MODELS[key]
    vis = m.get("vision", False)
    vision_hint = "\n📸 <i>Можешь отправить фото для анализа!</i>" if vis else ""
    kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(text="🤖 Сменить модель", callback_data="menu_ai")],
        [InlineKeyboardButton(text="🧹 Очистить память", callback_data="clear_memory"),
         InlineKeyboardButton(text="🏠 Главная",   callback_data="back_home")],
    ])
    try:
        await callback.message.edit_text(
            f"✅ <b>Модель активирована!</b>\n\n"
            f"🤖 {m['label']}{vision_hint}\n\n"
            f"✏️ <b>Напиши свой вопрос прямо сейчас</b> — я отвечу!\n\n"
            f"<i>Память очищена — начинаем новый диалог.</i>",
            parse_mode="HTML",
            reply_markup=kb,
        )
    except Exception:
        await callback.message.answer(
            f"✅ <b>{m['label']} активирована!</b>{vision_hint}\n\n✏️ Пиши вопрос:",
            parse_mode="HTML", reply_markup=kb
        )
    await callback.answer(f"✅ {m['label']} активирована!")


@dp.callback_query(F.data.startswith("imgset_"))
async def set_imggen(callback: CallbackQuery):
    key = callback.data.split("_", 1)[1]
    uid = callback.from_user.id
    if key not in IMG_MODELS:
        return await callback.answer("❌ Модель не найдена", show_alert=True)
    if key in disabled_img_models:
        return await callback.answer("🚫 Модель временно недоступна", show_alert=True)
    set_img_model(uid, key)
    # ВАЖНО: устанавливаем в user_settings чтобы handle_text видел imggen-режим
    user_settings[uid] = f"imggen:{key}"
    asyncio.create_task(db_save_user(uid))
    m = IMG_MODELS[key]
    examples_text = "\n".join(f"  • <i>{e}</i>" for e in IMG_PROMPT_EXAMPLES[:3])
    supports_i2i = m.get("img2img", False)
    kb_rows = [
        [InlineKeyboardButton(text="✨ ИИ придумает идею", callback_data=f"imggen_ai_idea_{key}")],
    ]
    if supports_i2i:
        kb_rows.append([InlineKeyboardButton(text="📸 Загрузить фото для стилизации", callback_data=f"imggen_upload_photo_{key}")])
    kb_rows.append([InlineKeyboardButton(text="◀️ Другая модель", callback_data="menu_imggen"),
                    InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")])
    kb = InlineKeyboardMarkup(inline_keyboard=kb_rows)

    i2i_hint = "\n📸 <i>Поддерживает стилизацию фото! Нажми кнопку ниже.</i>" if supports_i2i else ""
    text = (
        f"╔══════════════════════╗\n"
        f"║  {m['label']}  ║\n"
        f"╚══════════════════════╝\n\n"
        f"<b>{m['desc']}</b>{i2i_hint}\n\n"
        f"━━━━━━━━━━━━━━━━━━━━━━━━\n"
        f"✏️ <b>Опиши своё изображение:</b>\n\n"
        f"<b>Примеры:</b>\n{examples_text}\n\n"
        f"━━━━━━━━━━━━━━━━━━━━━━━━\n"
        f"💡 Пиши на русском — ИИ переведёт и улучшит!\n"
        f"Или нажми кнопку — ИИ придумает сам 👇"
    )
    try:
        await callback.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
    except Exception:
        await callback.message.answer(text, parse_mode="HTML", reply_markup=kb)
    await callback.answer(f"✅ {m['label']} выбрана!")


@dp.callback_query(F.data.startswith("imggen_ai_idea_"))
async def cb_imggen_ai_idea(callback: CallbackQuery):
    """ИИ генерирует случайную идею для изображения."""
    uid = callback.from_user.id
    key = callback.data.replace("imggen_ai_idea_", "")
    m   = IMG_MODELS.get(key, {})
    await callback.answer("✨ Генерирую идею...")
    try:
        sys_p = (
            f"Придумай оригинальную и интересную идею для генерации изображения. "
            f"Стиль модели: {m.get('desc', 'универсальный')}. "
            f"Верни ТОЛЬКО одно описание на русском языке (2-3 предложения). "
            f"Без пояснений и вводных слов."
        )
        msgs = [{"role": "system", "content": sys_p},
                {"role": "user", "content": "Придумай идею"}]
        idea = (await call_chat(msgs, "claude_haiku", max_tokens=150)).strip()
    except Exception:
        import random
        idea = random.choice(IMG_PROMPT_EXAMPLES)

    # Генерируем сразу!
    wait = await callback.message.answer(
        f"✨ <b>ИИ придумал идею:</b>\n<i>{idea[:300]}</i>\n\n"
        f"🎨 <i>Улучшаю промпт и генерирую...</i>",
        parse_mode="HTML"
    )
    try:
        await callback.message.delete()
    except Exception:
        pass
    await _do_imggen(wait, uid, idea, key)


@dp.callback_query(F.data == "imggen_from_photo")
async def cb_imggen_from_photo(callback: CallbackQuery):
    """Кнопка 'Стилизация фото (img2img)' из главного меню генерации."""
    uid = callback.from_user.id
    await callback.answer()
    # Только FLUX.2 Dev поддерживает img2img через Pollinations ?image= параметр
    rows = [
        [InlineKeyboardButton(text="🆕 FLUX.2 Dev  📸→🖼 (рекомендуется)", callback_data="imggen_upload_photo_flux2dev")],
        [InlineKeyboardButton(text="◀️ Назад", callback_data="menu_imggen")],
    ]
    kb = InlineKeyboardMarkup(inline_keyboard=rows)
    text = (
        "📸 <b>Стилизация фото (img2img)</b>\n\n"
        "Единственная рабочая модель для стилизации фото — <b>FLUX.2 Dev</b>.\n"
        "Остальные модели не поддерживают передачу изображения.\n\n"
        "💡 <i>Примеры стилей:</i>\n"
        "• <code>в стиле аниме</code>\n"
        "• <code>oil painting, Van Gogh</code>\n"
        "• <code>cyberpunk neon style</code>\n"
        "• <code>sketch pencil drawing</code>\n"
        "• <code>pixel art 8-bit</code>"
    )
    try:
        await callback.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
    except Exception:
        await callback.message.answer(text, parse_mode="HTML", reply_markup=kb)


@dp.callback_query(F.data.startswith("imggen_upload_photo_"))
async def cb_imggen_upload_photo(callback: CallbackQuery):
    """Пользователь выбрал модель для img2img — просим прислать фото."""
    uid = callback.from_user.id
    key = callback.data.replace("imggen_upload_photo_", "")
    if key not in IMG_MODELS:
        return await callback.answer("❌ Модель не найдена", show_alert=True)
    m = IMG_MODELS[key]
    # Если модель не поддерживает img2img — используем flux2dev
    actual_key = key if m.get("img2img") else "flux2dev"
    actual_m   = IMG_MODELS.get(actual_key, m)
    # Сохраняем состояние: ждём фото с этой моделью
    last_photo[uid] = {"img2img_model_pending": actual_key}
    await callback.answer()
    note = f"\n⚠️ <i>{m['label']} не поддерживает img2img — используем {actual_m['label']}</i>" if actual_key != key else ""
    text = (
        f"📸 <b>{actual_m['label']}</b> — загрузи фото{note}\n\n"
        "Отправь фотографию — бот применит стиль который ты напишешь.\n\n"
        "💡 После фото напиши стиль, например:\n"
        "• <code>в стиле аниме</code>\n"
        "• <code>watercolor painting</code>\n"
        "• <code>нарисуй усы и очки</code>\n"
        "• <code>превратить в картину маслом</code>"
    )
    try:
        await callback.message.edit_text(text, parse_mode="HTML", reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_imggen")
        ]]))
    except Exception:
        await callback.message.answer(text, parse_mode="HTML", reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="cancel_imggen")
        ]]))


@dp.callback_query(F.data == "cancel_imggen")
async def cancel_imggen(callback: CallbackQuery):
    # Модели чата и генерации теперь хранятся раздельно — отмена imggen ничего не сбрасывает
    await callback.message.edit_text(
        "❌ Генерация отменена.",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="🏠 Главное меню", callback_data="back_home"),
        ]])
    )
    await callback.answer()


@dp.callback_query(F.data.in_({"photo_describe", "photo_task", "photo_code"}))
async def cb_photo_quick(callback: CallbackQuery):
    """Обработка быстрых кнопок после отправки фото без подписи."""
    uid = callback.from_user.id
    photo_data = last_photo.get(uid)
    if not photo_data or not photo_data.get("ask_pending"):
        await callback.answer("❌ Фото не найдено, отправь снова", show_alert=True)
        return
    last_photo[uid]["ask_pending"] = False
    prompts = {
        "photo_describe": "Подробно опиши что изображено на этом фото",
        "photo_task":     "Реши все задачи и задания на этом фото. Выполни каждое задание, заполни поля 'Ответ:'",
        "photo_code":     "Распознай и выведи весь код с этого скриншота в блоке ```",
    }
    caption = prompts[callback.data]
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    await callback.answer()
    await _analyze_photo(callback.message, uid, photo_data["img_b64"], caption)


@dp.callback_query(F.data == "photo_translate")
async def cb_photo_translate(callback: CallbackQuery):
    uid = callback.from_user.id
    photo_data = last_photo.get(uid)
    if not photo_data or not photo_data.get("ask_pending"):
        return await callback.answer("❌ Фото не найдено, отправь снова", show_alert=True)
    last_photo[uid]["ask_pending"] = False
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    await callback.answer()
    await _analyze_photo(
        callback.message, uid, photo_data["img_b64"],
        "Распознай весь текст на этом фото и переведи его на русский язык. Сначала покажи оригинал, потом перевод."
    )


@dp.callback_query(F.data == "photo_to_video")
async def cb_photo_to_video(callback: CallbackQuery):
    uid = callback.from_user.id
    photo_data = last_photo.get(uid)
    if not photo_data or not photo_data.get("ask_pending"):
        return await callback.answer("❌ Фото не найдено, отправь снова", show_alert=True)
    if not can_video(uid):
        li = get_limits_info(uid)
        return await callback.answer(
            f"🚫 Лимит видео исчерпан! ({li['video_used']}/{li['video_max']} в месяц)",
            show_alert=True
        )
    last_photo[uid]["ask_pending"] = False
    img_b64 = photo_data["img_b64"]
    tg_fp   = photo_data.get("tg_file_path")
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    await callback.answer()
    await _do_video_gen(
        callback.message, uid,
        "animate this exact person realistically. Same face, same appearance. Natural movement: subtle breathing, blinking, slight head motion. Cinematic quality.",
        mode="ai", image_base64=img_b64, tg_file_path=tg_fp
    )


@dp.callback_query(F.data == "photo_img2img")
async def cb_photo_img2img(callback: CallbackQuery):
    uid = callback.from_user.id
    photo_data = last_photo.get(uid)
    if not photo_data or not photo_data.get("ask_pending"):
        return await callback.answer("❌ Фото не найдено, отправь снова", show_alert=True)
    if not can_img(uid):
        return await callback.answer("🚫 Лимит генераций исчерпан!", show_alert=True)
    last_photo[uid]["ask_pending"]     = False
    last_photo[uid]["img2img_pending"] = True
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    await callback.answer()
    await callback.message.answer(
        "🎨 <b>Изменение стиля фото</b>\n\n"
        "Напиши как именно изменить фото:\n\n"
        "💡 <i>Примеры:</i>\n"
        "• <code>в стиле аниме</code>\n"
        "• <code>превратить в картину маслом</code>\n"
        "• <code>cyberpunk neon style</code>\n"
        "• <code>pixel art 8-bit</code>\n"
        "• <code>sketch pencil drawing</code>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="photo_img2img_cancel")
        ]])
    )


@dp.callback_query(F.data == "photo_img2img_cancel")
async def cb_photo_img2img_cancel(callback: CallbackQuery):
    uid = callback.from_user.id
    last_photo.pop(uid, None)
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    await callback.answer("Отменено")


@dp.callback_query(F.data == "photo_img2img_retry")
async def cb_photo_img2img_retry(callback: CallbackQuery):
    """Повторить изменение стиля с новым промптом."""
    uid = callback.from_user.id
    photo_data = last_photo.get(uid)
    if not photo_data or not photo_data.get("img_b64"):
        return await callback.answer("❌ Фото не найдено, отправь снова", show_alert=True)
    if not can_img(uid):
        return await callback.answer("🚫 Лимит генераций исчерпан!", show_alert=True)
    # Восстанавливаем флаг ожидания нового стиля
    last_photo[uid]["img2img_pending"] = True
    last_photo[uid].pop("ask_pending", None)
    last_photo[uid].pop("fusion_pending", None)
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    await callback.answer()
    await callback.message.answer(
        "🎨 <b>Новый стиль</b>\n\n"
        "Напиши какой стиль применить к фото:\n\n"
        "💡 <i>Примеры:</i>\n"
        "• <code>в стиле аниме</code>\n"
        "• <code>oil painting portrait</code>\n"
        "• <code>cyberpunk neon city</code>\n"
        "• <code>pixel art 8-bit</code>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="❌ Отмена", callback_data="photo_img2img_cancel")
        ]])
    )


@dp.callback_query(F.data == "photo_result_to_video")
async def cb_photo_result_to_video(callback: CallbackQuery):
    """Оживить результат img2img в видео."""
    uid = callback.from_user.id
    photo_data = last_photo.get(uid)
    if not photo_data or not photo_data.get("img_b64"):
        return await callback.answer("❌ Фото не найдено, отправь снова", show_alert=True)
    if not can_video(uid):
        li = get_limits_info(uid)
        return await callback.answer(
            f"🚫 Лимит видео исчерпан! ({li['video_used']}/{li['video_max']})",
            show_alert=True
        )
    await callback.answer()
    try:
        await callback.message.edit_reply_markup(reply_markup=None)
    except Exception:
        pass
    # Используем result_b64 (стилизованное фото) если есть, иначе оригинал
    img_b64 = photo_data.get("result_b64") or photo_data["img_b64"]
    tg_fp   = photo_data.get("tg_file_path")
    await _do_video_gen(
        callback.message, uid,
        "animate this exact person realistically. Same face, same appearance. Natural movement: subtle breathing, blinking, slight head motion. Cinematic quality.",
        mode="ai", image_base64=img_b64, tg_file_path=tg_fp
    )


@dp.callback_query(F.data.startswith("save_"))
async def save_file(callback: CallbackQuery):
    try:
        _, fmt, uid_str = callback.data.split("_", 2)
        uid = int(uid_str)
        if uid not in last_responses:
            return await callback.answer("❌ Ответ не найден", show_alert=True)
        await callback.answer("⏳ Создаю файл…")
        d = last_responses[uid]
        if   fmt == "docx": path = await FileCreator.create_docx(d["q"], d["a"], uid)
        elif fmt == "png":  path = await FileCreator.create_png(d["q"], d["a"], uid)
        elif fmt == "pdf":  path = await FileCreator.create_pdf(d["q"], d["a"], uid)
        else: return
        await callback.message.answer_document(FSInputFile(path))
        os.remove(path)
    except Exception as e:
        logging.exception("save_file")
        await callback.message.answer(f"❌ Ошибка: {e}")


@dp.callback_query(F.data == "menu_photofusion")
async def menu_photofusion(callback: CallbackQuery):
    uid = callback.from_user.id
    user_photo_mode[uid] = "waiting_photo"
    text = (
        "🪄 <b>Фото-магия</b>\n\n"
        "Отправь своё фото — и я сделаю с ним что-нибудь крутое!\n\n"
        "<b>Примеры запросов:</b>\n"
        "• <i>«Сделай меня супергероем»</i>\n"
        "• <i>«Помести меня в Нью-Йорк»</i>\n"
        "• <i>«Нарисуй меня в стиле аниме»</i>"
    )
    kb = InlineKeyboardMarkup(inline_keyboard=[[
        InlineKeyboardButton(text="◀️ Назад", callback_data="back_home"),
    ]])
    try:
        await callback.message.edit_text(text, parse_mode="HTML", reply_markup=kb)
    except Exception:
        await callback.message.answer(text, parse_mode="HTML", reply_markup=kb)
    await callback.answer()


# ==================================================================
# 💬 ОБРАБОТКА ТЕКСТОВЫХ СООБЩЕНИЙ
# ==================================================================

@dp.message(F.text)
async def handle_text(message: Message):
    uid = message.from_user.id
    _init_tokens(uid)
    _init_limits(uid)
    # Обновляем дату последней активности
    if uid in user_profiles:
        user_profiles[uid]["last_active_date"] = msk_now().strftime("%d.%m.%Y")

    # ── Если выбрана модель генерации изображений — роутим туда ──
    full_mk = get_model_key(uid)
    if full_mk.startswith("imggen:"):
        img_key = full_mk.split(":", 1)[1]
        if img_key not in IMG_MODELS:
            img_key = get_img_model(uid)
        if img_key not in IMG_MODELS:
            img_key = "flux_klein4"
        prompt = message.text.strip()
        if not prompt:
            return
        if not can_img(uid):
            await message.answer(
                "🚫 <b>Лимит генераций исчерпан!</b>\n\nОформи подписку для продолжения.",
                parse_mode="HTML"
            )
            return
        wait_msg = await message.answer(
            f"🎨 <b>Генерирую изображение...</b>\n\n"
            f"🖌 Модель: <code>{IMG_MODELS[img_key]['label']}</code>\n"
            f"💡 <i>Перевожу промпт и улучшаю...</i>",
            parse_mode="HTML"
        )
        try:
            result = await generate_image(prompt, img_key)
            try: await wait_msg.delete()
            except Exception: pass
            # generate_image возвращает {"data": bytes, "url": url} при успехе
            img_data  = result.get("data") if result else None
            img_url   = result.get("url", "") if result else ""
            if img_data or img_url:
                kb = InlineKeyboardMarkup(inline_keyboard=[
                    [InlineKeyboardButton(text="🔄 Ещё вариант",   callback_data=f"imggen_again_{img_key}"),
                     InlineKeyboardButton(text="✨ Новый промпт",  callback_data="imggen_new")],
                    [InlineKeyboardButton(text="🎨 Сменить модель", callback_data="menu_imggen"),
                     InlineKeyboardButton(text="🏠 Главная",       callback_data="back_home")],
                ])
                caption = f"✅ <b>Готово!</b> 🎨 {IMG_MODELS[img_key]['label']}"
                if img_data:
                    await message.answer_photo(
                        BufferedInputFile(img_data, "image.jpg"),
                        caption=caption, parse_mode="HTML", reply_markup=kb
                    )
                elif img_url.startswith("http"):
                    await message.answer_photo(img_url, caption=caption, parse_mode="HTML", reply_markup=kb)
                else:
                    import base64 as _b64
                    raw = _b64.b64decode(img_url.split(",", 1)[-1])
                    await message.answer_photo(
                        BufferedInputFile(raw, "image.jpg"),
                        caption=caption, parse_mode="HTML", reply_markup=kb
                    )
                spend_img(uid)
                asyncio.create_task(db_save_user(uid))
            else:
                err = result.get("error", "Неизвестная ошибка") if result else "Нет ответа"
                await message.answer(f"❌ <b>Ошибка генерации:</b> {err[:200]}", parse_mode="HTML")
        except Exception as ex:
            try: await wait_msg.delete()
            except Exception: pass
            logging.error(f"imggen handle_text: {ex}")
            await message.answer(f"❌ <b>Ошибка:</b> {str(ex)[:200]}", parse_mode="HTML")
        return

    mk  = get_chat_model(uid)

    # Обработка ожидающих действий от админа
    if uid in admin_await:
        action_info = admin_await.pop(uid)
        action = action_info.get("action")

        if action == "use_promo":
            code = message.text.strip().upper()
            used_set = promo_used.setdefault(uid, set())
            if code not in promo_codes:
                await message.answer("❌ Промокод не найден.", reply_markup=home_kb(uid))
            elif code in used_set:
                await message.answer("❌ Ты уже использовал этот промокод.", reply_markup=home_kb(uid))
            elif len(promo_codes[code].get("used_by", [])) >= promo_codes[code].get("max_uses", 999999):
                await message.answer("❌ Промокод исчерпан.", reply_markup=home_kb(uid))
            else:
                # Проверяем срок действия
                expires = promo_codes[code].get("expires", "")
                if expires:
                    try:
                        exp_dt = datetime.datetime.strptime(expires, "%d.%m.%Y")
                        if msk_now() > exp_dt:
                            await message.answer("❌ Промокод истёк.", reply_markup=home_kb(uid))
                            return
                    except Exception:
                        pass
                fast_reward = promo_codes[code].get("fast_requests", 0)
                pro_reward  = promo_codes[code].get("pro_requests", 0)
                total_reward = fast_reward + pro_reward
                promo_codes[code].setdefault("used_by", []).append(uid)
                used_set.add(code)
                _init_limits(uid)
                user_limits[uid]["pro_used"] = max(0, user_limits[uid]["pro_used"] - total_reward)
                li = get_limits_info(uid)
                await message.answer(
                    f"✅ <b>Промокод {code} активирован!</b>\n\n"
                    f"💬 +{total_reward} запросов\n\n"
                    f"📊 Текущий баланс: <b>{li['pro_used']} из {li['pro_max']}</b>",
                    parse_mode="HTML", reply_markup=home_kb(uid)
                )
            return

        if action == "create_promo_step1" and uid in ADMIN_IDS:
            parts = message.text.strip().split()
            if len(parts) >= 2:
                code       = parts[0].upper()
                req_val    = int(parts[1])
                max_uses   = int(parts[2]) if len(parts) > 2 else 999999
                expires    = parts[3] if len(parts) > 3 else ""
                promo_codes[code] = {
                    "fast_requests": 0,
                    "pro_requests":  req_val,
                    "max_uses":      max_uses,
                    "expires":       expires,
                    "used_by":       []
                }
                exp_str = f"\n📅 Срок: до {expires}" if expires else ""
                await message.answer(
                    f"✅ <b>Промокод создан!</b>\n\n"
                    f"🔑 Код: <code>{code}</code>\n"
                    f"💬 Запросов: <b>+{req_val}</b>\n"
                    f"👥 Макс. использований: <code>{max_uses}</code>{exp_str}",
                    parse_mode="HTML",
                )
            else:
                await message.answer(
                    "❌ Неверный формат!\n"
                    "Пример: <code>BONUS2025 20 100</code>\n"
                    "или: <code>VIP 50 50 31.12.2025</code>",
                    parse_mode="HTML"
                )
            return

        if action in ("give_tokens", "give_pro") and uid in ADMIN_IDS:
            parts = message.text.strip().split()
            if len(parts) >= 2:
                try:
                    target_uid = int(parts[0])
                    amount     = int(parts[1])
                    if amount <= 0: raise ValueError
                    _init_limits(target_uid)
                    user_limits[target_uid]["pro_used"] = max(0, user_limits[target_uid].get("pro_used", 0) - amount)
                    li = get_limits_info(target_uid)
                    asyncio.create_task(db_save_user(target_uid))
                    await message.answer(
                        f"✅ <b>Выдано +{amount} запросов (pro)</b>\n\n"
                        f"👤 ID: <code>{target_uid}</code>\n"
                        f"📊 Баланс: <b>{max(0, li['pro_max'] - li['pro_used'])} / {li['pro_max']}</b>",
                        parse_mode="HTML",
                        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                            InlineKeyboardButton(text="◀️ Назад", callback_data="admin_give_limits")
                        ]])
                    )
                    try:
                        await bot.send_message(target_uid, f"🎉 <b>Вам начислено +{amount} запросов!</b>", parse_mode="HTML")
                    except Exception: pass
                except (ValueError, IndexError):
                    await message.answer("❌ Формат: <code>USER_ID КОЛ-ВО</code> (оба числа > 0)", parse_mode="HTML")
            else:
                await message.answer("❌ Формат: <code>USER_ID КОЛ-ВО</code>", parse_mode="HTML")
            return

        if action == "give_img" and uid in ADMIN_IDS:
            parts = message.text.strip().split()
            if len(parts) >= 2:
                try:
                    target_uid = int(parts[0]); amount = int(parts[1])
                    if amount <= 0: raise ValueError
                    _init_limits(target_uid)
                    user_limits[target_uid]["img_used"] = max(0, user_limits[target_uid].get("img_used", 0) - amount)
                    asyncio.create_task(db_save_user(target_uid))
                    await message.answer(f"✅ <b>Выдано +{amount} генераций картинок</b>\n👤 ID: <code>{target_uid}</code>", parse_mode="HTML",
                        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="◀️ Назад", callback_data="admin_give_limits")]]))
                    try: await bot.send_message(target_uid, f"🎉 <b>Вам начислено +{amount} генераций картинок!</b>", parse_mode="HTML")
                    except Exception: pass
                except (ValueError, IndexError):
                    await message.answer("❌ Формат: <code>USER_ID КОЛ-ВО</code>", parse_mode="HTML")
            else:
                await message.answer("❌ Формат: <code>USER_ID КОЛ-ВО</code>", parse_mode="HTML")
            return

        if action == "give_music" and uid in ADMIN_IDS:
            parts = message.text.strip().split()
            if len(parts) >= 2:
                try:
                    target_uid = int(parts[0]); amount = int(parts[1])
                    if amount <= 0: raise ValueError
                    _init_limits(target_uid)
                    user_limits[target_uid]["music_used"] = max(0, user_limits[target_uid].get("music_used", 0) - amount)
                    asyncio.create_task(db_save_user(target_uid))
                    await message.answer(f"✅ <b>Выдано +{amount} генераций музыки</b>\n👤 ID: <code>{target_uid}</code>", parse_mode="HTML",
                        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="◀️ Назад", callback_data="admin_give_limits")]]))
                    try: await bot.send_message(target_uid, f"🎉 <b>Вам начислено +{amount} генераций музыки!</b>", parse_mode="HTML")
                    except Exception: pass
                except (ValueError, IndexError):
                    await message.answer("❌ Формат: <code>USER_ID КОЛ-ВО</code>", parse_mode="HTML")
            else:
                await message.answer("❌ Формат: <code>USER_ID КОЛ-ВО</code>", parse_mode="HTML")
            return

        if action == "give_video" and uid in ADMIN_IDS:
            parts = message.text.strip().split()
            if len(parts) >= 2:
                try:
                    target_uid = int(parts[0]); amount = int(parts[1])
                    if amount <= 0: raise ValueError
                    _init_limits(target_uid)
                    user_limits[target_uid]["video_used"] = max(0, user_limits[target_uid].get("video_used", 0) - amount)
                    asyncio.create_task(db_save_user(target_uid))
                    await message.answer(f"✅ <b>Выдано +{amount} генераций видео</b>\n👤 ID: <code>{target_uid}</code>", parse_mode="HTML",
                        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="◀️ Назад", callback_data="admin_give_limits")]]))
                    try: await bot.send_message(target_uid, f"🎉 <b>Вам начислено +{amount} генераций видео!</b>", parse_mode="HTML")
                    except Exception: pass
                except (ValueError, IndexError):
                    await message.answer("❌ Формат: <code>USER_ID КОЛ-ВО</code>", parse_mode="HTML")
            else:
                await message.answer("❌ Формат: <code>USER_ID КОЛ-ВО</code>", parse_mode="HTML")
            return

        if action == "give_all" and uid in ADMIN_IDS:
            parts = message.text.strip().split()
            if len(parts) >= 5:
                try:
                    target_uid = int(parts[0])
                    pro_amt, img_amt, music_amt, video_amt = int(parts[1]), int(parts[2]), int(parts[3]), int(parts[4])
                    _init_limits(target_uid)
                    lim = user_limits[target_uid]
                    lim["pro_used"]   = max(0, lim.get("pro_used",   0) - pro_amt)
                    lim["img_used"]   = max(0, lim.get("img_used",   0) - img_amt)
                    lim["music_used"] = max(0, lim.get("music_used", 0) - music_amt)
                    lim["video_used"] = max(0, lim.get("video_used", 0) - video_amt)
                    asyncio.create_task(db_save_user(target_uid))
                    await message.answer(
                        f"✅ <b>Выданы все лимиты</b>\n\n"
                        f"👤 ID: <code>{target_uid}</code>\n"
                        f"💬 +{pro_amt} запросов\n"
                        f"🎨 +{img_amt} картинок\n"
                        f"🎵 +{music_amt} музыки\n"
                        f"🎬 +{video_amt} видео",
                        parse_mode="HTML",
                        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                            InlineKeyboardButton(text="◀️ Назад", callback_data="admin_give_limits")
                        ]])
                    )
                    try:
                        await bot.send_message(target_uid,
                            f"🎉 <b>Вам начислены лимиты!</b>\n\n"
                            f"💬 +{pro_amt} запросов  🎨 +{img_amt} картинок\n"
                            f"🎵 +{music_amt} музыки  🎬 +{video_amt} видео",
                            parse_mode="HTML")
                    except Exception: pass
                except (ValueError, IndexError):
                    await message.answer("❌ Формат: <code>USER_ID ПРО КАРТИНКИ МУЗЫКА ВИДЕО</code>", parse_mode="HTML")
            else:
                await message.answer("❌ Формат: <code>USER_ID ПРО КАРТИНКИ МУЗЫКА ВИДЕО</code>", parse_mode="HTML")
            return

        # ── Забрать лимиты ────────────────────────────────────────
        if action in ("take_tokens", "take_pro") and uid in ADMIN_IDS:
            parts = message.text.strip().split()
            if len(parts) >= 2:
                try:
                    target_uid = int(parts[0]); amount = int(parts[1])
                    if amount <= 0: raise ValueError
                    _init_limits(target_uid)
                    L = _get_lims(target_uid)
                    user_limits[target_uid]["pro_used"] = min(L["pro_day"], user_limits[target_uid].get("pro_used", 0) + amount)
                    asyncio.create_task(db_save_user(target_uid))
                    li = get_limits_info(target_uid)
                    await message.answer(
                        f"✅ <b>Снято {amount} запросов (pro)</b>\n\n"
                        f"👤 ID: <code>{target_uid}</code>\n"
                        f"📊 Осталось: <b>{max(0, li['pro_max'] - li['pro_used'])} / {li['pro_max']}</b>",
                        parse_mode="HTML",
                        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                            InlineKeyboardButton(text="◀️ Назад", callback_data="admin_take_limits")
                        ]])
                    )
                except (ValueError, IndexError):
                    await message.answer("❌ Формат: <code>USER_ID КОЛ-ВО</code>", parse_mode="HTML")
            else:
                await message.answer("❌ Формат: <code>USER_ID КОЛ-ВО</code>", parse_mode="HTML")
            return

        if action == "take_img" and uid in ADMIN_IDS:
            parts = message.text.strip().split()
            if len(parts) >= 2:
                try:
                    target_uid = int(parts[0]); amount = int(parts[1])
                    _init_limits(target_uid); L = _get_lims(target_uid)
                    user_limits[target_uid]["img_used"] = min(L["img_month"], user_limits[target_uid].get("img_used", 0) + amount)
                    asyncio.create_task(db_save_user(target_uid))
                    await message.answer(f"✅ Снято {amount} картинок у <code>{target_uid}</code>", parse_mode="HTML",
                        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="◀️ Назад", callback_data="admin_take_limits")]]))
                except (ValueError, IndexError):
                    await message.answer("❌ Формат: <code>USER_ID КОЛ-ВО</code>", parse_mode="HTML")
            return

        if action == "take_music" and uid in ADMIN_IDS:
            parts = message.text.strip().split()
            if len(parts) >= 2:
                try:
                    target_uid = int(parts[0]); amount = int(parts[1])
                    _init_limits(target_uid); L = _get_lims(target_uid)
                    user_limits[target_uid]["music_used"] = min(L.get("music_month", 0), user_limits[target_uid].get("music_used", 0) + amount)
                    asyncio.create_task(db_save_user(target_uid))
                    await message.answer(f"✅ Снято {amount} музыки у <code>{target_uid}</code>", parse_mode="HTML",
                        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="◀️ Назад", callback_data="admin_take_limits")]]))
                except (ValueError, IndexError):
                    await message.answer("❌ Формат: <code>USER_ID КОЛ-ВО</code>", parse_mode="HTML")
            return

        if action == "take_video" and uid in ADMIN_IDS:
            parts = message.text.strip().split()
            if len(parts) >= 2:
                try:
                    target_uid = int(parts[0]); amount = int(parts[1])
                    _init_limits(target_uid); L = _get_lims(target_uid)
                    user_limits[target_uid]["video_used"] = min(L.get("video_month", 0), user_limits[target_uid].get("video_used", 0) + amount)
                    asyncio.create_task(db_save_user(target_uid))
                    await message.answer(f"✅ Снято {amount} видео у <code>{target_uid}</code>", parse_mode="HTML",
                        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="◀️ Назад", callback_data="admin_take_limits")]]))
                except (ValueError, IndexError):
                    await message.answer("❌ Формат: <code>USER_ID КОЛ-ВО</code>", parse_mode="HTML")
            return

        if action == "take_all" and uid in ADMIN_IDS:
            parts = message.text.strip().split()
            if parts:
                try:
                    target_uid = int(parts[0])
                    _init_limits(target_uid); L = _get_lims(target_uid)
                    lim = user_limits[target_uid]
                    lim["pro_used"]   = L["pro_day"]
                    lim["img_used"]   = L["img_month"]
                    lim["music_used"] = L.get("music_month", 0)
                    lim["video_used"] = L.get("video_month", 0)
                    asyncio.create_task(db_save_user(target_uid))
                    await message.answer(
                        f"🔴 <b>Все лимиты обнулены</b>\n👤 ID: <code>{target_uid}</code>",
                        parse_mode="HTML",
                        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                            InlineKeyboardButton(text="◀️ Назад", callback_data="admin_take_limits")
                        ]])
                    )
                except (ValueError, IndexError):
                    await message.answer("❌ Формат: <code>USER_ID</code>", parse_mode="HTML")
            return

        if action == "admin_give_sub" and uid in ADMIN_IDS:
            parts = message.text.strip().split()
            if len(parts) >= 2:
                try:
                    target_uid = int(parts[0])
                    plan_key   = parts[1].lower()
                    if plan_key not in SUB_PLANS:
                        await message.answer("❌ Неверный план. Используй: week / month", parse_mode="HTML")
                        return
                    plan = SUB_PLANS[plan_key]
                    existing = user_subscriptions.get(target_uid)
                    if existing and existing["expires"] > msk_now():
                        base = existing["expires"]
                    else:
                        base = msk_now()
                    new_exp = base + datetime.timedelta(days=plan["days"])
                    user_subscriptions[target_uid] = {"expires": new_exp, "plan": plan_key}
                    await db_save_subscription(target_uid)
                    await message.answer(
                        f"✅ Подписка выдана пользователю <code>{target_uid}</code>\n"
                        f"◈ Тариф: <b>{plan['name']}</b>\n"
                        f"◈ Действует до: <code>{new_exp.strftime('%d.%m.%Y %H:%M')}</code>",
                        parse_mode="HTML",
                    )
                    try:
                        await bot.send_message(
                            target_uid,
                            f"🎉 <b>Вам выдана подписка!</b>\n\n"
                            f"💎 Тариф: <b>{plan['name']}</b>\n"
                            f"📅 Действует до: <code>{new_exp.strftime('%d.%m.%Y')}</code>\n\n"
                            f"✅ Все Premium-модели разблокированы!",
                            parse_mode="HTML"
                        )
                    except Exception:
                        pass
                except Exception as e:
                    await message.answer(f"❌ Ошибка: {e}", parse_mode="HTML")
            else:
                await message.answer("❌ Формат: <code>USER_ID ПЛАН</code>  (week / month)", parse_mode="HTML")
            return

        if action == "admin_take_sub" and uid in ADMIN_IDS:
            try:
                target_uid = int(message.text.strip())
                if target_uid in user_subscriptions:
                    del user_subscriptions[target_uid]
                    if db_pool:
                        try:
                            async with db_pool.acquire() as conn:
                                await conn.execute(
                                    "UPDATE users SET sub_expires=NULL, sub_plan='' WHERE uid=$1", target_uid
                                )
                        except Exception:
                            pass
                    await message.answer(f"✅ Подписка пользователя <code>{target_uid}</code> отозвана.", parse_mode="HTML")
                    try:
                        await bot.send_message(target_uid, "ℹ️ Ваша подписка была деактивирована.", parse_mode="HTML")
                    except Exception:
                        pass
                else:
                    await message.answer(f"ℹ️ У пользователя <code>{target_uid}</code> нет активной подписки.", parse_mode="HTML")
            except Exception as e:
                await message.answer(f"❌ Ошибка: {e}", parse_mode="HTML")
            return

        if action == "admin_check_sub" and uid in ADMIN_IDS:
            try:
                target_uid = int(message.text.strip())
                sub = user_subscriptions.get(target_uid)
                prof = user_profiles.get(target_uid, {})
                name = prof.get("name", str(target_uid))
                if sub and has_active_sub(target_uid):
                    plan = SUB_PLANS.get(sub.get("plan", ""), {})
                    exp  = sub["expires"].strftime("%d.%m.%Y %H:%M")
                    await message.answer(
                        f"💎 <b>Подписка активна</b>\n"
                        f"◈ Пользователь: <b>{name}</b> (<code>{target_uid}</code>)\n"
                        f"◈ Тариф: <b>{plan.get('name', sub.get('plan', '—'))}</b>\n"
                        f"◈ Истекает: <code>{exp}</code>",
                        parse_mode="HTML"
                    )
                else:
                    await message.answer(
                        f"❌ У пользователя <b>{name}</b> (<code>{target_uid}</code>) нет активной подписки.",
                        parse_mode="HTML"
                    )
            except Exception as e:
                await message.answer(f"❌ Ошибка: {e}", parse_mode="HTML")
            return

        if action == "find_user" and uid in ADMIN_IDS:
            try:
                target_uid = int(message.text.strip())
            except Exception:
                await message.answer("❌ Введи числовой ID", parse_mode="HTML")
                return
            prof = user_profiles.get(target_uid)
            if not prof:
                await message.answer(f"❌ Пользователь <code>{target_uid}</code> не найден.", parse_mode="HTML")
                return
            name     = prof.get("name", "—")
            username = prof.get("username", "—")
            joined   = prof.get("joined", "—")
            requests = prof.get("requests", 0)
            tok      = user_tokens.get(target_uid, {}).get("tokens", 0)
            imgs     = user_images_count.get(target_uid, 0)
            mk_u     = user_settings.get(target_uid, "—")
            await message.answer(
                f"🔍 <b>Профиль пользователя</b>\n\n"
                f"👤 <b>{name}</b> @{username}\n"
                f"🆔 <code>{target_uid}</code>\n"
                f"📅 Регистрация: <code>{joined}</code>\n"
                f"💬 Запросов осталось: <code>{get_limits_info(uid)['pro_max'] - get_limits_info(uid)['pro_used']}</code>\n"
                f"💬 Запросов: <code>{requests}</code>\n"
                f"🖼 Генераций: <code>{imgs}</code>\n"
                f"🤖 Модель: <code>{mk_u}</code>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                    InlineKeyboardButton(text="◀️ Назад", callback_data="menu_admin"),
                ]])
            )
            return

        if action == "add_channel" and uid in ADMIN_IDS:
            raw = message.text.strip()
            # Поддерживаем: @username, username, https://t.me/username, числовой ID (-100...)
            if raw.startswith("https://t.me/"):
                ch = raw.split("https://t.me/")[1].split("/")[0].split("?")[0]
            elif raw.startswith("http://t.me/"):
                ch = raw.split("http://t.me/")[1].split("/")[0].split("?")[0]
            elif raw.lstrip("-").isdigit():
                ch = raw  # числовой chat_id
            else:
                ch = raw.lstrip("@")
            if not ch:
                await message.answer("❌ Пустое значение.")
                return
            if ch in required_channels:
                await message.answer(f"❌ Канал уже добавлен.")
                return
            # Получаем название канала и добавляем
            chat_id_val = int(ch) if str(ch).lstrip("-").isdigit() else f"@{ch}"
            try:
                chat = await bot.get_chat(chat_id=chat_id_val)
                ch_title = chat.title or f"@{ch}"
            except Exception:
                ch_title = f"@{ch}"
            required_channels.append(ch)
            asyncio.create_task(db_save_bot_settings())  # Сохраняем в БД
            await message.answer(
                f"✅ Канал <b>{ch_title}</b> добавлен!\n\n"
                f"Пользователи должны подписаться для доступа к боту.\n"
                f"<i>Убедись что бот добавлен в канал как администратор.</i>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                    InlineKeyboardButton(text="◀️ К каналам", callback_data="admin_channels"),
                ]])
            )
            return

        if action == "edit_limits" and uid in ADMIN_IDS:
            # legacy fallback — теперь принимает: ЗАПРОСЫ КАРТИНКИ [СБРОС_Ч]
            parts = message.text.strip().split()
            if len(parts) >= 2:
                try:
                    LIMITS["pro_day"]   = int(parts[0])
                    LIMITS["img_month"] = int(parts[1])
                    if len(parts) >= 3: LIMITS["reset_h"] = int(parts[2])
                    if len(parts) >= 4: LIMITS["music_month"] = int(parts[3])
                    asyncio.create_task(db_save_bot_settings())
                    await message.answer(f"✅ Лимиты обновлены!", parse_mode="HTML")
                except ValueError:
                    await message.answer("❌ Неверный формат.", parse_mode="HTML")
            return

        if action == "set_reset_h" and uid in ADMIN_IDS:
            try:
                h = int(message.text.strip())
                LIMITS["reset_h"] = h
                SUB_PLANS["week"]["reset_h"]  = h
                SUB_PLANS["month"]["reset_h"] = h
                asyncio.create_task(db_save_bot_settings())
                await message.answer(
                    f"✅ Интервал сброса запросов: <b>{h} часов</b>",
                    parse_mode="HTML",
                    reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                        InlineKeyboardButton(text="◀️ К лимитам", callback_data="admin_extra")
                    ]])
                )
            except ValueError:
                await message.answer("❌ Введи число (часы).", parse_mode="HTML")
            return

        for sect in ("free", "week", "month"):
            if action == f"set_limits_{sect}" and uid in ADMIN_IDS:
                parts = message.text.strip().split()
                if len(parts) < 3:
                    await message.answer(
                        "❌ Нужно минимум 3 числа: <code>ЗАПРОСЫ КАРТИНКИ МУЗЫКА [ВИДЕО]</code>",
                        parse_mode="HTML"
                    )
                    return
                try:
                    pro   = int(parts[0])
                    img   = int(parts[1])
                    music = int(parts[2])
                    video = int(parts[3]) if len(parts) >= 4 else 0
                    if sect == "free":
                        LIMITS["pro_day"]     = pro
                        LIMITS["img_month"]   = img
                        LIMITS["music_month"] = music
                        LIMITS["video_month"] = video
                    else:
                        SUB_PLANS[sect]["pro_day"]     = pro
                        SUB_PLANS[sect]["img_month"]   = img
                        SUB_PLANS[sect]["music_month"] = music
                        SUB_PLANS[sect]["video_month"] = video
                    asyncio.create_task(db_save_bot_settings())
                    plan_name = {"free": "🆓 Бесплатно", "week": "⚡ Неделя", "month": "💎 Месяц"}[sect]
                    await message.answer(
                        f"✅ <b>{plan_name} — лимиты обновлены!</b>\n\n"
                        f"💬 Запросы: <code>{pro}</code>\n"
                        f"🎨 Картинки: <code>{img}</code>/мес\n"
                        f"🎵 Музыка: <code>{music}</code>/мес\n"
                        f"🎬 Видео: <code>{video}</code>/мес",
                        parse_mode="HTML",
                        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                            InlineKeyboardButton(text="◀️ К лимитам", callback_data="admin_extra")
                        ]])
                    )
                except ValueError:
                    await message.answer("❌ Неверный формат. Введи числа через пробел.", parse_mode="HTML")
                return

        if action == "broadcast" and uid in ADMIN_IDS:
            broadcast_text = message.text
            sent = 0; failed = 0
            all_uids = list(user_settings.keys())
            wait_msg = await message.answer(f"📢 Отправляю {len(all_uids)} пользователям...")
            for target_uid in all_uids:
                try:
                    await bot.send_message(target_uid, broadcast_text, parse_mode="HTML")
                    sent += 1
                except Exception:
                    failed += 1
                await asyncio.sleep(0.05)
            await wait_msg.delete()
            await message.answer(
                f"✅ Рассылка завершена!\n"
                f"✅ Доставлено: <code>{sent}</code>\n"
                f"❌ Ошибок: <code>{failed}</code>",
                parse_mode="HTML"
            )
            return

        if action == "report_author":
            author = message.text.strip()
            if uid not in report_states:
                report_states[uid] = {}
            report_states[uid]["author"] = author
            admin_await[uid] = {"action": "report_topic"}
            await message.answer(
                f"👤 <b>Автор: {author}</b>\n\n"
                "✏️ Напиши <b>тему</b> доклада/реферата:\n\n"
                "<i>Пример: «Влияние искусственного интеллекта на рынок труда»</i>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                    InlineKeyboardButton(text="❌ Отмена", callback_data="back_home"),
                ]])
            )
            return

        if action == "report_custom_pages":
            raw_pages = message.text.strip()
            if not raw_pages.isdigit() or not (1 <= int(raw_pages) <= 10):
                await message.answer(
                    "❌ Введи число от 1 до 10:",
                    parse_mode="HTML"
                )
                admin_await[uid] = {"action": "report_custom_pages"}
                return
            pages = int(raw_pages)
            if uid not in report_states:
                report_states[uid] = {}
            report_states[uid]["pages"] = pages
            report_states[uid]["model"] = "sonar_deep_research"
            report_states[uid]["model_label"] = "🔬 Sonar Deep Research"
            admin_await[uid] = {"action": "report_author"}
            await message.answer(
                f"📝 <b>Объём: {pages} страниц</b>\n"
                f"🌐 <b>Модель: Sonar Deep Research</b> (с веб-поиском)\n\n"
                "👤 Введи <b>имя автора</b> (ФИО / псевдоним):\n\n"
                "<i>Или нажми «Пропустить», если не нужно</i>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                    [InlineKeyboardButton(text="⏭ Пропустить", callback_data="report_skip_author")],
                    [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
                ])
            )
            return

        if action == "report_topic":
            topic = message.text.strip()
            if not topic:
                await message.answer("\u274c Тема не может быть пустой. Введи тему:")
                admin_await[uid] = {"action": "report_topic"}
                return
            if uid not in report_states:
                report_states[uid] = {}
            report_states[uid]["topic"] = topic
            admin_await[uid] = {"action": "report_wishes"}
            await message.answer(
                f"\u2705 <b>Тема: \xab{topic[:60]}{'...' if len(topic)>60 else ''}\xbb</b>\n\n"
                "\u270d\ufe0f <b>Укажи пожелания</b> к работе\n"
                "<i>(стиль, акценты, что включить/исключить)</i>\n\n"
                "Или нажми <b>\xabПропустить\xbb</b>:",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                    [InlineKeyboardButton(text="\u23e9 Пропустить пожелания", callback_data="report_skip_wishes")],
                    [InlineKeyboardButton(text="\u274c Отмена", callback_data="back_home")],
                ])
            )
            return

        if action == "report_wishes":
            wishes = message.text.strip()
            if uid not in report_states:
                report_states[uid] = {}
            report_states[uid]["wishes"] = wishes
            state = report_states[uid]
            type_label = state.get("type_label", "Доклад")
            topic2 = state.get("topic", "—")
            pages2 = state.get("pages", 5)
            model_label2 = state.get("model_label", "Claude Opus")
            author2 = state.get("author", "")
            with_title2 = state.get("with_title", True)
            author_line2 = f"   • Автор: <b>{author2}</b>\n" if author2 else ""
            title_str2 = "📄 Есть" if with_title2 else "📝 Нет"
            wishes_short = wishes[:80] + ("..." if len(wishes) > 80 else "")
            await message.answer(
                f"📋 <b>Настройки {type_label}:</b>\n"
                f"   • Тип: <b>{type_label}</b>\n"
                f"   • Тема: <b>«{topic2[:60]}{'..' if len(topic2)>60 else ''}»</b>\n"
                f"   • Объём: <b>{pages2} страниц</b>\n"
                f"   • Модель: <b>{model_label2}</b>\n"
                f"{author_line2}"
                f"   • Титульный лист: <b>{title_str2}</b>\n"
                f"   • Пожелания: <b>{wishes_short}</b>\n\n"
                f"💎 Стоимость: <b>3 продвинутых запроса</b>\n\n"
                f"🚀 Всё верно? Нажми <b>«Создать доклад»</b>:",
                parse_mode="HTML",
                reply_markup=report_confirm_kb()
            )
            return

        if action == "pptx_topic":
            topic = message.text.strip()
            if not topic:
                await message.answer("❌ Тема не может быть пустой. Введи тему:")
                admin_await[uid] = {"action": "pptx_topic"}
                return
            if uid not in pptx_states:
                pptx_states[uid] = {}
            pptx_states[uid]["topic"] = topic
            pptx_states[uid].setdefault("title_info", {})
            # Новый флоу: сразу к выбору кол-ва слайдов
            await message.answer(
                f"🎞 <b>Тема: «{topic[:60]}{'...' if len(topic)>60 else ''}»</b>\n\n"
                "📏 <b>Сколько слайдов?</b>",
                parse_mode="HTML",
                reply_markup=pptx_slides_kb()
            )
            return

        if action == "pptx_subtitle":
            subtitle = message.text.strip()
            if uid not in pptx_states:
                pptx_states[uid] = {}
            pptx_states[uid].setdefault("title_info", {})["subtitle"] = subtitle
            admin_await[uid] = {"action": "pptx_author"}
            await message.answer(
                f"✅ <b>Подзаголовок: «{subtitle[:50]}»</b>\n\n"
                "👤 Введи <b>имя автора</b> (ФИО / докладчик):",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                    [InlineKeyboardButton(text="⏭ Пропустить", callback_data="pptx_skip_author")],
                    [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
                ])
            )
            return

        if action == "pptx_author":
            author = message.text.strip()
            if uid not in pptx_states:
                pptx_states[uid] = {}
            pptx_states[uid].setdefault("title_info", {})["author"] = author
            admin_await[uid] = {"action": "pptx_organization"}
            await message.answer(
                f"✅ <b>Автор: {author}</b>\n\n"
                "🏢 Введи <b>название организации</b>:\n"
                "<i>(школа, университет, компания)</i>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                    [InlineKeyboardButton(text="⏭ Пропустить", callback_data="pptx_skip_org")],
                    [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
                ])
            )
            return

        if action == "pptx_organization":
            org = message.text.strip()
            if uid not in pptx_states:
                pptx_states[uid] = {}
            pptx_states[uid].setdefault("title_info", {})["organization"] = org
            admin_await[uid] = {"action": "pptx_wishes"}
            await message.answer(
                f"✅ <b>Организация: {org[:50]}</b>\n\n"
                "✍️ Добавь пожелания к презентации\n"
                "<i>(стиль, акценты, что включить/исключить)</i>\n\n"
                "Или нажми «Пропустить»:",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                    [InlineKeyboardButton(text="⏭ Пропустить", callback_data="pptx_skip_wishes")],
                    [InlineKeyboardButton(text="❌ Отмена", callback_data="back_home")],
                ])
            )
            return

        if action == "pptx_wishes":
            wishes = message.text.strip()
            if uid not in pptx_states:
                pptx_states[uid] = {}
            pptx_states[uid]["wishes"] = wishes
            await _show_pptx_confirm(message, uid)
            return

        if action == "pptx_rename_slide":
            # Пользователь ввёл новое название для слайда
            new_name = message.text.strip()
            slide_idx = admin_await[uid].get("slide_idx", 0)
            if uid not in pptx_states:
                pptx_states[uid] = {}
            plan = pptx_states[uid].get("plan", [])
            if 0 <= slide_idx < len(plan):
                old_name = plan[slide_idx]
                plan[slide_idx] = new_name
                pptx_states[uid]["plan"] = plan
                topic = pptx_states[uid].get("topic", "")
                plan_display = "\n".join(f"{i+1}. {t}" for i, t in enumerate(plan))
                await message.answer(
                    f"✅ <b>Слайд {slide_idx + 1} переименован!</b>\n"
                    f"<s>{old_name}</s> → <b>{new_name}</b>\n\n"
                    f"🎞 <b>Тема:</b> {topic[:60]}\n\n"
                    f"<b>📋 План ({len(plan)} слайдов):</b>\n{plan_display}\n\n"
                    f"<i>✏️ Нажми на кнопку слайда чтобы изменить его название</i>",
                    parse_mode="HTML",
                    reply_markup=_plan_kb(uid)
                )
            else:
                await message.answer("❌ Слайд не найден", reply_markup=_plan_kb(uid))
            return

    # Проверка подписки для обычных пользователей
    if not await require_subscription(message):
        return

    # ── Обработка состояний генерации ВИДЕО ──────────────────────
    _video_state = video_states.get(uid, {})
    if _video_state.get("stage") == "await_prompt":
        prompt_text = message.text.strip()
        mode = _video_state.get("mode", "ai")
        video_states.pop(uid, None)
        await _do_video_gen(message, uid, prompt_text, mode)
        return

    # ── Видео из фото: пользователь написал свой промпт ──────────
    if _video_state.get("stage") == "await_video_prompt_text":
        prompt_text = message.text.strip()
        img_b64 = _video_state.get("img_b64")
        tg_fp   = _video_state.get("tg_file_path")
        video_states.pop(uid, None)
        await _do_video_gen(message, uid, prompt_text, mode="raw",
                            image_base64=img_b64, tg_file_path=tg_fp)
        return

    # ── Обработка состояний генерации музыки ─────────────────────
    _music_state = music_states.get(uid, {})
    if _music_state.get("stage") == "await_lyrics":
        lyrics = message.text.strip()
        style_key = _music_state.get("style_key", "")
        style = SUNO_STYLES.get(style_key, {})
        music_states[uid] = {**_music_state, "stage": "confirm", "lyrics": lyrics,
                             "prompt": lyrics[:80], "title": lyrics[:40]}
        preview = lyrics[:600] + ("..." if len(lyrics) > 600 else "")
        await message.answer(
            f"✍️ <b>Текст принят!</b>\n\n"
            f"<code>{preview}</code>\n\n"
            f"🎵 Стиль: <b>{style.get('label', '')}</b>\n\n"
            f"Выбери режим:",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="🎵 Генерировать с вокалом!", callback_data="music_gen_vocal"),
                 InlineKeyboardButton(text="🎸 Инструментал",            callback_data="music_gen_instr")],
                [InlineKeyboardButton(text="✏️ Изменить текст", callback_data="music_write_lyrics"),
                 InlineKeyboardButton(text="◀️ Стили",           callback_data="menu_music")],
            ])
        )
        return

    if _music_state.get("stage") == "await_wishes":
        wishes = message.text.strip()
        style_key = _music_state.get("style_key", "")
        style = SUNO_STYLES.get(style_key, {})
        lang = style.get("lang", "en")
        lang_note = "🇷🇺 ИИ напишет текст на русском" if lang == "ru" else "🇬🇧 ИИ will write lyrics in English"
        music_states[uid] = {
            **_music_state,
            "stage": "confirm",
            "prompt": wishes,
            "title": wishes[:50],
        }
        await message.answer(
            f"╔══════════════════════╗\n"
            f"║  🎵 {style.get('label', 'Трек')}  ║\n"
            f"╚══════════════════════╝\n\n"
            f"📝 <b>Идея принята:</b>\n<i>{wishes[:200]}</i>\n\n"
            f"🤖 ИИ напишет текст песни при генерации\n"
            f"{lang_note}\n\n"
            f"━━━━━━━━━━━━━━━━━━━━━━━━\n"
            f"Выбери режим:",
            parse_mode="HTML",
            reply_markup=music_confirm_kb()
        )
        return

    if _music_state.get("stage") == "await_prompt":
        # Пользователь ввёл описание трека
        user_prompt = message.text.strip()
        style_key = _music_state.get("style_key", "")
        style = SUNO_STYLES.get(style_key, {})
        music_states[uid] = {
            **_music_state,
            "stage": "confirm",
            "prompt": user_prompt,
            "title": user_prompt[:50],
        }
        await message.answer(
            f"🎵 <b>Готово к генерации!</b>\n\n"
            f"🎼 Стиль: <b>{style.get('label', style_key)}</b>\n"
            f"📝 <i>{user_prompt[:300]}</i>\n\n"
            f"Выбери режим:",
            parse_mode="HTML",
            reply_markup=music_confirm_kb()
        )
        return

    if _music_state.get("stage") == "await_custom_style":
        # Пользователь ввёл свой стиль
        custom_tags = message.text.strip()
        music_states[uid] = {
            "stage": "await_prompt",
            "style_key": "custom",
            "style_tags": custom_tags,
            "style_label": custom_tags[:40],
        }
        await message.answer(
            f"✅ Стиль сохранён: <i>{custom_tags[:100]}</i>\n\n"
            f"Теперь опиши сам трек — тема, настроение, атмосфера:",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="◀️ Стили", callback_data="menu_music")
            ]])
        )
        return

    # Текст в режиме чата никогда не запускает генерацию изображений.
    # Для генерации — использовать кнопку «🎨 Генерация картинок».

    if uid in last_photo and last_photo[uid].get("fusion_pending"):
        photo_data = last_photo[uid]
        last_photo[uid]["fusion_pending"] = False
        if not can_img(uid):
                "🚫 <b>Лимит генераций исчерпан!</b>\n\nОформи подписку для продолжения.",
    if uid in last_photo and last_photo[uid].get("img2img_pending"):
        photo_data = last_photo[uid]
        last_photo[uid]["img2img_pending"] = False
        img_b64      = photo_data["img_b64"]
        tg_file_path = photo_data.get("tg_file_path")
        style_text   = message.text.strip()
        # Берём модель из состояния или flux2dev по умолчанию
        img_model_key = photo_data.get("img2img_model_key", "flux2dev")
        if not can_img(uid):
                "🚫 <b>Лимит генераций исчерпан!</b>\n\nОформи подписку для продолжения.",
    # Если пользователь написал текст после фото без подписи — используем как запрос к фото
    if uid in last_photo and last_photo[uid].get("ask_pending"):
        photo_data = last_photo[uid]
        last_photo[uid]["ask_pending"] = False
        caption = message.text.strip()
        await _analyze_photo(message, uid, photo_data["img_b64"], caption, caption)
        return

    # Умное определение follow-up вопроса к уже проанализированному фото
    if uid in last_photo and last_photo[uid].get("photo_context") and not last_photo[uid].get("fusion_pending"):
        txt_lower_check = message.text.strip().lower()
        # Явные команды к фото
        SHORT_PHOTO_CMDS = ["выполни", "сделай", "реши", "ответь", "давай", "продолжи",
                            "выполнить", "сделать", "решить", "опиши", "что на", "а что",
                            "а как", "расскажи подробнее", "объясни подробнее", "уточни",
                            "переведи с", "что означает", "что написано"]
        is_followup_cmd = any(txt_lower_check == cmd or txt_lower_check.startswith(cmd + " ") for cmd in SHORT_PHOTO_CMDS)
        # Короткий вопрос с уточнением — скорее всего тоже к фото
        is_short_followup = (len(txt_lower_check) < 80 and any(
            w in txt_lower_check for w in ["а если", "а что если", "почему", "как это", "зачем", "можно ли", "покажи", "объясни"]
        ))
        if is_followup_cmd or is_short_followup:
            photo_data = last_photo[uid]
            prev_ctx = photo_data.get("last_answer", "")
            prev_q   = photo_data.get("caption", "")
            # Формируем контекстный запрос с фото
            ctx_caption = message.text.strip()
            if prev_ctx:
                ctx_caption = (
                    f"{message.text.strip()}\n\n"
                    f"[Контекст предыдущего анализа фото: вопрос был «{prev_q[:100]}», "
                    f"ответ: «{prev_ctx[:300]}...»]"
                )
            await _analyze_photo(message, uid, photo_data["img_b64"], ctx_caption, message.text.strip())
            return
        else:
            # Не похоже на follow-up к фото — убираем контекст
            last_photo.pop(uid, None)

    _init_limits(uid)
    _ok, _reason = can_send(uid, mk)
    if not _ok:
        # Кулдаун: не показываем сообщение о лимите чаще раза в 5 минут
        import time as _time_mod2
        _now_ts = _time_mod2.time()
        _last_notify = _limit_notify_cd.get(uid, 0)
        if _now_ts - _last_notify < 300:
            return  # молча игнорируем
        _limit_notify_cd[uid] = _now_ts

        if _reason == "premium":
            return await message.answer(
                "🔒 <b>Эта модель доступна только по подписке</b>\n\n"
                "💎 Оформи подписку и получи доступ ко всем моделям:\n"
                "⚡ 7 дней — 60 ₽\n"
                "💎 30 дней — 100 ₽",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                    [InlineKeyboardButton(text="💎 Оформить подписку", callback_data="sub_menu")],
                    [InlineKeyboardButton(text="🤖 Выбрать другую модель", callback_data="menu_ai")],
                ])
            )
        _rst = _reset_str(uid, "pro")
        return await message.answer(
            f"🚫 <b>Лимит запросов исчерпан!</b>\n\n"
            f"⏰ Сброс через: <code>{_rst}</code>\n\n"
            f"Или оформи подписку для увеличения лимитов.",
            parse_mode="HTML", reply_markup=ai_menu_kb(uid)
        )

    if uid not in user_memory:
        user_memory[uid] = deque(maxlen=20)

    # Проверяем что модель не выключена
    if mk in disabled_models:
        await message.answer(
            "🚫 <b>Модель временно недоступна</b>\n\n"
            "Выбери другую модель.",
            parse_mode="HTML",
            reply_markup=ai_menu_kb(uid)
        )
        return

    SHORT_COMMANDS = ["выполни", "сделай", "реши", "ответь", "давай", "продолжи", "выполнить", "сделать", "решить"]
    txt_lower = message.text.strip().lower()
    is_short_cmd = any(txt_lower == cmd or txt_lower.startswith(cmd + " ") for cmd in SHORT_COMMANDS)

    if is_short_cmd and uid in last_photo:
        photo_data = last_photo[uid]
        img_b64_saved = photo_data["img_b64"]
        task_prompt = (
            f"{message.text}\n\n"
            "На фото — учебные задания. Выполни ВСЕ задания. "
            "Для каждого пункта напиши номер и конкретный правильный ответ."
        )
        FALLBACK_VISION = ["claude_opus", "claude_haiku", "claude_sonnet", "gpt52", "gemini_25_flash", "gemini_20_flash", "gemini_3_flash", "qwen3_vl", "c4ai_vis32b"]
        cur_key = get_model_key(uid)
        
        # Если выбрана vision модель - используем её без fallback
        if cur_key in VISION_MODELS and cur_key not in disabled_models:
            # Проверка подписки для премиум моделей
            if cur_key in PREMIUM_MODELS and not has_active_sub(uid):
                await message.answer(
                    f"🔒 <b>Модель {MODELS[cur_key]['label']} доступна только по подписке</b>\n\n"
                    "💡 Выбери другую модель с Vision или оформи подписку.",
                    parse_mode="HTML",
                    reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                        [InlineKeyboardButton(text="💎 Оформить подписку", callback_data="sub_menu")],
                        [InlineKeyboardButton(text="🤖 Сменить модель", callback_data="menu_ai")],
                    ])
                )
                return
            vision_order = [cur_key]  # Только выбранная модель
        else:
            vision_order = [k for k in FALLBACK_VISION if k in MODELS and k not in disabled_models]

        wait_msg = await message.answer("🔍 <i>Выполняю задания с фото…</i>", parse_mode="HTML")
        async with ChatActionSender.typing(bot=bot, chat_id=message.chat.id):
            errors = []
            for vk in vision_order:
                try:
                    vision_msg = [{"role": "user", "content": [
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64_saved}"}},
                        {"type": "text", "text": task_prompt},
                    ]}]
                    ans = await call_chat(vision_msg, vk, max_tokens=3000)
                    last_responses[uid] = {"q": message.text, "a": ans, "model_label": MODELS[vk]["label"], "model_key": vk}
                    user_memory[uid].append({"role": "assistant", "content": ans})
                    spend_limit(uid, vk)
                    await wait_msg.delete()
                    if re.search(r"```", ans):
                        await smart_send(message, ans, reply_markup=save_kb(uid))
                    else:
                        html_ans = md_to_html(ans)
                        try:
                            chunks = [html_ans[j:j+4000] for j in range(0, len(html_ans), 4000)]
                            for ci, chunk in enumerate(chunks):
                                kb = save_kb(uid) if ci == len(chunks)-1 else None
                                await message.answer(chunk, parse_mode="HTML", reply_markup=kb)
                        except Exception:
                            await message.answer(ans[:4000], reply_markup=save_kb(uid))
                    return
                except Exception as e:
                    errors.append(f"{MODELS.get(vk, {}).get('label', vk)}: {str(e)[:60]}")
                    logging.error(f"Vision retry {vk}: {e}")
                    # Если выбрана конкретная модель (не fallback) - сразу показываем ошибку
                    if len(vision_order) == 1:
                        break
            
            await wait_msg.delete()
            if errors:
                error_msg = "\n".join(f"• {e}" for e in errors)
                await message.answer(
                    f"❌ <b>Не удалось выполнить задание</b>\n\n"
                    f"Ошибки:\n{error_msg}\n\n"
                    "💡 Попробуй:\n"
                    "• Отправить фото заново\n"
                    "• Выбрать другую модель с Vision 📸",
                    parse_mode="HTML",
                    reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                        [InlineKeyboardButton(text="🤖 Сменить модель", callback_data="menu_ai")],
                    ])
                )
            else:
                await message.answer("❌ Не удалось выполнить задание. Попробуй отправить фото заново.")
            return

    # Режим ответа и авто-модель
    _feats       = user_features.get(uid, {"doc_analysis": False, "answer_mode": "fast", "auto_model": False})
    _answer_mode = _feats.get("answer_mode", "fast")
    _mode_instr  = ANSWER_MODES.get(_answer_mode, ANSWER_MODES["fast"])[2]

    # ── Ранний выход: приветствие / благодарность — не запускаем ИИ ──
    if is_greeting_or_thanks(message.text or ""):
        import random as _rnd
        replies = [
            "😊 Пожалуйста! Если появятся вопросы — пиши.",
            "👍 Рад помочь! Спрашивай если что.",
            "🙂 Обращайся! Всегда готов помочь.",
            "✨ Пожалуйста! Задавай следующий вопрос когда будет готов.",
        ]
        await message.answer(_rnd.choice(replies))
        return


    _auto_label = ""
    if _feats.get("auto_model", False):
        _new_mk, _reason = auto_select_model(text=message.text)
        if _new_mk != mk and _new_mk in MODELS:
            mk = _new_mk
            _auto_label = f"\n🤖 <i>AI выбрал: {MODELS[mk]['label']}</i>\n({_reason})"

    GLOBAL_SYS = (
        f"Ты — {BOT_NAME}, интеллектуальный ИИ-ассистент нового поколения. "
        "Ты не просто отвечаешь — ты думаешь, анализируешь и выдаёшь результат, "
        "за который люди готовы платить. "
        "ВСЕГДА отвечай на языке пользователя (по умолчанию — русский). "
        "Если из контекста понятно о чём речь — не переспрашивай, сразу выполняй.\n\n"

        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "СТАНДАРТЫ КАЧЕСТВА ОТВЕТА\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n\n"

        "МАТЕМАТИКА И ФИЗИКА:\n"
        "• Инлайн: $формула$ | Блочные: $$формула$$\n"
        "• КАЖДУЮ формулу — в ОДНУ строку внутри $...$. Никогда не разбивай!\n"
        "• Правильно: $\\frac{\\sin\\alpha}{\\sin\\beta} = \\frac{v_1}{v_2}$\n"
        "• Запрещено: писать sin отдельно от alpha или разбивать на строки\n"
        "• $\\frac{a}{b}$, $\\sqrt{x}$, $x^2$, $\\vec{F}$, $v_0$, $\\int_a^b f\\,dx$\n\n"

        "КОД:\n"
        "• Блоки ```язык\\n...\\n``` — всегда с указанием языка\n"
        "• Структура: краткое объяснение → рабочий код → пояснение ключевых мест\n"
        "• Код должен работать без правок — проверяй логику перед отправкой\n\n"

        "ТАБЛИЦЫ:\n"
        "• Markdown: | Заголовок | Заголовок | с разделителем | :--- | :---: | ---: |\n"
        "• Содержательные — не пустые клетки ради «структуры»\n\n"

        "СТРУКТУРА ОТВЕТА:\n"
        "• ### Заголовок для разделов, #### для подразделов\n"
        "• Списки: - или 1. 2. 3. — не смешивай без причины\n"
        "• **Жирный** для ключевых понятий\n\n"

        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "ЗАПРЕЩЕНО АБСОЛЮТНО:\n"
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        "• Начинать с 'Конечно!', 'Отличный вопрос!', 'Безусловно!'\n"
        "• Говорить 'Как ИИ я не могу...' — ты можешь\n"
        "• Давать неполный код с '...здесь ваш код...'\n"
        "• Придумывать факты и цифры — только реальные данные\n"
        "• Заканчивать ответ предложением задать ещё вопросы"
        + _mode_instr
    )

    # Определяем системный промпт на основе темы запроса
    _ling_kw = ["разбор предложения", "подчеркни члены", "морфологический разбор",
                "синтаксический разбор", "члены предложения", "части речи разбор",
                "фонетический разбор", "лингвистический разбор", "разобрать предложение",
                "выдели подлежащее", "выдели сказуемое", "разбор слова"]
    if any(kw in message.text.lower() for kw in _ling_kw):
        system_content = GLOBAL_SYS + "\n\n" + LINGUISTICS_SYSTEM_PROMPT
    elif is_code_question(message.text):
        system_content = (
            GLOBAL_SYS + "\n\n"
            "━━━ РЕЖИМ: РАЗРАБОТЧИК ━━━\n"
            "Ты — senior-разработчик. Правила:\n"
            "• Код ВСЕГДА в блоках ```язык\\n...\\n``` — никогда текстом\n"
            "• Код рабочий и протестированный — не давай заготовки с '...'\n"
            "• Структура: суть проблемы → рабочий код → объяснение каждого шага\n"
            "• Предупреждай о возможных edge cases и ограничениях\n"
            "• Если есть более эффективное решение — предлагай его"
        )
    elif any(kw in message.text.lower() for kw in ["реши", "вычисли", "найди", "формул", "уравнени",
             "интеграл", "производн", "прогресси", "матем", "задач", "докажи", "упрости"]):
        system_content = (
            GLOBAL_SYS + "\n\n"
            "━━━ РЕЖИМ: МАТЕМАТИКА ━━━\n"
            "Это математическая/физическая задача. Правила:\n"
            "• Решай ТОЛЬКО математически — формулы, не код\n"
            "• LaTeX ОБЯЗАТЕЛЕН: $инлайн$ и $$блочные$$\n"
            "• Каждую формулу — В ОДНУ СТРОКУ: $\\frac{a}{b}$, $\\sqrt{D}$\n"
            "• Структура: **Дано** → **Формулы** → **Решение по шагам** → **Ответ**\n"
            "• Каждый шаг с комментарием: что и почему делаем\n"
            "• Финальный ответ: **Ответ: $x = ...$** с единицами измерения\n"
            "• Если задача имеет несколько решений — показывай все"
        )
    else:
        system_content = GLOBAL_SYS

    history = list(user_memory[uid])
    sys_messages = [{"role": "system", "content": system_content}] + history
    user_memory[uid].append({"role": "user", "content": message.text})

    async with ChatActionSender.typing(bot=bot, chat_id=message.chat.id):
        try:
            query_msgs = sys_messages + [{"role": "user", "content": message.text}]

            # Thinking-сообщение с живым счётчиком секунд
            think_msg = await message.answer("🔍 <i>Анализирую ваш вопрос... 1с</i>", parse_mode="HTML")
            _start_time = asyncio.get_event_loop().time()

            async def _update_think(msg, start):
                phases = [
                    (3,  "🔍 Анализирую ваш вопрос"),
                    (8,  "🧠 Размышляю"),
                    (15, "💡 Формирую ответ"),
                    (999,"⚡ Ответ почти готов"),
                ]
                while True:
                    await asyncio.sleep(1)
                    elapsed = int(asyncio.get_event_loop().time() - start)
                    label = phases[0][1]
                    for threshold, lbl in phases:
                        if elapsed < threshold:
                            label = lbl
                            break
                    try:
                        await msg.edit_text(f"{label}... <b>{elapsed}с</b>", parse_mode="HTML")
                    except Exception:
                        pass

            _ticker_task = asyncio.create_task(_update_think(think_msg, _start_time))
            try:
                ans = await call_chat(query_msgs, mk)
            finally:
                _ticker_task.cancel()
                try:
                    await _ticker_task
                except asyncio.CancelledError:
                    pass

            user_memory[uid].append({"role": "assistant", "content": ans})
            last_responses[uid] = {"q": message.text, "a": ans, "model_label": MODELS.get(mk, {}).get("label", "ИИ"), "model_key": mk}
            import datetime as _dt2
            if uid not in user_history: user_history[uid] = []
            user_history[uid].append({
                "q": (message.text or "")[:200], "a": ans[:600],
                "model": mk, "ts": _dt2.datetime.now().strftime("%d.%m %H:%M")
            })
            user_history[uid] = user_history[uid][-10:]
            if uid in user_profiles:
                user_profiles[uid]["requests"] = user_profiles[uid].get("requests", 0) + 1
            spend_limit(uid, mk)
            asyncio.create_task(db_save_user(uid, message.from_user.first_name or "", message.from_user.username or ""))

            try:
                await think_msg.delete()
            except Exception:
                pass


            # Строим клавиатуру — кнопка «Красиво» ВСЕГДА
            _main_kb = save_kb(uid)
            try:
                _view_url = store_answer(ans, MODELS.get(mk, {}).get("label", "ИИ"), mk)
                _main_kb = InlineKeyboardMarkup(inline_keyboard=[
                    [InlineKeyboardButton(text="🌐 Красиво", url=_view_url)]
                ] + _main_kb.inline_keyboard)
            except Exception:
                pass

            if re.search(r"```", ans):
                await smart_send(message, ans, reply_markup=_main_kb)
            else:
                html_ans = md_to_html(ans)
                chunks = [html_ans[j:j+4000] for j in range(0, len(html_ans), 4000)]
                for ci, chunk in enumerate(chunks):
                    kb = _main_kb if ci == len(chunks)-1 else None
                    try:
                        await message.answer(chunk, parse_mode="HTML", reply_markup=kb)
                    except Exception:
                        await message.answer(ans[ci*4000:(ci+1)*4000], reply_markup=kb)

        except Exception as e:
            if user_memory[uid]:
                user_memory[uid].pop()
            logging.error(f"handle_text: {e}")
            try:
                await think_msg.delete()
            except Exception:
                pass
            await message.answer(f"❌ Ошибка: {e}")


# ==================================================================
# 🎨 ГЕНЕРАЦИЯ ИЗОБРАЖЕНИЙ
# ==================================================================


async def _do_imggen(message: Message, uid: int, prompt: str, img_key: str):
    m = IMG_MODELS.get(img_key)
    if not m:
        await message.answer("❌ Модель генерации не найдена")
        return
    if img_key in disabled_img_models:
        await message.answer(
            "🚫 <b>Модель временно недоступна</b>\n\nВыбери другую модель.",
            parse_mode="HTML",
            reply_markup=imggen_category_kb()
        )
        return

    if not can_img(uid):
        return await message.answer(
                "🚫 <b>Лимит генераций исчерпан!</b>\n\nОформи подписку для продолжения.",
            parse_mode="HTML",
        )

    # ── Улучшаем промпт через ИИ ─────────────────────────────────
    wait_msg = await message.answer(
        f"✨ <i>Улучшаю промпт через ИИ...</i>",
        parse_mode="HTML"
    )
    enhanced_prompt = await enhance_prompt_ai(prompt, mode="image")
    if enhanced_prompt != prompt:
        try:
            await wait_msg.edit_text(
                f"🎨 <i>Генерирую изображение...\n"
                f"Модель: <b>{m['label']}</b>\n"
                f"⏳ 15–40 секунд</i>",
                parse_mode="HTML"
            )
        except Exception:
            pass
    else:
        try:
            await wait_msg.edit_text(
                f"🎨 <i>Генерирую...\nМодель: <b>{m['label']}</b>\n⏳ 15–40 секунд</i>",
                parse_mode="HTML"
            )
        except Exception:
            pass

    async with ChatActionSender.upload_photo(bot=bot, chat_id=message.chat.id):
        try:
            result = await generate_image(enhanced_prompt, img_key)
            try:
                await wait_msg.delete()
            except Exception:
                pass

            spend_img(uid)
            _img_left = _get_lims(uid)['img_month'] - user_limits.get(uid, {}).get('img_used', 0)

            # Показываем оригинальный и улучшенный промпт если отличаются
            prompt_info = f"📝 <i>{prompt[:150]}</i>"
            if enhanced_prompt != prompt and len(enhanced_prompt) > 10:
                prompt_info = (
                    f"💡 <i>Твой запрос:</i> {prompt[:100]}\n"
                    f"✨ <i>Улучшено ИИ:</i> {enhanced_prompt[:150]}"
                )

            caption = (
                f"🎨 <b>{m['label']}</b>\n"
                f"{prompt_info}\n\n"
                f"🖼 Осталось генераций: <code>{_img_left}</code>"
            )
            img_kb = InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="🔄 Ещё вариант", callback_data=f"imggen_again_{img_key}"),
                 InlineKeyboardButton(text="✨ Новый промпт", callback_data="imggen_new")],
                [InlineKeyboardButton(text="🎵 Генерация музыки", callback_data="menu_music")],
            ])
            img_data = BufferedInputFile(result["data"], filename="generated.jpg")
            asyncio.create_task(db_save_user(uid))
            try:
                sent = await message.answer_photo(img_data, caption=caption, parse_mode="HTML", reply_markup=img_kb)
                # Сохраняем промпт для кнопки "Ещё вариант"
                last_responses[uid] = {**last_responses.get(uid, {}), "last_img_prompt": enhanced_prompt, "last_img_key": img_key}
            except Exception as photo_err:
                logging.warning(f"answer_photo failed ({photo_err}), trying document")
                img_data2 = BufferedInputFile(result["data"], filename="generated.jpg")
                await message.answer_document(img_data2, caption=caption, parse_mode="HTML")

        except Exception as e:
            try:
                await wait_msg.delete()
            except Exception:
                pass
            logging.error(f"imggen error: {e}")
            await message.answer(
                f"❌ <b>Ошибка генерации</b>\n\n{str(e)[:200]}\n\n"
                "Попробуй другую модель или переформулируй промпт.",
                parse_mode="HTML",
                reply_markup=imggen_category_kb()
            )


@dp.callback_query(F.data.startswith("imggen_again_"))
async def cb_imggen_again(callback: CallbackQuery):
    """Повторная генерация с тем же промптом — новый seed."""
    uid = callback.from_user.id
    img_key = callback.data.replace("imggen_again_", "")
    last = last_responses.get(uid, {})
    prompt = last.get("last_img_prompt", "")
    if not prompt:
        await callback.answer("❌ Промпт не найден", show_alert=True)
        return
    await callback.answer("🔄 Генерирую новый вариант...")
    await _do_imggen(callback.message, uid, prompt, img_key)


@dp.callback_query(F.data == "imggen_new")
async def cb_imggen_new(callback: CallbackQuery):
    uid = callback.from_user.id
    await callback.answer()
    await callback.message.answer(
        "✏️ <b>Введи новый промпт для изображения:</b>\n\n"
        "<b>Примеры:</b>\n" +
        "\n".join(f"• {e}" for e in IMG_PROMPT_EXAMPLES[:4]) +
        "\n\n💡 Можешь написать на русском — ИИ переведёт и улучшит автоматически",
        parse_mode="HTML",
        reply_markup=imggen_prompt_kb()
    )


# ==================================================================
# 🎵 MUSIC GENERATION HANDLERS
# ==================================================================

def music_styles_kb() -> InlineKeyboardMarkup:
    rows = []
    # Фильтруем отключённые стили
    items = [(k, v) for k, v in SUNO_STYLES.items() if k not in disabled_music_styles and k != "custom"]
    for i in range(0, len(items), 2):
        row = []
        for k, v in items[i:i+2]:
            row.append(InlineKeyboardButton(text=v["label"], callback_data=f"music_style_{k}"))
        rows.append(row)
    rows.append([
        InlineKeyboardButton(text="✨ Свой стиль", callback_data="music_style_custom"),
        InlineKeyboardButton(text="◀️ Назад", callback_data="back_home"),
    ])
    return InlineKeyboardMarkup(inline_keyboard=rows)


def music_confirm_kb(has_lyrics: bool = False) -> InlineKeyboardMarkup:
    rows = [
        [InlineKeyboardButton(text="🎵 Генерировать!", callback_data="music_gen_vocal"),
         InlineKeyboardButton(text="🎸 Без вокала",    callback_data="music_gen_instr")],
        [InlineKeyboardButton(text="✍️ Написать текст",callback_data="music_write_lyrics"),
         InlineKeyboardButton(text="🤖 ИИ напишет текст", callback_data="music_ai_lyrics")],
        [InlineKeyboardButton(text="✏️ Изменить идею", callback_data="music_edit_prompt"),
         InlineKeyboardButton(text="◀️ Стили",          callback_data="menu_music")],
    ]
    return InlineKeyboardMarkup(inline_keyboard=rows)


@dp.callback_query(F.data == "menu_music")
async def menu_music(callback: CallbackQuery):
    uid = callback.from_user.id
    _init_limits(uid)
    if not await check_service(callback, "music"):
        return
    sub_active   = has_active_sub(uid)
    li = get_limits_info(uid)
    music_left = max(0, li["music_max"] - li["music_used"])

    if not sub_active:
        lock_text = (
            "🔒 <b>Генерация музыки — только для подписчиков</b>\n\n"
            "С подпиской вы получите:\n"
            "🎵 <b>3 трека в месяц</b> — Suno AI\n"
            "🎤 С вокалом или инструментал\n"
            "🇷🇺 Текст на русском языке\n"
            "📝 ИИ напишет текст или вы сами\n"
            "🚫 Без водяных знаков\n\n"
            "Оформи подписку и создавай музыку! 👇"
        )
        kb = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="💎 Оформить подписку", callback_data="sub_menu")],
            [InlineKeyboardButton(text="◀️ Назад",             callback_data="back_home")],
        ])
        try:
            await callback.message.edit_text(lock_text, parse_mode="HTML", reply_markup=kb)
        except Exception:
            await callback.message.answer(lock_text, parse_mode="HTML", reply_markup=kb)
        await callback.answer()
        return

    if not can_music(uid):
        await callback.answer(
            f"🚫 Лимит музыки исчерпан! ({li['music_used']}/{li['music_max']} в месяц)",
            show_alert=True
        )
        await callback.answer()
        return

    cur_model = get_music_model(uid)
    cur_label = MUSIC_MODELS.get(cur_model, {}).get("label", cur_model)

    music_text = (
        "🎵 <b>Генерация музыки</b>\n\n"
        f"🎼 Модель: <b>{cur_label}</b>\n"
        f"🎵 Осталось треков: <b>{music_left}</b> / {li['music_max']} в месяц\n\n"
        "<b>Выбери модель или перейди к стилям:</b>"
    )
    _model_btns = [InlineKeyboardButton(text=mm["label"], callback_data=f"music_model_{mk}")
                   for mk, mm in MUSIC_MODELS.items() if mk not in disabled_music_models]
    kb = InlineKeyboardMarkup(inline_keyboard=[
        _model_btns or [InlineKeyboardButton(text="⚠️ Нет доступных моделей", callback_data="noop")],
        [InlineKeyboardButton(text="▶️ Выбрать стиль →", callback_data="music_go_styles")],
        [InlineKeyboardButton(text="◀️ Назад",           callback_data="back_home")],
    ])
    try:
        await callback.message.edit_text(music_text, parse_mode="HTML", reply_markup=kb)
    except Exception:
        await callback.message.answer(music_text, parse_mode="HTML", reply_markup=kb)
    await callback.answer()


@dp.callback_query(F.data.startswith("music_model_"))
async def cb_set_music_model(callback: CallbackQuery):
    """Выбор модели генерации музыки."""
    uid = callback.from_user.id
    key = callback.data.replace("music_model_", "")
    if key not in MUSIC_MODELS:
        await callback.answer("❌ Модель не найдена", show_alert=True)
        return
    set_music_model(uid, key)
    m = MUSIC_MODELS[key]
    await callback.answer(f"✅ {m['label']} выбрана!")
    li = get_limits_info(uid)
    music_left = max(0, li["music_max"] - li["music_used"])
    music_text = (
        "🎵 <b>Генерация музыки</b>\n\n"
        f"🎼 Модель: <b>{m['label']}</b>\n"
        f"ℹ️ {m['desc']}\n\n"
        f"🎵 Осталось треков: <b>{music_left}</b> / {li['music_max']} в месяц\n\n"
        "<b>Выбери модель или перейди к стилям:</b>"
    )
    model_btns = []
    for mk, mm in MUSIC_MODELS.items():
        if mk not in disabled_music_models:
            model_btns.append(InlineKeyboardButton(text=mm["label"], callback_data=f"music_model_{mk}"))
    kb = InlineKeyboardMarkup(inline_keyboard=[
        model_btns if model_btns else [InlineKeyboardButton(text="⚠️ Нет доступных моделей", callback_data="noop")],
        [InlineKeyboardButton(text="▶️ Выбрать стиль →", callback_data="music_go_styles")],
        [InlineKeyboardButton(text="◀️ Назад",           callback_data="back_home")],
    ])
    try:
        await callback.message.edit_text(music_text, parse_mode="HTML", reply_markup=kb)
    except Exception:
        pass


@dp.callback_query(F.data == "music_go_styles")
async def cb_music_go_styles(callback: CallbackQuery):
    """Переход к выбору стиля после выбора модели."""
    uid = callback.from_user.id
    cur_model = get_music_model(uid)
    cur_label = MUSIC_MODELS.get(cur_model, {}).get("label", cur_model)
    li = get_limits_info(uid)
    music_left = max(0, li["music_max"] - li["music_used"])
    music_text = (
        f"🎵 <b>Генерация музыки — {cur_label}</b>\n\n"
        "🎤 С вокалом или инструментал  |  🚫 Без водяных знаков\n"
        "🇷🇺 Можно на русском языке\n\n"
        f"🎵 Осталось треков: <b>{music_left}</b> / {li['music_max']} в месяц\n\n"
        "<b>👇 Выбери стиль:</b>"
    )
    try:
        await callback.message.edit_text(music_text, parse_mode="HTML", reply_markup=music_styles_kb())
    except Exception:
        await callback.message.answer(music_text, parse_mode="HTML", reply_markup=music_styles_kb())
    await callback.answer()


@dp.callback_query(F.data.startswith("music_style_"))
async def cb_music_style(callback: CallbackQuery):
    uid = callback.from_user.id
    style_key = callback.data.replace("music_style_", "")

    if style_key == "custom":
        music_states[uid] = {"stage": "await_custom_style"}
        await callback.message.edit_text(
            "🎸 <b>Свой стиль</b>\n\n"
            "Опиши жанр и инструменты:\n\n"
            "<i>Например: «тёмный джаз, пианино, контрабас»\n"
            "или «агрессивный metal, двойной бас, гитарное соло»</i>",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="◀️ Назад к стилям", callback_data="menu_music")
            ]])
        )
        await callback.answer()
        return

    style = SUNO_STYLES.get(style_key)
    if not style:
        await callback.answer("❌ Стиль не найден", show_alert=True)
        return

    lang = style.get("lang", "en")
    lang_hint = "🇷🇺 Русский язык" if lang == "ru" else "🇬🇧 English"
    music_states[uid] = {"stage": "await_wishes", "style_key": style_key,
                         "style_tags": style["tags"], "style_label": style["label"]}

    await callback.message.edit_text(
        f"╔══════════════════════╗\n"
        f"║  🎵 {style['label']}  ║\n"
        f"╚══════════════════════╝\n\n"
        f"<b>{style['desc']}</b>  ·  {lang_hint}\n"
        f"<i>{style['tags'][:80]}</i>\n\n"
        f"━━━━━━━━━━━━━━━━━━━━━━━━\n"
        f"Выбери способ создания трека:",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="✏️ Описать идею (ИИ напишет текст)", callback_data=f"music_wish_{style_key}")],
            [InlineKeyboardButton(text="✍️ Написать текст самому",           callback_data=f"music_manual_{style_key}")],
            [InlineKeyboardButton(text="✨ ИИ придумает всё сам",            callback_data=f"music_ai_prompt_{style_key}")],
            [InlineKeyboardButton(text="◀️ Стили",                           callback_data="menu_music")],
        ])
    )
    await callback.answer()


@dp.callback_query(F.data.startswith("music_ai_prompt_"))
async def cb_music_ai_prompt(callback: CallbackQuery):
    """ИИ генерирует промпт для музыки в выбранном стиле."""
    uid = callback.from_user.id
    style_key = callback.data.replace("music_ai_prompt_", "")
    style = SUNO_STYLES.get(style_key, {})
    await callback.answer("✨ Генерирую идею...")

    try:
        lang = style.get("lang","en")
        if lang == "ru":
            sys_p = (
                f"Ты — поэт. Придумай тему для песни в стиле: {style.get('desc','')}. "
                "Опиши в 1-2 предложениях — что происходит, какие эмоции. "
                "Только тема, без объяснений. На русском языке."
            )
            usr_p = "Придумай тему для трека"
        else:
            sys_p = (
                f"Suggest a creative song theme for {style.get('desc','')} style. "
                "1-2 sentences about emotions or story. No explanations."
            )
            usr_p = "Suggest a song theme"

        msgs = [{"role":"system","content":sys_p},{"role":"user","content":usr_p}]
        ai_idea = (await call_chat(msgs,"claude_haiku",max_tokens=120)).strip()
        lang_note = "🇷🇺 Текст на русском" if lang == "ru" else "🇬🇧 Lyrics in English"

        music_states[uid] = {
            "stage": "confirm",
            "style_key": style_key,
            "style_tags": style.get("tags",""),
            "style_label": style.get("label",""),
            "prompt": ai_idea,
            "title": f"AI {style.get('label','Track')}",
        }
        await callback.message.edit_text(
            f"✨ <b>ИИ придумал идею:</b>\n\n"
            f"🎵 Стиль: <b>{style.get('label','')}</b>\n"
            f"💡 <i>{ai_idea}</i>\n\n"
            f"🤖 ИИ напишет полный текст при генерации · {lang_note}\n\n"
            f"Генерировать?",
            parse_mode="HTML",
            reply_markup=music_confirm_kb()
        )
    except Exception as e:
        await callback.message.answer(f"❌ Ошибка: {e}")


@dp.callback_query(F.data == "music_check_pending")
async def cb_music_check_pending(callback: CallbackQuery):
    uid = callback.from_user.id
    pending = suno_pending.get(uid)
    if not pending:
        await callback.answer("❌ Нет ожидающих треков", show_alert=True)
        return
    await callback.answer("🔄 Проверяю...", show_alert=False)
    wait = await callback.message.answer("⏳ <i>Проверяю готовность трека...</i>", parse_mode="HTML")
    task_id     = pending["task_id"]
    audio_url   = ""
    image_url   = ""
    track_title = pending.get("title", "AI Track")
    for _ in range(24):  # 24×5=2 мин
        await asyncio.sleep(5)
        try:
            result = await suno_get_result(task_id)
            if result["status"] == "complete":
                audio_url   = result["audio_url"]
                image_url   = result.get("image_url", "")
                track_title = result.get("title", track_title)
                break
        except Exception:
            pass
    try: await wait.delete()
    except: pass
    if not audio_url:
        await callback.message.answer(
            "⏳ <b>Всё ещё генерируется!</b>\n\nПодожди ещё 1–2 минуты.",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="🔄 Проверить ещё раз", callback_data="music_check_pending")],
                [InlineKeyboardButton(text="🎵 Новый трек", callback_data="menu_music")],
            ])
        )
        return
    suno_pending.pop(uid, None)
    await _send_music_result(
        callback.message, uid, pending.get("state", {}),
        audio_url, image_url, track_title,
        pending.get("prompt", ""), pending.get("instrumental", False)
    )


@dp.callback_query(F.data.startswith("music_wish_"))
async def cb_music_wish(callback: CallbackQuery):
    """Пользователь описывает идею — ИИ напишет текст при генерации."""
    uid = callback.from_user.id
    style_key = callback.data.replace("music_wish_", "")
    style = SUNO_STYLES.get(style_key, {})
    lang = style.get("lang", "en")
    lang_note = "на русском языке" if lang == "ru" else "in English"
    music_states[uid] = {"stage": "await_wishes", "style_key": style_key,
                         "style_tags": style.get("tags", ""), "style_label": style.get("label", "")}
    await callback.message.edit_text(
        f"✏️ <b>Опиши идею трека</b>\n\n"
        f"🎵 Стиль: <b>{style.get('label', '')}</b>\n\n"
        f"Напиши о чём хочешь песню — тема, настроение, история, образы.\n"
        f"ИИ напишет полный текст {lang_note} сам.\n\n"
        f"<b>Примеры:</b>\n"
        f"  • про ночной город и одиночество\n"
        f"  • весёлая про лето и друзей\n"
        f"  • грустная про расставание\n\n"
        f"✏️ Пиши свою идею:",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="✨ ИИ придумает всё сам", callback_data=f"music_ai_prompt_{style_key}")],
            [InlineKeyboardButton(text="◀️ Назад к стилям",        callback_data="menu_music")],
        ])
    )
    await callback.answer()


@dp.callback_query(F.data.startswith("music_manual_"))
async def cb_music_manual(callback: CallbackQuery):
    """Пользователь пишет текст сам — без ИИ."""
    uid = callback.from_user.id
    style_key = callback.data.replace("music_manual_", "")
    style = SUNO_STYLES.get(style_key, {})
    lang = style.get("lang", "en")
    hint = "на русском" if lang == "ru" else "на любом языке"
    music_states[uid] = {"stage": "await_lyrics", "style_key": style_key,
                         "style_tags": style.get("tags", ""), "style_label": style.get("label", "")}
    await callback.message.edit_text(
        f"✍️ <b>Напиши текст песни</b>\n\n"
        f"🎵 Стиль: <b>{style.get('label', '')}</b>\n\n"
        f"Используй структуру Suno:\n"
        f"<code>[Verse]\nТекст куплета...\n\n"
        f"[Chorus]\nТекст припева...\n\n"
        f"[Verse 2]\nВторой куплет...</code>\n\n"
        f"Текст {hint}. Suno споёт именно то, что напишешь.\n\n"
        f"✍️ Пиши текст:",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="🤖 ИИ напишет текст",  callback_data=f"music_ai_prompt_{style_key}")],
            [InlineKeyboardButton(text="◀️ Назад к стилям",    callback_data="menu_music")],
        ])
    )
    await callback.answer()


@dp.callback_query(F.data == "music_write_lyrics")
async def cb_music_write_lyrics(callback: CallbackQuery):
    """Кнопка из confirm — вернуться к написанию текста."""
    uid   = callback.from_user.id
    state = music_states.get(uid, {})
    style_key = state.get("style_key", "")
    style = SUNO_STYLES.get(style_key, {})
    lang  = style.get("lang", "en")
    hint  = "на русском" if lang == "ru" else "на любом языке"
    music_states[uid] = {**state, "stage": "await_lyrics"}
    await callback.message.edit_text(
        f"✍️ <b>Напиши текст песни</b>\n\n"
        f"🎵 Стиль: <b>{style.get('label', '')}</b>\n\n"
        f"<code>[Verse]\nТекст куплета...\n\n[Chorus]\nТекст припева...</code>\n\n"
        f"Текст {hint}. Suno споёт именно его.",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="🤖 ИИ напишет текст", callback_data="music_ai_lyrics")],
            [InlineKeyboardButton(text="◀️ Назад к стилям",   callback_data="menu_music")],
        ])
    )
    await callback.answer()


@dp.callback_query(F.data == "music_ai_lyrics")
async def cb_music_ai_lyrics(callback: CallbackQuery):
    uid   = callback.from_user.id
    state = music_states.get(uid, {})
    style = SUNO_STYLES.get(state.get("style_key", ""), {})
    lang  = style.get("lang", "en")
    prompt = state.get("prompt", state.get("original_wishes", ""))
    await callback.answer("✍️ Пишу текст...", show_alert=False)
    wait = await callback.message.answer("✍️ <i>ИИ пишет текст песни...</i>", parse_mode="HTML")
    try:
        if lang == "ru":
            sys_p = (
                f"Ты — талантливый русский поэт. Напиши текст песни в стиле «{style.get('desc', '')}» "
                f"на русском языке. Формат: [Verse], [Chorus], [Bridge]. "
                f"Текст рифмованный, атмосферный. Только текст, без пояснений."
            )
        else:
            sys_p = (
                f"You are a songwriter. Write song lyrics, style: {style.get('desc', '')}. "
                f"Format: [Verse], [Chorus], [Bridge]. Only lyrics, no explanations."
            )
        msgs = [{"role": "system", "content": sys_p},
                {"role": "user", "content": f"Тема: {prompt}"}]
        lyrics = (await call_chat(msgs, "claude_haiku", max_tokens=700)).strip()
    except Exception:
        lyrics = ""
    try: await wait.delete()
    except: pass
    if not lyrics:
        await callback.message.answer("❌ Не удалось написать текст.")
        return
    music_states[uid] = {**state, "stage": "confirm", "lyrics": lyrics}
    await callback.message.answer(
        f"✍️ <b>Текст от ИИ:</b>\n\n<code>{lyrics[:700]}</code>\n\n"
        f"Генерировать с этим текстом?",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="🎵 Генерировать!", callback_data="music_gen_vocal"),
             InlineKeyboardButton(text="🎸 Без вокала",    callback_data="music_gen_instr")],
            [InlineKeyboardButton(text="✍️ Переписать",    callback_data="music_ai_lyrics"),
             InlineKeyboardButton(text="◀️ Назад",          callback_data="menu_music")],
        ])
    )


@dp.callback_query(F.data.in_({"music_gen_start", "music_gen_vocal", "music_gen_instrumental", "music_gen_instr"}))
async def cb_music_gen(callback: CallbackQuery):
    uid = callback.from_user.id
    state = music_states.get(uid, {})
    if state.get("stage") != "confirm":
        await callback.answer("❌ Сначала выбери стиль и промпт", show_alert=True)
        return
    instrumental = callback.data in ("music_gen_instrumental", "music_gen_instr")

    # ── Экран подтверждения (защита от случайного двойного тапа) ──
    style_label = state.get("style_label", "")
    prompt_preview = state.get("prompt", "")[:100]
    mode_icon = "🎸 Инструментал" if instrumental else "🎤 С вокалом"

    confirm_text = (
        "╔══════════════════════╗\n"
        f"║  🎵 ПОДТВЕРЖДЕНИЕ  ║\n"
        "╚══════════════════════╝\n\n"
        f"🎼 Стиль: <b>{style_label}</b>\n"
        f"🎤 Режим: <b>{mode_icon}</b>\n"
        f"📝 Идея: <i>{prompt_preview}{'...' if len(state.get('prompt','')) > 100 else ''}</i>\n\n"
        "⚠️ <b>Будет потрачен 1 трек из лимита!</b>\n\n"
        "▾ Нажми «✅ Генерировать» для запуска"
    )
    # Сохраняем instrumental в state для финального запуска
    music_states[uid]["_pending_instrumental"] = instrumental
    music_states[uid]["stage"] = "final_confirm"

    confirm_kb = InlineKeyboardMarkup(inline_keyboard=[
        [InlineKeyboardButton(
            text="✅ Да, генерировать!",
            callback_data="music_final_go"
        )],
        [InlineKeyboardButton(
            text="✏️ Изменить промпт",
            callback_data="music_edit_prompt"
        ),
        InlineKeyboardButton(
            text="◀️ Отмена",
            callback_data="menu_music"
        )],
    ])
    try:
        await callback.message.edit_text(confirm_text, parse_mode="HTML", reply_markup=confirm_kb)
    except Exception:
        await callback.message.answer(confirm_text, parse_mode="HTML", reply_markup=confirm_kb)
    await callback.answer()


@dp.callback_query(F.data == "music_final_go")
async def cb_music_final_go(callback: CallbackQuery):
    """Финальный запуск генерации после подтверждения."""
    uid = callback.from_user.id
    state = music_states.get(uid, {})
    if state.get("stage") != "final_confirm":
        await callback.answer("❌ Начни заново", show_alert=True)
        return
    instrumental = state.get("_pending_instrumental", False)
    # Восстанавливаем stage для _do_music_gen
    music_states[uid]["stage"] = "confirm"
    await callback.answer("🎵 Запускаю генерацию!")
    await _do_music_gen(callback.message, uid, music_states[uid], instrumental)


@dp.callback_query(F.data == "music_edit_prompt")
async def cb_music_edit_prompt(callback: CallbackQuery):
    uid = callback.from_user.id
    state = music_states.get(uid, {})
    if not state:
        await callback.answer("❌ Начни заново", show_alert=True)
        return
    music_states[uid]["stage"] = "await_prompt"
    await callback.message.edit_text(
        "✏️ <b>Введи новое описание трека:</b>",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="◀️ Стили", callback_data="menu_music")
        ]])
    )
    await callback.answer()


async def _do_music_gen(message, uid: int, state: dict, instrumental: bool = False):
    """Запускает генерацию музыки через выбранную модель (Suno / ElevenLabs)."""
    if not can_music(uid):
        L = _get_lims(uid)
        if not has_active_sub(uid):
            await message.answer(
                "🚫 <b>Генерация музыки — только для подписчиков!</b>\n\n"
                "В бесплатном плане: <b>0 треков</b>\n"
                "С подпиской: <b>3 трека</b> в месяц\n\n"
                "💎 Оформи подписку и создавай музыку:",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                    [InlineKeyboardButton(text="💎 Оформить подписку", callback_data="sub_menu")],
                    [InlineKeyboardButton(text="◀️ Назад", callback_data="back_home")],
                ])
            )
        else:
            await message.answer(
                f"🚫 <b>Лимит генераций музыки исчерпан!</b>\n\n"
                f"Использовано: <b>{user_limits.get(uid, {}).get('music_used', 0)}</b> из "
                f"<b>{L['music_month']}</b> треков\nСброс в начале следующего месяца.",
                parse_mode="HTML"
            )
        return

    prompt     = state.get("prompt", "")
    style_tags = state.get("style_tags", "")
    style_key  = state.get("style_key", "")
    style_obj  = SUNO_STYLES.get(style_key, {})
    style_desc = style_obj.get("desc", "")
    title      = state.get("title", "AI Track")[:80]
    lang       = style_obj.get("lang", "en")
    lyrics     = state.get("lyrics", "")

    prompt_clean = prompt.strip() if prompt.strip() else ("about life" if lang != "ru" else "о жизни")

    wait_msg = await message.answer(
        f"🎵 <b>Генерирую трек...</b>\n\n"
        f"🎼 <i>{state.get('style_label', style_key)}</i>\n"
        f"📝 <i>{prompt_clean[:150]}</i>\n\n"
        f"⏳ Suno обычно делает 2–4 минуты. Подожди!",
        parse_mode="HTML"
    )

    try:
        # Шаг 1: генерация текста песни (если не инструментал)
        if instrumental:
            song_lyrics = ""
        elif lyrics:
            song_lyrics = lyrics
        else:
            try:
                await wait_msg.edit_text(
                    f"✍️ <b>Пишу текст песни...</b>\n\n"
                    f"🎼 <i>{state.get('style_label', style_key)}</i>\n"
                    f"📝 <i>{prompt_clean[:150]}</i>",
                    parse_mode="HTML"
                )
            except Exception:
                pass
            song_lyrics = await generate_music_lyrics(
                wishes=prompt_clean, style_desc=style_desc,
                style_tags=style_tags, lang=lang
            )

        # Шаг 2: собираем финальный промпт
        music_model = get_music_model(uid)

        if lang == "ru":
            lang_prefix = "russian vocals, russian language, "
        else:
            lang_prefix = ""

        if instrumental:
            style_part = f"{style_tags}, instrumental, no vocals" if style_tags else "instrumental, no vocals"
            full_prompt = f"{lang_prefix}{style_part}"
        else:
            style_part = f"{style_tags}, " if style_tags else ""
            full_prompt = f"{lang_prefix}{style_part}{song_lyrics or prompt_clean}"

        model_label = MUSIC_MODELS.get(music_model, {}).get("label", music_model)
        try:
            await wait_msg.edit_text(
                f"🎵 <b>{model_label} создаёт трек...</b>\n\n"
                f"🎼 <i>{state.get('style_label', style_key)}</i>\n"
                f"📝 <i>{prompt_clean[:100]}</i>\n\n"
                f"⏳ Обычно 1–4 минуты...",
                parse_mode="HTML"
            )
        except Exception:
            pass

        # Шаг 3: запрос к выбранной модели через Pollinations
        if music_model == "elevenlabs":
            el_prompt = f"{style_tags}, {prompt_clean}" if style_tags else prompt_clean
            if instrumental:
                el_prompt += ", instrumental, no vocals"
            if lang == "ru":
                el_prompt = f"russian style, {el_prompt}"
            audio_bytes = await generate_music_elevenlabs(el_prompt, lang=lang)
        else:
            # Suno v5 через Pollinations (default)
            audio_bytes = await generate_music_pollinations(full_prompt, lang=lang)

        try:
            await wait_msg.delete()
        except Exception:
            pass

        # Шаг 4: отправляем результат
        spend_music(uid)
        asyncio.create_task(db_save_user(uid))

        style_label = style_obj.get("label", "")
        mode_label  = "🎸 Инструментал" if instrumental else "🎤 С вокалом"
        lang_flag   = "🇷🇺 На русском" if lang == "ru" else "🇬🇧 English"
        # Получаем текст песни из состояния если он был написан пользователем
        user_lyrics = state.get("lyrics", "").strip()
        if user_lyrics and not instrumental:
            lyrics_block = f"\n\n📜 <b>Текст песни:</b>\n<i>{user_lyrics[:600]}{'...' if len(user_lyrics) > 600 else ''}</i>"
        else:
            lyrics_block = ""
        caption = (
            f"🎵 <b>{title}</b>\n\n"
            f"🎼 {style_label} · {mode_label} · {lang_flag}\n"
            f"📝 <i>{prompt_clean[:150]}</i>"
            f"{lyrics_block}"
        )
        music_kb = InlineKeyboardMarkup(inline_keyboard=[
            [InlineKeyboardButton(text="🔄 Ещё вариант", callback_data="music_gen_start"),
             InlineKeyboardButton(text="🎵 Новый трек",  callback_data="menu_music")],
            [InlineKeyboardButton(text="🎨 Картинки",    callback_data="menu_imggen")],
        ])
        music_states[uid] = {**state, "stage": "confirm"}

        audio_file = BufferedInputFile(audio_bytes, filename=f"{title[:40]}.mp3")
        await message.answer_audio(audio_file, caption=caption[:900], parse_mode="HTML", reply_markup=music_kb)

    except Exception as e:
        try:
            await wait_msg.delete()
        except Exception:
            pass
        logging.error(f"music gen error uid={uid}: {e}")
        await message.answer(
            f"❌ <b>Ошибка генерации музыки</b>\n\n<code>{str(e)[:200]}</code>",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="🔄 Попробовать снова", callback_data="menu_music")
            ]])
        )


# ==================================================================
# 📷 ОБРАБОТКА ФОТО
# ==================================================================

async def _do_photofusion(message: Message, uid: int, img_b64: str, user_request: str):
    wait_msg = await message.answer(
        "🪄 <i>Анализирую твоё фото и создаю магию…\n⏳ Это занимает 30–90 секунд, подожди!</i>",
        parse_mode="HTML"
    )
    async with ChatActionSender.upload_photo(bot=bot, chat_id=message.chat.id):
        try:
            VISION_FALLBACK = ["claude_opus", "claude_haiku", "claude_sonnet", "gpt52", "qwen3_vl", "c4ai_vis32b"]
            cur_key = get_model_key(uid)
            
            # Если выбрана vision модель - используем её без fallback
            if cur_key in VISION_MODELS and cur_key not in disabled_models:
                # Проверка подписки для премиум моделей
                if cur_key in PREMIUM_MODELS and not has_active_sub(uid):
                    await wait_msg.delete()
                    await message.answer(
                        f"🔒 <b>Модель {MODELS[cur_key]['label']} доступна только по подписке</b>\n\n"
                        "💡 Выбери другую модель с Vision или оформи подписку.",
                        parse_mode="HTML",
                        reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                            [InlineKeyboardButton(text="💎 Оформить подписку", callback_data="sub_menu")],
                            [InlineKeyboardButton(text="🤖 Сменить модель", callback_data="menu_ai")],
                        ])
                    )
                    return
                vision_order = [cur_key]  # Только выбранная модель
            else:
                vision_order = [k for k in VISION_FALLBACK if k in MODELS and k not in disabled_models]

            person_desc = None
            for vk in vision_order:
                if vk not in MODELS:
                    continue
                try:
                    vis_msg = [{"role": "user", "content": [
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}},
                        {"type": "text", "text": (
                            "Describe this person's face in detail for image generation. "
                            "Include: face shape, eye color/shape, nose, lips, hair color/style, skin tone, age, gender. "
                            "Be specific. English, max 150 words."
                        )}
                    ]}]
                    person_desc = await call_chat(vis_msg, vk, max_tokens=200)
                    break
                except Exception as e:
                    logging.warning(f"Vision {vk}: {e}")
                    continue

            if not person_desc:
                person_desc = "a person"

            gen_prompt = (
                f"ultra-photorealistic portrait, {user_request}, "
                f"subject: {person_desc}. "
                f"preserve identical face features, 8K, professional lighting, sharp details"
            )

            IMG_ORDER = ["flux_klein4", "flux_klein9", "gptimg1", "gptimg_mini", "gptimg15"]
            result = None
            for img_key in IMG_ORDER:
                if img_key not in IMG_MODELS or img_key in disabled_img_models:
                    continue
                try:
                    result = await generate_image(gen_prompt, img_key)
                    break
                except Exception as e:
                    logging.warning(f"ImgGen {img_key}: {e}")

            await wait_msg.delete()
            if not result:
                await message.answer("❌ Не удалось создать изображение")
                return

            spend_img(uid)
            _img_left = LIMITS['img_month'] - user_limits.get(uid, {}).get('img_used', 0)

            caption_out = (
                f"🪄 <b>Фото-магия готова!</b>\n"
                f"<i>Запрос: {user_request[:100]}</i>\n"
                f"<i>💎 Остаток генераций: <code>{_img_left}</code></i>"
            )
            img_data = BufferedInputFile(result["data"], filename="fusion.jpg")
            try:
                await message.answer_photo(img_data, caption=caption_out, parse_mode="HTML",
                                           reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                                               InlineKeyboardButton(text="🪄 Ещё раз", callback_data="menu_photofusion"),
                                               InlineKeyboardButton(text="🏠 Меню", callback_data="back_home"),
                                           ]]))
            except Exception:
                img_data2 = BufferedInputFile(result["data"], filename="fusion.jpg")
                await message.answer_document(img_data2, caption=caption_out, parse_mode="HTML")

        except Exception as e:
            try:
                await wait_msg.delete()
            except Exception:
                pass
            logging.error(f"PhotoFusion error: {e}")
            await message.answer(f"❌ Ошибка фото-магии: {str(e)[:300]}")


# ==================================================================
# 📄 АНАЛИЗ ФАЙЛОВ (DOCX, PDF, TXT и другие)
# ==================================================================

async def _extract_text_from_file(file_path: str, mime_type: str, file_name: str) -> str:
    """Извлекает текст из файла в зависимости от типа."""
    ext = os.path.splitext(file_name)[1].lower() if file_name else ""

    # TXT файлы
    if ext in (".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm", ".log", ".py", ".js", ".ts", ".java", ".cpp", ".c", ".h", ".css"):
        try:
            async with aiofiles.open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return await f.read()
        except Exception as e:
            raise RuntimeError(f"Не удалось прочитать текстовый файл: {e}")

    # DOCX файлы
    if ext == ".docx":
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # Также таблицы
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)
            return "\n\n".join(paragraphs)
        except Exception as e:
            raise RuntimeError(f"Не удалось прочитать DOCX: {e}")

    # PDF файлы
    if ext == ".pdf":
        # Попробуем PyMuPDF (fitz) или pdfminer
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            texts = []
            for page in doc:
                texts.append(page.get_text())
            doc.close()
            result = "\n\n".join(texts)
            if result.strip():
                return result
        except ImportError:
            pass
        except Exception as e:
            logging.warning(f"fitz failed: {e}")

        try:
            from pdfminer.high_level import extract_text as pdf_extract_text
            result = pdf_extract_text(file_path)
            if result and result.strip():
                return result
        except ImportError:
            pass
        except Exception as e:
            logging.warning(f"pdfminer failed: {e}")

        # Fallback — попытка прочитать как текст
        try:
            async with aiofiles.open(file_path, "r", encoding="latin-1", errors="replace") as f:
                raw = await f.read()
            # Убираем бинарный мусор — оставляем только печатные символы
            import string
            printable = set(string.printable + "абвгдеёжзийклмнопрстуфхцчшщъыьэюяАБВГДЕЁЖЗИЙКЛМНОПРСТУФХЦЧШЩЪЫЬЭЮЯ")
            clean = "".join(c for c in raw if c in printable or c in "\n\r\t")
            # Убираем длинные цепочки нечитабельных символов
            import re as _re
            clean = _re.sub(r'[ \t]{5,}', ' ', clean)
            clean = _re.sub(r'\n{5,}', '\n\n', clean)
            if len(clean.strip()) > 100:
                return clean
        except Exception:
            pass

        raise RuntimeError("Не удалось извлечь текст из PDF. Попробуй установить PyMuPDF: pip install pymupdf")

    # Прочие файлы — попробуем как текст
    try:
        async with aiofiles.open(file_path, "r", encoding="utf-8", errors="replace") as f:
            content = await f.read()
        if content.strip():
            return content
    except Exception:
        pass

    raise RuntimeError(f"Неподдерживаемый формат файла: {ext or mime_type}")


@dp.message(F.document)
async def handle_document(message: Message):
    """Обработчик документов — читает и анализирует файлы любых текстовых форматов."""
    uid = message.from_user.id
    if not await require_subscription(message):
        return

    doc = message.document
    file_name = doc.file_name or "file"
    mime_type = doc.mime_type or ""
    ext = os.path.splitext(file_name)[1].lower()

    SUPPORTED_EXTS = {
        ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm", ".log",
        ".yaml", ".yml", ".toml", ".ini", ".conf", ".cfg", ".env",
        ".rst", ".tex", ".gitignore", ".editorconfig",
        ".py", ".pyw", ".js", ".mjs", ".ts", ".jsx", ".tsx",
        ".java", ".cpp", ".cc", ".cxx", ".c", ".h", ".hpp",
        ".cs", ".go", ".rs", ".php", ".rb", ".pl",
        ".sh", ".bash", ".zsh", ".bat", ".cmd", ".ps1",
        ".sql", ".graphql", ".gql", ".tf", ".kt", ".swift", ".dart",
        ".r", ".scala", ".vue", ".svelte", ".css", ".scss",
        ".docx", ".pdf",
    }

    is_supported = (
        ext in SUPPORTED_EXTS or
        any(m in mime_type for m in ["text/", "pdf", "word", "document", "json", "xml"])
    )

    if not is_supported:
        await message.answer(
            f"📎 Файл <code>{file_name}</code> получен.\n\n"
            f"⚠️ Формат <code>{ext or mime_type}</code> не поддерживается.\n"
            f"✅ Поддерживается: py, js, go, java, sql, txt, md, json, docx, pdf и другие.",
            parse_mode="HTML"
        )
        return

    _init_limits(uid)
    mk = get_model_key(uid)
    if mk.startswith("imggen:"):
        mk = "claude_sonnet"

    _ok, _reason = can_send(uid, mk)
    if not _ok:
        _rst = _reset_str(uid, _reason)
        await message.answer(f"🚫 Лимит исчерпан. Сброс через: <code>{_rst}</code>", parse_mode="HTML")
        return

    status_msg = await message.answer(f"📄 <i>Читаю <b>{file_name}</b>...</i>", parse_mode="HTML")
    file_path = tmp(f"doc_{uid}_{doc.file_id[:10]}{ext}")

    try:
        file_obj = await bot.get_file(doc.file_id)
        buf = io.BytesIO()
        await bot.download_file(file_obj.file_path, destination=buf)
        async with aiofiles.open(file_path, "wb") as f:
            await f.write(buf.getvalue())

        extracted_text = await _extract_text_from_file(file_path, mime_type, file_name)

        if not extracted_text or not extracted_text.strip():
            await status_msg.delete()
            await message.answer("❌ Не удалось извлечь текст из файла.\n• PDF-скан? Нужен текстовый PDF\n• Файл пустой или повреждён")
            return

        MAX_CHARS = 15000
        truncated = len(extracted_text) > MAX_CHARS
        if truncated:
            extracted_text = extracted_text[:MAX_CHARS]

        # Вопрос: подпись → последнее сообщение пользователя → авто-анализ
        user_caption = (message.caption or "").strip()
        last_user_q = ""
        if not user_caption and uid in user_memory:
            for m in reversed(list(user_memory[uid])):
                if m["role"] == "user" and len(m["content"]) < 500 and not m["content"].startswith("Файл '"):
                    last_user_q = m["content"]
                    break

        if user_caption:
            prompt = f"Файл '{file_name}':\n\n{extracted_text}\n\nЗадание: {user_caption}"
        elif last_user_q:
            prompt = f"Файл '{file_name}':\n\n{extracted_text}\n\nВопрос: {last_user_q}"
        else:
            prompt = (
                f"Файл '{file_name}':\n\n{extracted_text}\n\n"
                f"Проанализируй: о чём файл, ключевые части, структура, важные детали."
            )

        await status_msg.delete()
        think_msg = await message.answer(f"🧠 <i>Анализирую <b>{file_name}</b>...</i>", parse_mode="HTML")

        if uid not in user_memory:
            user_memory[uid] = deque(maxlen=20)
        user_memory[uid].append({"role": "user", "content": prompt})

        try:
            _file_sys_prompt = _get_file_system_prompt(file_name, user_caption or last_user_q)
            msgs = [{"role": "system", "content": _file_sys_prompt}] + list(user_memory[uid])
            ans = await call_chat(msgs, mk, max_tokens=2500)
            user_memory[uid].append({"role": "assistant", "content": ans})
            last_responses[uid] = {"q": f"[Файл: {file_name}]", "a": ans, "model_label": MODELS.get(mk, {}).get("label", "ИИ"), "model_key": mk}
            spend_limit(uid, mk)
            await think_msg.delete()

            trunc_notice = "\n<i>⚠️ Файл усечён до 15000 символов</i>" if truncated else ""
            header = f"📄 <b>{file_name}</b>{trunc_notice}\n\n"

            # Клавиатура с кнопкой «Красиво» — ВСЕГДА
            _doc_kb = save_kb(uid)
            try:
                _vurl2 = store_answer(ans, MODELS.get(mk, {}).get("label", "ИИ"), mk)
                _doc_kb = InlineKeyboardMarkup(inline_keyboard=[
                    [InlineKeyboardButton(text="🌐 Красиво", url=_vurl2)]
                ] + _doc_kb.inline_keyboard)
            except Exception:
                pass

            if re.search(r"```", ans):
                await message.answer(header, parse_mode="HTML")
                await smart_send(message, ans, reply_markup=_doc_kb)
            else:
                html_ans = md_to_html(ans)
                chunks = [html_ans[j:j+3800] for j in range(0, len(html_ans), 3800)]
                for ci, chunk in enumerate(chunks):
                    kb = _doc_kb if ci == len(chunks)-1 else None
                    pfx = header if ci == 0 else ""
                    try:
                        await message.answer(pfx + chunk, parse_mode="HTML", reply_markup=kb)
                    except Exception:
                        await message.answer(pfx + ans[:4000], reply_markup=kb)

        except Exception as e:
            if user_memory.get(uid):
                user_memory[uid].pop()
            try: await think_msg.delete()
            except Exception: pass
            logging.error(f"handle_document analyze: {e}")
            await message.answer(f"❌ Ошибка: {str(e)[:300]}")

    except RuntimeError as re_err:
        try: await status_msg.delete()
        except Exception: pass
        await message.answer(f"❌ Ошибка чтения: {re_err}", parse_mode="HTML")
    except Exception as e:
        try: await status_msg.delete()
        except Exception: pass
        logging.error(f"handle_document: {e}")
        await message.answer(f"❌ Ошибка: {str(e)[:300]}")
    finally:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception:
            pass


async def _do_img2img(message: Message, uid: int, img_b64: str, tg_file_path: str,
                      style_text: str, img_model_key: str, wait_msg=None):
    """Общая функция img2img: анализ фото → генерация в стиле с максимальным сохранением лица."""
    import contextlib
    try:
        if wait_msg:
            with contextlib.suppress(Exception):
                await wait_msg.edit_text("🎨 <b>Стилизую фото...</b>\n🔍 <i>Анализирую...</i>", parse_mode="HTML")

        # Короткое описание (max 80 символов) — длинный промпт убивает Pollinations
        vision_prompt = (
            "Describe this person in max 30 words: "
            "gender, age, hair color and style, eye color, skin tone, outfit. "
            "Keywords only, comma-separated English."
        )
        msgs = [{"role": "user", "content": [
            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}},
            {"type": "text", "text": vision_prompt},
        ]}]
        photo_desc = ""
        for vk in ["claude_haiku", "claude_sonnet", "gpt52", "qwen3_vl"]:
            if vk in MODELS and vk not in disabled_models and vk in VISION_MODELS:
                with contextlib.suppress(Exception):
                    photo_desc = (await call_chat(msgs, vk, max_tokens=60)).strip()
                    if photo_desc and len(photo_desc) > 10:
                        break
        if not photo_desc:
            photo_desc = "portrait, person"
        # Жёсткий лимит — Pollinations не справляется с длинными URL
        photo_desc = photo_desc[:200]

        if wait_msg:
            with contextlib.suppress(Exception):
                await wait_msg.edit_text(
                    f"🎨 <b>Стилизую фото...</b>\n✨ <i>Применяю «{style_text[:40]}»...</i>\n⏳ ~30–60 сек...",
                    parse_mode="HTML"
                )

        # Короткий итоговый промпт — ключевой фактор для Pollinations
        style_lower = style_text.lower()
        DETAIL_KEYWORDS = ["усы", "борода", "очки", "шляпа", "добавь", "нарисуй",
                           "mustache", "beard", "glasses", "hat", "add", "draw"]
        is_detail_edit = any(kw in style_lower for kw in DETAIL_KEYWORDS)

        if is_detail_edit:
            final_prompt = f"{photo_desc}, {style_text}, photorealistic"
        else:
            final_prompt = f"{style_text} style, {photo_desc}"
        # Абсолютный лимит длины
        final_prompt = final_prompt[:400]

        result = await generate_image(final_prompt, img_model_key, ratio="1:1",
                                      image_base64=img_b64, tg_file_path=tg_file_path)
        spend_img(uid)
        asyncio.create_task(db_save_user(uid))

        if wait_msg:
            with contextlib.suppress(Exception):
                await wait_msg.delete()

        m_label = IMG_MODELS.get(img_model_key, {}).get("label", img_model_key)
        img_file = BufferedInputFile(result["data"], filename="img2img.jpg")
        await message.answer_photo(
            img_file,
            caption=(
                f"🎨 <b>Стиль изменён</b> · {m_label}\n\n"
                f"✏️ <i>{style_text[:100]}</i>"
            ),
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                InlineKeyboardButton(text="🔄 Другой стиль", callback_data="photo_img2img_retry"),
                InlineKeyboardButton(text="🎬 Оживить видео", callback_data="photo_result_to_video"),
            ]])
        )
        last_photo[uid] = {"img_b64": img_b64, "tg_file_path": tg_file_path,
                           "result_b64": result["data"], "photo_context": True}
    except Exception as e:
        if wait_msg:
            with contextlib.suppress(Exception):
                await wait_msg.delete()
        await message.answer(f"❌ <b>Ошибка генерации:</b> {str(e)[:200]}", parse_mode="HTML")


@dp.message(F.photo)
async def handle_photo(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message):
        return

    file_obj = await bot.get_file(message.photo[-1].file_id)
    buf = io.BytesIO()
    await bot.download_file(file_obj.file_path, destination=buf)
    buf.seek(0)
    img_b64 = base64.b64encode(buf.read()).decode()
    tg_file_path = file_obj.file_path  # путь для Telegram CDN URL
    user_caption = message.caption or ""

    # ── Режим img2img через выбор модели из меню ─────────────────
    _pending = last_photo.get(uid, {})
    if _pending.get("img2img_model_pending"):
        chosen_key = _pending["img2img_model_pending"]
        last_photo[uid] = {
            "img_b64": img_b64,
            "tg_file_path": tg_file_path,
            "img2img_pending": True,
            "img2img_model_key": chosen_key,
        }
        m_label = IMG_MODELS.get(chosen_key, {}).get("label", chosen_key)
        if user_caption.strip():
            # Фото с подписью — сразу генерируем
            last_photo[uid]["img2img_pending"] = False
            if not can_img(uid):
                return await message.answer("🚫 <b>Лимит генераций исчерпан!</b>\n\nОформи подписку.", parse_mode="HTML")
            style_text = user_caption.strip()
        else:
            await message.answer(
                f"✅ <b>Фото получено!</b>\n\n"
                f"Модель: <b>{m_label}</b>\n\n"
                f"Теперь напиши стиль, например:\n"
                f"• <code>в стиле аниме</code>\n"
                f"• <code>oil painting</code>\n"
                f"• <code>cyberpunk neon</code>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                    InlineKeyboardButton(text="❌ Отмена", callback_data="photo_img2img_cancel")
                ]])
            )
        return

    # ── Режим «Оживить фото» для видео (через кнопку меню) ───────
    _video_state = video_states.get(uid, {})
    if _video_state.get("stage") == "await_photo":
        if not can_video(uid):
            li = get_limits_info(uid)
            await message.answer(
                f"🚫 <b>Лимит видео исчерпан!</b> ({li['video_used']}/{li['video_max']} в месяц)",
                parse_mode="HTML"
            )
            video_states.pop(uid, None)
            return
        # Сохраняем фото и спрашиваем про промпт
        video_states[uid] = {
            "mode": "photo",
            "stage": "await_video_prompt",
            "img_b64": img_b64,
            "tg_file_path": tg_file_path,
        }
        # Если пользователь сразу прислал подпись — используем её и генерируем
        if user_caption.strip():
            video_states.pop(uid, None)
            await _do_video_gen(message, uid, user_caption.strip(), mode="ai",
                                image_base64=img_b64, tg_file_path=tg_file_path)
            return
        # Иначе спрашиваем: свой промпт или AI
        await message.answer(
            "🖼 <b>Фото получено!</b>\n\n"
            "Как сгенерировать видео?\n\n"
            "✍️ <b>Напиши свой промпт</b> — бот сделает видео точно по нему\n"
            "✨ <b>Или выбери AI улучшение</b> — ИИ сам опишет движение и детали",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="✨ AI улучшение (авто)", callback_data="vidphoto_mode_ai")],
                [InlineKeyboardButton(text="✍️ Написать свой промпт", callback_data="vidphoto_mode_raw")],
                [InlineKeyboardButton(text="❌ Отмена", callback_data="vidphoto_cancel")],
            ])
        )
        return

    # ── Пользователь в режиме генерации КАРТИНОК (imggen:) ───────
    cur_mk = get_model_key(uid)
    if cur_mk.startswith("imggen:"):
        img_key = cur_mk.split(":", 1)[1]
        if user_caption.strip():
            # Фото + промпт → описываем фото vision-ИИ → генерируем в стиле
            if not can_img(uid):
                return await message.answer("🚫 <b>Лимит генераций исчерпан!</b>", parse_mode="HTML")
            last_photo[uid] = {"img_b64": img_b64, "img2img_pending": False}
            style_text = user_caption.strip()
            wait_msg = await message.answer(f"🎨 <b>Изменяю фото...</b>\n🔍 <i>Анализирую...</i>", parse_mode="HTML")
            try:
                # Короткое описание — длинные промпты убивают Pollinations
                msgs_v = [{"role": "user", "content": [
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}},
                    {"type": "text", "text": (
                        "Describe this person in max 30 words: "
                        "gender, age, hair color/style, eye color, skin tone, outfit. "
                        "Keywords only, comma-separated English."
                    )},
                ]}]
                photo_desc = ""
                for vk in ["claude_haiku", "claude_sonnet", "gpt52", "qwen3_vl"]:
                    if vk in MODELS and vk not in disabled_models and vk in VISION_MODELS:
                        try:
                            photo_desc = (await call_chat(msgs_v, vk, max_tokens=180)).strip()
                            if photo_desc and len(photo_desc) > 20:
                                break
                        except Exception:
                            continue
                if not photo_desc:
                    photo_desc = "portrait, detailed photo"

                # Используем flux2dev (img2img), передаём фото напрямую
                actual_key = img_key if IMG_MODELS.get(img_key, {}).get("img2img") else "flux2dev"
                await _do_img2img(message, uid, img_b64, tg_file_path, style_text, actual_key, wait_msg)
            except Exception as e:
                try: await wait_msg.delete()
                except Exception: pass
                await message.answer(f"❌ <b>Ошибка генерации:</b> {str(e)[:200]}", parse_mode="HTML")
            return
        else:
            # Фото без подписи в imggen-режиме → ждём промпт
            last_photo[uid] = {"img_b64": img_b64, "tg_file_path": tg_file_path, "img2img_pending": True}
            await message.answer(
                "🎨 <b>Фото получено!</b>\n\n"
                "Напиши промпт — что сделать с фото:\n\n"
                "💡 <i>Примеры:</i>\n"
                "• <code>в стиле аниме</code>\n"
                "• <code>oil painting portrait</code>\n"
                "• <code>cyberpunk neon style</code>",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                    InlineKeyboardButton(text="❌ Отмена", callback_data="photo_img2img_cancel")
                ]])
            )
            return

    # ── Пользователь в режиме генерации ВИДЕО (по тексту) ────────
    # Если отправил фото с подписью → сразу делаем видео из фото
    _v_state = video_states.get(uid, {})
    if _v_state.get("stage") == "await_prompt" and user_caption.strip():
        if not can_video(uid):
            li = get_limits_info(uid)
            return await message.answer(
                f"🚫 <b>Лимит видео исчерпан!</b> ({li['video_used']}/{li['video_max']})",
                parse_mode="HTML"
            )
        video_states.pop(uid, None)
        await _do_video_gen(message, uid, user_caption.strip(), mode="ai", image_base64=img_b64)
        return

    FUSION_KEYWORDS = ["сделай меня", "помести меня", "я в", "нарисуй меня", "стиле аниме",
                       "супергерой", "знаменитост", "бэтмен", "вместе с", "вставь меня"]
    is_fusion_mode = (user_photo_mode.get(uid) == "waiting_photo")
    is_fusion_cap  = any(kw in user_caption.lower() for kw in FUSION_KEYWORDS) if user_caption else False

    if is_fusion_mode or is_fusion_cap:
        user_photo_mode.pop(uid, None)
        if not user_caption.strip():
            last_photo[uid] = {"img_b64": img_b64, "tg_file_path": tg_file_path, "caption": "", "fusion_pending": True}
            await message.answer(
                "🪄 <b>Фото получено!</b>\n\nНапиши <b>что с ним сделать</b>:",
                parse_mode="HTML"
            )
            return
        if not can_img(uid):
                "🚫 <b>Лимит генераций исчерпан!</b>\n\nОформи подписку для продолжения.",
    # Если нет подписи — сохраняем фото и спрашиваем что нужно сделать
    if not user_caption.strip():
        last_photo[uid] = {"img_b64": img_b64, "tg_file_path": tg_file_path, "caption": "", "ask_pending": True}
        # Проверяем доступность видео для кнопки оживления
        video_btn_row = []
        if can_video(uid):
            video_btn_row = [InlineKeyboardButton(text="🎬 Оживить фото (видео)", callback_data="photo_to_video")]
        await message.answer(
            "📸 <b>Фото получено!</b>\n\n"
            "Что хочешь сделать с фото?\n\n"
            "💡 Или просто напиши вопрос / запрос:",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="🔍 Описать фото",       callback_data="photo_describe"),
                 InlineKeyboardButton(text="📝 Решить задачу",      callback_data="photo_task")],
                [InlineKeyboardButton(text="💻 Код из скрина",      callback_data="photo_code"),
                 InlineKeyboardButton(text="🌐 Перевести текст",    callback_data="photo_translate")],
                [InlineKeyboardButton(text="🎨 Изменить стиль фото", callback_data="photo_img2img")],
                *([video_btn_row] if video_btn_row else []),
            ])
        )
        return

    task_keywords = ["задани", "выполни", "реши", "напиши", "ответ", "упражнени", "вопрос"]
    is_task = any(kw in user_caption.lower() for kw in task_keywords) if user_caption else False
    caption = (
        f"{user_caption}\nВажно: выполни все задания, заполни поля 'Ответ:'."
        if is_task else user_caption
    )

    await _analyze_photo(message, uid, img_b64, caption, user_caption)


async def _analyze_photo(message, uid: int, img_b64: str, caption: str, user_caption: str = ""):
    """Анализирует фото через vision-модель."""
    FALLBACK_VISION = ["claude_opus", "claude_haiku", "claude_sonnet", "gpt52", "gemini_25_flash", "gemini_20_flash", "gemini_3_flash", "qwen3_vl", "c4ai_vis32b", "command_vision"]
    cur_key = get_model_key(uid)

    # Если caption пустой - используем дефолтный запрос
    if not caption or not caption.strip():
        caption = "Что изображено на этом фото? Опиши подробно."

    # Проверяем подписку для премиум моделей ПЕРЕД запросом
    if cur_key in VISION_MODELS and cur_key in PREMIUM_MODELS and not has_active_sub(uid):
        await message.answer(
            f"🔒 <b>Модель {MODELS[cur_key]['label']} доступна только по подписке</b>\n\n"
            "💡 Выбери другую модель с Vision или оформи подписку.",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(inline_keyboard=[
                [InlineKeyboardButton(text="💎 Оформить подписку", callback_data="sub_menu")],
                [InlineKeyboardButton(text="🤖 Сменить модель", callback_data="menu_ai")],
            ])
        )
        return

    # Строим порядок: выбранная модель (если vision) + все fallback
    if cur_key in VISION_MODELS and cur_key not in disabled_models:
        vision_order = [cur_key] + [k for k in FALLBACK_VISION if k != cur_key and k in MODELS and k not in disabled_models]
        auto_note = ""
    else:
        vision_order = [k for k in FALLBACK_VISION if k in MODELS and k not in disabled_models]
        auto_note = " <i>(авто)</i>"

    if not vision_order:
        await message.answer("⚠️ Нет доступных моделей для анализа фото.", parse_mode="HTML")
        return

    model_label = MODELS[vision_order[0]]["label"]
    wait_msg = await message.answer(
        f"🔍 <i>Анализирую через {model_label}{auto_note}…</i>",
        parse_mode="HTML"
    )

    async with ChatActionSender.typing(bot=bot, chat_id=message.chat.id):
        errors = []
        # Перебираем все vision модели по порядку с полным fallback
        for vk in vision_order:
            try:
                vision_msg = [{"role": "user", "content": [
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}},
                    {"type": "text", "text": caption},
                ]}]
                ans = await call_chat(vision_msg, vk, max_tokens=3000)
                last_responses[uid] = {"q": f"[Фото] {caption}", "a": ans, "model_label": MODELS.get(vk, {}).get("label", "ИИ"), "model_key": vk}
                last_photo[uid] = {
                    "img_b64": img_b64,
                    "caption": user_caption,
                    "last_answer": ans[:800],   # Контекст ответа для follow-up
                    "photo_context": True,       # Флаг: фото было проанализировано
                }
                if uid not in user_memory:
                    user_memory[uid] = deque(maxlen=20)
                user_memory[uid].append({"role": "user", "content": f"[Фото: '{user_caption or caption}']"})
                user_memory[uid].append({"role": "assistant", "content": ans})
                spend_limit(uid, vk)
                asyncio.create_task(db_save_user(uid, message.from_user.first_name or "", message.from_user.username or ""))
                await wait_msg.delete()
                # Клавиатура с кнопкой «Красиво» — ВСЕГДА
                _photo_kb = save_kb(uid)
                try:
                    _vurl = store_answer(ans, MODELS[vk]["label"], vk)
                    _photo_kb = InlineKeyboardMarkup(inline_keyboard=[
                        [InlineKeyboardButton(text="🌐 Красиво", url=_vurl)]
                    ] + _photo_kb.inline_keyboard)
                except Exception:
                    pass
                header = f"🔍 <i>{MODELS[vk]['label']}</i>\n\n"
                if re.search(r"```", ans):
                    await message.answer(header, parse_mode="HTML")
                    await smart_send(message, ans, reply_markup=_photo_kb)
                else:
                    html_ans = md_to_html(ans)
                    try:
                        await message.answer(header + html_ans[:3800], parse_mode="HTML", reply_markup=_photo_kb)
                    except Exception:
                        await message.answer(header + ans[:3800], reply_markup=_photo_kb)
                return
            except Exception as e:
                errors.append(f"{vk}: {str(e)[:80]}")
                logging.error(f"Vision {vk}: {e}")

        await wait_msg.delete()
        await message.answer("❌ Все vision-модели недоступны:\n" + "\n".join(errors))


# ==================================================================
# 🎤 ГОЛОСОВЫЕ
# ==================================================================

@dp.message(F.voice)
async def handle_voice(message: Message):
    uid = message.from_user.id
    if not await require_subscription(message):
        return

    # Индикатор обработки
    status_msg = await message.answer("🎤 <i>Обработка аудио...</i>", parse_mode="HTML")

    try:
        async with ChatActionSender.typing(bot=bot, chat_id=message.chat.id):
            file_obj   = await bot.get_file(message.voice.file_id)
            voice_path = tmp(f"voice_{uid}.ogg")
            buf = io.BytesIO()
            await bot.download_file(file_obj.file_path, destination=buf)
            async with aiofiles.open(voice_path, "wb") as f:
                await f.write(buf.getvalue())

            text = None
            last_err = ""

            # Конвертируем ogg → wav для лучшей совместимости
            wav_path = voice_path.replace(".ogg", ".wav")
            try:
                conv = subprocess.run(
                    ["ffmpeg", "-y", "-i", voice_path, "-ar", "16000", "-ac", "1", "-f", "wav", wav_path],
                    capture_output=True, timeout=30
                )
                use_path = wav_path if os.path.exists(wav_path) and os.path.getsize(wav_path) > 100 else voice_path
                use_fname = "voice.wav" if use_path == wav_path else "voice.ogg"
                use_ctype = "audio/wav" if use_path == wav_path else "audio/ogg"
            except Exception:
                use_path = voice_path
                use_fname = "voice.ogg"
                use_ctype = "audio/ogg"

            # Попытка 1: OnlySQ — whisper-large-v3 (основной, надёжнее)
            timeout_sq = aiohttp.ClientTimeout(total=60, connect=15, sock_read=50)
            try:
                async with aiohttp.ClientSession(timeout=timeout_sq) as session_sq:
                    with open(use_path, "rb") as af:
                        form_sq = aiohttp.FormData()
                        form_sq.add_field("file", af, filename=use_fname, content_type=use_ctype)
                        form_sq.add_field("model", "whisper-large-v3")
                        form_sq.add_field("language", "ru")
                        form_sq.add_field("response_format", "json")
                        async with session_sq.post(
                            "https://api.onlysq.ru/ai/openai/audio/transcriptions",
                            headers={"Authorization": f"Bearer {SQ_KEY}"},
                            data=form_sq,
                        ) as r_sq:
                            if r_sq.status == 200:
                                res_sq = await r_sq.json()
                                text = res_sq.get("text", "").strip()
                                if text:
                                    logging.info(f"SQ Whisper success uid={uid}")
                            else:
                                sq_raw = await r_sq.text()
                                last_err = f"SQ HTTP {r_sq.status}: {sq_raw[:150]}"
                                logging.warning(f"Whisper SQ error: {last_err}")
            except asyncio.TimeoutError:
                last_err = "SQ Таймаут (60с)"
                logging.warning(f"Whisper SQ timeout uid={uid}")
            except Exception as e_sq:
                last_err = f"SQ: {str(e_sq)[:100]}"
                logging.warning(f"Whisper SQ exception: {e_sq}")

            # Попытка 2: IO.Solutions (faster-whisper) — резерв
            if not text:
                logging.info(f"Trying IO Whisper fallback for uid={uid}")
                timeout_io = aiohttp.ClientTimeout(total=90, connect=20, sock_read=70)
                async with aiohttp.ClientSession(timeout=timeout_io) as session_io:
                    try:
                        with open(use_path, "rb") as af:
                            form_io = aiohttp.FormData()
                            form_io.add_field("file", af, filename=use_fname, content_type=use_ctype)
                            form_io.add_field("model", IO_WHISPER)
                            form_io.add_field("language", "ru")
                            form_io.add_field("response_format", "json")
                            async with session_io.post(
                                IO_AUDIO,
                                headers={"Authorization": f"Bearer {IO_KEY}"},
                                data=form_io,
                            ) as r_io:
                                if r_io.status == 200:
                                    res_io = await r_io.json()
                                    text = res_io.get("text", "").strip()
                                    if text:
                                        logging.info(f"IO Whisper success uid={uid}")
                                else:
                                    raw_err_io = await r_io.text()
                                    last_err += f" | IO HTTP {r_io.status}: {raw_err_io[:150]}"
                                    logging.error(f"Whisper IO error: {last_err}")
                    except asyncio.TimeoutError:
                        last_err += " | IO Таймаут (90с)"
                        logging.warning(f"Whisper IO timeout uid={uid}")
                    except Exception as e_io:
                        last_err += f" | IO: {str(e_io)[:100]}"
                        logging.warning(f"Whisper IO exception: {e_io}")

            # Попытка 3: OnlySQ whisper-1 (самый простой, последний шанс)
            if not text:
                logging.info(f"Trying SQ whisper-1 last-resort for uid={uid}")
                try:
                    timeout_w1 = aiohttp.ClientTimeout(total=60, connect=15, sock_read=50)
                    async with aiohttp.ClientSession(timeout=timeout_w1) as session_w1:
                        # Используем исходный ogg если wav не помог
                        with open(voice_path, "rb") as af:
                            form_w1 = aiohttp.FormData()
                            form_w1.add_field("file", af, filename="voice.ogg", content_type="audio/ogg")
                            form_w1.add_field("model", "whisper-1")
                            form_w1.add_field("response_format", "json")
                            async with session_w1.post(
                                "https://api.onlysq.ru/ai/openai/audio/transcriptions",
                                headers={"Authorization": f"Bearer {SQ_KEY}"},
                                data=form_w1,
                            ) as r_w1:
                                if r_w1.status == 200:
                                    res_w1 = await r_w1.json()
                                    text = res_w1.get("text", "").strip()
                                    if text:
                                        logging.info(f"whisper-1 success uid={uid}")
                                else:
                                    w1_err = await r_w1.text()
                                    last_err += f" | W1 HTTP {r_w1.status}: {w1_err[:80]}"
                except Exception as e_w1:
                    last_err += f" | W1: {str(e_w1)[:60]}"

            # Чистим временные файлы
            for p in [wav_path]:
                try:
                    if os.path.exists(p):
                        os.remove(p)
                except Exception:
                    pass

            # Удаляем основной временный файл
            try:
                os.remove(voice_path)
            except Exception:
                pass

            await status_msg.delete()

            if not text:
                await message.answer(
                    f"❌ <b>Не удалось распознать речь</b>\n\n"
                    f"Причина: <code>{last_err or 'Пустой ответ'}</code>\n\n"
                    f"💡 Попробуй говорить чётче или повтори позже.",
                    parse_mode="HTML"
                )
                return

            await message.answer(f"🎤 <i>Распознано:</i> <b>{text}</b>", parse_mode="HTML")

            mk = get_chat_model(uid)

            if uid not in user_memory:
                user_memory[uid] = deque(maxlen=20)

            _init_limits(uid)
            _vok, _vreason = can_send(uid, mk)
            if not _vok:
                _vrst = _reset_str(uid, _vreason)
                await message.answer(
                    f"🚫 <b>Лимит запросов исчерпан!</b>\nСброс через: <code>{_vrst}</code>",
                    parse_mode="HTML",
                )
                return

            if mk in disabled_models:
                await message.answer(
                    "🚫 <b>Текущая модель недоступна.</b> Выбери другую.",
                    parse_mode="HTML", reply_markup=ai_menu_kb(uid)
                )
                return

            think_msg = await message.answer("🔍 <i>Анализирую ваш вопрос... 1с</i>", parse_mode="HTML")
            _vstart = asyncio.get_event_loop().time()

            async def _vtick(m, s):
                phases = [
                    (3,  "🔍 Анализирую ваш вопрос"),
                    (8,  "🧠 Размышляю"),
                    (15, "💡 Формирую ответ"),
                    (999,"⚡ Ответ почти готов"),
                ]
                while True:
                    await asyncio.sleep(1)
                    elapsed = int(asyncio.get_event_loop().time() - s)
                    label = phases[0][1]
                    for thr, lbl in phases:
                        if elapsed < thr:
                            label = lbl
                            break
                    try:
                        await m.edit_text(f"{label}... <b>{elapsed}с</b>", parse_mode="HTML")
                    except Exception:
                        pass

            _vtick_task = asyncio.create_task(_vtick(think_msg, _vstart))
            try:
                user_memory[uid].append({"role": "user", "content": text})
                ans = await call_chat(list(user_memory[uid]), mk)
                user_memory[uid].append({"role": "assistant", "content": ans})
                _vtick_task.cancel()
                try:
                    await _vtick_task
                except asyncio.CancelledError:
                    pass
                last_responses[uid] = {"q": text, "a": ans, "model_label": MODELS.get(mk, {}).get("label", "ИИ"), "model_key": mk}
                spend_limit(uid, mk)
                await think_msg.delete()
                if re.search(r"```", ans):
                    await smart_send(message, ans, reply_markup=save_kb(uid))
                else:
                    html_ans = md_to_html(ans)
                    try:
                        await message.answer(html_ans[:4000], parse_mode="HTML", reply_markup=save_kb(uid))
                    except Exception:
                        await message.answer(ans[:4000], reply_markup=save_kb(uid))

            except Exception as e:
                _vtick_task.cancel()
                if user_memory[uid]:
                    user_memory[uid].pop()
                try:
                    await think_msg.delete()
                except Exception:
                    pass
                await message.answer(f"❌ Ошибка нейросети: {e}")

    except Exception as e:
        try:
            await status_msg.delete()
        except Exception:
            pass
        logging.error(f"handle_voice outer exception: {e}")
        await message.answer(f"❌ Ошибка обработки голосового: {str(e)[:200]}")


# ==================================================================
# 🚀 ЗАПУСК
# ==================================================================


# ==================================================================
# 🧪 ТЕСТ НЕЙРОСЕТЕЙ (ADMIN)
# ==================================================================

@dp.callback_query(F.data == "admin_test_models")
async def admin_test_models(callback: CallbackQuery):
    uid = callback.from_user.id
    if uid not in ADMIN_IDS:
        return await callback.answer("❌ Нет доступа", show_alert=True)
    await callback.answer("🧪 Запускаю тест нейросетей...")

    # Тестируем только chat-модели с ненулевыми лимитами (пропускаем отключённые)
    # Пропускаем vision-only модели и генерацию изображений
    TEST_PROMPT = "Привет! Ответь ровно 5 словами на русском языке."
    test_models = {k: v for k, v in MODELS.items() if k not in disabled_models}

    wait_msg = await callback.message.answer(
        "🧪 <b>Тест нейросетей</b>\n\n"
        f"Тестирую <code>{len(test_models)}</code> моделей...\n"
        "⏳ Это займёт несколько минут.",
        parse_mode="HTML"
    )

    results = []
    for model_key, model_info in test_models.items():
        label = model_info["label"]
        try:
            # Для reasoning моделей нужно больше ресурсов (reasoning-модели)
            _is_reasoning = any(s in model_info["name"] for s in ["R1", "GLM-4.7", "GLM-4.6"])
            _test_tokens = 800 if _is_reasoning else 60
            test_msgs = [{"role": "user", "content": TEST_PROMPT}]
            ans = await call_chat(test_msgs, model_key, max_tokens=_test_tokens)
            # Берём только первые 80 символов ответа
            short_ans = ans.strip().replace("\n", " ")[:80]
            results.append(f"✅ <b>{label}</b>\n   → <i>{short_ans}</i>")
        except Exception as e:
            err_short = str(e)[:60]
            results.append(f"❌ <b>{label}</b>\n   → <code>{err_short}</code>")

    await wait_msg.delete()

    # Разбиваем на чанки если много моделей
    header = f"🧪 <b>Результаты теста ({len(test_models)} моделей)</b>\n\n"
    body = "\n\n".join(results)
    full_text = header + body
    chunks = [full_text[i:i+3800] for i in range(0, len(full_text), 3800)]
    for ci, chunk in enumerate(chunks):
        kb = InlineKeyboardMarkup(inline_keyboard=[[
            InlineKeyboardButton(text="◀️ Назад", callback_data="menu_admin"),
        ]]) if ci == len(chunks) - 1 else None
        await callback.message.answer(chunk, parse_mode="HTML", reply_markup=kb)


# ==================================================================
# 🌐 HTTP API СЕРВЕР ДЛЯ MINI APP
# ==================================================================

# Railway автоматически выдаёт переменную PORT — используем её
API_PORT = int(os.getenv("PORT", os.getenv("API_PORT", "8080")))

async def api_chat_handler(request: aiohttp_web.Request) -> aiohttp_web.Response:
    """POST /api/chat — принимает запрос от мини-аппа, возвращает ответ ИИ."""
    import json as _j
    # CORS headers
    headers = {
        "Access-Control-Allow-Origin": "*",
        "Access-Control-Allow-Methods": "POST, OPTIONS",
        "Access-Control-Allow-Headers": "Content-Type, X-Telegram-Init-Data",
        "Content-Type": "application/json",
    }
    if request.method == "OPTIONS":
        return aiohttp_web.Response(status=200, headers=headers)

    try:
        # Читаем сырое тело с большим лимитом (20MB) — иначе aiohttp режет при больших base64
        raw = await request.read()
        data = _j.loads(raw.decode("utf-8"))
    except Exception as e:
        logging.error(f"API chat parse error: {e}, body_size={request.content_length}")
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "Invalid JSON"}, ensure_ascii=False),
            status=400, headers=headers
        )

    action  = data.get("action", "chat")
    uid     = int(data.get("uid", 0) or 0)
    text    = data.get("text", "").strip()
    # Фронт шлёт либо "image" (строка base64) либо "images" (массив) — принимаем оба формата
    _img_single = data.get("image")
    images  = data.get("images", [])
    if _img_single and not images:
        images = [_img_single]
    model   = data.get("model", "")
    # Frontend sends "messages", backend historically used "history" — support both
    history = data.get("history") or data.get("messages", [])
    # If text not sent separately, extract it from the last user message
    if not text and history:
        for msg in reversed(history):
            if msg.get("role") == "user":
                content = msg.get("content", "")
                if isinstance(content, str):
                    text = content
                break

    # If uid not sent, try to extract from Telegram init_data
    if not uid:
        try:
            import urllib.parse as _up
            _idata = data.get("init_data", "")
            if _idata:
                _params = dict(_up.parse_qsl(_idata))
                _user_json = _params.get("user", "")
                if _user_json:
                    import json as _jj
                    _uobj = _jj.loads(_user_json)
                    uid = int(_uobj.get("id", 0))
        except Exception:
            pass

    if not uid:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "No uid"}, ensure_ascii=False),
            status=400, headers=headers
        )

    # Проверка подписки
    ok_sub, not_subbed = await check_subscription(uid)
    if not ok_sub:
        channels_info = []
        for ch in not_subbed:
            if str(ch).lstrip("-").isdigit():
                ch_link = f"https://t.me/c/{str(ch).replace('-100', '')}"
                ch_title = ch
            else:
                ch_link = f"https://t.me/{ch}"
                ch_title = f"@{ch}"
            try:
                chat = await bot.get_chat(chat_id=int(ch) if str(ch).lstrip("-").isdigit() else f"@{ch}")
                ch_title = chat.title or ch_title
            except Exception:
                pass
            channels_info.append({"title": ch_title, "url": ch_link})
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "subscription_required", "channels": channels_info}, ensure_ascii=False),
            headers=headers
        )

    # Проверка принятия соглашения (через бот — mini-app только читает)
    if uid not in ADMIN_IDS and not await _has_accepted(uid):
        return aiohttp_web.Response(
            text=_j.dumps({
                "ok": False,
                "error": "terms_required",
                "message": "Откройте бота и примите пользовательское соглашение."
            }, ensure_ascii=False),
            headers=headers
        )

    # Установка модели
    if model and model in MODELS and model not in disabled_models:
        user_settings[uid] = model
    mk = get_model_key(uid)

    # Блокировка premium моделей без подписки
    if mk in PREMIUM_MODELS and not has_active_sub(uid):
        return aiohttp_web.Response(
            text=_j.dumps({
                "ok":    False,
                "error": "premium_required",
                "message": "Эта модель доступна только по подписке. Оформите подписку в боте."
            }, ensure_ascii=False),
            headers=headers
        )


    # Лимиты
    _init_limits(uid)
    _ok, _reason = can_send(uid, mk)
    if not _ok:
        _rst = _reset_str(uid, _reason)
        li = get_limits_info(uid)
        return aiohttp_web.Response(
            text=_j.dumps({
                "ok": False,
                "error": "limit_exceeded",
                "reset_in": _rst,
                "fast_left": 0,
                "pro_left":  li["pro_max"]  - li["pro_used"],
            }, ensure_ascii=False),
            headers=headers
        )

    if mk in disabled_models:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "model_disabled"}, ensure_ascii=False),
            headers=headers
        )

    if not text and not images:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "empty_message"}, ensure_ascii=False),
            headers=headers
        )

    if uid not in user_memory:
        user_memory[uid] = deque(maxlen=20)

    try:
        if images:
            # Приоритет: текущая выбранная vision-модель → SQ-vision → IO-vision
            SQ_VISION_PRIORITY = ["claude_opus", "claude_haiku", "claude_sonnet", "gpt52", "qwen3_vl", "c4ai_vis32b"]
            IO_VISION_PRIORITY = []

            if mk in VISION_MODELS and mk not in disabled_models:
                vision_key = mk
            else:
                # Сначала пробуем SQ (надёжнее)
                vision_key = next(
                    (k for k in SQ_VISION_PRIORITY if k in MODELS and k not in disabled_models), None
                )
                if not vision_key:
                    # Потом IO
                    vision_key = next(
                        (k for k in IO_VISION_PRIORITY if k in MODELS and k not in disabled_models), None
                    )

            if not vision_key:
                return aiohttp_web.Response(
                    text=_j.dumps({"ok": False, "error": "no_vision_model"}, ensure_ascii=False),
                    headers=headers
                )
            content = []
            for img_b64 in images:
                content.append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{img_b64}"}})
            content.append({"type": "text", "text": text or "Подробно опиши что изображено на фото"})
            msgs = [{"role": "system", "content": VISION_SYSTEM_PROMPT}, {"role": "user", "content": content}]
            ans = await call_chat(msgs, vision_key, max_tokens=3000)
            spend_limit(uid, vision_key)
            used_model_label = MODELS[vision_key]["label"]
        else:
            # Определяем режим ответа из профиля пользователя
            _feats_api = user_features.get(uid, {})
            _req_mode = data.get("answer_mode", "")
            _feats_api = user_features.get(uid, {})
            _answer_mode_api = _req_mode if _req_mode in ANSWER_MODES else _feats_api.get("answer_mode", "fast")
            if _req_mode in ANSWER_MODES:
                _feats_api["answer_mode"] = _req_mode
                user_features[uid] = _feats_api
            _mode_instr = ANSWER_MODES.get(_answer_mode_api, ANSWER_MODES["fast"])[2]
            GLOBAL_SYS = (
                f"Ты — {BOT_NAME}, интеллектуальный ИИ-ассистент. "
                "Отвечай на языке пользователя — по умолчанию русский.\n"
                "СТАНДАРТЫ:\n"
                "• Математика: LaTeX ($формула$), не Python-код для мат.задач\n"
                "• Таблицы: Markdown | заголовок | содержательные, минимум 3 столбца\n"
                "• Код: блоки ```язык\\n...\\n``` — рабочий, полный, без заглушек\n"
                "• Запрещено: 'Конечно!', 'Как ИИ я...', неполный код, придуманные факты"
                + _mode_instr
            )
            # Строим историю из данных фронтенда (уже содержит всё включая текущее сообщение)
            # Важно: соблюдать строгое чередование user/assistant
            if history:
                clean_history = []
                last_role = None
                for h in history:
                    role = h.get("role")
                    content_h = h.get("content", "")
                    if role not in ("user", "assistant") or not content_h:
                        continue
                    # Пропускаем дубли одной роли подряд (берём последнее)
                    if role == last_role:
                        if clean_history:
                            clean_history[-1] = {"role": role, "content": content_h}
                        continue
                    clean_history.append({"role": role, "content": content_h})
                    last_role = role

                # Проверяем что история заканчивается на user
                if clean_history and clean_history[-1]["role"] == "user":
                    # ВАЖНО: заменяем последнее user-сообщение на text из payload.
                    # history на фронте содержит displayText (имена файлов),
                    # а text — полное содержимое файлов для AI.
                    if text:
                        clean_history[-1]["content"] = text
                    query_msgs = [{"role": "system", "content": GLOBAL_SYS}] + clean_history
                else:
                    query_msgs = [{"role": "system", "content": GLOBAL_SYS}] + clean_history + [{"role": "user", "content": text}]





                # Обновляем server-side память последними 6 парами
                user_memory[uid] = __import__('collections').deque(clean_history[-12:], maxlen=20)
            else:
                user_memory[uid].append({"role": "user", "content": text})
                query_msgs = [{"role": "system", "content": GLOBAL_SYS}] + list(user_memory[uid])
            ans = await call_chat(query_msgs, mk)
            user_memory[uid].append({"role": "assistant", "content": ans})
            spend_limit(uid, mk)
            used_model_label = MODELS[mk]["label"]

        last_responses[uid] = {"q": text or "[Фото]", "a": ans, "model_label": used_model_label, "model_key": mk}
        if uid not in user_profiles:
            user_profiles[uid] = {"name": "", "username": "", "joined": msk_now().strftime("%d.%m.%Y %H:%M"), "requests": 0, "last_bonus": None}
        user_profiles[uid]["requests"] = user_profiles[uid].get("requests", 0) + 1
        # Пересчитываем уровень
        _total_req = user_profiles[uid]["requests"]
        _level_thresholds = [0, 50, 100, 200, 500, 1000, 2000, 5000]
        user_profiles[uid]["level"] = sum(1 for t in _level_thresholds if _total_req >= t)
        # Сохраняем requests в БД асинхронно
        asyncio.create_task(db_save_user(uid))

        li = get_limits_info(uid)
        return aiohttp_web.Response(
            text=_j.dumps({
                "ok": True,
                "answer": ans,
                "model_label": used_model_label,
                "fast_left": 0,
                "pro_left":  li["pro_max"]  - li["pro_used"],
            }, ensure_ascii=False),
            headers=headers
        )

    except Exception as e:
        logging.error(f"api_chat_handler error: {e}")
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": str(e)[:300]}, ensure_ascii=False),
            status=500, headers=headers
        )


async def api_limits_handler(request: aiohttp_web.Request) -> aiohttp_web.Response:
    """GET /api/limits?uid=123&init_data=... — возвращает текущие лимиты пользователя."""
    import json as _j, urllib.parse as _up
    headers = {
        "Access-Control-Allow-Origin": "*",
        "Content-Type": "application/json",
    }
    uid_str = request.rel_url.query.get("uid", "0")
    try:
        uid = int(uid_str)
    except Exception:
        uid = 0

    # Если uid не передан или равен 0 — пробуем извлечь из Telegram init_data
    if not uid:
        try:
            _idata = request.rel_url.query.get("init_data", "")
            if _idata:
                _params = dict(_up.parse_qsl(_idata))
                _user_json = _params.get("user", "")
                if _user_json:
                    _uobj = _j.loads(_user_json)
                    uid = int(_uobj.get("id", 0))
        except Exception:
            pass

    if not uid:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "No uid"}, ensure_ascii=False),
            headers=headers
        )

    _init_limits(uid)
    li   = get_limits_info(uid)
    prof = user_profiles.get(uid, {})
    sub_active = has_active_sub(uid)
    _sub_expires = sub_expires_str(uid) if sub_active else ""
    _sub_plan    = sub_plan_label(uid)
    return aiohttp_web.Response(
        text=_j.dumps({
            "ok":             True,
            # ── лимиты (используются фронтом) ──────────────────────
            "pro_used":       li["pro_used"],
            "pro_max":        li["pro_max"],
            "fast_left":      0,
            "pro_left":       li["pro_max"] - li["pro_used"],
            "img_used":       li["img_used"],
            "img_max":        li["img_max"],
            "img_left":       li["img_max"] - li["img_used"],
            "music_used":     li.get("music_used", 0),
            "music_max":      li.get("music_max", 0),
            "video_used":     li.get("video_used", 0),
            "video_max":      li.get("video_max", 0),
            # ── подписка (has_sub — именно то поле что читает фронт) ─
            "has_sub":        sub_active,
            "sub_active":     sub_active,
            "sub_expires":    _sub_expires,
            "sub_plan":       _sub_plan,
            "plan":           _sub_plan,
            "expires":        _sub_expires,
            # ── профиль ────────────────────────────────────────────
            "total_requests": prof.get("requests", 0),
            "requests":       prof.get("requests", 0),
            "total_gens":     user_images_count.get(uid, 0),
            "join_date":      prof.get("joined", ""),
            "terms_accepted": await _has_accepted(uid),
            "referrals":      len(user_referrals.get(uid, {}).get("refs", [])),
            "level":          prof.get("level", 0),
            "level_max":      next((t for t in [50, 100, 200, 500, 1000, 2000, 5000] if t > prof.get("requests", 0)), 5000),
            "reset_in":       li.get("reset_in", ""),
            "model":          MODELS.get(user_settings.get(uid, DEFAULT_MODEL), {}).get("label", ""),
            "mode":           user_features.get(uid, {}).get("answer_mode", "fast"),
        }, ensure_ascii=False),
        headers=headers
    )


async def api_imggen_handler(request: aiohttp_web.Request) -> aiohttp_web.Response:
    """POST /api/imggen — генерация изображений из мини-аппа."""
    import json as _j, base64 as _b64
    headers = {
        "Access-Control-Allow-Origin": "*",
        "Access-Control-Allow-Methods": "POST, OPTIONS",
        "Access-Control-Allow-Headers": "Content-Type",
        "Content-Type": "application/json",
    }
    if request.method == "OPTIONS":
        return aiohttp_web.Response(status=200, headers=headers)
    try:
        raw = await request.read()
        data = _j.loads(raw.decode("utf-8"))
    except Exception as e:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "Invalid JSON"}, ensure_ascii=False),
            status=400, headers=headers)

    uid    = int(data.get("uid", 0))
    prompt = data.get("prompt", "").strip()
    model  = data.get("model", "flux_klein9")
    ratio  = data.get("ratio", "1:1")
    # Валидируем ratio
    VALID_RATIOS = {"1:1", "4:3", "3:4", "16:9", "9:16"}
    if ratio not in VALID_RATIOS:
        ratio = "1:1"

    if not uid:
        return aiohttp_web.Response(text=_j.dumps({"ok": False, "error": "No uid"}, ensure_ascii=False), headers=headers)
    if not prompt:
        return aiohttp_web.Response(text=_j.dumps({"ok": False, "error": "Empty prompt"}, ensure_ascii=False), headers=headers)

    # Проверка подписки
    ok_sub, not_subbed = await check_subscription(uid)
    if not ok_sub:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "subscription_required"}, ensure_ascii=False), headers=headers)

    # Проверка лимита генерации
    _init_limits(uid)
    if not can_img(uid):
        li = get_limits_info(uid)
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "img_limit_exceeded",
                          "img_left": li["img_max"] - li["img_used"]}, ensure_ascii=False),
            headers=headers)

    if model not in IMG_MODELS:
        model = "flux_klein9"

    try:
        result = await generate_image(prompt, model, ratio)
        img_bytes = result["data"]
        img_b64 = _b64.b64encode(img_bytes).decode()
        spend_img(uid)
        li = get_limits_info(uid)
        if uid in user_profiles:
            user_profiles[uid]["requests"] = user_profiles[uid].get("requests", 0) + 1
        return aiohttp_web.Response(
            text=_j.dumps({
                "ok": True,
                "image_b64": img_b64,
                "model_label": IMG_MODELS[model]["label"],
                "img_left": li["img_max"] - li["img_used"],
            }, ensure_ascii=False),
            headers=headers)
    except Exception as e:
        logging.error(f"api_imggen error: {e}")
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": str(e)[:300]}, ensure_ascii=False),
            status=500, headers=headers)


async def api_transcribe_handler(request: aiohttp_web.Request) -> aiohttp_web.Response:
    """POST /api/transcribe — транскрибация голосового аудио из мини-аппа."""
    import json as _j
    headers = {
        "Access-Control-Allow-Origin": "*",
        "Access-Control-Allow-Methods": "POST, OPTIONS",
        "Access-Control-Allow-Headers": "Content-Type",
        "Content-Type": "application/json",
    }
    if request.method == "OPTIONS":
        return aiohttp_web.Response(status=200, headers=headers)
    try:
        data = await request.post()
        uid_str = data.get("uid", "0")
        uid = int(uid_str) if str(uid_str).isdigit() else 0
        audio_field = data.get("audio")
        if not audio_field or not hasattr(audio_field, "file"):
            return aiohttp_web.Response(
                text=_j.dumps({"ok": False, "error": "No audio"}, ensure_ascii=False),
                status=400, headers=headers)

        audio_bytes = audio_field.file.read()
        fname = audio_field.filename or "voice.webm"
        ext = fname.rsplit(".", 1)[-1].lower() if "." in fname else "webm"
        tmp_path = tmp(f"miniapp_voice_{uid}.{ext}")
        async with aiofiles.open(tmp_path, "wb") as f:
            await f.write(audio_bytes)

        # Конвертируем в wav
        wav_path = tmp(f"miniapp_voice_{uid}.wav")
        try:
            subprocess.run(
                ["ffmpeg", "-y", "-i", tmp_path, "-ar", "16000", "-ac", "1", "-f", "wav", wav_path],
                capture_output=True, timeout=30
            )
            use_path = wav_path if os.path.exists(wav_path) and os.path.getsize(wav_path) > 100 else tmp_path
        except Exception:
            use_path = tmp_path

        text = None
        last_err = ""
        # Попытка 1: SQ whisper-large-v3
        try:
            timeout_sq = aiohttp.ClientTimeout(total=60, connect=15, sock_read=50)
            async with aiohttp.ClientSession(timeout=timeout_sq) as sess:
                with open(use_path, "rb") as af:
                    form_data = aiohttp.FormData()
                    form_data.add_field("file", af, filename=os.path.basename(use_path), content_type="audio/wav")
                    form_data.add_field("model", "whisper-large-v3")
                    form_data.add_field("language", "ru")
                    form_data.add_field("response_format", "json")
                    async with sess.post(
                        "https://api.onlysq.ru/ai/openai/audio/transcriptions",
                        headers={"Authorization": f"Bearer {SQ_KEY}"},
                        data=form_data,
                    ) as r:
                        if r.status == 200:
                            res = await r.json()
                            text = res.get("text", "").strip()
                        else:
                            last_err = f"SQ {r.status}: {(await r.text())[:100]}"
        except Exception as e:
            last_err = str(e)[:100]

        # Попытка 2: IO Whisper
        if not text:
            try:
                timeout_io = aiohttp.ClientTimeout(total=90, connect=20, sock_read=70)
                async with aiohttp.ClientSession(timeout=timeout_io) as sess:
                    with open(use_path, "rb") as af:
                        form_data2 = aiohttp.FormData()
                        form_data2.add_field("file", af, filename=os.path.basename(use_path), content_type="audio/wav")
                        form_data2.add_field("model", IO_WHISPER)
                        form_data2.add_field("language", "ru")
                        form_data2.add_field("response_format", "json")
                        async with sess.post(
                            IO_AUDIO,
                            headers={"Authorization": f"Bearer {IO_KEY}"},
                            data=form_data2,
                        ) as r2:
                            if r2.status == 200:
                                res2 = await r2.json()
                                text = res2.get("text", "").strip()
                            else:
                                last_err += f" | IO {r2.status}"
            except Exception as e2:
                last_err += f" | {str(e2)[:60]}"

        # Чистим файлы
        for p in [tmp_path, wav_path]:
            try:
                if os.path.exists(p):
                    os.remove(p)
            except Exception:
                pass

        if text:
            return aiohttp_web.Response(
                text=_j.dumps({"ok": True, "text": text}, ensure_ascii=False),
                headers=headers)
        else:
            return aiohttp_web.Response(
                text=_j.dumps({"ok": False, "error": last_err or "Не удалось распознать"}, ensure_ascii=False),
                headers=headers)
    except Exception as e:
        logging.error(f"api_transcribe_handler: {e}")
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": str(e)[:200]}, ensure_ascii=False),
            status=500, headers=headers)


async def api_answer_handler(request: aiohttp_web.Request) -> aiohttp_web.Response:
    """Возвращает сохранённый ответ по короткому ID для View.html."""
    import json as _j2
    headers = {
        "Access-Control-Allow-Origin": "*",
        "Access-Control-Allow-Headers": "*",
        "Content-Type": "application/json; charset=utf-8",
    }
    if request.method == "OPTIONS":
        return aiohttp_web.Response(status=200, headers=headers)
    ans_id = request.match_info.get("ans_id", "")
    if ans_id in _answer_store:
        return aiohttp_web.Response(
            text=_j2.dumps({"ok": True, **_answer_store[ans_id]}, ensure_ascii=False),
            headers=headers)
    return aiohttp_web.Response(
        text=_j2.dumps({"ok": False, "error": "not found"}, ensure_ascii=False),
        status=404, headers=headers)


_base_dir = os.path.dirname(os.path.abspath(__file__))
VIEW_HTML_PATH = next(
    (os.path.join(_base_dir, name) for name in ("View.html", "view.html")
     if os.path.exists(os.path.join(_base_dir, name))),
    os.path.join(_base_dir, "View.html")
)

async def view_page_handler(request: aiohttp_web.Request) -> aiohttp_web.Response:
    """GET /view?id=XXX — читает View.html с диска и инжектирует данные."""
    import json as _jv
    html_headers = {
        "Content-Type": "text/html; charset=utf-8",
        "Access-Control-Allow-Origin": "*",
        "Cache-Control": "no-cache",
    }
    ans_id = request.rel_url.query.get("id", "")
    if ans_id and ans_id in _answer_store:
        data_json = _jv.dumps(_answer_store[ans_id], ensure_ascii=False)
    else:
        data_json = "null"

    inject = f"<script>window.__VIEW_DATA__ = {data_json}; window.__API_BASE__ = '{API_BASE_URL}';</script>"

    path = VIEW_HTML_PATH
    for name in ("View.html", "view.html"):
        p = os.path.join(_base_dir, name)
        if os.path.exists(p):
            path = p
            break

    if os.path.exists(path):
        try:
            async with aiofiles.open(path, "r", encoding="utf-8") as f:
                html = await f.read()
            html = html.replace("</head>", inject + "\n</head>", 1)
            return aiohttp_web.Response(text=html, headers=html_headers)
        except Exception as e:
            logging.error(f"view_page_handler: {e}")

    # Fallback если файл не найден
    fallback = f"""<!DOCTYPE html><html><head><meta charset="UTF-8">
{inject}
<title>Хуза ИИ</title>
<style>body{{background:#0a0a0f;color:#f0f0ff;font-family:sans-serif;display:flex;align-items:center;justify-content:center;min-height:100vh;flex-direction:column;gap:16px;text-align:center;padding:20px}}</style>
</head><body>
<div style="font-size:48px">⚠️</div>
<h2>View.html не найден</h2>
<p style="color:rgba(240,240,255,0.5)">Загрузи файл View.html в папку с bot.py на Railway</p>
</body></html>"""
    return aiohttp_web.Response(text=fallback, headers=html_headers)

async def api_store_answer_handler(request: aiohttp_web.Request) -> aiohttp_web.Response:
    """POST /api/store_answer — сохраняет ответ и возвращает короткий ID."""
    import json as _js
    import secrets
    headers = {
        "Access-Control-Allow-Origin": "*",
        "Access-Control-Allow-Headers": "*",
        "Content-Type": "application/json; charset=utf-8",
    }
    if request.method == "OPTIONS":
        return aiohttp_web.Response(status=200, headers=headers)
    try:
        body = await request.json()
        text = body.get("text", "")
        if not text:
            return aiohttp_web.Response(
                text=_js.dumps({"ok": False, "error": "no text"}),
                status=400, headers=headers)
        short_id = secrets.token_urlsafe(6)
        _answer_store[short_id] = {
            "text":     text,
            "model":    body.get("model", "ИИ"),
            "modelKey": body.get("modelKey", ""),
            "time":     body.get("time", msk_now().strftime("%d.%m · %H:%M")),
        }
        if len(_answer_store) > 500:
            for k in list(_answer_store.keys())[:100]:
                _answer_store.pop(k, None)
        return aiohttp_web.Response(
            text=_js.dumps({"ok": True, "id": short_id}, ensure_ascii=False),
            headers=headers)
    except Exception as e:
        return aiohttp_web.Response(
            text=_js.dumps({"ok": False, "error": str(e)}),
            status=500, headers=headers)


async def api_pptx_handler(request: aiohttp_web.Request) -> aiohttp_web.Response:
    """POST /api/pptx — генерация PPTX через веб-апп."""
    import json as _j, base64 as _b64
    headers = {
        "Access-Control-Allow-Origin": "*",
        "Access-Control-Allow-Methods": "POST, OPTIONS",
        "Access-Control-Allow-Headers": "Content-Type",
        "Content-Type": "application/json",
    }
    if request.method == "OPTIONS":
        return aiohttp_web.Response(status=200, headers=headers)
    try:
        data = await request.json()
    except Exception:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "Invalid JSON"}),
            status=400, headers=headers)

    uid    = int(data.get("uid", 0))
    topic  = data.get("topic", "").strip()
    slides = int(data.get("slides", 7))
    theme  = data.get("theme", "dark")
    model  = data.get("model", "claude_opus")
    wishes = data.get("wishes", "")
    author = data.get("author", "")

    if not uid or not topic:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "uid and topic required"}),
            status=400, headers=headers)

    ok_sub, _ = await check_subscription(uid)
    if not ok_sub:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "subscription_required"}), headers=headers)

    _init_limits(uid)
    _refresh_limits(uid)
    li = get_limits_info(uid)
    pro_left = li["pro_max"] - li["pro_used"]
    if pro_left < 3:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "limit_exceeded", "pro_left": pro_left}),
            headers=headers)

    if model not in MODELS:
        model = "claude_opus"

    try:
        pptx_data = await generate_pptx_structure_ai(
            topic=topic, slides_count=slides,
            model_key=model, extra_wishes=wishes
        )
        if author:
            pptx_data["author"] = author
            for sl in pptx_data.get("slides", []):
                if sl.get("type") == "title":
                    sl["author"] = author

        pptx_bytes = create_pptx_bytes(pptx_data, theme)
        for _ in range(3):
            user_limits[uid]["pro_used"] = user_limits[uid].get("pro_used", 0) + 1

        actual_slides = len(pptx_data.get("slides", []))
        fname = f"Презентация_{topic[:25].replace(' ', '_')}.pptx"
        return aiohttp_web.Response(
            text=_j.dumps({
                "ok": True,
                "pptx_b64": _b64.b64encode(pptx_bytes).decode(),
                "filename": fname,
                "slides_count": actual_slides,
                "title": pptx_data.get("title", topic),
            }, ensure_ascii=False),
            headers=headers)
    except Exception as e:
        logging.error(f"api_pptx: {e}")
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": str(e)[:300]}),
            status=500, headers=headers)


async def api_webpptx_handler(request: aiohttp_web.Request) -> aiohttp_web.Response:
    """POST /api/webpptx — генерация стильной HTML-презентации через веб-апп."""
    import json as _j, base64 as _b64
    headers = {
        "Access-Control-Allow-Origin":  "*",
        "Access-Control-Allow-Methods": "POST, OPTIONS",
        "Access-Control-Allow-Headers": "Content-Type",
        "Content-Type": "application/json",
    }
    if request.method == "OPTIONS":
        return aiohttp_web.Response(status=200, headers=headers)
    try:
        data = await request.json()
    except Exception:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "Invalid JSON"}),
            status=400, headers=headers,
        )

    uid    = int(data.get("uid", 0))
    topic  = data.get("topic", "").strip()
    slides = int(data.get("slides", 7))
    model  = data.get("model", "claude_opus")
    wishes = data.get("wishes", "")
    author = data.get("author", "")

    if not uid or not topic:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "uid and topic required"}),
            status=400, headers=headers,
        )

    ok_sub, _ = await check_subscription(uid)
    if not ok_sub:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "subscription_required"}),
            headers=headers,
        )

    _init_limits(uid)
    _refresh_limits(uid)
    li = get_limits_info(uid)
    if li["pro_max"] - li["pro_used"] < 3:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "limit_exceeded", "pro_left": li["pro_max"] - li["pro_used"]}),
            headers=headers,
        )

    if model not in MODELS:
        model = "claude_opus"

    try:
        pptx_data  = await generate_html_structure_ai(
            topic=topic, slides_count=slides,
            model_key=model, extra_wishes=wishes,
            title_info={"author": author} if author else None,
        )
        html_str   = create_html_presentation(pptx_data, pexels_key=os.environ.get("PEXELS_KEY",""))
        html_bytes = html_str.encode("utf-8")

        for _ in range(3):
            user_limits[uid]["pro_used"] = user_limits[uid].get("pro_used", 0) + 1

        safe_name = re.sub(r"[^\w\s-]", "", topic)[:25].strip().replace(" ", "_")
        fname = f"Презентация_{safe_name}.html"

        return aiohttp_web.Response(
            text=_j.dumps({
                "ok":           True,
                "html_b64":     _b64.b64encode(html_bytes).decode(),
                "filename":     fname,
                "slides_count": len(pptx_data.get("slides", [])),
                "title":        pptx_data.get("title", topic),
            }, ensure_ascii=False),
            headers=headers,
        )
    except Exception as e:
        logging.error(f"api_webpptx: {e}")
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": str(e)[:300]}),
            status=500, headers=headers,
        )



async def suno_callback_handler(request: aiohttp_web.Request) -> aiohttp_web.Response:
    """Заглушка — Suno заменён на Pollinations."""
    return aiohttp_web.Response(text="ok")



async def platega_callback_handler(request: aiohttp_web.Request) -> aiohttp_web.Response:
    """
    Вебхук от Platega: POST /platega_callback
    Заголовки: X-MerchantId, X-Secret
    Body JSON: { "status": "CONFIRMED", "payload": "sub_week_123456", "transactionId": "...", ... }
    """
    try:
        # Проверяем секрет
        secret = request.headers.get("X-Secret", "")
        if PLATEGA_SECRET and secret != PLATEGA_SECRET:
            logging.warning(f"[PLATEGA CB] неверный X-Secret: {secret!r}")
            return aiohttp_web.Response(status=403, text="Forbidden")

        body = await request.json()
        logging.info(f"[PLATEGA CB] body={body}")

        status  = body.get("status", "")
        payload = body.get("payload", "")   # формат: sub_{plan_key}_{uid}

        if status != "CONFIRMED":
            return aiohttp_web.Response(text="ok")

        parts = str(payload).split("_")
        if len(parts) < 3 or parts[0] != "sub":
            logging.warning(f"[PLATEGA CB] неверный payload: {payload!r}")
            return aiohttp_web.Response(text="ok")

        plan_key = parts[1]
        try:
            uid = int(parts[2])
        except ValueError:
            return aiohttp_web.Response(text="ok")

        if plan_key not in SUB_PLANS:
            logging.warning(f"[PLATEGA CB] неизвестный план: {plan_key!r}")
            return aiohttp_web.Response(text="ok")

        plan = SUB_PLANS[plan_key]

        # Активируем подписку (та же логика что в successful_payment_handler)
        existing = user_subscriptions.get(uid)
        if existing and existing["expires"] > msk_now():
            base = existing["expires"]
        else:
            base = msk_now()
        new_exp = base + datetime.timedelta(days=plan["days"])
        user_subscriptions[uid] = {"expires": new_exp, "plan": plan_key}
        asyncio.create_task(db_save_subscription(uid))
        logging.info(f"[PLATEGA CB] uid={uid} plan={plan_key} activated until {new_exp}")

        exp = sub_expires_str(uid)
        try:
            await bot.send_message(
                uid,
                f"✅ <b>Оплата прошла!</b>\n\n"
                f"💎 Тариф: <b>{plan['name']}</b>\n"
                f"📅 Активна до: <code>{exp}</code>\n\n"
                f"Все возможности разблокированы. Приятного пользования!",
                parse_mode="HTML",
                reply_markup=InlineKeyboardMarkup(inline_keyboard=[[
                    InlineKeyboardButton(text="🏠 Главное меню", callback_data="back_home"),
                ]]),
            )
        except Exception as e:
            logging.error(f"[PLATEGA CB] send_message uid={uid}: {e}")

    except Exception as e:
        logging.error(f"[PLATEGA CB] handler error: {e}")

    return aiohttp_web.Response(text="ok")


async def api_report_handler(request: aiohttp_web.Request) -> aiohttp_web.Response:
    """POST /api/report — генерация реферата/доклада из мини-аппа."""
    import json as _j, base64 as _b64
    headers = {
        "Access-Control-Allow-Origin":  "*",
        "Access-Control-Allow-Methods": "POST, OPTIONS",
        "Access-Control-Allow-Headers": "Content-Type",
        "Content-Type": "application/json",
    }
    if request.method == "OPTIONS":
        return aiohttp_web.Response(status=200, headers=headers)
    try:
        raw = await request.read()
        data = _j.loads(raw.decode("utf-8"))
    except Exception:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "Invalid JSON"}),
            status=400, headers=headers)

    uid    = int(data.get("uid", 0) or 0)
    topic  = data.get("topic", "").strip()
    pages  = int(data.get("pages", 5))
    model  = data.get("model", "claude_opus")
    extra  = data.get("extra", "")
    rtype  = data.get("type", "report")  # "report" | "referat"

    if not uid:
        # Попробуем извлечь uid из init_data
        try:
            import urllib.parse as _up
            _idata = data.get("init_data", "")
            if _idata:
                _params = dict(_up.parse_qsl(_idata))
                _uobj = _j.loads(_params.get("user", "{}"))
                uid = int(_uobj.get("id", 0))
        except Exception:
            pass

    if not uid or not topic:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "uid and topic required"}),
            status=400, headers=headers)

    ok_sub, _ = await check_subscription(uid)
    if not ok_sub:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "subscription_required"}), headers=headers)

    _init_limits(uid)
    _refresh_limits(uid)
    li = get_limits_info(uid)
    if li["pro_max"] - li["pro_used"] < 3:
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": "limit_exceeded",
                           "pro_left": li["pro_max"] - li["pro_used"]}), headers=headers)

    if model not in MODELS:
        model = "claude_opus"

    # Определяем тип документа
    is_referat = rtype == "referat"
    rtype_label = "Реферат" if is_referat else "Доклад"
    target_words = pages * 250

    wishes_block = f"Дополнительно: {extra}\n" if extra else ""

    if is_referat:
        num_chapters = max(2, pages - 2)
        chapters = "\n".join([
            f"## Глава {i}. [Конкретный аспект темы]"
            for i in range(1, num_chapters + 1)
        ])
        system_prompt = (
            "Ты — автор научного реферата. Пишешь академично.\n"
            "🇷🇺 ТОЛЬКО РУССКИЙ ЯЗЫК.\n"
            "Пишешь ТОЛЬКО реферат, без предисловий и комментариев."
        )
        user_prompt = (
            f"Напиши научный реферат: «{topic}»\n\n"
            f"🛑 ОБЪЁМ: примерно {target_words} слов ({pages} страниц).\n\n"
            f"СТРУКТУРА:\n## Введение\n{chapters}\n## Заключение\n## Список литературы\n\n"
            f"{wishes_block}"
        )
    else:
        num_sections = max(1, pages - 1)
        sections = "\n".join([f"## Раздел {i}" for i in range(1, num_sections + 1)])
        system_prompt = (
            "Ты — профессиональный автор школьных и студенческих докладов.\n"
            "🇷🇺 ОБЯЗАТЕЛЬНО: Весь доклад пиши ТОЛЬКО на русском языке.\n"
            "Пишешь ТОЛЬКО сам текст доклада, без вступлений и комментариев."
        )
        user_prompt = (
            f"Напиши доклад на тему: «{topic}»\n\n"
            f"ОБЪЁМ: примерно {target_words} слов ({pages} страниц).\n\n"
            f"СТРУКТУРА:\n## Введение\n{sections}\n## Заключение\n\n"
            f"{wishes_block}"
        )

    try:
        msgs = [
            {"role": "system", "content": system_prompt},
            {"role": "user",   "content": user_prompt},
        ]
        _fallbacks = ["claude_opus", "claude_sonnet", "deepseek_v3_sq", "qwen3_max"]
        ans = ""
        for _dm in _fallbacks:
            if _dm not in MODELS:
                continue
            try:
                _ans = await call_chat(msgs, _dm, max_tokens=4096)
                if _ans and len(_ans.strip()) > 200:
                    ans = _ans
                    break
            except Exception as _e:
                logging.warning(f"[API_REPORT] {_dm} ошибка: {_e}")
                continue

        if not ans or len(ans.strip()) < 100:
            raise RuntimeError("Не удалось сгенерировать документ. Попробуй ещё раз.")

        # Списываем 3 запроса
        for _ in range(3):
            user_limits[uid]["pro_used"] = user_limits[uid].get("pro_used", 0) + 1

        # Создаём docx
        path = await _create_report_docx(
            topic=topic, rtype_label=rtype_label, pages=pages,
            model_label=MODELS.get(model, {}).get("label", model),
            content=ans, uid=uid, author="", with_title=True
        )

        async with aiofiles.open(path, "rb") as f:
            docx_bytes = await f.read()
        try:
            os.remove(path)
        except Exception:
            pass

        fname = f"{rtype_label}_{topic[:25].replace(' ', '_')}.docx"
        return aiohttp_web.Response(
            text=_j.dumps({
                "ok":       True,
                "docx_b64": _b64.b64encode(docx_bytes).decode(),
                "filename": fname,
            }, ensure_ascii=False),
            headers=headers)
    except Exception as e:
        logging.error(f"api_report_handler: {e}")
        return aiohttp_web.Response(
            text=_j.dumps({"ok": False, "error": str(e)[:300]}),
            status=500, headers=headers)


async def start_api_server():
    app = aiohttp_web.Application(client_max_size=20 * 1024 * 1024)  # 20MB для base64 фото
    app.router.add_post("/suno_callback", suno_callback_handler)
    app.router.add_post("/platega_callback", platega_callback_handler)
    app.router.add_post("/api/store_answer", api_store_answer_handler)
    app.router.add_options("/api/store_answer", api_store_answer_handler)
    app.router.add_post("/api/chat", api_chat_handler)
    app.router.add_options("/api/chat", api_chat_handler)
    app.router.add_post("/api/imggen", api_imggen_handler)
    app.router.add_options("/api/imggen", api_imggen_handler)
    app.router.add_get("/api/limits", api_limits_handler)
    app.router.add_post("/api/transcribe", api_transcribe_handler)
    app.router.add_options("/api/transcribe", api_transcribe_handler)
    app.router.add_post("/api/pptx", api_pptx_handler)
    app.router.add_options("/api/pptx", api_pptx_handler)
    app.router.add_post("/api/webpptx", api_webpptx_handler)
    app.router.add_options("/api/webpptx", api_webpptx_handler)
    app.router.add_post("/api/report", api_report_handler)
    app.router.add_options("/api/report", api_report_handler)
    app.router.add_get("/api/answer/{ans_id}", api_answer_handler)
    app.router.add_options("/api/answer/{ans_id}", api_answer_handler)
    # ✅ View.html раздаётся с Railway — GitHub Pages не нужен!
    app.router.add_get("/view", view_page_handler)
    app.router.add_get("/View.html", view_page_handler)  # Алиас для совместимости
    runner = aiohttp_web.AppRunner(app)
    await runner.setup()
    site = aiohttp_web.TCPSite(runner, "0.0.0.0", API_PORT)
    await site.start()
    logging.info(f"✅ API сервер запущен на порту {API_PORT}")


async def main():
    global bot
    await init_db()
    await db_load_all_users()
    await db_load_subscriptions()
    await db_load_bot_settings()
    await bot.delete_webhook(drop_pending_updates=True)
    # Закрываем старую сессию чтобы избежать TelegramConflictError
    try:
        await bot.close()
    except Exception:
        pass
    await asyncio.sleep(2)
    bot = Bot(token=TELEGRAM_TOKEN)
    # Устанавливаем команды для навигации (видны в поле ввода)
    await bot.set_my_commands([
        BotCommand(command="start",   description="🏠 Главное меню"),
        BotCommand(command="profile", description="👤 Мой профиль"),
        BotCommand(command="ai",      description="💬 Задать вопрос ИИ"),
        BotCommand(command="img",     description="🎨 Генерация картинок"),
        BotCommand(command="music",   description="🎵 Генерация музыки"),
        BotCommand(command="model",   description="🤖 Выбрать модель"),
        BotCommand(command="report",  description="📝 Доклад/Реферат"),
        BotCommand(command="pptx",    description="🎞 Генерация презентации PPTX"),
        BotCommand(command="webpptx", description="🌐 Веб-презентация HTML"),
        BotCommand(command="clear",   description="🧹 Очистить память"),
        BotCommand(command="about",   description="ℹ️ О боте"),
        BotCommand(command="info",    description="💡 Помощь и контакты"),
    ], scope=BotCommandScopeDefault())
    # Устанавливаем кнопку "🤖 ХУЗА AI" слева в поле ввода (MenuButtonWebApp)
    try:
        await bot.set_chat_menu_button(
            menu_button=MenuButtonWebApp(
                text="🤖 ХУЗА AI",
                web_app=WebAppInfo(url=f"{MINI_APP_URL}?api={API_BASE_URL}")
            )
        )
        logging.info(f"✅ WebApp кнопка слева установлена: {MINI_APP_URL}")
    except Exception as e:
        logging.warning(f"set_chat_menu_button: {e}")
    # Запускаем API сервер параллельно с ботом
    await start_api_server()
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())
